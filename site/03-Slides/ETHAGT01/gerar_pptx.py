#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT01: Arquitetura Cognitiva de Agentes LLM
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Paleta Etho ───
PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ───

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj=""):
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT01"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

T = 62  # Total

# ═══════════════════════════════════════
# SEÇÃO A — Abertura (1-6)
# ═══════════════════════════════════════

s = title_slide(
    "Arquitetura Cognitiva de Agentes LLM",
    "Universidade Etho · Especialização em Programação Agêntica\nFase A — Fundamentos Agênticos · 25 h",
    "ETHAGT01"
)
add_notes(s, "📖 Bem-vindos à primeira aula. Hoje estabelecemos o bloco fundamental — Augmented LLM em loop. Não vamos falar de frameworks; vamos falar de arquitetura.\n💡 Analogia: aprender a cozinhar vs aprender a usar uma panela. Hoje vocês aprendem a cozinhar.\n❓ 'Quantos já usaram framework de agentes em produção?'\n⚠️ Alunos chegam querendo LangGraph. Redirecionar: arquitetura > framework.\n➡️ Vamos aos objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: estabelecer o bloco fundamental — Augmented LLM em loop",
    "",
    "Objetivos específicos:",
    "1. Explicar a transição de 'LLM como oráculo' para 'controlador cognitivo'",
    "2. Decompor um agente em Perception · Brain · Planning · Action · Tool Use",
    "3. Implementar um agente ReAct do zero (sem framework)",
    "4. Diferenciar workflows de agentes e justificar quando usar cada",
    "5. Adicionar observabilidade mínima desde a primeira implementação",
    "6. Avaliar criticamente 3 frameworks sob a lente dos princípios",
], "📖 Cada objetivo é mensurável: explicar, decompor, implementar, diferenciar, adicionar, avaliar.\n❓ 'Qual objetivo vocês acham mais desafiador?' (geralmente #4 ou #6)\n⚠️ Alunos confundem 'aprender ReAct' com 'aprender LangGraph'. ReAct é o padrão; LangGraph é instanciação.\n➡️ Competências.", 2, T, "Objetivos mensuráveis")

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → I (Intermediário)",
    "C3 MCP & Tool Use → B (Básico)",
    "C4 Agent Memory → B (Básico)",
    "C5 AgentOps & Avaliação → B (Básico)",
    "",
    "C1 atinge Intermediário: constrói agente funcional e justifica escolhas",
    "C3, C4, C5 em Básico: fundação para módulos posteriores",
], "📖 O Framework Etho tem 6 competências em 3 níveis. Este módulo toca 4 delas. C1 atinge Intermediário.\n💡 Analogia: sair do estacionamento (Básico) e dirigir na cidade (Intermediário).\n⚠️ 'Básico' não significa 'superficial' — significa fundação.\n➡️ Agenda.", 3, T, "Conectar ao Framework Etho")

s = content_slide("Agenda da Aula", [
    "Bloco 1 (45 min):",
    "  Abertura (8 min) — objetivos, motivação, contexto",
    "  Fundamentos (22 min) — Augmented LLM, ReAct, DEMO",
    "  Workflows vs Agentes (10 min) — 5 padrões, decisão",
    "  Intervalo (5 min)",
    "",
    "Bloco 2 (45 min):",
    "  Frameworks (15 min) — Python puro vs LangGraph vs OpenAI SDK",
    "  Observabilidade (8 min) — traces, logs, custo",
    "  Fechamento (12 min) — boas práticas, quiz, Q&A",
], "📖 Dois blocos. Primeiro: teoria + demo. Segundo: prática + fechamento. Quiz no final.\n➡️ Por que estamos aqui?", 4, T, "Roteiro da aula")

s = content_slide("O Problema do Oráculo", [
    "Cenário: 'Reserve voo + hotel + carro para a conferência'",
    "",
    "LLM puro (oráculo):",
    "  • Responde uma vez e acabou",
    "  • Não pode agir (sem tools)",
    "  • Não pode verificar (sem environment)",
    "  • Não pode iterar (sem loop)",
    "",
    "Pergunta: Quantos já tentaram tarefa multi-step com ChatGPT e falharam?",
], "📖 LLM puro é como um oráculo — pergunta, responde, acabou. Tarefas reais precisam de loop, tools e memória.\n💡 Analogia: consultor de viagens que só pode falar uma vez — não pode pesquisar, ligar, verificar.\n❓ Mãos levantadas: 'Quantos já falharam com ChatGPT multi-step?'\n⚠️ Alunos acham que 'prompt melhor' resolve. Não resolve — falta loop, tools, memória.\n➡️ Por que só agora é viável?", 5, T, "Criar tensão — LLM puro falha em multi-step")

s = content_slide("Por Que Agora", [
    "Linha do tempo:",
    "  2020: GPT-3 (linguagem, sem tools)",
    "  2022: Chain-of-Thought (reasoning emergente)",
    "  Mar/2023: ReAct (Yao et al.)",
    "  Jun/2023: Function calling (OpenAI API)",
    "  Dez/2024: Anthropic 'Building Effective Agents'",
    "  2025: Agent frameworks amadurecem",
    "",
    "Confluência de 4 fatores:",
    "  1. Reasoning (CoT, ToT, Reflexion)",
    "  2. Tool calling estruturado",
    "  3. Context window expandida (128k+)",
    "  4. Redução de custo (gpt-4o-mini, Haiku)",
], "📖 Agentes não são ideia nova — existem desde os anos 90. Mas 4 fatores convergiram: reasoning, tools, context, custo. A Anthropic em dez/2024 sistematizou o que funciona.\n💡 Analogia: avião — conceito existia há séculos, mas precisou de motor leve + aerodinâmica + materiais.\n❓ 'Qual fator foi o gatilho mais recente?' (Resposta: custo)\n⚠️ Alunos acham que agentes são 'novos'. ReAct é de 2022. Novo é a viabilidade econômica.\n➡️ Vamos ao bloco fundamental.", 6, T, "Confluência histórica")

# ═══════════════════════════════════════
# SEÇÃO B — Fundamentos (7-22)
# ═══════════════════════════════════════

s = section_slide(1, "Do LLM ao Agente")
add_notes(s, "Início do bloco de fundamentos. O que muda quando transformamos um LLM em agente?")

s = content_slide("O Que É um Agente?", [
    "Definição ampla: sistema autônomo, longo prazo, multi-tool",
    "Definição restritiva: segue workflow predefinido",
    "Recuo pragmático (Anthropic): 'agentic systems' = workflows + agentes",
    "",
    "Workflows: LLMs e tools orquestrados via CÓDIGO PREDEFINIDO",
    "Agentes: LLM DIRECIONA seu próprio processo e tool usage",
    "",
    "Pergunta: Qual definição sua equipe usa?",
], "📖 A palavra 'agente' é problemática — cada um define diferente. Anthropic cunhou 'agentic systems' como guarda-chuva.\n💡 Analogia: trem (workflow — trilhos predefinidos) vs táxi (agente — motorista decide o caminho).\n❓ 'Qual definição sua equipe usa?'\n⚠️ Usar 'agente' para tudo infla expectativas e confunde stakeholders.\n➡️ Vamos decompor um agente em partes.", 8, T, "Definições concorrentes")

s = content_slide("Taxonomia Unificada de Agentes", [
    "Fonte: arXiv:2601.12560 (Arunkumar et al., 2026)",
    "",
    "6 componentes:",
    "  1. Perception — como o agente percebe o mundo (input, tool results)",
    "  2. Brain — o LLM (motor cognitivo, reasoning)",
    "  3. Planning — decompor objetivo em sub-tarefas",
    "  4. Action — executar (tool calls, API calls)",
    "  5. Tool Use — selecionar e usar ferramentas",
    "  6. Collaboration — comunicação com outros agentes/humanos",
], "📖 O survey propôs uma taxonomia que cobre todo sistema agêntico. Brain = LLM. Perception = input + tool results. Planning = decompor (ETHAGT04). Action = executar. Tool Use = ACI (ETHAGT02). Collaboration = multi-agente (ETHAGT09).\n💡 Analogia: cérebro (Brain), sentidos (Perception), pensar 'primeiro A depois B' (Planning), mover mãos (Action), pegar chave de fenda (Tool Use), pedir ajuda (Collaboration).\n❓ 'Qual componente é mais fraco hoje?' (Planning)\n⚠️ Alunos acham que 'Brain' é tudo. Sem Perception, agente é cego. Sem Action, é oráculo.\n➡️ A transição histórica.", 9, T, "Taxonomia de 6 componentes")

s = content_slide("De Geração Única para Controle Cognitivo", [
    "Antes (Geração única):",
    "  prompt → response (1 step)",
    "  Sem estado, sem memória, sem iteração",
    "  Custo previsível, latência baixa",
    "",
    "Agora (Controle cognitivo):",
    "  prompt → plan → act → observe → adapt → ... → response (N steps)",
    "  Estado persistente, memória entre steps",
    "  Custo variável, latência alta",
    "  Erro composto (falha em step N propaga)",
], "📖 A transição não é 'mais chamadas' — é mudança de paradigma com implicações arquiteturais: estado, memória, orçamento, observabilidade, error handling.\n💡 Analogia: pedir prato no restaurante (1 step) vs cozinhar (N steps — planejar, comprar, cortar, cozinhar, provar, ajustar).\n⚠️ Alunos subestimam custo. 10 steps × 2k tokens × 1000 usuários/dia = $6k/mês.\n➡️ O bloco fundamental.", 10, T, "Mudança de paradigma")

s = section_slide(2, "O Augmented LLM: O Bloco Fundamental")
add_notes(s, "Este é o slide mais importante. Se só lembrarem de uma coisa: o Augmented LLM é o bloco fundamental de tudo.")

s = content_slide("O Augmented LLM", [
    "Definição (Anthropic): LLM enhanced with retrieval, tools, and memory",
    "",
    "4 componentes:",
    "  • LLM — motor cognitivo (reasoning, geração)",
    "  • Retrieval — o modelo gera suas próprias queries",
    "  • Tools — extensão de ação (APIs, DB, execução)",
    "  • Memory — working (context window) vs persistente",
    "",
    "Princípio: o MODELO decide quando recuperar, qual tool chamar, o que reter",
    "Interface bem documentada = regra de ouro (ACI)",
], "📖 O Augmented LLM não é 'LLM + alguma coisa'. É mudança qualitativa: o modelo é ATIVO na decisão.\n💡 Analogia: cérebro em jarra (LLM puro) vs cérebro com olhos (retrieval), mãos (tools) e memória. A diferença é autonomia.\n❓ 'Qual componente é mais negligenciado em produção?' (Memory)\n⚠️ 'LLM + tools' não é Augmented LLM. Faltam retrieval e memory.\n➡️ Vamos aprofundar cada componente.", 12, T, "O bloco fundamental de Anthropic")

s = content_slide("Retrieval In-Loop", [
    "RAG tradicional (fixo):",
    "  query → retrieve → generate (pipeline predefinido)",
    "  Sempre recupera, mesmo para 'oi, tudo bem?'",
    "",
    "RAG agêntico (in-loop):",
    "  Modelo decide SE e QUANDO recuperar",
    "  Modelo gera suas próprias search queries",
    "  Pode recuperar múltiplas vezes com queries diferentes",
    "",
    "Vantagem: eficiência + flexibilidade",
    "Aprofundamento: ETHAGT06 — RAG Agêntico",
], "📖 A diferença é a DECISÃO. RAG tradicional sempre recupera — desperdiça tokens. RAG agêntico: o modelo avalia se precisa.\n💡 Analogia: estudante que sempre consulta o livro (fixo) vs que sabe quando consultar (agêntico).\n❓ 'Em que cenário RAG fixo é melhor?' (quando todas as queries precisam de retrieval)\n⚠️ Sem tool de busca explicitamente chamável, o modelo não pode decidir.\n➡️ Tools.", 13, T, "Retrieval como decisão do modelo")

s = content_slide("Tools como Extensão de Ação", [
    "Tools = capacidades que o modelo não tem (API, DB, execução, busca)",
    "",
    "Tool calling estruturado:",
    "  • JSON Schema define nome, descrição, parâmetros",
    "  • Modelo gera tool_call com argumentos",
    "  • Sistema executa e retorna observation",
    "  • Modelo consome a observation e continua",
    "",
    "A descrição é o que o modelo usa para decidir QUANDO chamar",
    "Descrição ruim = tool errada chamada",
    "Aprofundamento: ETHAGT02 — ACI",
], "📖 Tools são a forma do agente AGIR. Sem tools, só fala. A descrição é crítica — é o que o modelo usa para decidir.\n💡 Analogia: manual de ferramentas para um funcionário. Se o manual for confuso, pega a ferramenta errada.\n❓ 'O que acontece se duas tools tiverem descrições sobrepostas?' (modelo confunde)\n⚠️ Menos tools bem descritas > muitas tools mal descritas.\n➡️ Memory.", 14, T, "Tool calling estruturado")

s = content_slide("Memory: Working vs Persistente", [
    "Working memory (context window):",
    "  • Tudo que o modelo 'vê' no prompt atual",
    "  • Efêmera — some quando a sessão acaba",
    "  • Limitada (128k-200k tokens, mas custo cresce)",
    "",
    "Persistent memory (checkpointer):",
    "  • Estado serializado entre sessões",
    "  • Postgres, Redis, SQLite",
    "  • Permite resumir e retomar",
    "",
    "Tensão: mais memória = melhor decisão mas mais custo",
    "Estratégias: summarization, eviction, entity memory (ETHAGT05)",
], "📖 Memory é o componente mais subestimado. Working = context window. Persistent = checkpointer. A tensão é: mais contexto = melhor decisão mas mais custo.\n💡 Analogia: memória de curto prazo (reunião) vs caderno de anotações (consultar depois).\n⚠️ Tentar colocar todo o histórico na context window funciona por um tempo, mas custo explode.\n➡️ A regra de ouro.", 15, T, "Working vs persistente")

s = content_slide("A Regra de Ouro: ACI", [
    "Princípio (Anthropic): 'Invista tanto esforço em ACI quanto em HCI'",
    "",
    "ACI = Agent-Computer Interface",
    "  • Como o agente interage com tools",
    "  • Descrição de tool = docstring para um júnior",
    "  • Poka-yoke: fazer difícil errar",
    "",
    "Caso real (SWE-bench):",
    "  Problema: agente usava paths relativos errados",
    "  Solução: mudaram a tool para exigir paths absolutos",
    "  Resultado: erro desapareceu, sem mexer no prompt",
    "",
    "Regra: tempo em tools > tempo no prompt principal",
], "📖 A Anthropic gastou MAIS tempo nas tools do que no prompt principal. Por quê? A tool é o ponto de contato entre modelo e mundo. Se confusa, o modelo erra.\n💡 Analogia: design de UI. Se usuários clicam no botão errado, a solução não é 'treinar' — é mover o botão.\n❓ 'Já tiveram bug onde a solução era melhorar o prompt mas o problema era a tool?'\n⚠️ Passar horas no prompt e 5 min na tool é invertê-lo.\n➡️ O Agent Loop.", 16, T, "Princípio ACI")

s = section_slide(3, "O Agent Loop: ReAct")
add_notes(s, "O Augmented LLM é o bloco. O Agent Loop é como o colocamos em movimento. ReAct é o padrão fundacional.")

s = content_slide("ReAct: O Padrão Fundacional", [
    "Padrão: Thought → Action → Observation em loop",
    "Fonte: Yao et al., ICLR 2023 (arXiv:2210.03629)",
    "",
    "Por que funciona:",
    "  • Thought força reasoning explícito (não pular para ação)",
    "  • Action executa no ambiente (grounding)",
    "  • Observation traz fato real (não alucinável)",
    "  • Loop permite adaptação (se inesperado, repensar)",
    "",
    "Guardrails:",
    "  • max_steps — limite de iterações",
    "  • max_cost — orçamento de tokens",
    "  • timeout — limite de tempo",
], "📖 ReAct unifica reasoning (CoT) e acting (tool use). Thought → Action → Observation. A Observation é o grounding — não é alucinável. max_steps é não-negociável.\n💡 Analogia: cientista fazendo experimento. Thought = hipótese. Action = experimento. Observation = resultado. Se inesperado, repensa.\n❓ 'O que acontece se removermos o Thought?' (modelo chama tools aleatoriamente sem plano)\n⚠️ Sem max_steps, loop infinito. Always set max_steps.\n➡️ Trace real.", 18, T, "Padrão ReAct")

s = code_slide("Trace Real de Execução", """[STEP 0]
  Thought: "Preciso saber a previsão do tempo no RJ"
  Action: search("previsão tempo Rio de Janeiro amanhã")
  Observation: "25°C, possibilidade de chuva à tarde"

[STEP 1]
  Thought: "Agora posso responder sobre o clima"
  Action: (sem tool_call — resposta final)
  Answer: "Leve guarda-chuva. Amanhã no RJ: 25°C com
          possibilidade de chuva à tarde.\"""",
    "📖 Trace real de ReAct. Note: (1) cada step tem Thought, Action, Observation; (2) Observation é fato; (3) no step 1, modelo decide parar.\n💡 Analogia: 'pensar em voz alta' — você vê cada passo do raciocínio.\n❓ 'Como saberia se está funcionando sem o trace?' (não saberia)\n⚠️ Não logar o Thought = trace opaco. Sem Thought, não sabe POR QUE escolheu aquela action.\n➡️ Limitações.", 19, T, "Trace real de ReAct")

s = content_slide("Limitações do ReAct", [
    "1. Loops infinitos — sem max_steps, repete indefinidamente",
    "2. Custo cumulativo — cada iteração adiciona tokens ao contexto",
    "3. Alucinação em ação — modelo chama tool inexistente ou args errados",
    "4. Latência serial — N steps = N × latência",
    "5. Erro composto — falha em step N propaga para N+1",
    "6. Context exhaustion — em conversas longas, context window enche",
    "",
    "Mitigações: max_steps, summarization, validação de args, retry, traces",
], "📖 ReAct não é bala de prata. Cada limitação tem mitigação, mas nenhuma é perfeita. ETHAGT04 mostra evoluções: Reflexion, ToT, LATS.\n💡 Analogia: limitações de um carro — velocidade, consumo, desgaste. Nenhuma é fatal, mas precisa gerenciar.\n❓ 'Qual limitação é mais crítica em produção?' (erro composto — mais difícil de debugar)\n⚠️ Sem mitigações, funciona na demo mas não em produção.\n➡️ Demo ao vivo.", 20, T, "Limitações honestas")

s = code_slide("DEMO: ReAct em 50 Linhas", """import json, os
from openai import OpenAI

client = OpenAI(api_key=os.getenv("LLM_API_KEY"))

TOOLS = [{
    "type": "function",
    "function": {
        "name": "calculator",
        "description": "Executa operação aritmética",
        "parameters": {
            "type": "object",
            "properties": {
                "a": {"type": "number"},
                "b": {"type": "number"},
                "op": {"type": "string", "enum": ["+","-","*","/"]}
            },
            "required": ["a", "b", "op"]
        }
    }
}]

def calculator(a, b, op):
    return {"+":a+b, "-":a-b, "*":a*b, "/":a/b if b else "erro"}[op]

messages = [
    {"role": "system", "content": "Você é um agente ReAct."},
    {"role": "user", "content": "Quanto é 123 * 456 + 789?"}
]

max_steps = 5
for step in range(max_steps):
    response = client.chat.completions.create(
        model="gpt-4o-mini", messages=messages, tools=TOOLS, tool_choice="auto")
    msg = response.choices[0].message
    messages.append(msg)
    if not msg.tool_calls:
        print(f"Resposta: {msg.content}")
        break
    for tc in msg.tool_calls:
        args = json.loads(tc.function.arguments)
        result = calculator(**args)
        print(f"Passo {step}: {tc.function.name}({args}) = {result}")
        messages.append({"role":"tool","tool_call_id":tc.id,"content":str(result)})""",
    "📖 Demo ao vivo. Prestem atenção em: (a) como o modelo decide qual tool chamar; (b) como o resultado volta; (c) como o modelo decide parar. A query exige 2 calls: 123*456, depois resultado+789.\n💡 Analogia: ver o motor com o capô aberto. Depois, LangGraph/CrewAI fará sentido.\n❓ Após rodar: 'Viram quantos steps? Por que não 1?' (calculator só faz uma operação por vez)\n⚠️ Se API falhar, tenho screenshot do trace.\n➡️ Pergunta da demo.", 21, T, "Demo ao vivo — ReAct em Python puro")

s = exercise_slide("Pergunta da Demo", [
    "Em duplas (2 min):",
    "",
    "1. O que acontece se a tool retornar erro? Quem trata?",
    "2. E se o modelo não chamar tool nenhuma?",
    "3. E se o modelo chamar uma tool que não existe?",
    "",
    "Discussão em duplas + compartilhar",
], "📖 Respostas: (1) Sem try/except, exceção quebra o loop. Solução: capturar e retornar como Observation estruturada. (2) Código entra no 'if not tool_calls' e imprime resposta — mesmo que errada. (3) calculator(**args) falha com KeyError.\n💡 Analogia: perguntar 'o que se o pneu furar' antes de dirigir.\n⚠️ Modelo não 'sabe' tratar erros — só vê a Observation. Se for traceback de Python, fica confuso.\n➡️ Workflows vs Agentes.", 22, T)

# ═══════════════════════════════════════
# SEÇÃO C — Workflows vs Agentes (23-28)
# ═══════════════════════════════════════

s = section_slide(4, "Workflows vs Agentes")
add_notes(s, "Já temos o Augmented LLM e o loop. Agora: QUANDO usar agente vs workflow?")

s = content_slide("Workflows vs Agentes", [
    "Workflows: LLMs e tools orquestrados via CÓDIGO PREDEFINIDO",
    "  → Previsível, testável, barato",
    "  → Limitação: inflexível",
    "",
    "Agentes: LLM DIRECIONA seu próprio processo e tool usage",
    "  → Flexível, adaptativo, poderoso",
    "  → Limitação: imprevisível, caro, difícil de testar",
    "",
    "Fonte: Anthropic, Building Effective Agents (2024)",
    "",
    "Critério: comece com workflow; só use agente quando necessário",
], "📖 A distinção é sobre QUEM controla o fluxo. Workflow = seu código. Agente = o modelo. Workflow = trilhos de trem. Agente = táxi.\n💡 Analogia: receita (workflow) vs cozinhar livre (agente). Receita é previsível. Livre é criativo mas varia.\n❓ 'Em que situação usariam agente em vez de workflow?' (passos imprevisíveis, input muito variado)\n⚠️ Começar com agente porque é 'mais legal' = sistema caro e imprevisível.\n➡️ Os 5 padrões.", 24, T, "Distinção canônica de Anthropic")

s = content_slide("Os 5 Workflows Canônicos", [
    "1. Prompt Chaining — sequência de steps (gerar → revisar → traduzir)",
    "2. Routing — classificar e dispatch (tipo de query → handler)",
    "3. Parallelization — sectioning (dividir) ou voting (múltiplas tentativas)",
    "4. Orchestrator-Workers — orquestrador delega dinamicamente",
    "5. Evaluator-Optimizer — loop de refinamento (gerar → avaliar → melhorar)",
    "",
    "Aprofundamento: ETHAGT03 — Padrões de Workflow",
], "📖 5 padrões canônicos da Anthropic. Cada um resolve um problema. Prompt Chaining para sequência linear. Routing para handlers especializados. Parallelization para velocidade/confiabilidade. Orchestrator-Workers para subtasks imprevisíveis. Evaluator-Optimizer para refinamento.\n💡 Analogia: padrões de design do GoF — você não usa todos, mas precisa conhecer.\n❓ 'Qual já usaram?' (Prompt Chaining e Routing mais comuns)\n⚠️ Orchestrator-Workers NÃO é agente — o orquestrador segue fluxo predefinido.\n➡️ Árvore de decisão.", 25, T, "5 padrões canônicos")

s = content_slide("Árvore de Decisão: Workflow ou Agente?", [
    "Q1: Consigo listar os passos a priori?",
    "  Sim → WORKFLOW (comece pelo padrão mais simples)",
    "  Não → Q2",
    "",
    "Q2: O número de passos é previsível?",
    "  Sim, mas variável → Q3",
    "  Não → AGENTE",
    "",
    "Q3: Posso confiar no modelo sem HITL?",
    "  Sim → AGENTE (autônomo)",
    "  Não → AGENTE + HITL forte (checkpoints críticos)",
], "📖 A primeira pergunta é a mais importante: 'Consigo listar os passos?' Se sim, workflow. Se não, agente.\n💡 Analogia: triagem médica. Resfriado (protocolo fixo = workflow). Diagnóstico complexo (iterativo = agente). Cirurgia (agente + HITL).\n❓ 'Pensem no sistema de vocês: workflow ou agente?'\n⚠️ 80% dos casos = workflow. Agente é para os 20% que precisam de flexibilidade.\n➡️ Escalada de complexidade.", 26, T, "Critério de decisão")

s = content_slide("A Escalada de Complexidade", [
    "Nível 0: Single LLM call + retrieval + examples",
    "  → 90% dos casos. Antes de tudo, tente isto.",
    "Nível 1: Workflow simples (prompt chaining, routing)",
    "Nível 2: Workflow complexo (orchestrator-workers, evaluator-optimizer)",
    "Nível 3: Agente autônomo (modelo controla o fluxo)",
    "Nível 4: Multi-agente (múltiplos agentes colaborando)",
    "",
    "Regra: só suba um nível com EVIDÊNCIA de que o atual é insuficiente",
], "📖 O princípio mais importante: COMECE SIMPLES. A Anthropic é explícita. Cada nível adiciona custo, latência, imprevisibilidade. 90% dos casos = Nível 0.\n💡 Analogia: escalada de remédios. Não começa com quimioterapia para dor de cabeça. Começa com paracetamol.\n❓ 'Em que nível está o sistema de vocês?'\n⚠️ Pular do Nível 0 para o 3 ou 4 porque 'agentes são o futuro' = sistema 10x mais caro.\n➡️ Exercício rápido.", 27, T, "Princípio: comece simples")

s = exercise_slide("Exercício: Workflow ou Agente?", [
    "Votação rápida (mãos levantadas):",
    "",
    "1. 'Traduzir documentação técnica EN→PT' → W ou A?",
    "2. 'Resolver issue de GitHub arbitrário' → W ou A?",
    "3. 'Classificar ticket de suporte' → W ou A?",
    "4. 'Pesquisar e sintetizar tema novo' → W ou A?",
    "5. 'Revisar código em busca de bugs' → W ou A?",
    "6. 'Negociar contrato com fornecedor' → W ou A?",
], "📖 Respostas: (1) W — chaining: traduzir → revisar → ajustar. (2) A — cada issue é diferente. (3) W — routing. (4) A — busca iterativa. (5) W — parallelization/voting. (6) A + HITL — alta stakes.\n⚠️ Se votam 'agente' em tudo, lembrar: passos previsíveis = workflow.\n➡️ Intervalo. Voltamos para frameworks.", 28, T)

# ═══════════════════════════════════════
# SEÇÃO D — Frameworks (29-36)
# ═══════════════════════════════════════

s = section_slide(5, "Implementação: Do Zero vs Framework")
add_notes(s, "Já entendemos o que é um agente e quando usar. Agora COMO implementar: 3 abordagens.")

s = content_slide("O Mesmo Agente, 3 Implementações", [
    "Versão 1: Python puro + OpenAI SDK (~50 linhas)",
    "  Controle total, transparência total, mais código",
    "",
    "Versão 2: LangGraph (~20 linhas)",
    "  StateGraph + checkpointer embutido, abstrai loop",
    "",
    "Versão 3: OpenAI Agents SDK (~15 linhas)",
    "  Agent(model, tools, instructions) → run()",
    "  Máxima produtividade, mínima transparência",
    "",
    "Menos código ≠ melhor — depende do CONTROLE que você precisa",
], "📖 O mesmo agente ReAct em 50, 20 ou 15 linhas. A diferença é o que você vê e controla.\n💡 Analogia: carro manual vs automático vs autônomo. Manual: controla tudo. Automático: mais fácil. Autônomo: só diz o destino.\n❓ 'Qual para POC de 1 semana? E para produção com 10k usuários?' (POC: SDK. Produção: LangGraph ou Python puro)\n⚠️ Escolher SDK para produção porque 'é mais fácil' = problema quando o agente falha e você não sabe por quê.\n➡️ Python puro.", 30, T, "3 implementações do mesmo agente")

s = code_slide("Python Puro: Controle Total", """# Loop principal — você controla tudo
max_steps = 5
for step in range(max_steps):
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=messages,
        tools=TOOLS,
        tool_choice="auto"
    )
    msg = response.choices[0].message
    messages.append(msg)

    if not msg.tool_calls:
        print(f"Resposta: {msg.content}")
        break

    for tc in msg.tool_calls:
        args = json.loads(tc.function.arguments)
        result = calculator(**args)
        messages.append({
            "role": "tool",
            "tool_call_id": tc.id,
            "content": str(result)
        })""",
    "📖 Python puro: você é o mestre. Prompt efetivo é EXATAMENTE o que você escreve. Loop é seu. Estado é a lista messages.\n💡 Analogia: cozinhar do zero vs mistura de bolo. Do zero controla cada ingrediente.\n❓ 'Quem já implementou agente em Python puro?'\n⚠️ 'Muito trabalho' — o loop são 15 linhas. O trabalho está nas tools, prompt, error handling. Existe independentemente do framework.\n➡️ LangGraph.", 31, T, "Python puro — controle total")

s = code_slide("LangGraph: Abstração Controlada", """from langgraph.prebuilt import create_react_agent

# 2 linhas e você tem um agente ReAct
agent = create_react_agent(
    model="gpt-4o-mini",
    tools=TOOLS
)

result = agent.invoke({
    "messages": [
        {"role": "user", "content": "Quanto é 123*456+789?"}
    ]
})

# O que LangGraph abstrai:
#   - Loop (StateGraph com nodes e edges)
#   - Estado (Checkpointer embutido)
#   - Tools (ToolNode wrapper)
#
# O que esconde:
#   - Prompt efetivo (framework adiciona system prompt)
#   - Formato exato das messages""",
    "📖 LangGraph é o meio-termo. Abstrai loop, estado e tools. Mas você ainda vê o grafo. O que esconde: prompt efetivo.\n💡 Analogia: framework web (Flask/FastAPI) vs HTTP puro. Abstrai o repetitivo, deixa ver o importante.\n❓ 'Desvantagem de create_react_agent vs StateGraph custom?' (menos controle do loop)\n⚠️ Usar create_react_agent sem ler o código fonte = problema quando algo dá errado.\n➡️ OpenAI SDK.", 32, T, "LangGraph — meio-termo")

s = code_slide("OpenAI Agents SDK: Máxima Produtividade", """from agents import Agent, Runner

# 4 linhas — máximo de abstração
agent = Agent(
    name="Calculator Agent",
    instructions="Você é um agente ReAct.",
    tools=[calculator_tool]
)

result = Runner.run_sync(
    agent,
    "Quanto é 123 * 456 + 789?"
)

# O que abstrai: TUDO
#   - Loop, estado, tools, handoffs,
#     guardrails, tracing nativo
#
# O que esconde: QUASE TUDO
#   - Prompt efetivo (caixa preta)
#   - Lógica de loop (interna)
#   - Serialização de estado (opaca)
#
# Quando usar: POC rápido, demos
# Quando NÃO usar: produção com auditoria""",
    "📖 OpenAI SDK é a camada mais alta. Máxima produtividade, mínima transparência. Em 5 linhas: agente funcional com tracing.\n💡 Analogia: Wix vs codar site. Wix é produtivo. Mas se precisa de controle granular, código do zero.\n⚠️ Em produção com compliance: 'o SDK decidiu' não é resposta aceitável.\n➡️ Comparação estrutural.", 33, T, "OpenAI SDK — máxima abstração")

s = comparison_slide("Comparação Estrutural: 3 Abordagens",
    "Python Puro", [
        "Controle do prompt: ✅ Total",
        "Transparência: ✅ Total",
        "Estado: ❌ Manual",
        "Tools: ✅ Você valida",
        "Linhas: ~50",
        "Debug: ✅ Granular",
        "Multi-agente: ❌ Manual",
    ],
    "LangGraph / SDK", [
        "Controle do prompt: ⚠️ Parcial",
        "Transparência: ✅/❌ Alto/Baixo",
        "Estado: ✅ Embutido",
        "Tools: ⚠️ Wrapper",
        "Linhas: ~20/~15",
        "Debug: ✅/❌ Limitado",
        "Multi-agente: ✅ Nativo",
    ],
    "📖 Não há 'vencedor'. Python puro ganha em controle. SDK em produtividade. LangGraph é meio-termo. A escolha depende do contexto.\n❓ 'Qual critério é mais importante para vocês?'\n⚠️ Focar em 'linhas de código' é o critério MENOS importante.\n➡️ Quando reduzir abstração.", 34, T, "Tabela comparativa")

s = content_slide("Quando Reduzir Camadas de Abstração", [
    "Sinais de alerta (se você disser qualquer um, reduza abstração):",
    "",
    "  'Não entendo por que o agente fez X'",
    "  'O prompt efetivo é diferente do que escrevi'",
    "  'Custo inesperado — não sei onde os tokens vão'",
    "  'Latência inexplicável — não sei qual step está lento'",
    "  'Não consigo reproduzir o bug localmente'",
    "",
    "Regra da Anthropic: 'Ensure you understand the underlying code'",
    "Em produção: a opacidade se torna problema de auditoria",
], "📖 Frameworks são ótimos para começar. Em produção, a opacidade vira problema. Se não consegue explicar POR QUE o agente decidiu X, tem problema de auditoria.\n💡 Analogia: dirigir sem saber trocar pneu. Na cidade, tudo bem. Na estrada remota, precisa saber.\n❓ 'Já disseram alguma dessas frases?'\n⚠️ Docs não são código. Ler o código fonte do framework pelo menos uma vez.\n➡️ Discussão.", 35, T, "Quando reduzir abstração")

s = exercise_slide("Discussão: Qual Framework?", [
    "Perguntas para discussão aberta (3 min):",
    "",
    "1. Qual framework dá mais controle? E mais produtividade?",
    "2. Em que cenário você NÃO usaria framework?",
    "3. Como convencer um PM de que Python puro é melhor?",
    "",
    "Síntese: não há 'melhor' — há o certo para o contexto",
], "📖 Síntese: (1) Controle: Python puro > LangGraph > SDK. Produtividade: inverso. (2) Produção com auditoria. (3) Custo de manutenção vs desenvolvimento; risco de caixa preta.\n⚠️ Transformar em 'meu framework é melhor' — redirecionar: depende do contexto.\n➡️ Observabilidade.", 36, T)

# ═══════════════════════════════════════
# SEÇÃO E — Observabilidade (37-41)
# ═══════════════════════════════════════

s = section_slide(6, "Observabilidade Desde o Dia 1")
add_notes(s, "Observabilidade não é 'fase de produção'. É algo que você adiciona na PRIMEIRA linha de código.")

s = content_slide("Por Que Desde o Dia 1?", [
    "Sem traces: melhoria é adivinhação",
    "  'O agente errou' — por quê? Qual step? Qual tool?",
    "  'Custo subiu' — onde? Qual step consumiu mais?",
    "  'Latência alta' — em qual step?",
    "",
    "Com traces: você vê o pensamento do agente",
    "  Cada Thought, Action, Observation logado",
    "  Custo e latência por step",
    "  Possível reproduzir e debugar",
    "",
    "Em produção: debugging sem traces = impossível",
], "📖 Sem traces, você só sabe que errou. Com traces, sabe por quê, onde e como corrigir. Em produção, é a diferença entre hotfix de 5 min e 3h de debug.\n💡 Analogia: caixa preta de avião. Não espera cair para instalar. Está lá desde o primeiro voo.\n❓ 'Têm traces nos sistemas de LLM de vocês hoje?' (maioria não tem)\n⚠️ Adicionar logs APÓS o primeiro bug = já perdeu os dados.\n➡️ Padrão mínimo.", 38, T, "Por que observabilidade desde o início")

s = code_slide("Logging Estruturado", """import json, time

def log_step(step, thought, action, args, obs, usage):
    log = {
        "timestamp": time.isoformat(),
        "step": step,
        "thought": thought,
        "action": action,
        "args": args,
        "observation": obs,
        "tokens_in": usage.prompt_tokens,
        "tokens_out": usage.completion_tokens,
        "cost_usd": (usage.prompt_tokens * 0.15
                    + usage.completion_tokens * 0.60) / 1_000_000,
        "latency_ms": elapsed
    }
    print(json.dumps(log, ensure_ascii=False))

# Campos mínimos:
#   timestamp, step, thought, action, args,
#   observation, tokens_in, tokens_out, cost, latency

# Evolução:
#   print → structured logs → LangSmith / Phoenix""",
    "📖 JSON estruturado com campos padronizados. O thought é a parte mais importante — sem ele, você sabe O QUE fez mas não POR QUÊ.\n💡 Analogia: extrato de banco. Cada transação tem data, valor, descrição. Sem extrato, não sabe para onde o dinheiro foi.\n⚠️ Logar só Action e Observation, esquecendo Thought = trace opaco.\n➡️ Traces.", 39, T, "Padrão mínimo de logging")

s = content_slide("Traces: O Que São", [
    "Trace = árvore de spans que representa a execução completa",
    "",
    "Span = unidade atômica:",
    "  • Nome (ex: 'LLM call', 'tool: calculator')",
    "  • Duração (start → end)",
    "  • Input e Output",
    "  • Metadata (tokens, cost, model)",
    "",
    "Estrutura hierárquica:",
    "  Agent Run (root span)",
    "  ├── LLM Call #1 → Tool: calculator",
    "  ├── LLM Call #2 → Tool: calculator",
    "  └── LLM Call #3 (answer)",
    "",
    "Ferramentas: LangSmith, Phoenix, Langfuse, OpenTelemetry",
], "📖 Trace é como um profile de execução. Cada chamada é um span. A visualização em timeline mostra o que aconteceu, QUANDO e QUANTO TEMPO levou.\n💡 Analogia: DevTools do navegador mostrando waterfall. Cada recurso é uma barra. O trace de agente é igual, mas para LLM e tools.\n❓ 'Quem já usou LangSmith ou Phoenix?'\n⚠️ Log ≠ trace. Log é uma linha. Trace é árvore estruturada com timing e hierarquia.\n➡️ Custo e latência.", 40, T, "Conceito de trace")

s = content_slide("Custo e Latência como Métricas de Primeira Classe", [
    "Custo por step: tokens_in × price_in + tokens_out × price_out",
    "  Ex: gpt-4o-mini = $0.15/1M in + $0.60/1M out",
    "  10 steps × 2k tokens = ~$0.006 por execução",
    "  1000 execuções/dia = $6/dia = $180/mês",
    "",
    "Latência cumulativa (serial):",
    "  Cada step espera o anterior → latência_total = Σ latência_step",
    "  10 steps × 2s = 20s total",
    "",
    "Dashboard mínimo: step, tokens, custo, latência p50/p95/p99",
    "Em produção: orçamento por execução (max_cost), alertas",
], "📖 Custo e latência não são 'detalhes de produção' — são métricas de primeira classe. Sem medir, não sabe se custa $0.01 ou $1.00 por execução.\n💡 Analogia: conta de energia. Não espera o fim do mês. Tem medidor em tempo real.\n⚠️ Calcular custo só com a primeira chamada. Esquece que o loop multiplica: 10 steps × 2k = 20k tokens.\n➡️ Fechamento.", 41, T, "Custo e latência")

# ═══════════════════════════════════════
# SEÇÃO F — Fechamento (42-62)
# ═══════════════════════════════════════

s = section_slide(7, "Boas Práticas e Fechamento")
add_notes(s, "Vamos consolidar tudo em boas práticas, anti-patterns, caso de estudo e quiz.")

s = content_slide("Boas Práticas (DO)", [
    "✅ Comece simples — single LLM call antes de agente",
    "✅ Documente tools como se fosse para um júnior",
    "✅ Defina max_steps desde o dia 1",
    "✅ Log estruturado desde a primeira linha",
    "✅ Sandbox para execução de tools",
    "✅ HITL em pontos críticos (ações irreversíveis)",
    "✅ Otimize o prompt efetivo, não apenas o que você escreve",
    "✅ Entenda o framework — leia o código fonte",
    "✅ Meça custo e latência por step desde o início",
    "✅ Teste em sandbox antes de produção",
], "📖 Cada item é uma lição aprendida da forma difícil em produção. 'Comece simples' é o #1 — Anthropic repete em todo documento.\n➡️ Anti-patterns.", 43, T, "Checklist de boas práticas")

s = content_slide("Anti-Patterns (DON'T)", [
    "❌ Começar com agente quando workflow basta",
    "❌ Adicionar complexidade sem evidência",
    "❌ Confiar no framework cegamente ('caixa preta')",
    "❌ Não definir max_steps (loop infinito)",
    "❌ Não logar (debug cego)",
    "❌ Tool mal documentada (ACI fraca)",
    "❌ Prompt gigante sem estrutura",
    "❌ Não testar em sandbox",
    "❌ Não medir custo ('descobri no fim do mês')",
    "❌ Esquecer error handling (tool quebra → agente morre)",
    "❌ Não ter fallback (se API cai, agente trava)",
], "📖 Cada anti-pattern é um bug que vi em produção. 'Começar com agente' é o mais comum. 'Sem max_steps' causa loops que custam milhares.\n❓ 'Qual já cometeram?'\n⚠️ Mesmo times sênior cometem sob pressão de prazo.\n➡️ Caso de estudo.", 44, T, "Checklist de anti-patterns")

s = content_slide("Caso: Anthropic Coding Agent no SWE-bench", [
    "SWE-bench: benchmark que mede capacidade de resolver issues de GitHub",
    "  Issues de Django, scikit-learn, sympy, etc.",
    "",
    "O coding agent da Anthropic:",
    "  Input: texto da issue (PR description)",
    "  Output: patch (diff) que passa nos testes",
    "  Arquitetura: Augmented LLM + loop + ACI + sandbox",
    "  Não há 'magia' — é o que vimos hoje",
    "",
    "Resultados: Claude 3.5 Sonnet ~49% no SWE-bench Verified",
    "",
    "Por que é agente (não workflow):",
    "  Cada issue é diferente — não dá para predefinir passos",
], "📖 Este caso prova que os conceitos de hoje são reais. A arquitetura é simples: Augmented LLM + loop + ACI + sandbox. 49% significa que em metade das issues, o agente produz um patch que passa nos testes — sem intervenção humana.\n💡 Analogia: ver um chef cozinhar. A receita é simples — os mesmos ingredientes. A diferença está na execução.\n❓ 'Se a arquitetura é simples, por que nem todo mundo consegue 49%?' (qualidade das tools, prompt, sandbox, modelo)\n➡️ A lição ACI.", 45, T, "Caso real — SWE-bench")

s = content_slide("A Lição ACI do SWE-bench", [
    "Problema: agente saía do diretório raiz e usava paths relativos errados",
    "",
    "Solução ingênua: 'melhorar o prompt para dizer use paths absolutos'",
    "  → Frágil: modelo pode esquecer",
    "  → Não escalável: um novo prompt para cada edge case",
    "",
    "Solução robusta: mudar a TOOL para exigir paths absolutos",
    "  → Poka-yoke: erro impossível pela design",
    "  → Escalável: funciona independentemente do prompt",
    "",
    "Citação: 'We spent more time optimizing our tools than the overall prompt'",
    "  — Anthropic",
], "📖 A solução robusta foi mudar a tool, não o prompt. Aí o modelo NÃO PODE errar, mesmo que queira. Poka-yoke.\n💡 Analogia: proteção de tomada. Não ensina a criança (prompt). Coloca a proteção (poka-yoke).\n❓ 'Têm algum bug hoje que poderia ser resolvido redesenhando a tool?'\n➡️ Exercício de debug.", 46, T, "Lição ACI — SWE-bench")

s = exercise_slide("Exercício: Debug de Trace Quebrado", [
    "Analise o seguinte trace:",
    "",
    '[STEP 0] Thought: "Vou buscar o preço do produto X"',
    '         Action: search_price("Produto X")',
    '         Observation: "Produto X: R$ 99,90"',
    '[STEP 1] Thought: "Vou buscar o preço do produto X"',
    '         Action: search_price("Produto X")',
    '         Observation: "Produto X: R$ 99,90"',
    '[STEP 2-4] (mesmo padrão repetido)',
    '[STEP 5] max_steps reached. Abortando.',
    "",
    "Em duplas: (1) identificar o loop (2) propor 2 correções",
], "📖 Problema: agente não processa a Observation. Causas: prompt fraco, contexto não inclui observations, modelo fraco.\nCorreções: (1) prompt: 'Use informações já obtidas antes de buscar'; (2) código: detecção de repetição (se action == anterior 3x, forçar resposta).\n⚠️ 'Aumentar max_steps' não é correção — só faz gastar mais tokens antes de falhar.\n➡️ Resumo.", 47, T)

s = content_slide("Resumo da Aula", [
    "1. Augmented LLM = bloco fundamental (LLM + retrieval + tools + memory)",
    "2. ReAct = padrão fundacional (Thought → Action → Observation)",
    "3. Workflow vs Agente = decisão arquitetural (previsibilidade vs flexibilidade)",
    "4. Comece simples — só aumente complexidade com evidência",
    "5. ACI é onde mora a confiabilidade (tempo em tools > tempo no prompt)",
    "6. Observabilidade desde o dia 1 (traces, custo, latência)",
], "📖 Se só lembrarem de 6 coisas: Augmented LLM é o bloco. ReAct é o padrão. Workflow vs agente é a decisão. Comece simples. ACI é a confiabilidade. Observabilidade é a operação.\n➡️ Checklist.", 48, T, "6 pontos-chave")

s = content_slide("Checklist da Aula", [
    "[ ] Definiu agente vs workflow",
    "[ ] Explicou Augmented LLM (4 componentes)",
    "[ ] Apresentou taxonomia unificada (6 componentes)",
    "[ ] Implementou ReAct do zero (demo)",
    "[ ] Comparou 3 frameworks",
    "[ ] Discutiu quando reduzir abstração",
    "[ ] Adicionou logging estruturado",
    "[ ] Introduziu traces e custo/latência",
    "[ ] Discutiu boas práticas e anti-patterns",
    "[ ] Analisou caso real (SWE-bench)",
    "[ ] Fez exercício de debug de trace",
    "[ ] Conectou com próximos módulos",
], "📖 Confirmar que cobri tudo que prometi. Auto-avaliação: se algum item não está claro, é hora de perguntar.\n➡️ Quiz.", 49, T, "Confirmar cobertura")

# Quiz
s = content_slide("Quiz 1/5", [
    "Qual é o bloco fundamental de qualquer sistema agêntico?",
    "",
    "A) Um framework como LangGraph",
    "B) O Augmented LLM (LLM + retrieval + tools + memory)",
    "C) Um prompt bem escrito",
    "D) Um modelo grande o suficiente",
    "",
    "Resposta: B",
], "📖 Se marcou A, ainda confunde framework com arquitetura.\n➡️ Próxima.", 50, T, "Quiz 1")

s = content_slide("Quiz 2/5", [
    "Em ReAct, o que força o grounding do modelo?",
    "",
    "A) O prompt do sistema",
    "B) A temperatura baixa",
    "C) A Observation (resultado real da tool)",
    "D) O max_steps",
    "",
    "Resposta: C",
], "📖 Grounding = ancorar em fato real. Observation é o único componente não-alucinável.\n➡️ Próxima.", 51, T, "Quiz 2")

s = content_slide("Quiz 3/5", [
    "Quando DEVERIA usar um agente em vez de workflow?",
    "",
    "A) Quando a tarefa tem passos predefinidos",
    "B) Quando precisa de previsibilidade total",
    "C) Quando os passos são imprevisíveis e exigem flexibilidade",
    "D) Quando o custo precisa ser mínimo",
    "",
    "Resposta: C",
], "📖 Agentes trocam previsibilidade e custo por flexibilidade.\n➡️ Próxima.", 52, T, "Quiz 3")

s = content_slide("Quiz 4/5", [
    "O que é ACI?",
    "",
    "A) Agent Communication Interface",
    "B) Agent-Computer Interface (design de tools)",
    "C) Automated Code Injection",
    "D) Agent Control Infrastructure",
    "",
    "Resposta: B",
], "📖 ACI = Agent-Computer Interface. Design das tools. Poka-yoke.\n➡️ Próxima.", 53, T, "Quiz 4")

s = content_slide("Quiz 5/5", [
    "Qual é a PRIMEIRA coisa a adicionar a um agente em produção?",
    "",
    "A) Um framework",
    "B) Múltiplas tools",
    "C) Observabilidade (logs estruturados + traces)",
    "D) HITL",
    "",
    "Resposta: C",
], "📖 Sem observabilidade, não pode debugar, medir custo, ou reproduzir bugs.\n➡️ Perguntas de discussão.", 54, T, "Quiz 5")

s = content_slide("Perguntas para Discussão", [
    "1. 'Toda aplicação de LLM deveria ser um agente?' (V/F justificado)",
    "2. 'Em que cenário Python puro é melhor que LangGraph?'",
    "3. 'Como convencer um PM de que workflow basta e agente é overkill?'",
    "4. 'Qual o maior risco de começar com framework sem entender o básico?'",
], "📖 Respostas: (1) Falso — Anthropic: 'consider not building agentic systems at all'. (2) Produção com auditoria. (3) Custo 5-10x, latência 10-30s vs 2-5s, workflow é testável. (4) Caixa preta — quando falha, não sabe por quê.\n➡️ Conexão.", 55, T, "Discussão profunda")

s = content_slide("Conexão com a Especialização", [
    "ETHAGT02 — Tool Calling e ACI (aprofunda ferramentas)",
    "ETHAGT03 — Padrões de Workflow (os 5 padrões em detalhe)",
    "ETHAGT04 — Reasoning & Planning (evolução do ReAct)",
    "ETHAGT05 — Memória de Agentes (working vs persistente)",
    "ETHAGT06 — RAG Agêntico (retrieval in-loop aprofundado)",
    "ETHAGT08 — MCP (tool use como padrão de servidor)",
    "ETHAGT09-10 — Multi-Agentes (Collaboration)",
    "ETHAGT12 — AgentOps (observabilidade em profundidade)",
    "ETHAGT13 — Segurança (tools como vetor de ataque)",
], "📖 ETHAGT01 é a fundação. Todo módulo expande um aspecto. A especialização é um zoom progressivo em cada componente do Augmented LLM.\n➡️ Leitura.", 56, T, "Mapa da especialização")

s = content_slide("Leitura Recomendada", [
    "Obrigatório (antes de ETHAGT02):",
    "  • Anthropic, Building Effective Agents (2024)",
    "    Foco: Apêndice 2 'Prompt Engineering your Tools'",
    "  • Yao et al., ReAct (ICLR 2023) — arXiv:2210.03629",
    "    Foco: Seção 3 (method)",
    "",
    "Recomendado:",
    "  • arXiv:2601.12560 — survey de arquiteturas",
    "  • Schluntz & Albert (YouTube) — Building more effective AI agents",
    "",
    "Opcional:",
    "  • Karpathy, Software 3.0 (talks)",
    "  • HuggingFace Agents Course",
], "📖 O Apêndice 2 da Anthropic prepara para ETHAGT02. O paper do ReAct dá a base teórica da demo.\n➡️ Referências.", 57, T, "Leitura obrigatória")

s = content_slide("Referências Completas", [
    "1. Anthropic. Building Effective Agents. 2024. anthropic.com/engineering/building-effective-agents",
    "2. Yao et al. ReAct. ICLR 2023. arXiv:2210.03629",
    "3. Arunkumar et al. Agentic AI. arXiv:2601.12560, 2026",
    "4. Schick et al. Toolformer. NeurIPS 2023. arXiv:2302.04761",
    "5. Jimenez et al. SWE-bench. arXiv:2310.06770",
    "6. Shinn et al. Reflexion. NeurIPS 2023. arXiv:2303.11366",
    "7. Wang et al. Survey on LLM Agents. arXiv:2308.11432",
    "8. OpenAI. Practical Guide to Building Agents. 2024",
    "9. LangGraph docs. github.com/langchain-ai/langgraph",
    "10. Anthropic. Effective Context Engineering for AI Agents. 2025",
], "📖 Todas as fontes em 20-Research/ETHAGT01-pesquisa.md. As 3 canônicas (1, 2, 3) são leitura obrigatória.\n➡️ Projeto.", 58, T, "Referências")

s = content_slide("Projeto do Módulo", [
    "Implementar 'research assistant' em 3 versões:",
    "  1. Python puro",
    "  2. LangGraph",
    "  3. Framework à escolha (OpenAI SDK / CrewAI / PydanticAI)",
    "",
    "O agente decide: responder direto, recuperar de base local, ou usar tool de busca",
    "",
    "Entrega: repo + README comparativo + traces",
    "Critério: justificar escolha com ≥3 critérios (controle, transparência, esforço)",
    "",
    "Prazo: 2 semanas",
    "Avaliação: Técnico 40% / Consultivo 30% / Comportamental 20% / Prático 10%",
], "📖 A chave não é 'funcionar' — é a COMPARAÇÃO. Justifiquem com dados qual abordagem escolheriam em produção.\n➡️ Labs.", 59, T, "Projeto do módulo")

s = content_slide("Próximos Passos", [
    "1. Ler: Anthropic Building Effective Agents (15 min)",
    "   Foco no Apêndice 2 (ACI)",
    "2. Ler: ReAct paper — arXiv:2210.03629 (20 min)",
    "3. Rodar: exemplo ReAct do repositório",
    "   19-Examples/ETHAGT01/exemplo-01-react-loop/main.py",
    "4. Iniciar: Lab 1 — 'ReAct em 50 linhas'",
    "5. Próxima aula: ETHAGT02 — Tool Calling e ACI",
], "📖 O mais importante é ler o Apêndice 2 — é a base da próxima aula.\n➡️ Q&A.", 60, T, "Próximos passos")

# Slides 61-62 (labs + Q&A)
s = content_slide("Labs do Módulo", [
    "Lab 1 (4h): 'ReAct em 50 linhas'",
    "  Implementar agent loop em Python puro com tool de cálculo",
    "  Entrega: repo + trace de exemplo",
    "",
    "Lab 2 (4h): 'Augmented LLM'",
    "  Adicionar retrieval (ChromaDB/Qdrant) + 2 tools ao Lab 1",
    "  O agente decide quando recuperar vs usar tool",
    "  Entrega: agente funcional + traces",
    "",
    "Ambos individuais. Usar 19-Examples/ETHAGT01/ como referência.",
], "📖 Lab 1 é obrigatório antes do projeto. Lab 2 adiciona retrieval e mais tools.\n➡️ Encerramento.", 61, T, "Labs")

s = title_slide(
    "Perguntas?",
    "ETHAGT01 — Arquitetura Cognitiva de Agentes LLM\nPróxima aula: ETHAGT02 — Tool Calling e ACI",
    "Universidade Etho"
)
add_notes(s, "📖 Abrir para Q&A. Se não houver perguntas: 'Qual parte foi menos clara?'\nMensagem final: 'O bloco fundamental é o Augmented LLM. Tudo o que vem depois é composição. Se entenderam isso, a especialização vai fluir.'\n➡️ Fim da aula.")

# ─── Salvar ───
output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT01-Apresentacao.pptx")
prs.save(output)
print(f"PPTX gerado: {output}")
print(f"Total de slides: {len(prs.slides)}")
