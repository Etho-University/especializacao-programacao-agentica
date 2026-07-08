#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT02: Tool Calling e Agent-Computer Interface (ACI)
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Paleta Etho ───
ETHO_PRIMARY   = RGBColor(0x1B, 0x3A, 0x5E)
ETHO_ACCENT    = RGBColor(0xE8, 0x5D, 0x2F)
ETHO_DARK      = RGBColor(0x0F, 0x1E, 0x2D)
ETHO_LIGHT     = RGBColor(0xF7, 0xF8, 0xFA)
ETHO_MUTED     = RGBColor(0x6B, 0x7A, 0x8F)
ETHO_SUCCESS   = RGBColor(0x2D, 0x86, 0x59)
ETHO_DANGER    = RGBColor(0xC0, 0x39, 0x2B)
ETHO_INFO      = RGBColor(0x29, 0x80, 0xB9)
ETHO_WARNING   = RGBColor(0xD4, 0xA0, 0x17)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
BLACK          = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ───

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=ETHO_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=16, color=ETHO_DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

def add_footer(slide, slide_num, total, objetivo=""):
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(6), Inches(0.4),
                objetivo, font_size=10, color=ETHO_MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4),
                f"Slide {slide_num} / {total}", font_size=10, color=ETHO_MUTED, alignment=PP_ALIGN.RIGHT)

def add_header_bar(slide, module_code="ETHAGT02"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ETHO_PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3),
                "Universidade Etho", font_size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3),
                module_code, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, module_code):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, ETHO_DARK)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                title, font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                subtitle, font_size=20, color=ETHO_MUTED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                module_code, font_size=14, color=ETHO_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    return slide

def section_slide(number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ETHO_PRIMARY)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(2), Inches(2),
                str(number), font_size=80, color=ETHO_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5),
                title, font_size=32, color=WHITE, bold=True)
    return slide

def content_slide(title, bullets, notes, slide_num, total, objetivo="", subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ETHO_LIGHT)
    add_header_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
                title, font_size=32, color=ETHO_PRIMARY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
                    subtitle, font_size=18, color=ETHO_MUTED)
    add_bullets(slide, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5),
                bullets, font_size=16, color=ETHO_DARK)
    add_footer(slide, slide_num, total, objetivo)
    add_notes(slide, notes)
    return slide

def code_slide(title, code_text, notes, slide_num, total, objetivo=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ETHO_DARK)
    add_header_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
                title, font_size=28, color=WHITE, bold=True)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(slide, slide_num, total, objetivo)
    add_notes(slide, notes)
    return slide

def comparison_slide(title, left_title, left_items, right_title, right_items, notes, slide_num, total, objetivo=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ETHO_LIGHT)
    add_header_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
                title, font_size=28, color=ETHO_PRIMARY, bold=True)
    # Left
    add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5),
                left_title, font_size=20, color=ETHO_DANGER, bold=True)
    add_bullets(slide, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5),
                left_items, font_size=15, color=ETHO_DARK)
    # Right
    add_textbox(slide, Inches(7), Inches(1.5), Inches(6), Inches(0.5),
                right_title, font_size=20, color=ETHO_SUCCESS, bold=True)
    add_bullets(slide, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5),
                right_items, font_size=15, color=ETHO_DARK)
    add_footer(slide, slide_num, total, objetivo)
    add_notes(slide, notes)
    return slide

def exercise_slide(title, enunciado, notes, slide_num, total, objetivo="Exercício"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ETHO_LIGHT)
    add_header_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7),
                title, font_size=28, color=ETHO_ACCENT, bold=True)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = ETHO_WARNING
    box.line.width = Pt(2)
    add_bullets(slide, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5),
                enunciado, font_size=16, color=ETHO_DARK)
    add_footer(slide, slide_num, total, objetivo)
    add_notes(slide, notes)
    return slide

# ─── Total de slides (atualizado dinamicamente) ───
TOTAL = 57

# ═══════════════════════════════════════════════════
# SEÇÃO A — Abertura (Slides 1-6)
# ═══════════════════════════════════════════════════

# Slide 1 — Capa
s = title_slide(
    "Tool Calling e Agent-Computer Interface (ACI)",
    "Universidade Etho · Especialização em Programação Agêntica\nFase A — Fundamentos Agênticos · 25 h",
    "ETHAGT02"
)
add_notes(s, """📖 EXPLICAÇÃO: Bem-vindos à segunda aula. Hoje vamos aprofundar o componente mais crítico do Augmented LLM: as tools. Na aula passada vimos que o Augmented LLM = LLM + retrieval + tools + memory. Hoje é o dia das tools.
💡 ANALOGIA: ETHAGT01 foi aprender a cozinhar. ETHAGT02 é aprender a afiar as facas — sem facas boas, o melhor chef falha.
❓ PERGUNTA: "Quantos de vocês leram o Apêndice 2 da Anthropic?" (verificar homework)
➡️ TRANSIÇÃO: Vamos aos objetivos.""")

# Slide 2 — Objetivos
s = content_slide(
    "Objetivos do Módulo",
    [
        "Objetivo geral: projetar, documentar e validar tools que tornam um agente confiável",
        "",
        "Objetivos específicos:",
        "1. Dominar o mecanismo de tool calling (function calling, JSON Schema, structured outputs)",
        "2. Aplicar princípios ACI: poka-yoke, exemplos, paths absolutos, formato natural",
        "3. Projetar tools idempotentes com tratamento de erro e timeouts",
        "4. Distinguir tools seguras de destrutivas; aplicar HITL onde necessário",
        "5. Avaliar empiricamente o uso de tools (workbench, iterar)",
    ],
    """📖 EXPLICAÇÃO: Cada objetivo é mensurável. O mais importante é #2 — ACI. Se você só lembrar de uma coisa desta aula, deve ser: invista em ACI tanto quanto em HCI.
❓ PERGUNTA: "Qual objetivo vocês acham mais difícil?" (geralmente #5 — avaliar empiricamente)
➡️ TRANSIÇÃO: Competências.""",
    2, TOTAL, "Estabelecer objetivos mensuráveis"
)

# Slide 3 — Competências
s = content_slide(
    "Competências Desenvolvidas",
    [
        "C1 Programação Agêntica → I (Intermediário)",
        "C3 MCP & Tool Use → I (Intermediário) — núcleo do módulo",
        "C5 AgentOps & Avaliação → B (Básico)",
        "C6 Agent Security → B (Básico) — introdução via matriz de risco",
        "",
        "Conexão com ETHAGT01: C1 saiu de B → I; C3 era B, agora atinge I",
    ],
    """📖 EXPLICAÇÃO: Este módulo eleva C3 (MCP & Tool Use) de Básico para Intermediário. É o núcleo da aula. C6 (Security) é introduzido como Básico via a matriz de risco e HITL — o aprofundamento vem em ETHAGT13.
➡️ TRANSIÇÃO: Agenda.""",
    3, TOTAL, "Conectar ao Framework Etho"
)

# Slide 4 — Agenda
s = content_slide(
    "Agenda da Aula",
    [
        "Bloco 1 (45 min):",
        "  • Abertura (8 min) — motivação, contexto ACI",
        "  • Mecanismo do Tool Calling (15 min) — function calling, JSON Schema",
        "  • ACI como Disciplina (15 min) — 5 princípios, caso SWE-bench",
        "  • Intervalo (5 min)",
        "",
        "Bloco 2 (45 min):",
        "  • Engenharia de Tools (12 min) — schemas, idempotência, tipologia",
        "  • Tools Perigosas e HITL (8 min) — matriz de risco",
        "  • Erros Comuns e Avaliação (10 min) — workbench, métricas",
        "  • Fechamento (15 min) — boas práticas, quiz, Q&A",
    ],
    """📖 EXPLICAÇÃO: O bloco 1 é teoria + princípios. O bloco 2 é prática + engenharia. O quiz final tem 5 perguntas.
➡️ TRANSIÇÃO: Por que ACI importa?""",
    4, TOTAL, "Mostrar roteiro da aula"
)

# Slide 5 — Motivação
s = content_slide(
    "O Problema: Tools Mal Desenhadas",
    [
        "Tools mal desenhadas fazem o agente falhar SILENTIOSAMENTE",
        "",
        "Cenário real (Anthropic SWE-bench):",
        "  • Agente saía do diretório raiz → usava paths relativos errados",
        "  • Resultado: editava arquivo errado ou falha 'arquivo não encontrado'",
        "  • O modelo não estava 'errado' — a tool permitia o erro",
        "",
        "Pergunta: O que acontece quando uma tool devolve erro não tratado?",
        "",
        "Resposta: o agente alucina uma solução ou entra em loop",
    ],
    """📖 EXPLICAÇÃO: A falha silenciosa é o pior tipo de bug. O agente não crasha — ele continua operando com informação errada. No caso da Anthropic, o agente editava o arquivo errado e seguia em frente. Sem traces, ninguém percebe.
💡 ANALOGIA: É como um mecânico com uma chave de boca mal calibrada. Ele aperta o parafuso achando que está certo, mas a chave escorrega. O parafuso fica frouxo, o carro anda, e o problema só aparece 1000 km depois.
❓ PERGUNTA: "Vocês já tiveram um bug onde a 'solução' era melhorar o prompt, mas o real problema era a tool?"
➡️ TRANSIÇÃO: A solução é tratar ACI como disciplina.""",
    5, TOTAL, "Criar tensão — tools mal desenhadas = falha silenciosa"
)

# Slide 6 — Contexto
s = content_slide(
    "ACI como Disciplina de Design",
    [
        "ACI = Agent-Computer Interface (análogo a HCI — Human-Computer Interface)",
        "",
        "Princípio da Anthropic:",
        "  \"Invista tanto esforço em ACI quanto em HCI\"",
        "",
        "Na prática:",
        "  • Tempo em tools > tempo no prompt principal",
        "  • Descrição de tool = docstring para dev júnior",
        "  • Poka-yoke: tornar o erro estruturalmente impossível",
        "  • Testar com N inputs antes de confiar",
        "",
        "Resultado: agente confiável sem 'prompt engineering' heroico",
    ],
    """📖 EXPLICAÇÃO: ACI é o conceito central de hoje. A Anthropic passou MAIS tempo nas tools do que no prompt principal do SWE-bench. Isso inverte a intuição: a maioria das pessoas passa 90% do tempo no prompt e 10% nas tools. O certo é o inverso.
💡 ANALOGIA: É como design de UI. Se usuários clicam no botão errado, a solução não é 'treinar os usuários' — é mover o botão. Com agentes: se o modelo chama a tool errada, a solução não é 'prompt melhor' — é redesenhar a tool.
➡️ TRANSIÇÃO: Vamos entender o mecanismo antes dos princípios.""",
    6, TOTAL, "Introduzir ACI como disciplina"
)

# ═══════════════════════════════════════════════════
# SEÇÃO B — O Mecanismo do Tool Calling (Slides 7-15)
# ═══════════════════════════════════════════════════

s = section_slide(1, "O Mecanismo do Tool Calling")
add_notes(s, "Início do bloco de mecanismo. Como o tool calling funciona tecnicamente.")

# Slide 8 — Function Calling
s = content_slide(
    "Function Calling: Do Prompt para JSON",
    [
        "Antes: parsear texto livre do LLM (frágil, quebra com formatação)",
        "Agora: JSON Schema define o contrato",
        "",
        "Fluxo:",
        "  1. Sistema registra tools com nome, descrição, schema",
        "  2. LLM recebe prompt + lista de tools",
        "  3. LLM decide: resposta direta OU tool_call (JSON com args)",
        "  4. Sistema executa a função",
        "  5. Resultado volta como 'tool' message",
        "  6. LLM consome o resultado e continua",
        "",
        "O JSON é estruturado — não é texto livre parseado",
    ],
    """📖 EXPLICAÇÃO: O function calling foi o que tornou tool use confiável. Antes, você pedia pro LLM "responda em JSON" e parseava — frágil. Agora, o modelo gera um tool_call com schema definido. O sistema executa e retorna o resultado. O modelo consome.
💡 ANALOGIA: É como a diferença entre pedir para alguém "me passe o parafuso de 10mm" (texto livre — pode trazer o errado) vs ter um catálogo com código de barras (JSON Schema — não tem ambiguidade).
❓ PERGUNTA: "Quem aqui já usou function calling?" (calibrar)
➡️ TRANSIÇÃO: O JSON Schema é o contrato.""",
    8, TOTAL, "Explicar o mecanismo de function calling"
)

# Slide 9 — JSON Schema
s = code_slide(
    "JSON Schema como Contrato",
    '''{
  "type": "function",
  "function": {
    "name": "search_product",
    "description": "Busca produto por nome no catálogo. Retorna preço, estoque e SKU. Use quando o usuário perguntar sobre preço ou disponibilidade. NÃO use para listar categorias — use list_categories.",
    "parameters": {
      "type": "object",
      "properties": {
        "query": {
          "type": "string",
          "description": "Nome do produto ou parte. Ex: 'iPhone 15' ou 'notebook Dell'"
        },
        "category": {
          "type": "string",
          "enum": ["electronics", "clothing", "food", "books"],
          "description": "Filtro de categoria. Omita se não souber."
        }
      },
      "required": ["query"]
    }
  }
}''',
    """📖 EXPLICAÇÃO: Note três coisas neste schema: (1) a descrição diz quando USAR e quando NÃO USAR; (2) category é enum — o modelo não pode inventar categorias; (3) query tem exemplo na descrição. Isso é ACI aplicada.
💡 Poka-yoke aplicados: enum em category (não pode errar), descrição com exemplo, regra de fronteira (não usar para listar categorias).
➡️ TRANSIÇÃO: Structured outputs.""",
    9, TOTAL, "Mostrar JSON Schema como contrato"
)

# Slide 10 — Structured Outputs
s = content_slide(
    "Structured Outputs / Constrained Decoding",
    [
        "O que é: forçar o LLM a gerar apenas JSON válido conforme schema",
        "",
        "Como funciona (visão):",
        "  • Constrained decoding: restringe tokens durante geração",
        "  • Garante 100% de conformidade com o schema",
        "  • Elimina erros de parse",
        "",
        "Vantagem: zero erros de formatação",
        "Custo: pode limitar criatividade do modelo em edge cases",
        "",
        "OpenAI: response_format = { type: 'json_schema', json_schema: {...} }",
        "Anthropic: tool_choice = { type: 'tool', name: '...' }",
    ],
    """📖 EXPLICAÇÃO: Structured outputs é o próximo nível. Em vez de o modelo gerar texto e torcer que o JSON é válido, o decoding é constrangido — o modelo SÓ pode gerar tokens que produzem JSON válido. Isso elimina 100% dos erros de parse. O custo é que em edge cases raros, o modelo pode ficar "preso" no schema e não conseguir expressar algo.
➡️ TRANSIÇÃO: E se o modelo precisa chamar múltiplas tools?""",
    10, TOTAL, "Introduzir structured outputs"
)

# Slide 11 — Multi-tool Calls
s = content_slide(
    "Multi-tool Calls e Paralelismo",
    [
        "O modelo pode chamar múltiplas tools em UMA resposta",
        "",
        "Exemplo: 'Qual o preço do iPhone E do Samsung?'",
        "  → tool_call: search_product('iPhone')",
        "  → tool_call: search_product('Samsung')",
        "  (executadas em paralelo)",
        "",
        "Vantagem: latência reduzida (paralelo vs serial)",
        "Risco: ordem importa se há dependência",
        "",
        "Boa prática: documentar se a tool é independente ou sequencial",
    ],
    """📖 EXPLICAÇÃO: Múltiplas tools em paralelo reduz latência. Mas cuidado: se tool B depende do resultado de tool A, não pode paralelizar. O modelo precisa saber — inclua na descrição se a tool tem dependências.
➡️ TRANSIÇÃO: Quanto custa isso em tokens?""",
    11, TOTAL, "Explicar multi-tool calls paralelas"
)

# Slide 12 — Custo de Tokens
s = content_slide(
    "Custo: Tokens de Descrição vs Benefício",
    [
        "Cada tool adiciona tokens ao prompt (descrição + schema)",
        "",
        "Estimativa por tool bem descrita:",
        "  • ~150-300 tokens (nome + descrição + schema + exemplos)",
        "  • 10 tools = ~2000-3000 tokens extras por chamada",
        "",
        "Trade-off:",
        "  • Mais tools = mais capacidades mas mais custo",
        "  • Descrições ricas = melhor uso mas mais tokens",
        "",
        "Regra: 5-10 tools bem descritas > 50 tools vagas",
        "Se precisar de 50+ tools: usar tool retrieval (ETAGT08 — MCP)",
    ],
    """📖 EXPLICAÇÃO: Cada tool adiciona ~200 tokens ao prompt. 10 tools = 2k tokens extras. Em 10k chamadas/dia, são 20M tokens extras = $3-30/dia dependendo do modelo. Gorilla (arXiv:2305.15334) mostrou que com 1700+ APIs, a performance cai drasticamente. A solução para muitas tools é tool retrieval — o modelo busca apenas as tools relevantes. Isso é MCP (ETHAGT08).
❓ PERGUNTA: "Quantas tools vocês têm no sistema de vocês?" (calibrar)
➡️ TRANSIÇÃO: Vamos ver um exemplo comparativo.""",
    12, TOTAL, "Mostrar custo de tokens por tool"
)

# Slide 13 — Exemplo: Bem vs Mal Descrita
s = comparison_slide(
    "Tool Bem vs Mal Descrita",
    "❌ Mal Descrita",
    [
        'name: "calc"',
        'description: "Faz cálculo"',
        'parameters:',
        '  a: number',
        '  b: number',
        '  op: string',
        '',
        "Problemas:",
        "• Nome não descritivo",
        "• Descrição vaga",
        "• op é string livre (não enum)",
        "• Sem exemplos",
        "• Sem fronteiras",
    ],
    "✅ Bem Descrita",
    [
        'name: "calculator"',
        'description: "Executa +, -, *, / em dois números."',
        'parameters:',
        '  a: number (ex: 123)',
        '  b: number (ex: 456)',
        '  op: enum ["+", "-", "*", "/"]',
        '',
        "Melhorias:",
        "• Nome descritivo",
        "• Descrição com operadores",
        "• op é enum (poka-yoke)",
        "• Exemplos nos params",
    ],
    """📖 EXPLICAÇÃO: A diferença é clara. A tool mal descrita permite o modelo passar op="multiply" ou op="soma" — vai quebrar. A bem descrita com enum só permite os 4 operadores. Isso é poka-yoke: tornar o erro impossível.
➡️ TRANSIÇÃO: Demo ao vivo.""",
    13, TOTAL, "Comparar tool bem vs mal descrita"
)

# Slide 14 — DEMO
s = code_slide(
    "DEMO: Tool Calling na Prática",
    '''from openai import OpenAI
import json

client = OpenAI()

TOOLS = [{
    "type": "function",
    "function": {
        "name": "calculator",
        "description": "Executa +, -, *, / em dois números.",
        "parameters": {
            "type": "object",
            "properties": {
                "a": {"type": "number", "description": "Primeiro número. Ex: 123"},
                "b": {"type": "number", "description": "Segundo número. Ex: 456"},
                "op": {"type": "string", "enum": ["+","-","*","/"],
                        "description": "Operação"}
            },
            "required": ["a", "b", "op"]
        }
    }
}]

messages = [
    {"role": "system", "content": "Você é um agente. Use tools quando necessário."},
    {"role": "user", "content": "Quanto é 123 * 456?"}
]

response = client.chat.completions.create(
    model="gpt-4o-mini", messages=messages, tools=TOOLS, tool_choice="auto"
)
msg = response.choices[0].message

if msg.tool_calls:
    for tc in msg.tool_calls:
        args = json.loads(tc.function.arguments)
        result = eval(f"{args['a']} {args['op']} {args['b']}")
        print(f"Tool: {tc.function.name}({args}) = {result}")''',
    """📖 EXPLICAÇÃO: Demo ao vivo. Mostrar: (1) como o modelo escolhe a tool; (2) como os args vêm estruturados; (3) como o resultado volta. Se API falhar, tenho screenshot.
❓ PERGUNTA: "Viram como o enum garantiu que op foi '*' e não 'multiply'?"
➡️ TRANSIÇÃO: Pergunta da demo.""",
    14, TOTAL, "Demo ao vivo de tool calling"
)

# Slide 15 — Pergunta da Demo
s = exercise_slide(
    "Pergunta da Demo",
    [
        "Em duplas (2 min):",
        "",
        "1. O que acontece se o modelo passar op='×' (x multiplicação) em vez de '*'?",
        "   Como o enum ajuda aqui?",
        "",
        "2. O que acontece se a=123 e b=0 e op='/'?",
        "   A tool deve tratar? Ou o modelo deve saber?",
        "",
        "3. Como você adicionaria uma tool de 'potência' sem quebrar a existente?",
    ],
    """📖 EXPLICAÇÃO: Respostas: (1) O enum rejeita '×' — o modelo SÓ pode gerar '+', '-', '*', '/'. Poka-yoke. (2) A tool DEVE tratar divisão por zero e retornar {"error": "division by zero"} como Observation. O modelo decide o que fazer. (3) Adicionar power(a, b) como nova tool, ou estender calculator com op="^" — mas cuidado: "^" pode confundir com XOR.
➡️ TRANSIÇÃO: Agora os princípios de ACI.""",
    15, TOTAL
)

# ═══════════════════════════════════════════════════
# SEÇÃO C — ACI como Disciplina (Slides 16-26)
# ═══════════════════════════════════════════════════

s = section_slide(2, "ACI como Disciplina de Design")
add_notes(s, "Transição para os princípios de ACI.")

# Slide 17 — ACI :: HCI
s = content_slide(
    "A Analogia: ACI :: HCI",
    [
        "HCI (Human-Computer Interface):",
        "  • Como humanos interagem com computadores",
        "  • Décadas de pesquisa: usabilidade, affordance, feedback",
        "  • Investimento óbvio — todo produto digital tem UI/UX",
        "",
        "ACI (Agent-Computer Interface):",
        "  • Como agentes (LLMs) interagem com computadores (tools)",
        "  • Mesmos princípios: usabilidade, affordance, feedback",
        "  • Investimento NÃO óbvio — a maioria esquece",
        "",
        "Princípio: invista em ACI TANTO quanto em HCI",
        "  → Se sua UI tem designer, suas tools precisam também",
    ],
    """📖 EXPLICAÇÃO: A analogia é poderosa. Tudo que você sabe sobre HCI se aplica a ACI. Affordance: um botão parece clicável → uma tool parece chamável (descrição clara). Feedback: UI mostra loading → tool retorna Observation estruturada. Consistência: mesma ação = mesmo botão → mesma operação = mesma tool.
💡 ANALOGIA: Se você contratasse um funcionário e desse ferramentas sem manual, sem treinamento, seria irresponsável. Com agentes é igual: tools sem descrição = ferramentas sem manual.
➡️ TRANSIÇÃO: Primeiro princípio.""",
    17, TOTAL, "Estabelecer a analogia ACI :: HCI"
)

# Slide 18 — Princípio 1
s = content_slide(
    "Princípio 1: Ponha-se no Lugar do Modelo",
    [
        "Pergunta-chave: \"É óbvio como usar esta tool apenas pela descrição?\"",
        "",
        "Se você precisa pensar muito para entender → o modelo também precisa",
        "",
        "Teste prático:",
        "  1. Leia apenas o nome e a descrição da tool",
        "  2. Tente pensar: quando eu chamaria isso?",
        "  3. Se a resposta for 'depende' → a descrição está incompleta",
        "",
        "Uma boa tool definition inclui:",
        "  • Exemplo de uso",
        "  • Edge cases conhecidos",
        "  • Requisitos de formato",
        "  • Fronteira com outras tools (quando usar esta vs aquela)",
    ],
    """📖 EXPLICAÇÃO: O teste é simples: leia a descrição como se fosse um júnior. Se você ficar confuso, o modelo também ficará. A diferença é que o júnior pode perguntar; o modelo não pode — ele adivinha. E quando adivinha, erra.
➡️ TRANSIÇÃO: Segundo princípio.""",
    18, TOTAL, "Princípio 1 — empatia com o modelo"
)

# Slide 19 — Princípio 2
s = content_slide(
    "Princípio 2: Descrições Ricas (Docstring de Júnior)",
    [
        "Descrição de tool = docstring para um dev júnior",
        "",
        "Deve incluir:",
        "  • O que a tool faz (não o que É, o que FAZ)",
        "  • Quando usar (gatilhos)",
        "  • Quando NÃO usar (fronteiras)",
        "  • Formato esperado dos parâmetros",
        "  • Exemplos",
        "  • Edge cases",
        "",
        "Anti-pattern: 'Faz busca' → não diz o que busca, nem quando, nem como",
        "Boa prática: 'Busca produto por nome no catálogo. Retorna preço e estoque.'",
    ],
    """📖 EXPLICAÇÃO: 'Faz busca' é a descrição mais comum e mais inútil. Busca o quê? Onde? Retorna o quê? Quando usar vs quando não usar? A boa descrição responde todas essas perguntas em 1-2 frases.
➡️ TRANSIÇÃO: Terceiro princípio.""",
    19, TOTAL, "Princípio 2 — descrições ricas"
)

# Slide 20 — Princípio 3
s = content_slide(
    "Princípio 3: Formato Próximo ao Natural",
    [
        "O modelo escreve melhor em formatos que vê naturalmente na web",
        "",
        "❌ Ruim: pedir para o modelo gerar diffs com headers de chunk",
        "  (exige contar linhas antes de escrever o código)",
        "",
        "✅ Bom: pedir para o modelo reescrever o arquivo inteiro",
        "  (formato natural — texto plano)",
        "",
        "❌ Ruim: JSON com escaping de newlines e quotes",
        "✅ Bom: Markdown com code blocks",
        "",
        "Regra: se o formato tem 'overhead' (contar linhas, escaping), mude",
    ],
    """📖 EXPLICAÇÃO: A Anthropic descobriu que o formato importa tanto quanto a descrição. Diffs exigem que o modelo saiba quantas linhas mudaram ANTES de escrever o código — isso é cognitivamente caro. Reescrever o arquivo inteiro é mais natural. JSON com escaping é difícil; markdown é fácil.
💡 ANALOGIA: É como pedir para alguém escrever uma carta. Se você pedir em formato de diff (com @@ -1,3 +1,5 @@), a pessoa precisa contar linhas. Se pedir em texto plano, flui natural.
➡️ TRANSIÇÃO: Quarto princípio — poka-yoke.""",
    20, TOTAL, "Princípio 3 — formato natural"
)

# Slide 21 — Princípio 4
s = content_slide(
    "Princípio 4: Poka-yoke (Tornar o Erro Impossível)",
    [
        "Poka-yoke = design que previne erros (conceito japonês de manufatura)",
        "",
        "Em tools:",
        "  • Enum em vez de string livre → não pode inventar valor",
        "  • Paths absolutos em vez de relativos → não pode errar caminho",
        "  • required fields → não pode omitir críticos",
        "  • pattern regex → valida formato antes de executar",
        "  • default values → reduz carga cognitiva do modelo",
        "",
        "Caso SWE-bench:",
        "  Antes: aceitava paths relativos → modelo errava após mudar de diretório",
        "  Depois: exige paths absolutos → erro estruturalmente impossível",
    ],
    """📖 EXPLICAÇÃO: Poka-yoke é o princípio mais poderoso. Em vez de confiar que o modelo 'lembre' de usar paths absolutos (frágil), mude a tool para EXIGIR paths absolutos (robusto). O modelo não pode errar mesmo que queira. Isso é design defensivo.
💡 ANALOGIA: É como a proteção da tomada: você não ensina a criança a não colocar o dedo (prompt); você coloca a proteção (poka-yoke).
➡️ TRANSIÇÃO: Quinto princípio.""",
    21, TOTAL, "Princípio 4 — poka-yoke"
)

# Slide 22 — Princípio 5
s = content_slide(
    "Princípio 5: Tokens para Pensar",
    [
        "\"Give the model enough tokens to think before it writes itself into a corner\"",
        "— Anthropic",
        "",
        "O que significa:",
        "  • Não restrinja a resposta a poucos tokens se a tarefa é complexa",
        "  • Permita que o modelo 'raciocine' antes de produzir o output final",
        "  • max_tokens generoso > max_tokens restritivo",
        "",
        "Na prática:",
        "  • Tool de geração de código: max_tokens=4096, não 512",
        "  • Tool de análise: permita reasoning antes do JSON",
        "  • Custo: mais tokens = mais caro, mas qualidade é melhor",
    ],
    """📖 EXPLICAÇÃO: Se você cortar tokens cedo demais, o modelo 'se encurrala' — começa a escrever algo e não consegue terminar. O resultado é truncado ou errado. Dar espaço para pensar é investir em qualidade.
➡️ TRANSIÇÃO: Vamos praticar identificar o que falta.""",
    22, TOTAL, "Princípio 5 — tokens para pensar"
)

# Slide 23 — Exercício
s = exercise_slide(
    "Exercício: O Que Está Faltando?",
    [
        "Leia esta descrição de tool:",
        "",
        '  name: "send_email"',
        '  description: "Envia email"',
        '  parameters: { to: string, subject: string, body: string }',
        "",
        "Em trios (3 min):",
        "1. Liste 5 coisas que estão faltando",
        "2. Reescreva a descrição aplicando os 5 princípios",
        "3. Quais poka-yokes você adicionaria?",
    ],
    """📖 EXPLICAÇÃO: Faltam: (1) descrição rica — 'envia email' não diz para quem, com que formato, nem limitações; (2) validação de email (pattern regex); (3) max length no body; (4) não diz se é HTML ou texto plano; (5) não há HITL (enviar email é ação destrutiva); (6) não há rate limiting; (7) não diz o que retorna (success? error? message_id?). Poka-yokes: pattern de email, enum de formato, required fields, dry-run opcional.
➡️ TRANSIÇÃO: Vamos ver o antes e depois de uma refatoração real.""",
    23, TOTAL
)

# Slide 24 — Antes vs Depois
s = comparison_slide(
    "Antes vs Depois: Refatoração ACI",
    "❌ Antes",
    [
        'name: "send_email"',
        'description: "Envia email"',
        'parameters:',
        '  to: string',
        '  subject: string',
        '  body: string',
        '',
        "Problemas:",
        "• Descrição vaga",
        "• Sem validação de email",
        "• Sem formato (HTML/text)",
        "• Sem HITL",
        "• Sem rate limit",
    ],
    "✅ Depois",
    [
        'name: "send_email"',
        'description: "Envia email texto ou HTML."',
        '  Use para comunicação com usuários.',
        '  NÃO use para emails em massa — use bulk_email.',
        'parameters:',
        '  to: string (pattern: email regex)',
        '  subject: string (maxLength: 200)',
        '  body: string (maxLength: 10000)',
        '  format: enum ["text", "html"]',
        '  dry_run: boolean (default: true)',
        '',
        "Melhorias: poka-yoke em tudo",
    ],
    """📖 EXPLICAÇÃO: A diferença é dramática. A versão refatorada tem: enum em formato, pattern em email, maxLength em subject e body, dry_run default true (HITL), fronteira com bulk_email. O modelo não pode mais enviar email inválido ou sem confirmação.
➡️ TRANSIÇÃO: Vamos ver o caso real da Anthropic.""",
    24, TOTAL, "Comparar antes vs depois de refatoração ACI"
)

# Slide 25 — Caso SWE-bench
s = content_slide(
    "Caso Real: Anthropic SWE-bench — Redesign de Tools",
    [
        "Contexto: coding agent para resolver issues de GitHub",
        "",
        "Problema:",
        "  • Agente navegava fora do diretório raiz",
        "  • Começava a usar paths relativos incorretos",
        "  • Editava arquivo errado ou falhava 'não encontrado'",
        "",
        "Tentação: ajustar o prompt ('lembre-se de usar paths absolutos')",
        "  → Frágil: modelo esquece em conversas longas",
        "",
        "Solução: mudar a TOOL para exigir paths absolutos sempre",
        "  → Robusto: erro estruturalmente impossível",
        "  → Resultado: 'o modelo usou a tool perfeitamente' (Schluntz)",
    ],
    """📖 EXPLICAÇÃO: Este é o caso canônico de ACI. A lição: erros sistemáticos são sinais de design de tool, não de prompt. Antes de reescrever o prompt pela 3ª vez, pergunte: posso redesenhar a tool?
💡 ANALOGIA: Se um usuário clica sempre no botão errado, a solução não é 'treinar o usuário'. É mover o botão.
➡️ TRANSIÇÃO: A lição consolidada.""",
    25, TOTAL, "Caso real de redesign de tool"
)

# Slide 26 — A Lição
s = content_slide(
    "A Lição: Tempo em Tools > Tempo no Prompt",
    [
        "\"We spent more time optimizing our tools than the overall prompt\"",
        "— Anthropic (Building Effective Agents, Appendix 2)",
        "",
        "Mudança de mentalidade:",
        "  ❌ Antes: 90% prompt, 10% tools",
        "  ✅ Depois: 30% prompt, 70% tools",
        "",
        "Checklist mental antes de mexer no prompt:",
        "  1. A tool está bem descrita?",
        "  2. O schema tem poka-yoke?",
        "  3. Há exemplos na descrição?",
        "  4. Há fronteiras entre tools similares?",
        "  5. O erro é sistemático ou ocasional?",
        "     (Sistemático = design de tool; Ocasional = talvez prompt)",
    ],
    """📖 EXPLICAÇÃO: Esta é a inversão mental que a aula inteira tenta provocar. Se o erro é sistemático (acontece sempre da mesma forma), é design de tool. Se é ocasional (às vezes acontece), pode ser prompt. Antes de reescrever o prompt, passe pelo checklist.
➡️ TRANSIÇÃO: Intervalo. Voltamos para engenharia de tools.""",
    26, TOTAL, "Fixar a lição principal"
)

# ═══════════════════════════════════════════════════
# SEÇÃO D — Engenharia de Tools (Slides 27-36)
# ═══════════════════════════════════════════════════

s = section_slide(3, "Engenharia de Tools")
add_notes(s, "Transição para engenharia de tools — schemas, idempotência, tipologia.")

# Slide 27 — Schemas Claros
s = content_slide(
    "Schemas Claros: Pydantic, TypedDict, Zod",
    [
        "Pydantic (Python): validação automática + serialização",
        "TypedDict (Python): tipo estático sem runtime",
        "Zod (TypeScript): validação com inferência de tipo",
        "",
        "Vantagem Pydantic:",
        "  • Valida args antes de executar a tool",
        "  • Mensagens de erro estruturadas",
        "  • Compatível com OpenAI function calling",
        "",
        "Exemplo:",
        "  class SearchProduct(BaseModel):",
        "      query: str = Field(..., min_length=1, max_length=200)",
        "      category: Literal['electronics','clothing','food'] | None = None",
    ],
    """📖 EXPLICAÇÃO: Pydantic é o padrão de facto em Python. Define o schema como classe, valida automaticamente, e gera o JSON Schema para a API. Se o modelo passar args inválidos, Pydantic rejeita antes de executar — mais um poka-yoke.
➡️ TRANSIÇÃO: Idempotência.""",
    27, TOTAL, "Schemas com Pydantic/TypedDict/Zod"
)

# Slide 28 — Idempotência
s = content_slide(
    "Idempotência",
    [
        "Idempotência: chamar a mesma tool com os mesmos args N vezes = mesmo resultado",
        "",
        "ONDE IMPORTA:",
        "  ✅ Cobrança (não cobrar 2x se retry)",
        "  ✅ Envio de email (não enviar 2x)",
        "  ✅ DB writes (não criar 2 registros)",
        "  ✅ Deploy (não fazer deploy 2x)",
        "",
        "COMO GARANTIR:",
        "  • request_id (UUID) — rejeita duplicatas",
        "  • Idempotency key header",
        "  • Check antes de executar (já existe?)",
        "  • Transactions com unique constraint",
    ],
    """📖 EXPLICAÇÃO: Idempotência é crítica em agentes porque retries acontecem. O agente chama a tool, o timeout dispara, o agente tenta de novo — se a tool não é idempotente, você cobra 2x ou envia 2 emails. request_id é a solução mais simples: a tool guarda o ID e rejeita duplicatas.
💡 ANALOGIA: É como um co-piloto que confirma: "Já fizemos isso?" antes de repetir.
➡️ TRANSIÇÃO: Timeouts e retries.""",
    28, TOTAL, "Idempotência em tools"
)

# Slide 29 — Timeouts, Retries, Fallbacks
s = content_slide(
    "Timeouts, Retries, Fallbacks",
    [
        "Toda tool que chama serviço externo PRECISA de:",
        "",
        "1. TIMEOUT — não deixe o agente esperar indefinidamente",
        "   • Default: 30s para APIs, 5s para DB",
        "   • Retornar como Observation: {error: 'timeout after 30s'}",
        "",
        "2. RETRY — falhas transientes existem",
        "   • Max 3 retries com backoff exponencial (1s, 2s, 4s)",
        "   • NÃO retryar erros 4xx (são permanentes)",
        "",
        "3. FALLBACK — o que fazer se tudo falha",
        "   • Retornar erro estruturado para o modelo",
        "   • NÃO deixar exceção propagar (quebra o loop)",
    ],
    """📖 EXPLICAÇÃO: Sem timeout, uma API lenta pode travar o agente por minutos. Sem retry, uma falha transiente (500 temporário) mata a execução. Sem fallback, uma exceção não tratada quebra o loop. Cada um desses é um poka-yoke de produção.
➡️ TRANSIÇÃO: Tipologia de tools.""",
    29, TOTAL, "Timeouts, retries, fallbacks"
)

# Slide 30 — Tipologia
s = content_slide(
    "Tipologia: Leitura / Escrita / Destrutiva / External",
    [
        "1. LEITURA (segura): search, get, list, query",
        "   → Sem side effects, idempotente, sem HITL",
        "",
        "2. ESCRITA (modifica estado): create, update, set",
        "   → Side effect, precisa idempotência, HITL opcional",
        "",
        "3. DESTRUTIVA (irreversível): delete, drop, reset",
        "   → HITL OBRIGATÓRIO, dry-run antes, log de auditoria",
        "",
        "4. EXTERNAL SIDE-EFFECT (mundo real): email, deploy, transfer",
        "   → HITL OBRIGATÓRIO, confirmação explícita, allowlist",
        "",
        "A tipologia determina o nível de proteção necessário",
    ],
    """📖 EXPLICAÇÃO: A tipologia guia as decisões de segurança. Leitura é segura — pode rodar sem proteção. Escrita modifica estado — precisa idempotência. Destrutiva é irreversível — HITL obrigatório. External afeta o mundo real — máxima proteção.
➡️ TRANSIÇÃO: Matriz de risco e HITL.""",
    30, TOTAL, "Tipologia de tools"
)

# Slide 31 — Versionamento
s = content_slide(
    "Versionamento de Tools",
    [
        "Problema: mudar uma tool em produção pode quebrar agentes existentes",
        "",
        "Estratégias:",
        "1. Versionamento explícito: search_product_v1, search_product_v2",
        "   → Simples, mas polui o namespace",
        "",
        "2. Versionamento semântico: search_product (v1.2.0)",
        "   → Melhor, mas o modelo precisa entender versões",
        "",
        "3. Compatibilidade retroativa:",
        "   • Adicionar campos opcionais = OK",
        "   • Remover campos = QUEBRA",
        "   • Mudar tipo de campo = QUEBRA",
        "   • Adicionar enum values = OK",
        "",
        "Regra: NUNCA quebrar uma tool em produção. Sempre adicionar nova.",
    ],
    """📖 EXPLICAÇÃO: Versionamento de tools é como versionamento de API. Adicionar campos opcionais é seguro. Remover ou mudar tipo quebra. Se precisar mudar, crie uma nova tool (v2) e deprecie a antiga gradualmente.
➡️ TRANSIÇÃO: Quando usar 1 tool com mode vs 2 tools?""",
    31, TOTAL, "Versionamento de tools"
)

# Slide 32 — 1 tool vs 2 tools
s = content_slide(
    "Quando 1 Tool com mode vs 2 Tools Separadas?",
    [
        "1 tool com parâmetro mode:",
        "  search(mode='product' | 'category' | 'user')",
        "  ✅ Menos tools no namespace (economiza tokens)",
        "  ❌ Modelo precisa escolher o mode correto",
        "  ❌ Schema mais complexo",
        "",
        "2 tools separadas:",
        "  search_product() + search_category() + search_user()",
        "  ✅ Cada tool tem schema simples e focado",
        "  ✅ Modelo escolhe naturalmente pela descrição",
        "  ❌ Mais tools = mais tokens de descrição",
        "",
        "Regra: se as tools compartilham 80% do schema → 1 tool com mode",
        "        se são conceitualmente diferentes → tools separadas",
    ],
    """📖 EXPLICAÇÃO: A regra de 80% é um guia prático. Se search_product e search_user compartilham quase todo o schema (query, limit, offset), vale consolidar em uma com mode. Se são conceitualmente diferentes (buscar produto vs buscar usuário), separe.
➡️ TRANSIÇÃO: Ferramentas perigosas e HITL.""",
    32, TOTAL, "1 tool com mode vs 2 separadas"
)

# Slide 33 — [SEÇÃO] HITL
s = section_slide(4, "Tools Perigosas e HITL")
add_notes(s, "Transição para segurança e HITL.")

# Slide 34 — Matriz de Risco
s = content_slide(
    "Matriz de Risco: Irreversível × Impactante",
    [
        "          | Baixo impacto    | Alto impacto",
        "----------|------------------|------------------",
        "Reversível| search           | create_order",
        "          | get_status       | update_profile",
        "          | → Sem HITL       | → HITL opcional",
        "----------|------------------|------------------",
        "Irrevers. | delete_cache     | delete_account",
        "          | reset_config     | deploy, transfer",
        "          | → HITL recomend. | → HITL OBRIGAT.",
        "",
        "Quadrante crítico: irreversível + alto impacto = HITL sempre",
        "  Exemplos: delete, deploy, transfer, email em massa",
    ],
    """📖 EXPLICAÇÃO: A matriz classifica tools por dois eixos: reversibilidade (posso desfazer?) e impacto (quanto dano causa?). O quadrante crítico é irreversível + alto impacto — HITL é não-negociável. Exemplos: deletar conta de usuário, fazer deploy, transferir dinheiro, enviar email em massa.
➡️ TRANSIÇÃO: Como implementar HITL.""",
    34, TOTAL, "Matriz de risco"
)

# Slide 35 — HITL
s = content_slide(
    "HITL: Confirmação, Dry-run, Simulação",
    [
        "1. CONFIRMAÇÃO EXPLÍCITA",
        "   • Antes de executar: 'Confirmar envio para X? [y/n]'",
        "   • Agente pausa, humano aprova, executa",
        "",
        "2. DRY-RUN",
        "   • Parâmetro dry_run=true: simula sem executar",
        "   • Retorna o que FARIA, não faz",
        "   • Exemplo: 'enviar_email(dry_run=true) → preview do email'",
        "",
        "3. SIMULAÇÃO / SANDBOX",
        "   • Executa em ambiente isolado (container, DB de teste)",
        "   • Resultado é real mas não afeta produção",
        "",
        "4. ALLOWLIST",
        "   • Só executa actions pré-aprovadas",
        "   • Dinâmico: modelo propõe, humano aprova uma vez",
    ],
    """📖 EXPLICAÇÃO: HITL tem níveis. Confirmação é o mais simples (y/n). Dry-run é mais rico (mostra o que faria). Simulação é o mais completo (executa de verdade em sandbox). Allowlist é para produção (ações pré-aprovadas).
➡️ TRANSIÇÃO: Erros comuns.""",
    35, TOTAL, "Implementação de HITL"
)

# Slide 36 — Exercício HITL
s = exercise_slide(
    "Exercício: Classifique as Tools",
    [
        "Para cada tool, classifique:",
        "  (a) Tipo: leitura / escrita / destrutiva / external",
        "  (b) HITL: nenhum / opcional / obrigatório",
        "  (c) Idempotente: sim / não",
        "",
        "1. search_orders(query) — busca pedidos no DB",
        "2. update_order_status(id, status) — muda status",
        "3. delete_user(id) — remove conta",
        "4. send_invoice(email, pdf) — envia fatura por email",
        "5. refund_payment(id, amount) — estorna pagamento",
    ],
    """📖 EXPLICAÇÃO: Respostas: (1) Leitura, sem HITL, idempotente. (2) Escrita, HITL opcional, idempotente (mesmo status = no-op). (3) Destrutiva, HITL obrigatório, NÃO idempotente. (4) External, HITL obrigatório, NÃO idempotente (sem request_id). (5) External+Destrutiva, HITL obrigatório, precisa idempotência (request_id para não estornar 2x).
➡️ TRANSIÇÃO: Erros comuns e correções.""",
    36, TOTAL
)

# ═══════════════════════════════════════════════════
# SEÇÃO E — Erros Comuns e Avaliação (Slides 37-44)
# ═══════════════════════════════════════════════════

s = section_slide(5, "Erros Comuns e Avaliação")
add_notes(s, "Transição para erros comuns e avaliação de tools.")

# Slide 37 — Os 5 Erros
s = content_slide(
    "Os 5 Erros Mais Comuns em Design de Tools",
    [
        "1. Paths relativos → absolutos (caso Anthropic SWE-bench)",
        "   Sintoma: agente erra caminho após mudar de diretório",
        "",
        "2. Muitas tools similares → consolidar ou renomear",
        "   Sintoma: modelo confunde search_product com search_items",
        "",
        "3. Descrições vagas → reescrever como docstring de júnior",
        "   Sintoma: modelo chama tool errada ou não chama",
        "",
        "4. Schema frouxo → apertar com enum/patterns",
        "   Sintoma: modelo passa valores inválidos",
        "",
        "5. Falta de erro tratado → propagar com mensagem útil",
        "   Sintoma: exceção quebra o loop do agente",
    ],
    """📖 EXPLICAÇÃO: Cada erro tem um sintoma característico. Aprender a diagnosticar pelo sintoma é a habilidade mais valiosa desta aula. Se o modelo erra caminho → paths. Se confunde tools → consolidar. Se passa valores inválidos → schema. Se quebra o loop → error handling.
❓ PERGUNTA: "Qual desses erros vocês já cometeram?"
➡️ TRANSIÇÃO: Vamos aprofundar cada um.""",
    37, TOTAL, "Os 5 erros mais comuns"
)

# Slide 38 — Erro detalhado 1
s = comparison_slide(
    "Erro 1: Paths Relativos → Absolutos",
    "❌ Antes",
    [
        'name: "edit_file"',
        'description: "Edita um arquivo"',
        'parameters:',
        '  path: string',
        '    (descrição: caminho do arquivo)',
        '',
        "Problema:",
        "• 'path' é ambíguo — relativo ou absoluto?",
        "• Modelo perde track do CWD",
        "• Após navegar, usa path relativo errado",
    ],
    "✅ Depois",
    [
        'name: "edit_file"',
        'description: "Edita arquivo."',
        'parameters:',
        '  path: string',
        '    description: "Path ABSOLUTO."',
        '    pattern: "^/.*"',
        '    (ex: "/src/main.py")',
        '',
        "Solução:",
        "• Descrição diz ABSOLUTO",
        "• Regex pattern força / inicial",
        "• Exemplo na descrição",
        "• Poka-yoke: impossível usar relativo",
    ],
    """📖 EXPLICAÇÃO: O pattern "^/.*" é um poka-yoke regex — rejeita qualquer path que não comece com /. O modelo nem consegue gerar um path relativo válido. Isso é design defensivo aplicado.
➡️ TRANSIÇÃO: Erro 2.""",
    38, TOTAL, "Erro 1 — paths relativos"
)

# Slide 39 — Erro 2 e 3
s = content_slide(
    "Erros 2 e 3: Tools Similares e Descrições Vagas",
    [
        "ERRO 2: MUITAS TOOLS SIMILARES",
        "  Sintoma: search_product vs search_item vs search_products",
        "  Solução: consolidar em search(query, type='product'|'item')",
        "  Ou: renomear para evitar confusão (get_product vs list_products)",
        "",
        "ERRO 3: DESCRIÇÕES VAGAS",
        "  ❌ 'Busca dados' — busca o quê? onde? retorna o quê?",
        "  ✅ 'Busca produto por nome no catálogo. Retorna preço e estoque.'",
        "",
        "Checklist da boa descrição:",
        "  [ ] Diz o que FAZ (não o que É)",
        "  [ ] Diz quando USAR",
        "  [ ] Diz quando NÃO USAR (fronteira)",
        "  [ ] Tem exemplo de parâmetro",
        "  [ ] Diz o que RETORNA",
    ],
    """📖 EXPLICAÇÃO: Erro 2 é sobre namespace pollution. Se duas tools têm nomes similares, o modelo confunde. Solução: consolidar ou renomear. Erro 3 é sobre clareza. A boa descrição responde 5 perguntas: o que faz, quando usar, quando não usar, exemplo, retorno.
➡️ TRANSIÇÃO: Avaliação de tools.""",
    39, TOTAL, "Erros 2 e 3"
)

# Slide 40 — [SEÇÃO] Avaliando Tools
s = section_slide(6, "Avaliando Tools")
add_notes(s, "Transição para avaliação empírica de tools.")

# Slide 41 — Workbench
s = content_slide(
    "Workbench: Rodar, Observar, Iterar",
    [
        "Workbench = ferramenta para testar tools com N inputs",
        "",
        "Fluxo:",
        "  1. Definir conjunto de casos de teste (20+ inputs)",
        "  2. Rodar o agente com cada input",
        "  3. Observar: tool chamada corretamente? Args corretos?",
        "  4. Medir: taxa de uso correto, custo, latência",
        "  5. Identificar padrões de erro",
        "  6. Refatorar tools (aplicar ACI)",
        "  7. Re-rodar e comparar (antes vs depois)",
        "",
        "Não é opcional — é como você sabe que suas tools funcionam",
    ],
    """📖 EXPLICAÇÃO: Workbench é o teste de regressão para tools. Sem ele, você 'acha' que as tools funcionam. Com workbench, você PROVA que funcionam — e mede a melhoria após refatoração. O Lab 1 replica exatamente este fluxo.
➡️ TRANSIÇÃO: Métricas.""",
    41, TOTAL, "Workbench para avaliar tools"
)

# Slide 42 — Métricas
s = content_slide(
    "Métricas: Taxa de Uso, Custo, Latência",
    [
        "1. TAXA DE USO CORRETO",
        "   • % de inputs onde a tool correta foi chamada com args corretos",
        "   • Meta: ≥85% (referência do projeto do módulo)",
        "",
        "2. CUSTO POR CHAMADA",
        "   • Tokens consumidos pela descrição + args + resultado",
        "   • Comparar antes vs depois da refatoração",
        "",
        "3. LATÊNCIA",
        "   • Tempo médio de execução da tool",
        "   • p50, p95, p99",
        "",
        "4. ERROS TRATADOS",
        "   • % de chamadas que retornam erro estruturado vs crash",
        "   • Meta: 100% de erros tratados (0 crashes)",
    ],
    """📖 EXPLICAÇÃO: As 4 métricas formam o dashboard de quality das tools. Taxa de uso correto é a principal. As outras três são de suporte. Se a taxa é <85%, refatore. Se custo subiu após refatoração (descrições mais longas), avalie se o benefício justifica.
➡️ TRANSIÇÃO: Fechamento.""",
    42, TOTAL, "Métricas de avaliação"
)

# Slide 43 — Erros 4 e 5
s = content_slide(
    "Erros 4 e 5: Schema Frouxo e Erro Não Tratado",
    [
        "ERRO 4: SCHEMA FROUXO",
        "  ❌ op: string (modelo pode passar 'multiply', 'soma', '×')",
        "  ✅ op: enum ['+', '-', '*', '/'] (só 4 valores possíveis)",
        "  ❌ email: string (qualquer texto)",
        "  ✅ email: string with pattern regex",
        "",
        "ERRO 5: ERRO NÃO TRATADO",
        "  ❌ Exceção propaga → loop do agente quebra",
        "  ✅ try/except → retorna {error: '...', suggestion: '...'}",
        "  O modelo consome o erro e decide o que fazer",
        "",
        "Regra: toda tool tem try/except. SEMPRE.",
    ],
    """📖 EXPLICAÇÃO: Schema frouxo permite valores inválidos. Schema apertado (enum, pattern, min/max) rejeita antes de executar. Erro não tratado quebra o loop. Erro tratado retorna mensagem útil que o modelo pode usar para se recuperar.
➡️ TRANSIÇÃO: Boas práticas e anti-patterns.""",
    43, TOTAL, "Erros 4 e 5"
)

# Slide 44 — Testes de Regressão
s = content_slide(
    "Testes de Regressão para Tools",
    [
        "Conjunto de testes que roda a cada mudança de tool",
        "",
        "Estrutura:",
        "  • 20+ casos de teste cobrindo: happy path, edge cases, erros",
        "  • Cada caso: input + tool esperada + args esperados",
        "  • Rodar antes e depois de qualquer mudança",
        "",
        "Exemplo de caso:",
        "  input: 'Quanto é 10 / 0?'",
        "  tool esperada: calculator(a=10, b=0, op='/')",
        "  resultado esperado: {error: 'division by zero'}",
        "  (não deve crashar — deve retornar erro estruturado)",
        "",
        "CI/CD: rodar testes de regressão a cada commit",
    ],
    """📖 EXPLICAÇÃO: Testes de regressão garantem que uma refatoração não quebra o que funcionava. Sem eles, você melhora a tool para um caso e quebra para outro. O conjunto de 20+ casos cobre happy path (funcionamento normal), edge cases (valores extremos), e erros (entradas inválidas).
➡️ TRANSIÇÃO: Fechamento da aula.""",
    44, TOTAL, "Testes de regressão para tools"
)

# ═══════════════════════════════════════════════════
# SEÇÃO F — Fechamento (Slides 45-57)
# ═══════════════════════════════════════════════════

s = section_slide(7, "Boas Práticas e Fechamento")
add_notes(s, "Transição para o fechamento.")

# Slide 45 — Boas Práticas
s = content_slide(
    "Boas Práticas (DO)",
    [
        "✅ Descrição rica = docstring para júnior (exemplos, edge cases)",
        "✅ Poka-yoke: enum, pattern, required, default",
        "✅ Paths absolutos (não relativos)",
        "✅ Formato próximo ao natural (markdown > JSON com escaping)",
        "✅ Idempotência com request_id para ações destrutivas",
        "✅ try/except em TODA tool — sempre retornar erro estruturado",
        "✅ HITL para tools destrutivas e external side-effects",
        "✅ dry_run como parâmetro default true",
        "✅ Versionar tools sem quebrar (adicionar, não remover)",
        "✅ Workbench com 20+ casos antes de produção",
        "✅ max_tokens generoso (dar espaço para pensar)",
    ],
    """📖 EXPLICAÇÃO: Cada item é uma lição que alguém aprendu da forma difícil. A maioria vem direto da Anthropic. O resto vem de experiência de produção.
➡️ TRANSIÇÃO: E o que NÃO fazer.""",
    45, TOTAL, "Boas práticas DO"
)

# Slide 46 — Anti-patterns
s = content_slide(
    "Anti-Patterns (DON'T)",
    [
        "❌ Descrição de 3 palavras ('Faz busca')",
        "❌ Schema sem enum/pattern (string livre para tudo)",
        "❌ Paths relativos sem documentar",
        "❌ Sem try/except (exceção quebra o loop)",
        "❌ Sem max_tokens (modelo se encurrala)",
        "❌ Sem timeout (API lenta trava o agente)",
        "❌ 50+ tools no namespace (performance cai)",
        "❌ Tools com nomes similares (search_product vs search_item)",
        "❌ Sem HITL para ações destrutivas (delete, deploy, email)",
        "❌ Não testar com workbench ('acho que funciona')",
        "❌ Mudar schema de tool em produção sem versionar",
    ],
    """📖 EXPLICAÇÃO: Cada anti-pattern aqui é um bug que eu já vi em produção. 'Faz busca' é o mais comum. Sem try/except é o mais perigoso. Sem HITL para destrutivas é o mais caro.
➡️ TRANSIÇÃO: Caso de estudo.""",
    46, TOTAL, "Anti-patterns DON'T"
)

# Slide 47 — Caso de Estudo Consolidado
s = content_slide(
    "Caso Consolidado: Anthropic Tool Redesign",
    [
        "O que aconteceu:",
        "  • Coding agent para SWE-bench",
        "  • Agente saía do diretório raiz → paths relativos errados",
        "  • Edita arquivo errado ou falha 'não encontrado'",
        "",
        "O que NÃO funcionou:",
        "  • Ajustar prompt ('lembre-se de usar paths absolutos')",
        "  → Modelo esquecia em conversas longas",
        "",
        "O que FUNCIONOU:",
        "  • Mudar a tool para exigir paths absolutos (pattern: ^/.*)",
        "  → Erro estruturalmente impossível",
        "  → 'O modelo usou a tool perfeitamente' (Schluntz)",
        "",
        "Lição: erros sistemáticos = design de tool, não de prompt",
    ],
    """📖 EXPLICAÇÃO: Este caso consolida toda a aula. A lição é: quando o erro é sistemático (sempre da mesma forma), a causa é o design da tool, não o prompt. A solução é redesenhar a tool com poka-yoke.
➡️ TRANSIÇÃO: Exercício final.""",
    47, TOTAL, "Caso consolidado"
)

# Slide 48 — Exercício Final
s = exercise_slide(
    "Exercício: Anti-patterns em Schema",
    [
        "Em trios (4 min):",
        "",
        "Dado este schema, identifique 3 anti-patterns e reescreva:",
        "",
        '  name: "db"',
        '  description: "Faz operação no banco"',
        '  parameters: {',
        '    table: string,',
        '    action: string,',
        '    data: string',
        '  }',
        "",
        "Dica: aplique enum, pattern, descrições ricas, poka-yoke",
    ],
    """📖 EXPLICAÇÃO: Anti-patterns: (1) nome 'db' é vago; (2) descrição não diz qual operação, qual banco, o que retorna; (3) action é string livre (deveria ser enum); (4) data é string (deveria ser object); (5) sem required fields; (6) sem HITL para ações destrutivas; (7) table sem validação. Refatoração: rename para execute_db_query, enum em action (select/insert/update/delete), data como object, HITL para delete/drop.
➡️ TRANSIÇÃO: Resumo.""",
    48, TOTAL
)

# Slide 49 — Resumo
s = content_slide(
    "Resumo da Aula",
    [
        "1. Tool calling = JSON Schema como contrato (não texto livre)",
        "2. ACI :: HCI — invista tanto quanto investiria em UI",
        "3. 5 princípios ACI: empatia, descrição rica, formato natural, poka-yoke, tokens para pensar",
        "4. Engenharia: schemas claros, idempotência, timeouts, tipologia",
        "5. HITL obrigatório para destrutivas e external",
        "6. 5 erros comuns: paths, tools similares, descrições vagas, schema frouxo, erro não tratado",
        "7. Workbench: rodar 20+ casos, medir, iterar",
        "8. Lição SWE-bench: tempo em tools > tempo no prompt",
    ],
    """📖 EXPLICAÇÃO: Se vocês só lembrarem de 3 coisas: (1) ACI é disciplina de design; (2) poka-yoke > prompt; (3) teste com workbench antes de confiar.
➡️ TRANSIÇÃO: Quiz.""",
    49, TOTAL, "Resumo dos 8 pontos"
)

# Slide 50 — Quiz 1
s = content_slide(
    "Quiz 1/5",
    [
        "O que significa ACI?",
        "",
        "A) Agent Communication Interface",
        "B) Agent-Computer Interface (design de tools)",
        "C) Automated Code Injection",
        "D) Agent Control Infrastructure",
        "",
        "Resposta: B",
    ],
    """📖 EXPLICAÇÃO: ACI = Agent-Computer Interface, análogo a HCI. É o design das tools.
➡️ TRANSIÇÃO: Próxima.""",
    50, TOTAL, "Quiz 1"
)

# Slide 51 — Quiz 2
s = content_slide(
    "Quiz 2/5",
    [
        "Qual é o poka-yoke aplicado no caso do SWE-bench?",
        "",
        "A) Aumentar o max_steps",
        "B) Adicionar mais exemplos no prompt",
        "C) Mudar a tool para exigir paths absolutos",
        "D) Usar um modelo maior",
        "",
        "Resposta: C",
    ],
    """📖 EXPLICAÇÃO: A mudança da tool para exigir paths absolutos tornou o erro estruturalmente impossível. Poka-yoke.
➡️ TRANSIÇÃO: Próxima.""",
    51, TOTAL, "Quiz 2"
)

# Slide 52 — Quiz 3
s = content_slide(
    "Quiz 3/5",
    [
        "Quando HITL é OBRIGATÓRIO?",
        "",
        "A) Para toda tool de leitura",
        "B) Apenas para tools com mais de 3 parâmetros",
        "C) Para tools destrutivas e external side-effects",
        "D) HITL é sempre opcional",
        "",
        "Resposta: C",
    ],
    """📖 EXPLICAÇÃO: HITL é obrigatório para ações irreversíveis (delete, deploy, transfer, email em massa).
➡️ TRANSIÇÃO: Próxima.""",
    52, TOTAL, "Quiz 3"
)

# Slide 53 — Quiz 4
s = content_slide(
    "Quiz 4/5",
    [
        "O que fazer quando uma tool retorna erro?",
        "",
        "A) Deixar a exceção propagar e quebrar o loop",
        "B) Retornar erro estruturado {error: '...', suggestion: '...'}",
        "C) Silenciosamente ignorar e continuar",
        "D) Reiniciar o agente",
        "",
        "Resposta: B",
    ],
    """📖 EXPLICAÇÃO: try/except em toda tool. Retornar erro estruturado permite que o modelo decida o que fazer (tentar de novo, usar outra tool, ou informar o usuário).
➡️ TRANSIÇÃO: Próxima.""",
    53, TOTAL, "Quiz 4"
)

# Slide 54 — Quiz 5
s = content_slide(
    "Quiz 5/5",
    [
        "Qual é a regra de ouro da Anthropic sobre tools?",
        "",
        "A) Use sempre o maior modelo disponível",
        "B) Invista em ACI tanto quanto em HCI",
        "C) Quanto mais tools, melhor",
        "D) Descrições curtas economizam tokens",
        "",
        "Resposta: B",
    ],
    """📖 EXPLICAÇÃO: A Anthropic é explícita: invista tanto esforço em ACI quanto em HCI. Tempo em tools > tempo no prompt.
➡️ TRANSIÇÃO: Conexão.""",
    54, TOTAL, "Quiz 5"
)

# Slide 55 — Conexão
s = content_slide(
    "Conexão com Próximos Módulos",
    [
        "ETHAGT03 — Padrões de Workflow: como tools se compõem",
        "ETHAGT04 — Reasoning: o Thought antes da Action",
        "ETHAGT05 — Memória: tools de persistência",
        "ETHAGT08 — MCP: tool use como padrão de servidor (escala)",
        "ETHAGT12 — AgentOps: avaliando tools em produção",
        "ETHAGT13 — Segurança: tools como vetor de ataque (prompt injection via tool)",
    ],
    """📖 EXPLICAÇÃO: ETHAGT08 (MCP) é a continuação natural — escala tool use para centenas de tools. ETHAGT13 (Security) mostra como tools podem ser vetores de ataque.
➡️ TRANSIÇÃO: Leitura.""",
    55, TOTAL, "Conexão com próximos módulos"
)

# Slide 56 — Leitura e Referências
s = content_slide(
    "Leitura Recomendada e Referências",
    [
        "Obrigatório (antes de ETHAGT03):",
        "  • Anthropic, Building Effective Agents — Appendix 2 (15 min)",
        "  • OpenAI Function Calling Guide (10 min)",
        "",
        "Recomendado:",
        "  • Gorilla paper (arXiv:2305.15334) — escala de tools",
        "  • ToolLLM paper (arXiv:2307.16789) — benchmark",
        "  • Pydantic AI docs — padrões de schema",
        "",
        "Referências completas: 20-Research/ETHAGT02-pesquisa.md",
    ],
    """📖 EXPLICAÇÃO: O Apêndice 2 da Anthropic é a leitura mais importante. Se já leram (homework), parabéns. Se não, ler antes de ETHAGT03.
➡️ TRANSIÇÃO: Projeto e próximos passos.""",
    56, TOTAL, "Leitura recomendada"
)

# Slide 57 — Q&A / Encerramento
s = title_slide(
    "Perguntas?",
    "ETHAGT02 — Tool Calling e Agent-Computer Interface\nPróxima aula: ETHAGT03 — Padrões de Workflow",
    "Universidade Etho"
)
add_notes(s, """📖 EXPLICAÇÃO: Abrir para Q&A. Se não houver perguntas, fazer: "Qual parte da aula foi menos clara?"
Mensagem final: "ACI é onde mora a confiabilidade. Invistam em tools tanto quanto investiriam em UI. Poka-yoke > prompt. Testem com workbench."
➡️ TRANSIÇÃO: Fim da aula.""")

# ─── Salvar ───
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT02-Apresentacao.pptx")
prs.save(output_path)
print(f"PPTX gerado: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
