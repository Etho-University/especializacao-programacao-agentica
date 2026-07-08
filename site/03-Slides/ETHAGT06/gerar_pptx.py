#!/usr/bin/env python3
"""Gerador de PPTX — ETHAGT06: RAG Agêntico. Universidade Etho · Especialização em Programação Agêntica."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj=""):
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT06"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

T = 85

s = title_slide(
    "RAG Agêntico",
    "Universidade Etho · Especialização em Programação Agêntica\nFase B — RAG Avançado · 30 h",
    "ETHAGT06"
)
add_notes(s, "📖 Bem-vindos. Vamos evoluir do RAG ingênuo — que funciona na demo e quebra em produção — para o RAG agêntico. 4 arquiteturas: Adaptive, CRAG, Self-RAG, Agentic.\n💡 Analogia: estagiário que sempre busca no mesmo arquivo vs pesquisador sênior que decide quando consultar, avalia a fonte, reflete sobre a resposta.\n❓ 'Quantos já usaram RAG em produção?'\n⚠️ Alunos querem pular para Agentic. Redirecionar: comece com Adaptive.\n➡️ Vamos aos objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: evoluir do RAG 'ingênuo' para RAG agêntico, onde o agente decide quando/o quê/como recuperar",
    "",
    "Objetivos específicos:",
    "1. Diagnosticar limites do RAG ingênuo em produção",
    "2. Implementar Adaptive RAG, CRAG, Self-RAG e Agentic RAG",
    "3. Aplicar técnicas de qualidade: chunking, re-ranking, query rewriting, hybrid search",
    "4. Construir pipeline de avaliação de RAG (faithfulness, relevance, context recall/precision)",
    "5. Produzir um sistema RAG multi-tenant com segurança"
], "📖 Cada objetivo é mensurável: diagnosticar, implementar, aplicar, construir, produzir.\n💡 Analogia: checklist de pré-voo.\n❓ 'Qual objetivo é mais desafiador?' (costuma ser #2 ou #4)\n⚠️ 'Fazer RAG agêntico' ≠ 'usar LangGraph'. Arquitetura > framework.\n➡️ Vamos às competências.", 2, T)

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → A (Avançado)",
    "C3 MCP & Tool Use → B (Básico)",
    "C4 Agent Memory → I (Intermediário)",
    "C5 AgentOps & Avaliação → I (Intermediário)",
    "",
    "C1 atinge Avançado: projetar e justificar arquiteturas RAG completas",
    "C5 atinge Intermediário: pipeline de avaliação automatizado (Ragas)",
    "C4 atinge Intermediário: estado entre hops (memória de trabalho)"
], "📖 C1 Avançado = escolher arquitetura certa e justificar trade-offs.\n💡 Analogia: ETHAGT01 foi sair do estacionamento; ETHAGT04 foi a cidade; hoje é rodovia com trânsito intenso.\n⚠️ 'Avançado' ≠ 'saber LangGraph avançado'. Significa justificar com trade-offs.\n➡️ Vamos à agenda.", 3, T)

s = content_slide("Agenda da Aula", [
    "Bloco 1 (~58 min):",
    "  Abertura (8 min) — objetivos, motivação, contexto",
    "  RAG Ingênuo Falha (10 min) — 4 tipos de falha",
    "  Adaptive RAG (12 min) — quando recuperar, routing",
    "  Corrective RAG (15 min) — avaliador, 3 caminhos, web fallback",
    "  Self-RAG (13 min) — tokens de reflexão",
    "",
    "Bloco 2 (~52 min):",
    "  Agentic RAG (15 min) — multi-hop, GraphRAG, DEMO",
    "  Engenharia de Qualidade (15 min) — chunking, re-rank, HyDE",
    "  Avaliação de RAG (10 min) — métricas, Ragas",
    "  Fechamento (12 min) — boas práticas, quiz, projeto, Q&A"
], "📖 Bloco 1 é a escalada (Adaptive → CRAG → Self-RAG). Bloco 2 é consolidação.\n💡 Analogia: escalada. Bloco 1 = subida; Bloco 2 = cume e descida.\n⚠️ Avisar: Seção B (RAG Ingênuo Falha) motiva tudo.\n➡️ Vamos conectar com ETHAGT04.", 4, T)

s = content_slide("Pré-requisitos: ETHAGT04 e o Que Você Já Sabe", [
    "ETHAGT04 cobriu: RAG básico, vector DBs, embeddings, pipeline retrieve→generate",
    "Você já deveria saber: cosine similarity, top-k retrieval, LangChain básico",
    "",
    "O que esta aula adiciona: o agente entra no loop de recuperação",
    "Você NÃO precisa ter feito ETHAGT04 se já pratica RAG em produção"
], "📖 ETHAGT04 = pipeline fixo. Hoje: agente decide se/o quê/como recuperar.\n💡 Analogia: ETHAGT04 = cozinhar com receita; hoje = cozinhar sem receita, adaptando.\n⚠️ Sem ETHAGT04? Recap rápido de embeddings e top-k.\n➡️ Mas por que estamos falando disso?", 5, T)

comparison_slide(
    "RAG 'Funciona' Até Parar de Funcionar",
    "Na DEMA (verde)", [
        "Usuário feliz",
        "Resposta parece correta",
        "Confiança alta"
    ],
    "Em PRODUÇÃO (vermelho)", [
        "Recupera chunks errados",
        "Responde com info genérica",
        "Confiança alta, grounding baixo",
        "Usuário age sobre info falsa"
    ],
    "📖 RAG ingênuo falha SILENCIOSAMENTE: sem erro 500, sem crash. Responde com confiança alta mesmo errado.\n💡 Analogia: médico que sempre dá diagnóstico confiante, mesmo sem saber.\n❓ 'Quantos já usaram RAG e obtiveram resposta errada com alta confiança?'\n⚠️ 'Prompt melhor' não resolve — problema é sistêmico.\n➡️ Como evoluímos?", 6, T
)

s = content_slide("Do RAG Fixo ao RAG In-Loop", [
    "Linha do tempo:",
    "  2020: RAG original (Lewis et al., arXiv:2005.11401)",
    "  2023: Self-RAG (Asai et al., arXiv:2310.11511)",
    "  2024: CRAG (Yan et al., arXiv:2401.15884)",
    "  2024: GraphRAG (Edge et al., arXiv:2404.16130)",
    "  2025: Agentic RAG",
    "",
    "RAG fixo: query → retrieve → generate (rígido)",
    "RAG agêntico: agente decide se/quando/o quê/como, com correção iterativa",
    "4 arquiteturas: Adaptive → CRAG → Self-RAG → Agentic"
], "📖 A evolução: cada arquitetura adiciona um checkpoint de decisão.\n💡 Analogia: piloto automático. Fixo = segue rota; Adaptive = decide se desvia; CRAG = avalia rota; Self-RAG = reflete se voa bem; Agentic = pilota tudo.\n❓ 'Qual desses papers vocês já leram?'\n⚠️ Confundem Self-RAG com CRAG. Self-RAG = modelo reflete; CRAG = pipeline avalia.\n➡️ Vamos diagnosticar as falhas.", 7, T)

s = section_slide("1", "Por que o RAG Ingênuo Falha")
add_notes(s, "📖 Transição: vamos desconstruir o pipeline ingênuo e identificar 4 classes de falha.\n➡️ Primeiro: o pipeline e seus pontos de falha.")

s = content_slide("O Pipeline do RAG Ingênuo e Seus Pontos de Falha", [
    "Pipeline: query → embed → vector search (top-k) → stuff into prompt → generate",
    "",
    "Pontos de falha:",
    "  Embed: sinônimos não capturados, multilingual distante",
    "  Search: top-k traz irrelevante junto com relevante",
    "  Stuff: contexto quebrado, janela estourada, lost in the middle",
    "  Generate: alucinação sobre docs ruins",
    "",
    "Fonte: arXiv:2005.11401 (RAG original, Lewis et al.)"
], "📖 Pipeline sem NENHUMA reflexão sobre qualidade. Cada etapa confia na anterior.\n💡 Analogia: linha de montagem sem controle de qualidade.\n❓ 'Em qual etapa vocês tiveram mais problemas?' (search ou generate)\n⚠️ Focar só no generate é erro — retrieval é o gargalo mais subestimado.\n➡️ Vamos detalhar as 4 falhas.", 9, T, obj="Diagrama D2")

s = content_slide("Falha 1: Chunking Quebra Contexto", [
    "Chunking por tamanho fixo (512 tokens) corta no meio de uma ideia",
    "Exemplo: política de férias dividida entre dois chunks — nenhum faz sentido sozinho",
    "Sobreposição (overlap) ajuda mas não resolve",
    "Consequência: embedding captura fragmento sem contexto → recuperação irrelevante",
    "",
    "Solução (preview Seção G): chunking semântico, hierárquico, late-chunking"
], "📖 O chunk é o átomo do RAG. Átomo ruim = molécula ruim.\n💡 Analogia: cortar livro a cada 500 palavras sem olhar capítulos.\n❓ 'Qual tamanho de chunk vocês usam?'\n⚠️ Overlap de 50 tokens resolve pouco se o conceito tem 800 tokens.\n➡️ Mesmo com chunking bom, embedding tem limites.", 10, T)

s = content_slide("Falha 2: Embedding Não Captura Semântica", [
    "Embedding captura similaridade lexical, não raciocínio",
    "Exemplo: 'férias de estagiário' vs 'licença de aprendiz' — equivalentes, lexicalmente diferentes",
    "Multilingual: mesma ideia em PT vs EN → embeddings distantes",
    "Dados tabulares: embeddings de texto não capturam estrutura de tabela",
    "",
    "Solução (preview Seção G): query rewriting, HyDE, hybrid search"
], "📖 Embeddings têm cegueira semântica: aprendem co-ocorrência, não significado.\n💡 Analogia: duas pessoas descrevendo a mesma coisa com palavras diferentes.\n❓ 'Já tiveram problema de sinônimos no RAG?'\n⚠️ 'Embedding melhor' resolve pouco — query rewriting e hybrid são mais impactantes.\n➡️ Mesmo com embedding bom, top-k traz lixo.", 11, T)

s = content_slide("Falha 3: Sem Re-Rank, Top-k Traz Lixo", [
    "Top-k retorna os k mais similares — nem todos são relevantes",
    "Sem re-rank: o chunk #1 pode ser irrelevante e 'afogar' os bons",
    "LLM recebe contexto poluído → alucinação ou confusão",
    "Exemplo: top-5 traz 1 chunk bom e 4 irrelevantes",
    "",
    "Solução (preview Seção G): re-ranking com Cohere, bge, Jina"
], "📖 'Mais próximo' ≠ 'mais relevante'. Top-k é aproximação.\n💡 Analogia: pedir 5 livros por título parecido e receber 1 bom + 4 inúteis.\n❓ 'Vocês usam re-rank hoje?' (provavelmente minoria)\n⚠️ Aumentar k PIORA sem re-rank. Mais contexto ≠ melhor.\n➡️ Mas a pior falha é não saber que está falhando.", 12, T)

s = content_slide("Falha 4: Sem Avaliação, Você Está Cego", [
    "'Funciona na minha máquina' ≠ funciona em produção",
    "Sem métricas: melhoria é adivinhação",
    "Métricas que importam: faithfulness, answer relevance, context precision/recall",
    "Sem dataset de holdout: não há como detectar regressão",
    "",
    "Solução (preview Seção H): Ragas, TruLens, DeepEval"
], "📖 Sem métricas, não sabe se melhorou ou piorou. Sem holdout, sem regressão.\n💡 Analogia: melhorar receita sem provar o resultado.\n❓ 'Vocês têm dataset de holdout?' (provavelmente <30%)\n⚠️ 'A resposta parece certa' é subjetivo. Precisamos de métricas automáticas.\n➡️ Essas 4 falhas alimentam um mito.", 13, T)

comparison_slide(
    "O Mito 'Vector DB Resolve Tudo'",
    "O MITO", [
        "'Só joga no vector DB e funciona'",
        "Vector DB é apresentado como solução universal"
    ],
    "A REALIDADE", [
        "Vector DB é uma PEÇA, não a solução",
        "Dados tabulares: embeddings não capturam estrutura → SQL",
        "Multilingual: PT vs EN → distância vetorial alta",
        "Multi-modal: texto + imagem + tabela → estratégias diferentes"
    ],
    "📖 'Vector DB resolve tudo' é o mito mais caro da indústria.\n💡 Analogia: 'Excel resolve tudo'. Ótimo para planilhas, péssimo para edição de vídeo.\n❓ 'Pior falha: não responder ou responder errado com confiança?' (responder errado — usuário age sobre info falsa)\n⚠️ Embedar tabelas como texto funciona mal. SQL retrieval é melhor.\n➡️ Agora que sabemos onde falha, vamos construir: Adaptive RAG.", 14, T
)

s = section_slide("2", "Adaptive RAG: Decidir Quando Recuperar")
add_notes(s, "📖 Primeira evolução: em vez de SEMPRE recuperar, decide SE recuperar.\n➡️ A ideia central.")

comparison_slide(
    "Decidir Quando Recuperar",
    "RAG Ingênuo (sempre recupera)", [
        "Recupera MESMO para 'Quem é o presidente?'",
        "Latência e custo desperdiçados",
        "Risco: docs irrelevantes confundem o modelo"
    ],
    "Adaptive RAG (decide)", [
        "Classificador decide se precisa recuperar",
        "Não precisa → responder direto (LLM puro)",
        "Pergunta simples → recuperar top-3",
        "Pergunta complexa → recuperar + query rewrite"
    ],
    "📖 MUITAS perguntas não precisam de retrieval. 'Quem é o presidente?' o GPT-4 sabe.\n💡 Analogia: médico triagista. Nem todo paciente precisa de exames.\n❓ 'Pensem numa pergunta que o modelo sabe sem retrieval.'\n⚠️ 'Sempre recuperar' NÃO é mais seguro — recuperação irrelevante é ruído.\n➡️ Mas quando decidimos recuperar, quanto?", 16, T
)

s = content_slide("Decidir Quanto Recuperar", [
    "Top-k fixo (sempre 5) é subótimo",
    "Pergunta factual: 1-2 chunks bastam",
    "Pergunta analítica: 5-10 chunks podem ser necessários",
    "Estratégia: classificar complexidade → ajustar k dinamicamente",
    "",
    "Trade-off: mais chunks = mais contexto = mais custo + mais ruído"
], "📖 Top-k fixo é preguiça arquitetural. k deve ser dinâmico.\n💡 Analogia: pedir referências. Pergunta simples = 1 fonte; tese = 50.\n⚠️ Aumentar k sem re-rank: top-20 com 3 bons e 17 ruins. Mais contexto ≠ melhor.\n➡️ Como classificamos? Estratégias de routing.", 17, T)

s = content_slide("Estratégias de Routing por Complexidade", [
    "Routing por regras: keywords, comprimento da pergunta",
    "Routing por LLM: classificador zero-shot ('precisa de retrieval?')",
    "Routing por complexidade: simples / média / complexa",
    "Routing por fonte: base local vs web vs KG",
    "",
    "LangGraph: RouteQuery node com conditional edges"
], "📖 Sweet spot: LLM leve (GPT-4o-mini, Haiku) com structured output. Regras para óbvios.\n💡 Analogia: triagem de central. Regras para casos óbvios + humano para ambíguos.\n⚠️ Usar LLM para perguntas óbvias é desperdício. Regras bastam para 'quem é X?'.\n➡️ Vamos ver o diagrama.", 18, T)

s = content_slide("Adaptive RAG — Diagrama", [
    "Fluxo: pergunta → classificador → (direto | simples | complexa) → gerar → resposta",
    "Classificador decide: responder direto, retrieve top-3, ou retrieve + query rewrite",
    "",
    "Diagrama: 12-Diagrams/ETHAGT06/adaptive-rag.mmd",
    "Fonte: LangGraph examples — adaptive_rag"
], "📖 Coração é o classificador (nó de decisão). 3 caminhos convergem para 'gerar resposta'.\n💡 Analogia: restaurante com 3 cardápios. Garçom (classificador) decide.\n❓ 'Qual caminho será mais usado?' (depende, geralmente 'simples')\n⚠️ Classificador muito complexo = erro. Comece com 3 classes.\n➡️ Vamos concretizar com exemplos.", 19, T, obj="Diagrama D5")

s = content_slide("Exemplo Prático: Quando Responder Direto", [
    "'Quem é o presidente do Brasil?' → responder direto (paramétrico)",
    "'Qual a política de férias para estagiários?' → recuperar (específico)",
    "'Compare as políticas de 2023 e 2024' → recuperar + query rewrite",
    "",
    "O classificador pode errar — por isso CRAG existe"
], "📖 Vamos praticar a decisão. 3 exemplos, 3 caminhos diferentes.\n💡 Analogia: GPS que decide rota. Às vezes escolhe beco sem saída. CRAG recalcula.\n⚠️ Confiar 100% no classificador é erro. ~85% de acerto é bom; 15% precisam de correção.\n➡️ Como implementar? Código.", 20, T)

code_slide(
    "Implementação: Adaptive RAG em LangGraph",
    "from langgraph.graph import StateGraph\n"
    "from pydantic import BaseModel\n\n"
    "class RouteQuery(BaseModel):\n"
    "    route: str  # 'direct' | 'simple' | 'complex'\n\n"
    "graph = StateGraph(State)\n"
    "graph.add_node('route_query', route_query_node)\n"
    "graph.add_node('retrieve', retrieve_node)\n"
    "graph.add_node('generate', generate_node)\n\n"
    "graph.add_conditional_edges('route_query', lambda s: s['route'],\n"
    "    {'direct': 'generate', 'simple': 'retrieve', 'complex': 'retrieve'})\n"
    "graph.add_edge('retrieve', 'generate')\n"
    "graph.set_entry_point('route_query')\n"
    "compiled = graph.compile()",
    "📖 StateGraph com nós route_query, retrieve, generate. Conditional edge direciona.\n💡 Analogia: circuito com interruptores.\n⚠️ Não coloque lógica de classificação no prompt do generate. Separe em nó próprio.\n➡️ Limitações do Adaptive.",
    21, T, obj="Diagrama D6"
)

s = content_slide("Limitações do Adaptive RAG", [
    "Classificador pode errar (falso negativo: não recupera quando deveria)",
    "Não avalia a qualidade dos docs recuperados — confia cegamente",
    "Não tem fallback se a base local não tem a resposta",
    "Não corrige: se recuperou lixo, gera com lixo",
    "",
    "→ Motivação para CRAG (próxima seção)"
], "📖 Adaptive resolve 'recuperar quando não precisa', mas NÃO 'recuperar lixo'.\n💡 Analogia: médico que sabe QUANDO pedir exames, mas não interpreta resultados ruins.\n➡️ Vamos praticar a decisão.", 22, T)

exercise_slide(
    "Exercício: Direto ou Recuperar?",
    [
        "Para cada pergunta, vote: Direto (D) ou Recuperar (R):",
        "",
        "1. 'O que é Python?'",
        "2. 'Qual a política de home office da Etho?'",
        "3. 'Quem escreveu Dom Casmurro?'",
        "4. 'Compare as APIs de dois produtos internos'",
        "5. 'Qual a temperatura de ebulição da água?'",
        "",
        "Votação rápida (mãos levantadas)"
    ],
    "📖 Gabarito: 1=D, 2=R, 3=D, 4=R+rewrite, 5=D.\n❓ Votar em cada uma. Anotar divergências.\n⚠️ Pergunta 4 gera dúvida: precisa de duas buscas + rewrite.\n➡️ Vocês fizeram o trabalho do classificador. Agora, e se o retrieval trouxer lixo? CRAG.",
    23, T
)

s = section_slide("3", "Corrective RAG (CRAG): Avaliar Antes de Confiar")
add_notes(s, "📖 Segunda evolução: CRAG adiciona avaliador de relevância. Controle de qualidade que faltava.\n➡️ A ideia central.")

comparison_slide(
    "Avaliar Antes de Confiar",
    "Adaptive RAG (confia)", [
        "Decide SE recuperar",
        "Depois confia cegamente nos docs",
        "Se recuperou lixo, gera com lixo"
    ],
    "CRAG (avalia)", [
        "Recupera, AVALIA relevância dos docs",
        "3 caminhos: usar / refinar / web fallback",
        "Fonte: Yan et al., arXiv:2401.15884"
    ],
    "📖 CRAG adiciona avaliação ANTES de gerar. Se docs são bons, usa; se parciais, refina; se ruins, web.\n💡 Analogia: cozinheiro que prova ingredientes antes de cozinhar.\n❓ 'Em qual caso de vocês o CRAG seria útil?'\n⚠️ CRAG não é 'adaptive com web'. O coração é o AVALIADOR, não o web fallback.\n➡️ Como funciona o avaliador?", 25, T
)

s = content_slide("O Retrieval Evaluator (Grau de Relevância)", [
    "Avaliador: LLM ou modelo leve classifica cada doc recuperado",
    "Três classes de relevância: correto / ambíguo / incorreto",
    "Score agregado decide o caminho:",
    "  Maioria correta → usar",
    "  Mistura → refinar",
    "  Maioria incorreta → web search",
    "",
    "Implementação: prompt estruturado ou modelo fine-tuned"
], "📖 Para cada doc, atribui classe: correto/ambíguo/incorreto. Agregação decide caminho.\n💡 Analogia: jurado que vota culpado/inocente/dúvida.\n⚠️ Threshold binário (relevante/não) é menos informativo. 3 classes permitem caminho 'refinar'.\n➡️ Vamos ver os 3 caminhos.", 26, T)

s = content_slide("Três Caminhos: Usar / Corrigir / Web", [
    "Caminho 1 (relevante): usar docs diretamente → gerar resposta",
    "Caminho 2 (ambíguo): knowledge refinement — extrair apenas partes relevantes",
    "Caminho 3 (irrelevante): web search como fallback → gerar com web results",
    "",
    "Trigger de fallback: quando a base local não tem a resposta"
], "📖 Caminho 1 (verde): usar direto. Caminho 2 (amarelo): refinement. Caminho 3 (azul): web.\n💡 Analogia: pesquisador. Achou referência perfeita → usa. Artigo parcial → extrai útil. Nada na biblioteca → Google Scholar.\n➡️ Diagrama completo.", 27, T)

s = content_slide("CRAG Flow — Diagrama", [
    "Fluxo: query → retrieve (KB local) → avaliador → (sim/ambíguo/não) → gerar → resposta",
    "O avaliador é o coração do CRAG",
    "Web search só é acionado quando docs locais são irrelevantes",
    "",
    "Diagrama: 12-Diagrams/ETHAGT06/crag-flow.mmd",
    "Fonte: arXiv:2401.15884; LangGraph example crag"
], "📖 Avaliador é camada de controle de qualidade entre retrieve e generate.\n💡 Analogia: esteira de fábrica com estação de inspeção.\n❓ 'Qual caminho será mais comum?' (sim, se base é boa)\n⚠️ CRAG NÃO busca sempre na web. Web é ÚLTIMO recurso.\n➡️ Detalhando cada caminho.", 28, T, obj="Diagrama D7")

s = content_slide("Caminho 1: Usar (Docs Relevantes)", [
    "Avaliador classifica como 'correto' → docs vão direto para geração",
    "Equivale ao RAG ingênuo, mas com confirmação de qualidade",
    "Caminho mais comum quando a base é boa e a pergunta é clara",
    "Latência: mínima (só overhead do avaliador)"
], "📖 Caminho feliz. ~80% dos casos em base curada. Overhead: 1 chamada de LLM (avaliador).\n💡 Analogia: controle de qualidade que aprova 95%. Parece desperdício, mas evita recall caro.\n➡️ E quando os docs são parcialmente bons?", 29, T)

s = content_slide("Caminho 2: Corrigir (Knowledge Refinement)", [
    "Avaliador classifica como 'ambíguo' → docs têm partes relevantes e irrelevantes",
    "Knowledge refinement: extrair apenas as partes relevantes de cada doc",
    "Técnicas: decompor em sentenças, re-avaliar, filtrar",
    "Resultado: contexto limpo para o gerador",
    "",
    "Trade-off: mais um step de LLM = mais latência e custo"
], "📖 Decompor doc em sentenças, re-avaliar cada uma, ficar com as relevantes. Contexto limpo.\n💡 Analogia: peneirar areia. Pá com pedras → areia fina.\n⚠️ Refinement sem re-avaliar é só heurística (tamanho), não relevância.\n➡️ E quando os docs são completamente irrelevantes?", 30, T)

s = content_slide("Caminho 3: Web Search (Fallback)", [
    "Avaliador classifica como 'incorreto' → base local não tem a resposta",
    "Web search: Tavily, SerpAPI, ou busca nativa do modelo",
    "Combinar web results com qualquer fragmento útil dos docs locais",
    "Quando ativa: base desatualizada, pergunta sobre evento recente, gap de cobertura",
    "",
    "Pergunta: Quando CRAG decide buscar na web? Quando o avaliador diz 'incorreto'"
], "📖 Web é plano B, não padrão. Ativa quando base local falha completamente.\n💡 Analogia: enciclopédia física sem verbete → internet.\n❓ 'Quando CRAG decide buscar na web?' (quando avaliador diz 'incorreto')\n⚠️ Ativar web sempre é erro. Latência, custo, risco de web não confiável.\n➡️ Vamos ver o código do avaliador.", 31, T)

code_slide(
    "Implementação: O Avaliador",
    "from pydantic import BaseModel\n\n"
    "class GradeDocuments(BaseModel):\n"
    "    binary_score: str  # 'yes' | 'no'\n\n"
    "def grade_documents(state):\n"
    "    filtered = []\n"
    "    for doc in state['documents']:\n"
    "        score = grader.invoke({'question': state['question'], 'doc': doc})\n"
    "        if score.binary_score == 'yes':\n"
    "            filtered.append(doc)\n"
    "    return {'documents': filtered}\n\n"
    "graph.add_conditional_edges('grade_documents',\n"
    "    lambda s: 'generate' if s['documents'] else 'web_search')",
    "📖 Função grade_documents itera sobre docs, classifica, acumula aprovados. Conditional edge decide.\n💡 Analogia: filtro de spam de email.\n⚠️ Threshold contínuo é menos simples que binário (yes/no). Para 3 classes, use correto/ambíguo/incorreto.\n➡️ Comparação com RAG ingênuo.",
    32, T
)

comparison_slide(
    "RAG Ingênuo vs CRAG",
    "RAG Ingênuo", [
        "Conflia cegamente nos docs",
        "Sem fallback",
        "Sem avaliação",
        "Barato, rápido"
    ],
    "CRAG", [
        "Avalia relevância dos docs",
        "Refina se parcial, web se irrelevante",
        "Custo: +1-2 chamadas de LLM (~$0.005/query)",
        "Grounding muito maior"
    ],
    "📖 CRAG adiciona ~$0.005/query. Em 10k/dia = $50/dia. Geralmente vale.\n💡 Analogia: comprar sem marca (ingênuo) vs com certificação (CRAG).\n⚠️ CRAG em latência crítica (<500ms) pode não caber.\n➡️ Quando CRAG brilha e quando falha?",
    33, T
)

s = content_slide("CRAG — Quando Brilha e Quando Falha", [
    "Brilha: base com cobertura parcial, perguntas diversas, dados que envelhecem",
    "Falha: avaliador erra (falso positivo de relevância), web search retorna lixo",
    "Melhoria: combinar CRAG com re-ranking (Seção G)",
    "",
    "Limitação: não reflete sobre a PRÓPRIA resposta — isso é Self-RAG"
], "📖 CRAG avalia DOCS, não RESPOSTA. Se resposta alucina sobre docs bons, CRAG não percebe.\n💡 Analogia: inspetor de matéria-prima. Garante ingredientes, mas não prova o prato. Self-RAG prova.\n➡️ Vamos ao próximo nível: Self-RAG.", 34, T)

s = section_slide("4", "Self-RAG: Modelo que Reflete sobre Recuperação")
add_notes(s, "📖 Terceira evolução: Self-RAG avalia TUDO — docs, resposta, processo. Auto-avaliação.\n➡️ A ideia central.")

comparison_slide(
    "Modelo que Reflete sobre Recuperação",
    "CRAG (avalia docs)", [
        "Avalia relevância dos docs ANTES de gerar",
        "Não avalia a resposta gerada",
        "Pipeline com gate de decisão"
    ],
    "Self-RAG (avalia tudo)", [
        "Reflete sobre TODO o processo",
        "Julga se a resposta é suportada pelos docs",
        "Modelo treinado para emitir tokens de reflexão",
        "Fonte: Asai et al., arXiv:2310.11511"
    ],
    "📖 CRAG avalia docs (antes de gerar). Self-RAG avalia resposta (depois de gerar). Se não suportada, regenera.\n💡 Analogia: cozinheiro que prova os ingredientes vs cozinheiro que prova o prato final e refaz.\n❓ 'Em qual caso vocês precisariam de Self-RAG?' (jurídico, médico, financeiro)\n⚠️ Self-RAG não é 'CRAG com mais chamadas'. Adiciona REFLEXÃO sobre a resposta.\n➡️ Como o modelo reflete? Tokens de reflexão.",
    36, T
)

s = content_slide("Tokens de Reflexão do Self-RAG", [
    "[Retrieve] — decidir se precisa recuperar",
    "[Retrieve(target)] — decidir qual fonte recuperar",
    "[Relevant] / [Irrelevant] — julgar relevância do doc",
    "[Fully supported] / [Partly supported] / [No support] — julgar se resposta é suportada",
    "[Utility] — julgar utilidade da resposta final",
    "",
    "Cada token controla uma decisão no pipeline"
], "📖 Self-RAG original é FINE-TUNED para emitir tokens especiais. Reflexão embutida, sem prompts extras.\n💡 Analogia: cirurgião que murmura checklist durante a operação.\n⚠️ Tokens não funcionam com GPT-4 (não foi treinado). Para não-treinados, prompting.\n➡️ Vamos ver o fluxo.", 37, T)

s = content_slide("Self-RAG — Fluxo com Tokens de Reflexão", [
    "Fluxo: pergunta → [Retrieve?] → recuperar → [Relevant?] → [Fully supported?] → resposta",
    "Se [No support]: regenerar ou recuperar mais",
    "O modelo decide em cada etapa via tokens de reflexão",
    "",
    "Fonte: arXiv:2310.11511"
], "📖 Decision points: [Retrieve?] (como Adaptive), [Relevant?] (como CRAG), [Fully supported?] (NOVO!).\n💡 Analogia: estudante com rubrica. A cada questão se pergunta: preciso consultar? página relevante? resposta fundamentada?\n❓ 'Qual token é mais importante?' ([Fully supported] — diferencia Self-RAG de CRAG)\n⚠️ Tokens são opcionais. Pode usar só [Retrieve] e [Fully supported].\n➡️ Como funciona o modelo treinado?", 38, T, obj="Diagrama D8")

s = content_slide("Self-RAG Original — Modelo Treinado", [
    "Self-RAG original: LLM fine-tuned para emitir tokens de reflexão",
    "Treinamento: dados com anotações de reflexão (relevância, suporte, utilidade)",
    "Vantagem: reflexão embutida no modelo, sem prompts extras",
    "Desvantagem: requer modelo específico — não funciona com GPT-4, Claude",
    "Por isso: adaptação para modelos não-treinados via prompting"
], "📖 Original é Llama-Self-RAG fine-tuned. Elegante (uma passada) mas preso àquele modelo.\n💡 Analogia: músico de jazz (treinado) improvisa natural; músico clássico precisa de partitura (prompting).\n⚠️ Baixar 'Self-RAG model' existe mas é pequeno. Para produção, prompting em modelo forte é melhor.\n➡️ Como adaptar para qualquer modelo? Prompting.", 39, T)

s = content_slide("Self-RAG via Prompting", [
    "Replicar os tokens de reflexão via prompting estruturado",
    "Prompt pede ao LLM para:",
    "  1. Decidir se precisa recuperar",
    "  2. Avaliar relevância de cada doc",
    "  3. Avaliar se a resposta é suportada",
    "  4. Regenerar se não suportada",
    "Custo: múltiplas chamadas de LLM",
    "Benefício: funciona com qualquer modelo (GPT-4, Claude, Llama)"
], "📖 Para GPT-4/Claude: prompt estruturado pede reflexão. JSON {supported: fully|partly|no}. Se 'no', regenera.\n💡 Analogia: traduzir receita de chef para iniciante. Mesmo resultado, método diferente.\n⚠️ Um prompt gigante para todas reflexões é ruim. Separe em chamadas (nós do grafo).\n➡️ Código.", 40, T)

code_slide(
    "Implementação: Self-RAG via Prompt",
    "from pydantic import BaseModel\n\n"
    "class GradeDocuments(BaseModel):\n"
    "    binary_score: str  # 'yes' | 'no'\n\n"
    "class GradeHallucinations(BaseModel):\n"
    "    binary_score: str  # 'yes' (supported) | 'no'\n\n"
    "class GradeAnswer(BaseModel):\n"
    "    binary_score: str  # 'yes' (useful) | 'no'\n\n"
    "graph.add_node('grade_documents', grade_docs_node)\n"
    "graph.add_node('generate', generate_node)\n"
    "graph.add_node('grade_hallucination', grade_halluc_node)\n"
    "graph.add_edge('generate', 'grade_hallucination')\n"
    "graph.add_conditional_edges('grade_hallucination',\n"
    "    lambda s: 'useful' if supported else 'regenerate')",
    "📖 Pydantic models para saídas estruturadas. Nós: grade_documents, generate, grade_hallucination, grade_answer. Loop com max_retries.\n💡 Analogia: redação com rubrica. Cada parágrafo passa por critérios.\n⚠️ Sem max_retries na regeneração = loop infinito. Sempre use guardrail.\n➡️ Vamos sistematizar as 3 arquiteturas.",
    41, T
)

comparison_slide(
    "Adaptive vs CRAG vs Self-RAG",
    "Checkpoints", [
    ],
    "Tabela comparativa", [
        "Adaptive: decide SE recuperar ✅",
        "  Avalia docs ❌ | Avalia resposta ❌",
        "CRAG: decide SE recuperar ❌",
        "  Avalia docs ✅ | Avalia resposta ❌",
        "Self-RAG: decide SE recuperar ✅",
        "  Avalia docs ✅ | Avalia resposta ✅",
        "",
        "Complexidade/custo: Adaptive < CRAG < Self-RAG"
    ],
    "📖 Cada nível adiciona reflexão. Adaptive (1 checkpoint), CRAG (2), Self-RAG (3).\n💡 Analogia: níveis de inspeção. Visual → instrumentos → auditoria.\n❓ 'Qual vocês usariam no projeto?' (justificarem)\n⚠️ Ir direto para Self-RAG 'para ser seguro' é overkill. Adaptive + CRAG resolvem 80%.\n➡️ Quando usar Self-RAG?",
    42, T, obj="Diagrama D9"
)

s = content_slide("Quando Usar Self-RAG", [
    "Use Self-RAG quando:",
    "  Alucinação é inaceitável (jurídico, médico, financeiro)",
    "  Você pode pagar o custo de múltiplas chamadas",
    "  Precisa de garantia de que a resposta é suportada",
    "",
    "Não use quando:",
    "  Latência é crítica",
    "  Base é pequena e confiável",
    "  Caso de uso tolera imprecisão"
], "📖 Self-RAG vale quando alucinação é inaceitável E pode pagar o custo.\n💡 Analogia: auditoria forense. Vale para fraudes milionárias; overkill para conta do café.\n➡️ Vamos praticar.", 43, T)

exercise_slide(
    "Exercício — Adaptive vs Self-RAG",
    [
        "Cenário: sistema de FAQ jurídico com 10.000 documentos",
        "",
        "Em duplas (5 min):",
        "1. Quando usar Adaptive RAG vs Self-RAG?",
        "2. Justifique com 3 critérios: tipo de pergunta, custo, qualidade",
        "3. Escreva o prompt de reflexão para um doc recuperado",
        "",
        "3 min discussão + 2 min compartilhar"
    ],
    "📖 Adaptive: perguntas simples, latência crítica. Self-RAG: jurídico, alucinação inaceitável.\n❓ Em duplas, 3 min. Depois 2 duplas compartilham.\n⚠️ Escolher Self-RAG 'por segurança' em FAQ pode ser overkill. Avaliar risco regulatório.\n➡️ Vamos ao próximo bloco: Agentic RAG.",
    44, T
)

s = section_slide("5", "Agentic RAG: O Agente Dirige Todo o Processo")
add_notes(s, "📖 Topo da escalada. Não há pipeline predefinido. Agente decide tudo. Conexão com ETHAGT01: Augmented LLM com retrieval in-loop.\n➡️ A ideia central.")

comparison_slide(
    "O Agente Dirige Todo o Processo",
    "Adaptive/CRAG/Self-RAG (fixo)", [
        "Pipeline com gates de decisão",
        "Estrutura predefinida",
        "Determinístico na estrutura"
    ],
    "Agentic RAG (loop)", [
        "O AGENTE decide tudo",
        "Planeja busca, refina queries, decide parar",
        "Tools de busca como ferramentas do agente",
        "Loop do agente É o pipeline"
    ],
    "📖 Adaptive/CRAG/Self-RAG são receitas (passos predefinidos). Agentic é chef criando prato novo.\n💡 Analogia: receita vs chef.\n❓ 'Qual o risco de Agentic?' (loop infinito, custo imprevisível — max_steps é essencial)\n⚠️ Pular para Agentic sem entender o básico é anti-pattern. Comece simples.\n➡️ O que o agente faz exatamente?",
    46, T
)

s = content_slide("Capacidades do Agente RAG", [
    "Planeja: 'Preciso buscar X, depois Y, depois combinar'",
    "Refina queries: 'A busca por X retornou pouco — vou reformular'",
    "Decide parar: 'Tenho informação suficiente' ou 'Preciso de mais um hop'",
    "max_steps como guardrail (lição de ETHAGT01)",
    "",
    "O agente raciocina sobre o QUE recuperou e o QUE ainda falta"
], "📖 3 capacidades únicas: planeja, refina, decide parar. Emergem do loop ReAct (ETHAGT01).\n💡 Analogia: detetive investigando. Planeja interrogatórios, reformula perguntas, decide quando tem provas.\n⚠️ Sem critério de parada no prompt = agente não sabe quando parar. Ensine.\n➡️ A capacidade mais poderosa: multi-hop.",
    47, T
)

s = content_slide("Multi-Hop: Cadeias de Recuperação", [
    "Single-hop: uma busca → resposta",
    "Multi-hop: múltiplas buscas encadeadas, cada uma dependendo da anterior",
    "",
    "Exemplo: 'Quem fundou a empresa que criou o ChatGPT?'",
    "  Hop 1: buscar 'quem criou o ChatGPT?' → OpenAI",
    "  Hop 2: buscar 'quem fundou a OpenAI?' → Sam Altman, Elon Musk",
    "  Resposta: sintetizar",
    "",
    "Sem agente: difícil de orquestrar"
], "📖 Multi-hop: resposta da 1ª busca é INPUT para a 2ª. Sem agente, muito difícil.\n💡 Analogia: investigação em camadas. Suspeito → passado → aliases. Cada hop depende do anterior.\n❓ 'Pensem em uma pergunta multi-hop do sistema de vocês.'\n⚠️ Multi-hop NÃO é buscas paralelas. É SEQUENCIAL e DEPENDENTE.\n➡️ Vamos ver um trace real.", 48, T, obj="Diagrama D10")

s = content_slide("Multi-Hop — Trace Detalhado", [
    "Pergunta: 'Qual a diferença entre a política de férias da Etho e a CLT?'",
    "",
    "Thought: 'Preciso de duas fontes'",
    "Action: search_etho_docs('política de férias')",
    "Observation: '30 dias após 1 ano...'",
    "Action: search_web('CLT férias Brasil')",
    "Observation: '30 dias corridos após 12 meses...'",
    "Thought: 'Agora posso comparar'",
    "Answer: 'A Etho segue a CLT, com diferença em...'"
], "📖 Trace real: pergunta exige duas fontes (interna + web). Agente raciocina, busca, combina, responde.\n💡 Analogia: jornalista investigativo. Arquivo interno + internet + escreve matéria.\n⚠️ Agente não 'sabe' que precisa de duas fontes — raciocina a partir do prompt.\n➡️ As tools de busca são especiais?", 49, T)

s = content_slide("Tools de Busca como Ferramentas do Agente", [
    "Agente tem tools: search_internal(query), search_web(query), search_kg(query)",
    "ACI importa: tool bem documentada > prompt melhor (ETHAGT01/ETHAGT02)",
    "Agente escolhe qual tool usar baseado na pergunta",
    "Exemplo: 'política interna?' → search_internal; 'notícia de hoje?' → search_web",
    "",
    "Tool return format: lista de docs com score, não texto bruto"
], "📖 Tools de busca são tools como qualquer outra. Lição de ACI aplica totalmente.\n💡 Analogia: caixa de ferramentas. Agente é mecânico que escolhe a ferramenta. Etiqueta ruim = erro.\n❓ 'Quantas tools de busca vocês têm? Documentadas como?'\n⚠️ Retornar texto bruto é ruim. Melhor: lista estruturada de docs com score.\n➡️ Fontes.", 50, T)

s = content_slide("Múltiplas Fontes: Web, Interno, KG", [
    "Web search: Tavily, SerpAPI, Brave — dados externos e recentes",
    "Interno: vector DB (Qdrant, Milvus), SQL DB — dados proprietários",
    "Knowledge Graph: Neo4j, GraphRAG — dados relacionais estruturados",
    "",
    "Agente combina fontes: web para contexto + interno para detalhe",
    "Agente decide quando uma fonte é insuficiente"
], "📖 Diversidade de fontes é superpoder. Web (recente), interno (proprietário), KG (relacional).\n💡 Analogia: pesquisador com 3 bibliotecas. Sabe qual consultar para cada pergunta.\n⚠️ Usar só vector DB é limitante. Web e KG resolvem classes inteiras de perguntas.\n➡️ Uma fonte especial: Knowledge Graph via GraphRAG.", 51, T)

s = content_slide("GraphRAG: Do Local ao Global", [
    "RAG tradicional: recupera chunks locais (vizinhança vetorial)",
    "GraphRAG: constrói grafo de conhecimento e recupera subgrafos",
    "Vantagem: responde perguntas globais ('quais os temas deste corpus?')",
    "Pipeline: extrair entidades/relações → comunidades → sumarizar → recuperar",
    "",
    "Fonte: Edge et al., Microsoft, arXiv:2404.16130",
    "No Agentic RAG: GraphRAG como uma tool"
], "📖 GraphRAG resolve perguntas GLOBAIS que vector RAG não consegue.\n💡 Analogia: vector RAG = olhar páginas; GraphRAG = olhar índice e sumário.\n❓ 'Qual pergunta global vocês fariam no corpus de vocês?'\n⚠️ GraphRAG para perguntas locais é overkill. Vector RAG é mais rápido.\n➡️ DEMO.", 52, T, obj="Diagrama D11")

s = content_slide("DEMO — Agentic RAG Multi-Hop", [
    "Código: 05-Labs/ETHAGT06/Lab2-Agentic-RAG-Multi-Hop",
    "Agente com tools: search_internal, search_web",
    "Pergunta multi-hop: 'Compare a política de férias da Etho com a CLT'",
    "",
    "Trace: planejamento → busca interna → busca web → síntese",
    "max_steps guardrail em ação"
], "📖 Demo ao vivo. Se API falhar, tenho screenshot do trace.\n💡 Analogia: assistir detetive trabalhar. Cada thought/action/observation visível.\n⚠️ Se API falhar, NÃO debugue ao vivo. Use screenshot e siga.\n➡️ Pergunta sobre a demo.", 53, T, obj="DEMO")

exercise_slide(
    "Pergunta da DEMO",
    [
        "'O que acontece se o agente não encontrar nada na busca interna?'",
        "'Como o agente decide que tem informação suficiente para parar?'",
        "'E se o max_steps for muito baixo?'",
        "",
        "Discussão em duplas (2 min)"
    ],
    "📖 (1) Depende do prompt — se instruído a combinar, tenta web. (2) Critério no prompt: 'pare quando tiver informação suficiente'. (3) max_steps baixo = resposta incompleta.\n❓ Em duplas, 2 min. Depois 2 duplas compartilham.\n➡️ Vamos sistematizar a escalada.",
    54, T
)

s = content_slide("A Escalada — Adaptive → CRAG → Self-RAG → Agentic", [
    "Adaptive: decide SE recuperar",
    "CRAG: avalia docs recuperados",
    "Self-RAG: avalia docs + resposta",
    "Agentic: agente dirige todo o processo",
    "",
    "Trade-off: mais controle = mais custo e complexidade",
    "Regra: comece com Adaptive, só suba com evidência de insuficiência"
], "📖 Cada degrau adiciona reflexão. Adaptive < CRAG < Self-RAG < Agentic em custo/complexidade.\n💡 Analogia: aprender a dirigir. Cidade → rodovia → noite → corrida.\n❓ 'Em qual degrau está o RAG de vocês hoje?' (maioria em Adaptive ou nenhum)\n⚠️ Querer Agentic porque é 'mais avançado' é anti-pattern. Adaptive resolve 60-70%.\n➡️ Agora, engenharia de qualidade.", 55, T, obj="Diagrama D12")

s = section_slide("6", "Engenharia de Qualidade: O RAG é Tão Bom Quanto Seus Chunks")
add_notes(s, "📖 Arquiteturas são a estrutura. Engenharia de qualidade é o conteúdo. Chunking, re-rank, rewriting, hybrid.\n➡️ Começando pelo fundamento: chunking.")

s = content_slide("Chunking: O Fundamento Esquecido", [
    "'Garbage in, garbage out' — chunking ruim contamina todo o pipeline",
    "Chunking fixo (512 tokens): simples, mas quebra contexto",
    "Chunking por sentença/parágrafo: preserva semântica",
    "Chunking semântico: agrupa por tópicos",
    "",
    "O chunk é o átomo do RAG — tudo depende dele"
], "📖 Chunking é a decisão mais impactante e mais subestimada. Átomo do RAG.\n💡 Analogia: cortar ingredientes. Corte ruim = prato ruim, não importa o tempero.\n❓ 'Qual estratégia de chunking vocês usam?' (maioria: fixo)\n⚠️ Otimizar embeddings antes de chunking é ordem errada. Chunking > embeddings > re-rank.\n➡️ Duas estratégias avançadas.", 57, T)

s = content_slide("Chunking Semântico e Hierárquico", [
    "Chunking semântico:",
    "  Detectar mudanças de tópico (embeddings de sentença → clustering)",
    "  Agrupar sentenças relacionadas",
    "",
    "Chunking hierárquico:",
    "  Múltiplos níveis: seção → parágrafo → sentença",
    "  Parent-child: recuperar filho, retornar pai para contexto",
    "",
    "Semântico para docs densos; hierárquico para estruturados"
], "📖 Semântico: embeddings de sentenças detectam mudanças de tópico. Hierárquico: parent-child (filho match preciso, pai contexto rico).\n💡 Analogia: semântico = agrupar conversas por assunto; hierárquico = livro com capítulos/seções/parágrafos.\n⚠️ Semântico em docs curtos não adiciona. Hierárquico em docs planos não ajuda.\n➡️ E a novidade de 2024?", 58, T)

s = content_slide("Late Chunking (Contextual Retrieval)", [
    "Problema: chunk isolado perde contexto do documento",
    "Late chunking (Anthropic Contextual Retrieval, 2024):",
    "  Processar documento inteiro ANTES de chunkar",
    "  Cada chunk recebe contexto do documento",
    "  Embeddings capturam contexto, não só fragmento",
    "Implementação: doc completo → chunkar → embedar",
    "Reduz falhas de recuperação em 30-50% (Anthropic)"
], "📖 Novidade mais impactante de 2024. Processa doc inteiro antes de chunkar. Embeddings muito mais ricos. -30-50% falhas.\n💡 Analogia: ler capítulo inteiro antes de destacar frase. Com contexto, destaque faz sentido.\n❓ 'Conheciam Contextual Retrieval?' (provavelmente minoria)\n⚠️ Late chunking ≠ 'chunking maior'. É processar doc antes.\n➡️ Re-ranking.", 59, T)

s = content_slide("Re-Ranking: Por Que e Como", [
    "Recuperação (BM25/densa): rápida, barata, mas imprecisa",
    "Re-ranking: lento, caro, mas preciso",
    "Pipeline: retrieve top-50 → re-rank → top-5 para o LLM",
    "Re-ranker avalia relevância query-doc, não só similaridade",
    "",
    "Reduz ruído no contexto do LLM"
], "📖 Retrieve MUITOS (top-50) com método barato, re-rankeie com cross-encoder PRECISO, fique com top-5.\n💡 Analogia: contratar. 100 currículos → triagem rápida → entrevista 5 melhores.\n⚠️ Re-rank de TODA a base é caro. Só dos top-50 já recuperados.\n➡️ Quais re-rankers?", 60, T)

s = content_slide("Re-Rankers: Cohere, bge, Jina", [
    "Cohere Rerank: API comercial, alta qualidade, paga por requisição",
    "bge-reranker (BAAI): open-source, roda local, bom custo-benefício",
    "Jina Reranker: API e self-hosted, foco em multilingual",
    "",
    "Critério: latência, custo, qualidade, privacidade",
    "Trade-off: cross-encoder é mais preciso mas mais lento"
], "📖 Cohere (qualidade, pago), bge (local, grátis), Jina (multilingual). Escolha conforme caso.\n💡 Analogia: tradutor. Profissional (Cohere), estagiário fluente (bge), poliglota (Jina).\n❓ 'Qual vocês usariam no projeto?'\n➡️ Query transformation.", 61, T)

s = content_slide("Query Rewriting e HyDE", [
    "Query rewriting: LLM reescreve a pergunta para melhorar recuperação",
    "  'férias estagiário' → 'política de licença para aprendizes e estagiários'",
    "",
    "HyDE (Hypothetical Document Embeddings):",
    "  Gerar resposta HIPOTÉTICA à pergunta",
    "  Embedar a resposta hipotética (não a pergunta)",
    "  Buscar por similaridade com a resposta hipotética",
    "Por que funciona: resposta hipotética é mais próxima dos docs que a pergunta"
], "📖 Query rewriting: reformular para casar melhor com docs. HyDE: gerar resposta hipotética e embedar ELA.\n💡 Analogia: rewriting = reformular busca no Google. HyDE = imaginar resposta ideal e buscar por algo parecido.\n❓ 'Query rewriting: sempre reescrever? Quando não?' (não para diretas; sim para ambíguas/sinônimos)\n⚠️ HyDE sempre adiciona latência. Vale para abstratas; rewriting basta para diretas.\n➡️ Hybrid search.", 62, T, obj="Diagrama D14")

s = content_slide("Hybrid Search: BM25 + Densa", [
    "BM25 (lexical): captura keywords exatas, termos técnicos, nomes próprios",
    "Densa (embedding): captura semântica, sinônimos, paráfrase",
    "Hybrid: combinar scores com reciprocal rank fusion (RRF)",
    "",
    "Resultado: melhor recall que qualquer método isolado",
    "Implementação: Qdrant e Milvus suportam hybrid nativamente"
], "📖 Hybrid é o PADRÃO de produção. BM25 + densa com RRF. Melhor recall.\n💡 Analogia: BM25 = buscar por número de telefone; densa = buscar por 'restaurante italiano perto do parque'. Precisa dos dois.\n❓ 'Usam hybrid search hoje?' (provavelmente minoria)\n⚠️ Usar SÓ densa porque 'é moderna' é erro. BM25 vence em termos técnicos.\n➡️ Multi-vector.", 63, T, obj="Diagrama D15")

s = content_slide("Multi-Vector: ColBERT", [
    "Embedding tradicional: 1 vetor por documento",
    "ColBERT: 1 vetor POR TOKEN do documento",
    "Recuperação: max-similarity entre tokens da query e tokens do doc",
    "Vantagem: granularidade muito maior, captura match a nível de token",
    "Desvantagem: armazenamento e computação muito maiores",
    "Quando usar: corpora técnicos com terminologia específica"
], "📖 ColBERT: N vetores por doc (um por token). Max-similarity granular. Caro mas preciso.\n💡 Analogia: tradicional = resumir livro em frase; ColBERT = manter cada palavra indexada.\n⚠️ ColBERT para tudo é overkill. Reserve para corpora técnicos específicos.\n➡️ Multimodal.", 64, T)

s = content_slide("Multimodal RAG: Texto, Imagem, Tabela", [
    "Texto: embedding de texto (padrão)",
    "Imagem: CLIP, vision embeddings — buscar por descrição ou imagem",
    "Tabela: converter para texto estruturado ou embedar como tabular",
    "Multimodal: combinar texto + imagem + tabela no mesmo espaço",
    "Modelos: GPT-4V, Claude Vision, LLaVA",
    "Desafio: chunking de imagem (regiões), alinhamento modal"
], "📖 RAG multimodal integra texto, imagem, tabela. CLIP para imagens, conversão para tabelas.\n💡 Analogia: enciclopédia ilustrada. Busca por texto, imagem, ou tabela.\n⚠️ Embedar imagens com embeddings de texto NÃO funciona. Precisa de CLIP/vision.\n➡️ Vamos praticar a escolha.", 65, T)

exercise_slide(
    "Exercício — Escolhendo a Estratégia",
    [
        "3 cenários — em duplas, escolher estratégia:",
        "",
        "1. Documentação técnica em inglês (10k páginas)",
        "2. FAQ jurídico em português (5k docs)",
        "3. Catálogo de produtos com imagens (50k items)",
        "",
        "Justificar com 2 critérios cada",
        "2 min discussão"
    ],
    "📖 (1) Hybrid + re-rank. (2) Query rewriting + semantic chunking. (3) Multimodal (CLIP) + ColBERT.\n❓ Em duplas, 2 min. Depois 2 duplas compartilham.\n➡️ Agora, avaliação.",
    66, T
)

s = section_slide("7", "Avaliação de RAG: Sem Métricas, Você Está Cego")
add_notes(s, "📖 Seção mais subestimada. Sem métricas, não sabe se melhorou. Sem holdout, sem regressão. Separa hobby de produção.\n➡️ Por que avaliar RAG é diferente?")

s = content_slide("Por Que Avaliar RAG é Diferente", [
    "LLM puro: a resposta está certa ou errada",
    "RAG: 3 componentes independentes:",
    "  1. Retrieval: docs recuperados são relevantes?",
    "  2. Generation: resposta é fiel aos docs?",
    "  3. End-to-end: resposta responde à pergunta?",
    "Cada componente precisa de métricas próprias",
    "Sem isso: não sabe se problema é retrieval ou generation"
], "📖 3 componentes falham independentemente. Sem medir cada um, não sabe onde otimizar.\n💡 Analogia: carro que não anda. Motor? Combustível? Transmissão? Sem medir, troca peças às cegas.\n❓ 'Se faithfulness 0.95 mas relevance 0.40, qual o problema?' (generation fiel mas não responde — docs não cobrem)\n➡️ Métricas de geração.", 68, T)

s = content_slide("Métricas — Faithfulness e Answer Relevance", [
    "Faithfulness: a resposta é fiel aos docs? (sem alucinação)",
    "  Decompor resposta em claims → verificar cada claim contra os docs",
    "  Score = claims suportadas / total de claims",
    "",
    "Answer relevance: a resposta responde à pergunta?",
    "  Gerar perguntas a partir da resposta → comparar com pergunta original",
    "  Score = similaridade semântica",
    "",
    "Tensão: alta faithfulness pode baixar relevance"
], "📖 Faithfulness decompõe em claims, verifica cada uma contra docs. Relevance gera perguntas da resposta, compara com original.\n💡 Analogia: faithfulness = jornalista não inventa fatos; relevance = responde ao tema da pauta.\n⚠️ Otimizar SÓ faithfulness = respostas fiéis mas inúteis. Precisamos das duas.\n➡️ Métricas de retrieval.", 69, T)

s = content_slide("Métricas — Context Precision e Context Recall", [
    "Context precision: dos docs recuperados, quantos são relevantes?",
    "  Score = relevantes no top-k / k",
    "  Penaliza ruído",
    "",
    "Context recall: dos docs necessários, quantos foram recuperados?",
    "  Score = recuperados / necessários",
    "  Penaliza gaps",
    "",
    "Tensão: mais docs → mais recall, menos precision",
    "Re-ranking melhora precision sem sacrificar recall"
], "📖 Precision: quanto do que achou é valioso. Recall: quanto do total você encontrou.\n💡 Analogia: precision = proporção de ouro no cascalho; recall = quanto do ouro total você achou. Re-rank = bateia.\n⚠️ Otimizar SÓ recall = LLM recebe muito ruído. Precisamos das duas.\n➡️ Frameworks.", 70, T)

s = content_slide("Frameworks: Ragas, TruLens, DeepEval", [
    "Ragas: foco em RAG, métricas padronizadas, open-source, integra LangChain",
    "TruLens: tracing + avaliação, dashboard interativo, foco em produção",
    "DeepEval: testes unitários para LLM, integra CI/CD, pytest-style",
    "",
    "Critério: stack existente, necessidade de tracing, CI/CD",
    "Todos usam LLM-as-judge internamente"
], "📖 Ragas (RAG puro), TruLens (tracing + dashboard), DeepEval (pytest + CI/CD). Todos LLM-as-judge.\n💡 Analogia: Ragas = controle de qualidade; TruLens = monitor de produção; DeepEval = testes no pipeline.\n❓ 'Qual stack vocês usam hoje?'\n➡️ Pipeline de avaliação.", 71, T)

s = content_slide("Eval Pipeline — Diagrama", [
    "Pipeline: pergunta → retrieve → generate → resposta",
    "Eval coleta: pergunta, docs recuperados, resposta gerada",
    "Métricas: faithfulness, answer_relevance, context_precision, context_recall",
    "Dataset de holdout: perguntas + ground truth",
    "Execução: batch eval em CI/CD para detectar regressão",
    "",
    "Diagrama: 12-Diagrams/ETHAGT06/eval-pipeline.mmd"
], "📖 Eval corre em paralelo ao pipeline RAG. Coleta 3 coisas, computa 4 métricas. Em produção: batch + CI/CD.\n💡 Analogia: linha de produção com estação de qualidade. Amostra, mede, bloqueia lote se falhar.\n❓ 'Quantas perguntas tem o holdout de vocês?' (provavelmente zero)\n⚠️ Eval manual é subjetivo. Automático é obrigatório em produção.\n➡️ LLM-as-judge.", 72, T, obj="Diagrama D16")

s = content_slide("LLM-as-Judge e Dataset de Holdout", [
    "LLM-as-judge: usar LLM (GPT-4, Claude) para avaliar respostas",
    "Vieses a mitigar:",
    "  Position bias → randomizar ordem",
    "  Self-preference → usar modelo diferente do gerador",
    "  Length bias → normalizar",
    "Dataset de holdout: perguntas curadas com ground truth",
    "Regressão: faithfulness cai 0.90→0.82 → bloquear deploy",
    "Critério do projeto: faithfulness ≥ 0.85, context recall ≥ 0.80"
], "📖 LLM-as-judge tem vieses. Mitigue: randomize ordem, use modelo diferente, normalize tamanho. Holdout para regressão.\n💡 Analogia: jurado de concurso com preferências. Mitigue com randomização e jurados diversos.\n❓ 'Se faithfulness cair 5%, o que fazem?' (investigar, reverter, NÃO deployar)\n⚠️ Mesmo modelo para gerar e julgar = self-preference. Use modelos diferentes.\n➡️ Fechamento.", 73, T)

s = section_slide("8", "Boas Práticas, Anti-Patterns e Próximos Passos")
add_notes(s, "📖 Fechamento. Consolidar: boas práticas, anti-patterns, resumo, quiz, conexão, projeto, Q&A.\n➡️ O que fazer.")

s = content_slide("Boas Práticas (DO)", [
    "Comece com Adaptive RAG antes de pular para Agentic",
    "Avalie docs recuperados antes de gerar (CRAG)",
    "Use re-ranking sempre que possível",
    "Chunking semântico > chunking fixo",
    "Avaliação automatizada desde o dia 1 (Ragas)",
    "Dataset de holdout para regressão",
    "Hybrid search (BM25 + densa) como padrão",
    "max_steps em agentes RAG (guardrail)"
], "📖 8 boas práticas. Fazer 3 já coloca à frente de 80% dos RAGs em produção.\n💡 Analogia: fundamentos de esporte. Sem fundamento, jogada fancy falha.\n❓ 'Quantas dessas 8 vocês já fazem?' (levantar mãos por número)\n➡️ E o que NÃO fazer?", 75, T)

s = content_slide("Anti-Patterns (DON'T)", [
    "Começar com Agentic RAG sem entender Adaptive/CRAG",
    "Chunking fixo sem pensar no conteúdo",
    "Confiar cegamente no top-k sem re-rank",
    "Sem avaliação — 'funciona na demo'",
    "Sem fallback web quando a base local é insuficiente",
    "Adicionar mais docs sem re-rank (mais contexto ≠ melhor)",
    "Sem dataset de holdout — sem como detectar regressão",
    "Query rewriting automático sem validação"
], "📖 8 anti-patterns. Cada um é erro caro em produção.\n💡 Analogia: erros clássicos de culinária. Cada um arruína o prato.\n❓ 'Qual desses vocês já cometeram?' (confessar gera risos e aprendizado)\n➡️ Resumo.", 76, T)

s = content_slide("Resumo da Aula", [
    "RAG ingênuo falha silenciosamente — 4 tipos de falha",
    "Adaptive RAG: decide QUANDO recuperar",
    "CRAG: AVALIA docs antes de usar, com fallback web",
    "Self-RAG: reflete sobre docs E resposta (hallucination check)",
    "Agentic RAG: agente dirige todo o processo (multi-hop)",
    "Qualidade: chunking + re-rank + hybrid + query rewriting",
    "Avaliação: faithfulness, relevance, context precision/recall (Ragas)",
    "Escalada: Adaptive → CRAG → Self-RAG → Agentic (comece simples)"
], "📖 8 pontos-chave. Se saírem com esses, a aula cumpriu o objetivo.\n💡 Analogia: decálogo do RAG. 8 mandamentos para RAG de produção.\n➡️ Checklist.", 77, T)

s = content_slide("Checklist da Aula", [
    "[ ] Diagnosticou 4 tipos de falha do RAG ingênuo",
    "[ ] Explicou Adaptive RAG (decidir quando recuperar)",
    "[ ] Implementou CRAG (avaliar + 3 caminhos)",
    "[ ] Diferenciou Self-RAG (reflexão sobre resposta)",
    "[ ] Descreveu Agentic RAG (agente dirige, multi-hop)",
    "[ ] Listou 3+ técnicas de qualidade (chunking, re-rank, hybrid)",
    "[ ] Definiu 4 métricas de avaliação (faithfulness, relevance, precision, recall)"
], "📖 Vamos fazer juntos. Para cada item, levantem a mão se dominam. Se não estiver claro, falem agora.\n❓ 'Algum item que NÃO está claro?' (pausar — refazer se houver)\n➡️ Quiz.", 78, T)

exercise_slide(
    "Quiz — Pergunta 1",
    [
        "'Qual é a diferença fundamental entre Adaptive RAG e Self-RAG?'",
        "",
        "A) Adaptive usa vector DB, Self-RAG usa grafo",
        "B) Adaptive decide quando recuperar, Self-RAG também avalia a resposta",
        "C) Adaptive é mais caro que Self-RAG",
        "D) Não há diferença significativa",
        "",
        "Resposta: B"
    ],
    "📖 Quiz individual, sem consulta. Resposta: B. Adaptive decide quando; Self-RAG também avalia resposta.\n❓ Votar A/B/C/D. Anotar distribuição.\n➡️ Pergunta 2.",
    79, T, obj="Quiz"
)

exercise_slide(
    "Quiz — Pergunta 2",
    [
        "'Em CRAG, quando o sistema decide buscar na web?'",
        "",
        "A) Sempre, como primeiro passo",
        "B) Quando o avaliador classifica os docs como irrelevantes",
        "C) Quando a resposta é muito curta",
        "D) Quando o usuário pede explicitamente",
        "",
        "Resposta: B"
    ],
    "📖 Resposta: B. Quando avaliador classifica como irrelevantes. Web é fallback, não padrão.\n❓ Votação.\n➡️ Pergunta 3.",
    80, T, obj="Quiz"
)

exercise_slide(
    "Quiz — Pergunta 3",
    [
        "'Qual métrica mede se a resposta é fiel aos documentos recuperados?'",
        "",
        "A) Context precision",
        "B) Answer relevance",
        "C) Faithfulness",
        "D) Context recall",
        "",
        "Resposta: C"
    ],
    "📖 Resposta: C. Faithfulness decompõe em claims, verifica cada uma contra docs.\n❓ Votação. Anotar acertos — ≥2 = compreensão básica.\n➡️ Conexão com próximos módulos.",
    81, T, obj="Quiz"
)

s = content_slide("Conexão com Próximos Módulos", [
    "ETHAGT07 — Knowledge Graphs & Vector DBs (GraphRAG em profundidade)",
    "ETHAGT90 — Projeto final (aplicar Agentic RAG)",
    "Conexão com ETHAGT01: Agentic RAG = Augmented LLM com retrieval in-loop",
    "Conexão com ETHAGT05: memória do agente RAG (estado entre hops)",
    "Conexão com ETHAGT12: AgentOps — traces de agentes RAG"
], "📖 ETHAGT06 é pivô. ETHAGT07 aprofunda KG/Vector DBs. ETHAGT90 aplica Agentic RAG.\n➡️ Referências.", 82, T)

s = content_slide("Referências Completas", [
    "1. Lewis, P. et al. Retrieval-Augmented Generation. arXiv:2005.11401, 2020",
    "2. Asai, A. et al. Self-RAG. arXiv:2310.11511, 2023",
    "3. Yan, S. et al. Corrective RAG (CRAG). arXiv:2401.15884, 2024",
    "4. Edge, D. et al. GraphRAG: From Local to Global. arXiv:2404.16130, 2024",
    "5. Anthropic. Contextual Retrieval. 2024",
    "6. LangGraph examples: adaptive_rag, crag, self_rag, agentic_rag",
    "7. Cohere. Rerank Documentation",
    "8. Ragas: Evaluation framework for RAG"
], "📖 4 canônicas são leitura obrigatória. Anthropic Contextual Retrieval é o último grito. LangGraph = implementação de referência.\n➡️ Projeto.", 83, T)

s = content_slide("Projeto do Módulo + Labs", [
    "Projeto: sistema RAG de produção sobre corpus técnico",
    "  Pipeline agêntico (Adaptive/CRAG/Agentic)",
    "  Eval automatizado com Ragas",
    "  Entrega: sistema + eval report + ADR",
    "  Critério: faithfulness ≥ 0.85 e context recall ≥ 0.80",
    "",
    "Lab 1 (4h): 'Diagnosticando falhas' — RAG ingênuo em corpus problemático",
    "Lab 2 (5h): 'Agentic RAG multi-hop' — agente que refina queries e combina fontes"
], "📖 Projeto = consolidação. ~20h de trabalho. Lab 1 diagnostica; Lab 2 implementa Agentic.\n💡 Analogia: prova de piloto. Não só dirige — planeja rota, escolhe carro, relata.\n❓ 'Alguma dúvida sobre o projeto?' (esclarecer prazos)\n➡️ Q&A.", 84, T)

s = title_slide(
    "Perguntas?",
    "Contato do professor\nNa próxima aula: ETHAGT07 — Knowledge Graphs & Vector DBs\nLeitura: Asai et al. Self-RAG (arXiv:2310.11511) antes da próxima aula",
    "ETHAGT06"
)
add_notes(s, "📖 Chegamos ao fim. Perguntas? Se não houver, façam a pergunta inversa: 'qual parte foi menos clara?'\n❓ 'Perguntas? Ou qual parte foi menos clara?'\n⚠️ Sempre reserve 2-3 min para Q&A.\n➡️ Fim da aula.")

prs.save("ETHAGT06-RAG-Agentico.pptx")
print(f"Gerado: ETHAGT06-RAG-Agentico.pptx ({len(prs.slides)} slides)")
