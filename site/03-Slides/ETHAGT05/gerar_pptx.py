#!/usr/bin/env python3
"""Gerador de PPTX — ETHAGT05: Memoria de Agentes
Universidade Etho · Especializacao em Programacao Agentica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj=""):
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT05"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def table_slide(title, headers, rows, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    rows_count = len(rows) + 1
    cols_count = len(headers)
    table = s.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.5), Inches(12), Inches(5)).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercicio"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

T = 70

s = title_slide(
    "Memoria de Agentes",
    "Universidade Etho · Especializacao em Programacao Agentica\nFase B — Memoria, Contexto e Persistencia · 25 h",
    "ETHAGT05"
)
add_notes(s, "Bem-vindos. Esta e a aula de memoria de agentes — o modulo que transforma um agente amnesico em um agente com contexto acumulado, aprendizado e persistencia.\n💡 Analogia: a diferenca entre um hospede de hotel (esquece cada conversa) e um amigo de longa data (lembra de tudo).\n❓ 'Quantos ja ficaram frustrados com o ChatGPT esquecendo algo no meio de uma conversa longa?'\n⚠️ Alunos chegam achando que memoria e so vector DB. Nao — memoria e uma arquitetura com 4 camadas.\n➡️ Vamos aos objetivos.")

s = content_slide("Objetivos do Modulo", [
    "Objetivo geral: arquitetar sistemas de memoria que dao a agentes persistencia,",
    "  contexto acumulado e aprendizado — alem da context window",
    "",
    "Objetivos especificos:",
    "1. Distinguir tipos de memoria (working, episodica, semantica, procedural)",
    "2. Implementar persistencia via checkpointer (Postgres/SQLite/Redis)",
    "3. Gerenciar a janela de contexto: sumarizacao, eviction, janela deslizante",
    "4. Construir memoria vetorial para recall episodico",
    "5. Lidar com consistencia, privacidade e custo de memoria",
], "📖 Cada objetivo e mensuravel: 'distinguir', 'implementar', 'gerenciar', 'construir', 'lidar'.\n❓ 'Qual desses objetivos voces acham mais desafiador?' (costuma ser #3 ou #5)\n⚠️ Alunos confundem 'memoria' com 'vector DB'. O objetivo #1 e entender que ha 4 camadas.\n➡️ Vamos ver onde estamos no mapa da especializacao.", 2, T, "Objetivos")

s = table_slide("Competencias Desenvolvidas", ["Competencia", "Nivel", "Proximo nivel em"], [
    ["C1 Programacao Agentica", "A (Avancado)", "ETHAGT07, ETHAGT09"],
    ["C4 Agent Memory", "I (Intermediario)", "ETHAGT14, ETHAGT15"],
    ["C5 AgentOps & Avaliacao", "B (Basico)", "ETHAGT12"],
    ["C6 Agent Security", "B (Basico)", "ETHAGT13"],
], "📖 C1 atinge Avancado aqui. C4 — Agent Memory — atinge Intermediario, o foco da aula.\n⚠️ 'Basico' em C6 nao significa que seguranca nao importa — PII em memoria e critico.\n➡️ Agenda.", 3, T, "Competencias")

s = content_slide("Agenda da Aula", [
    "BLOCO 1 (45 min):",
    "  Abertura (8 min) — objetivos, motivacao, contexto",
    "  Tipos de Memoria (12 min) — as 4 camadas + MemGPT",
    "  Checkpointer (15 min) — estado persistente, backends, DEMO",
    "  Gerenciamento de Contexto (10 min) — custo, sumarizacao, eviction",
    "  Intervalo (5 min)",
    "",
    "BLOCO 2 (45 min):",
    "  Memoria Vetorial (12 min) — embedding, recall, re-ranking",
    "  Semantica e Grafos (8 min) — consolidacao, KG, Generative Agents",
    "  Producao (8 min) — consistencia, PII, esquecimento, custo",
    "  Fechamento (12 min) — boas praticas, quiz, projeto, Q&A",
], "📖 A aula tem dois blocos. DEMO ao vivo no Slide 28. Quiz final com 3 perguntas.\n⚠️ Alunos chegam atrasados e perdem a Seccao B (4 camadas) — e a fundacao.\n➡️ Motivacao.", 4, T, "Agenda")

s = comparison_slide("O Agente Amnesico",
    "Sem Memoria", [
        "Assistente: 'lembra do projeto X?' — nao sabe",
        "Suporte: pergunta a mesma coisa a cada ticket",
        "Coding agent: nao lembra convencoes do repo",
        "NPC: esquece que voce ja o ajudou ontem",
        "Nao aprende, nao melhora, nao personaliza",
    ],
    "Com Memoria", [
        "Lembra preferencias entre sessoes",
        "Acumula contexto de interacoes passadas",
        "Aprende com sucessos e falhas",
        "Personaliza por usuario/projeto",
        "Evolui ao longo do tempo",
    ],
    "📖 Agente sem memoria e tragicamente limitado. Otimo em uma conversa isolada, mas nao acumula nada.\n💡 Analogia: medico que esquece seu historico a cada visita — voce explica tudo de novo.\n❓ 'Quanto contexto voce acha que um assistente pessoal precisa reter?' (meses a anos)\n⚠️ Context window maior nao resolve — 200k tokens nao basta para 1 ano de uso diario.\n➡️ Por que memoria agora?", 5, T, "Motivacao")

s = content_slide("Por Que Memoria Agora", [
    "Linha do tempo:",
    "  2020: GPT-3 (2k tokens, sem tools, sem persistencia)",
    "  2022: 4k-8k tokens (GPT-3.5, ChatGPT)",
    "  2023: 16k-32k + tool calling + LangGraph com checkpointer",
    "  2024: 128k-200k + frameworks com persistencia nativa",
    "  2025: MemGPT, Zep, A-MEM, Letta (memoria como 1a classe)",
    "",
    "Confluencia de 4 fatores:",
    "  1. Context windows maiores (128k+)",
    "  2. Custo menor (GPT-4o-mini, Claude Haiku)",
    "  3. Frameworks com checkpointer (LangGraph)",
    "  4. Vector DBs acessiveis (Qdrant, Chroma, Pinecone)",
    "",
    "Marco: MemGPT (arXiv:2310.08560) — LLMs como sistemas operacionais",
], "📖 Memoria de agentes nao e ideia nova — sistemas especialistas dos anos 80 ja tinham 'memoria'. Mas context windows cresceram, custo baixou, frameworks trouxeram checkpointer.\n❓ 'Qual fator foi o gatilho mais recente?' (Resposta: frameworks com checkpointer)\n➡️ As 4 camadas de memoria.", 6, T, "Contexto")

s = section_slide(1, "Tipos de Memoria")
add_notes(s, "Inicio do bloco de fundamentos. Quais tipos de memoria um agente precisa? Sao quatro: working, episodica, semantica, procedural.\n➡️ Working memory — a context window.")

s = content_slide("Working Memory — A Context Window", [
    "Working memory = context window do LLM (tokens visiveis no prompt)",
    "Efemera: desaparece ao fim da sessao",
    "Limitada: 4k a 200k tokens (modelo-dependente)",
    "",
    "Estrategias:",
    "  Token budget (orcamento fixo por sessao)",
    "  Sliding window (manter ultimas N mensagens)",
    "  System prompt + ultimas N mensagens",
    "",
    "Custo: cada token no contexto e processado a cada chamada (input cost)",
], "📖 Context window e a memoria de trabalho — o que o LLM 've' no momento. Efemera e custa por uso.\n💡 Analogia: mesa de trabalho. So cabe certa quantidade de papeis. Se enche, voce arquiva (evict) ou resume.\n❓ 'Se um modelo tem 200k tokens, isso e muito ou pouco?' (1 sessao: muito. 1 ano de uso diario: quase nada)\n⚠️ Context window nao e 'memoria' — e working memory. Memoria de longo prazo e outra coisa.\n➡️ Limites da context window.", 8, T)

s = content_slide("Limites da Context Window", [
    "'Lost in the Middle' (arXiv:2307.03172):",
    "  LLMs ignoram informacao no MEIO do contexto",
    "  Recall accuracy em curva U: alta no inicio, baixa no meio, alta no fim",
    "",
    "Custo cresce com o tamanho do contexto (linear em API, mas significativo)",
    "Latencia aumenta com contexto maior",
    "200k tokens != 200k tokens uteis (ruido, redundancia, irrelevancia)",
    "",
    "Conclusao: context window e working memory, nao memoria de longo prazo",
], "📖 O paper 'Lost in the Middle' mostrou que LLMs perdem precisao no meio de contextos longos. Destroi a ilusao de que 'context window grande resolve tudo'.\n💡 Analogia: biblioteca sem catalogacao. 10 milhoes de livros 'la' mas voce se perde.\n❓ 'Ja notaram o ChatGPT esquecendo algo no meio de uma conversa longa?'\n⚠️ Mito: 'custo cresce quadraticamente'. Em API, custo e linear. O que cresce e latencia e degradacao de qualidade.\n➡️ Episodica.", 9, T)

s = content_slide("Memoria Episodica — Eventos com Timestamp", [
    "Episodica = registro de eventos com timestamp (o que aconteceu, quando)",
    "Exemplo: 'Usuario perguntou sobre React hooks em 12/mar' → recall por similaridade",
    "Armazenamento tipico: vector DB (embedding + metadata de tempo)",
    "Recall: 'o que aconteceu antes que se pareca com X?'",
    "Diferenca de log: episodica e recuperavel por SIGNIFICADO, nao so query estruturada",
    "Inspirado em: Tulving (1972) — memoria episodica cognitiva",
], "📖 Episodica registra eventos com timestamp. Recuperacao por similaridade semantica — busca por significado, nao por chave estruturada.\n💡 Analogia: diario pessoal. Voce le entradas similares, nao busca por 'ID 1234'.\n❓ 'Diferenca entre episodica e log?' (Episodica e fuzzy, recuperada por similaridade; log e exato, por SQL)\n⚠️ Episodica nao e log. Sao complementares, nao iguais.\n➡️ Semantica.", 10, T)

s = content_slide("Memoria Semantica — Fatos e Conhecimento", [
    "Semantica = fatos, conhecimento sobre o mundo/usuario (o que e verdade)",
    "Exemplo: 'Joao e desenvolvedor Python' (fato, nao evento)",
    "Armazenamento tipico: KB estruturada, knowledge graph, tabela de perfil",
    "Como surge: consolidacao da episodica (multiplos eventos → um fato)",
    "Diferenca de episodica: semantica nao tem 'quando', tem 'o que e'",
    "Inspirado em: Tulving (1972) — memoria semantica cognitiva",
], "📖 Semantica armazena fatos — 'Joao e dev Python'. Nao tem timestamp relevante. Surge por consolidacao de eventos.\n💡 Analogia: episodios de serie (episodica) vs biblia da serie (semantica).\n❓ \"'O usuario prefere portugues' — episodica ou semantica?\" (Semantica)\n⚠️ Regra: se tem 'quando', e episodica. Se e 'o que e verdade', e semantica.\n➡️ Procedural.", 11, T)

s = content_slide("Memoria Procedural — Skills Aprendidas", [
    "Procedural = sequencias de actions para objetivos recorrentes (skills)",
    "Exemplo: 'para deploy: testes → build → push → deploy' (aprendido, nao hardcoded)",
    "Armazenamento tipico: biblioteca de playbooks, prompt templates, tool sequences",
    "Como surge: reflexao sobre sucessos/falhas (estilo Reflexion, arXiv:2303.11366)",
    "Diferenca: procedural e sobre COMO AGIR, nao sobre o que aconteceu ou o que e verdade",
    "Inspirado em: Squire (memoria de habilidades); Reflexion (Shinn et al.)",
], "📖 Procedural armazena 'como fazer' — sequencias de actions. Surge por reflexao: tenta, falha, reflete, extrai licao.\n💡 Analogia: aprender a dirigir. Primeiro pensa em cada passo (working); depois fica automatico (procedural).\n⚠️ Procedural nao e 'tools'. Tools sao fixas; procedural e a sequencia aprendida de QUANDO e COMO usar tools.\n➡️ As 4 camadas integradas.", 12, T)

s = table_slide("As 4 Camadas de Memoria", ["Camada", "O que armazena", "Onde", "Como recupera", "Quando usar"], [
    ["Working", "Mensagens atuais", "Context window", "No prompt", "Sempre (base)"],
    ["Episodica", "Eventos (timestamp)", "Vector DB", "Similaridade + metadata", "Recall de contexto"],
    ["Semantica", "Fatos (KB/KG)", "KB / KG", "Query estruturada", "Verdade estavel"],
    ["Procedural", "Skills (playbooks)", "Playbook library", "Match de objetivo", "Tarefas repetidas"],
], "📖 As 4 camadas sao complementares, nao alternativas. Chatbot simples: so working. Assistente pessoal: working + episodica + semantica.\n💡 Regra: comece com working + checkpointer. Adicione as outras conforme necessidade.\n❓ 'Para um chatbot de FAQ, quais camadas?' (Working + semantica)\n⚠️ Over-engineering de memoria e anti-pattern. Comece simples.\n➡️ Diagrama integrado.", 13, T, "Comparacao")

s = content_slide("Arquitetura Integrada de Memoria", [
    "Diagrama: 12-Diagrams/ETHAGT05/memory-layers.mmd",
    "",
    "Agente no centro, conectado as 4 camadas:",
    "  Working (context window, volatil)",
    "  Episodica (vector DB, recall por similaridade)",
    "  Semantica (KB/KG, fatos e relacoes)",
    "  Procedural (skills, sequencias de actions)",
    "",
    "Checkpointer (Postgres/Redis) persiste o ESTADO do agente entre sessoes",
    "  (Nao e uma das 4 camadas — e infraestrutura de persistencia)",
], "📖 Slide-chave do bloco. Memorizem este diagrama. O agente acessa 4 camadas de memoria, cada uma com storage e padrao de recuperacao distintos.\n💡 Analogia: arquitetura de um cerebro digital. Working = cortex pre-frontal. Episodica = hipocampo. Semantica = cortex temporal. Procedural = ganglios da base.\n⚠️ Checkpointer NAO e uma camada de memoria. E infraestrutura que persiste o estado do grafo do agente.\n➡️ Inspiracao cognitiva — sem literalismo.", 14, T, "Diagrama")

s = content_slide("Inspiracao Cognitiva sem Literalismo", [
    "'A memoria humana e inspiracao, nao especificacao.'",
    "",
    "Working ↔ memoria de trabalho (Baddeley, 1974)",
    "Episodica/Semantica ↔ Tulving (1972)",
    "Procedural ↔ memoria de habilidades (Squire)",
    "",
    "MAS: LLMs nao consolidam durante o sono, nao tem hipocampo,",
    "     nao esquecem por razoes biologicas",
    "",
    "Aprendizado: use a taxonomia como FRAMEWORK DE DESIGN,",
    "             nao como modelo neural",
], "📖 A taxonomia vem da psicologia cognitiva. Util como framework, mas LLMs nao sao cerebros.\n💡 Analogia: asa do aviao inspirada na asa do passaro, mas nao e asa de passaro — nao bate, e fixa, tem turbines.\n❓ 'Inspiracao cognitiva: framework de design ou modelo neural?' (Framework de design)\n⚠️ Nao replique fenomenos cognitivos sem beneficio pratico.\n➡️ MemGPT — a melhor instanciacao dessa taxonomia.", 15, T)

s = content_slide("MemGPT — LLMs como Sistemas Operacionais", [
    "MemGPT (Packer et al., arXiv:2310.08560): LLM como SO",
    "",
    "Context window = RAM (limitada, rapida)",
    "Memoria persistente = disco (ilimitada, lenta)",
    "O modelo gerencia sua propria memoria: page-in / page-out",
    "Self-editing memory: o modelo atualiza suas proprias notas/perfil",
    "",
    "Inspirou: Zep, A-MEM, Letta",
    "",
    "Diagrama: Analogia SO (RAM ↔ Disco ↔ MMU) vs MemGPT (Context ↔ Persistent ↔ Self-editing)",
], "📖 MemGPT e o paper canonico de memoria de agentes. Trata o LLM como SO: context window como RAM, memoria persistente como disco, self-editing memory.\n💡 Analogia: gerenciar abas do navegador. Fecha as menos usadas (page-out), reabre quando precisa (page-in).\n❓ 'Se o modelo decide page-in/out, nao ha risco de decidir mal?' (Sim, problema ativo de pesquisa)\n⚠️ MemGPT nao e vector DB — e uma arquitetura de gestao de memoria.\n➡️ Checkpointer — persistencia entre sessoes.", 16, T)

s = section_slide(2, "Checkpointer e Estado Persistente")
add_notes(s, "Inicio do bloco de checkpointer. Por que precisamos de persistencia? Como o LangGraph modela? Quais backends? Resume/replay/branching. DEMO no final.\n➡️ O problema.")

s = content_slide("Estados Efemeros", [
    "Sem checkpointer: agente perde todo estado ao reiniciar",
    "",
    "Casos que quebram:",
    "  Servidor reinicia → conversa perdida",
    "  HITL demora dias → contexto evaporou",
    "  Debug: nao da para reproduzir execucao passada",
    "  A/B testing: nao da para voltar no tempo",
    "",
    "Solucao: serializar estado em storage duravel (a cada step)",
], "📖 Sem checkpointer, o agente e efemero. Perde tudo ao restart. Quebra casos reais: deploy, HITL longo, debug, A/B testing.\n💡 Analogia: videogame sem save — morre e volta ao inicio. Checkpointer e o 'save'.\n⚠️ Checkpointer nao e so 'salvar mensagens' — serializa o estado completo do grafo.\n➡️ Conceito do LangGraph.", 18, T)

s = content_slide("LangGraph Checkpointer — Conceito", [
    "Checkpointer = camada que serializa o estado do grafo a cada step",
    "Estado = qualquer objeto serializavel (TypedDict, Pydantic, dataclass)",
    "Cada node execution → checkpoint gravado",
    "",
    "thread_id identifica a conversa/sessao",
    "checkpoint_id identifica o momento (versionado)",
    "",
    "Profundidade em ETHAGT03 (StateGraph)",
], "📖 LangGraph serializa o estado a cada node execution. thread_id agrupa checkpoints de uma sessao; checkpoint_id identifica o momento.\n💡 Analogia: save de videogame com slots. thread_id = slot (cada jogador); checkpoint_id = momento do save.\n⚠️ Nao confunda thread_id com session_id. Use thread_id como chave de agrupamento.\n➡️ Estado serializavel.", 19, T)

s = code_slide("Estado Serializavel",
    """from typing import TypedDict, Annotated
from langgraph.graph import StateGraph
from langgraph.checkpoint.postgres import PostgresSaver
from langgraph.graph.message import add_messages

DB_URI = "postgresql://user:pass@localhost:5432/etho"

class AgentState(TypedDict):
    messages: Annotated[list, add_messages]
    intermediate_results: dict
    metadata: dict
    schema_version: int

graph = StateGraph(AgentState)
checkpointer = PostgresSaver.from_conn_string(DB_URI)
app = graph.compile(checkpointer=checkpointer)""",
    "📖 Estado deve ser JSON-serializable. Padrao: TypedDict com messages, intermediate_results, metadata, schema_version.\n💡 Analogia: salvar jogo. Salva posicao/inventario/HP (serializavel), nao conexao com servidor (efemero).\n⚠️ Nao coloque objetos nao-serializaveis no estado (engine DB, file handles). Se nao vira JSON, nao vai no estado.\n➡️ Backends.", 20, T, "Codigo")

s = content_slide("Backends de Checkpointer", [
    "Postgres: producao, multi-tenant, ACID, escalavel",
    "SQLite: desenvolvimento local, single-node, zero-config",
    "Redis: baixa latencia, TTL automatico, sessoes curtas",
    "",
    "Outros: MongoDB, DynamoDB (comunidade)",
    "",
    "Criterio: durabilidade vs latencia vs escala",
], "📖 Tres backends oficiais. Postgres (producao), SQLite (dev), Redis (cache de sessao).\n💡 Analogia: SQLite = bicicleta (simples, local). Postgres = carro (robusto, multiuso). Redis = foguete (rapido mas volatile).\n❓ 'Qual backend usaria para multi-tenant em producao?' (Postgres)\n➡️ Comparacao detalhada.", 21, T)

s = table_slide("Comparacao de Backends", ["Eixo", "Postgres", "SQLite", "Redis"], [
    ["Durabilidade", "Alta", "Alta (local)", "Media"],
    ["Latencia", "Media", "Baixa", "Baixissima"],
    ["Multi-tenant", "Sim", "Nao", "Sim"],
    ["ACID", "Sim", "Sim", "Parcial"],
    ["TTL nativo", "Nao", "Nao", "Sim"],
    ["Custo", "Medio", "Zero", "Medio"],
    ["Operacao", "Requer DBA", "Zero-config", "Requer ops"],
    ["Uso tipico", "Producao", "Desenvolvimento", "Cache de sessao"],
], "📖 Postgres ganha em durabilidade e multi-tenant (producao). SQLite em simplicidade (dev). Redis em latencia (cache).\n💡 Regra: comece com SQLite (dev) → Postgres (prod) → Redis (cache de sessao ativa).\n⚠️ Nao use Redis como fonte de verdade — persistencia e limitada.\n➡️ Resume.", 22, T, "Comparacao")

s = code_slide("Resume — Retomando Execucao",
    """# Sessao 1: agente inicia tarefa
config = {"configurable": {"thread_id": "demo-001"}}
result = app.invoke(
    {"messages": [HumanMessage("Analise o relatorio Q1")]},
    config=config
)

# ... tempo passa (dias) ...

# Sessao 2: retoma com mesmo thread_id
result = app.invoke(None, config=config)

# O agente retoma EXATAMENTE de onde parou
# (mensagens, contexto, steps preservados)""",
    "📖 Resume carrega o ultimo checkpoint do thread_id e continua. O agente nao 'sabe' que foi pausado.\n💡 Analogia: pausar filme. Voce para no minuto 45, dorme, e no dia seguinte continua do 45 — nao do inicio.\n⚠️ Pegadinha: se o modelo mudou (GPT-4-turbo → GPT-4o), comportamento pode diferir. Se schema mudou, pode quebrar.\n➡️ Diagrama do resume.", 23, T, "Codigo")

s = content_slide("Checkpointer Resume — Pause e Retomada", [
    "Diagrama: 12-Diagrams/ETHAGT05/checkpointer-resume.mmd",
    "",
    "Sessao 1: agente executa → estado serializado → Checkpointer (Postgres)",
    "Sessao 2 (dias depois): thread_id → carrega checkpoint → continua",
    "Sessao 3 (semana depois): mesmo thread_id → mesma logica",
    "HITL: agente interrompe → humano aprova → resume",
    "",
    "Tudo transparente para o agente",
], "📖 O fluxo: sessao 1 grava, sessao 2 carrega com mesmo thread_id, e o agente continua.\n💡 Analogia: serie da Netflix — 'continuar de onde parou?' exatamente no minuto.\n➡️ Replay.", 24, T, "Diagrama")

s = content_slide("Replay — Debugando o Passado", [
    "Replay: carregar checkpoint de um momento especifico e RE-EXECUTAR",
    "",
    "Casos: 'por que o agente tomou decisao X?' → replay do checkpoint anterior",
    "Diferenca de resume: replay re-executa, resume continua",
    "Util para: debugging, A/B testing de prompts, regressao",
    "",
    "Ferramentas: LangSmith + checkpointer = replay visual",
], "📖 Replay e diferente de resume. Resume continua. Replay re-executa do passado. Essencial para debug e A/B testing.\n💡 Analogia: 'instant replay' nos esportes — voce volta no tempo e re-ve.\n⚠️ Nao confunda: resume = continuar; replay = re-executar.\n➡️ Branching.", 25, T)

s = content_slide("Branching — Time Travel", [
    "Branching: partir de um checkpoint e seguir caminho diferente",
    "Casos: 'e se o agente tivesse usado tool Y em vez de X?'",
    "Como: fork do checkpoint_id → nova execucao",
    "Analogo a: git checkout <commit> + git checkout -b <novo-branch>",
    "",
    "Pergunta: Time travel — ferramenta de debug ou risco de seguranca?",
], "📖 Branching e o time travel do checkpointer. Pega um checkpoint antigo, faz fork, executa diferente. Como save scumming em jogos.\n❓ 'Time travel e debug ou risco?' (Ambos — em dev e ferramenta; em prod, precisa de controle de acesso)\n⚠️ Branching nao desfaz a execucao original — cria um RAMO novo, como no git.\n➡️ Versionamento de schema.", 26, T)

s = content_slide("Versionamento de Schema de Estado", [
    "Problema: estado v1 (sem 'priority') vs estado v2 (com 'priority')",
    "Checkpoints antigos tem schema diferente → quebra ao retomar",
    "",
    "Estrategias:",
    "  Migracao lazy: converter no load",
    "  Version field: schema_version: 2",
    "  Backward compatibility: campos novos com default",
    "",
    "Analogia: migracao de DB, mas para estado de agente",
], "📖 O schema do estado evolui. Sem versionamento, checkpoints antigos quebram ao retomar. Com schema_version + migracao lazy, converte no load.\n💡 Analogia: atualizar app no celular. App v2 espera campo que v1 nao salvava. Migracao automatica resolve.\n⚠️ Versione o schema desde o dia 1: schema_version: 1 no primeiro commit.\n➡️ DEMO!", 27, T)

s = code_slide("DEMO — Checkpointer em Postgres",
    """# Passo 1: agente inicia tarefa
$ python main.py --thread-id demo-001
[INFO] Agente iniciado. thread_id=demo-001
[STEP 0] Thought: vou analisar o relatorio
[STEP 1] Action: read_file("Q1.pdf")
[STEP 2] Observation: 42 paginas carregadas
^C  # <-- Ctrl+C! Processo morto!

# ... dias depois ...

# Passo 2: retoma com mesmo thread_id
$ python main.py --thread-id demo-001
[INFO] Checkpoint carregado. Continuando do STEP 2.
[STEP 3] Thought: agora vou extrair KPIs...
# Estado preservado! Como se nunca tivesse parado.""",
    "DEMO AO VIVO.\n📖 Vou iniciar agente com thread_id=demo-001. Ele comeca tarefa. No meio, Ctrl+C — processo morre. Depois reinicio com mesmo thread_id e ele RETOMA exatamente de onde parou.\n💡 Analogia: truque de magica de 'serrar a pessoa ao meio'. Mato o agente (Ctrl+C) e ele revive intacto. A magica e o checkpointer.\n⚠️ Se API ou Postgres cair, tenho screenshot pre-gravado. Plano B.\n➡️ Pergunta da demo.", 28, T, "DEMO")

s = exercise_slide("Pergunta da DEMO", [
    "1. O que acontece se o modelo foi atualizado entre sessao 1 e sessao 2?",
    "",
    "2. E se o schema do estado mudou?",
    "",
    "3. Como voce testaria que o resume funciona corretamente?",
    "",
    "Discutam em duplas (2 min)",
], "📖 Discutam em duplas, 2 min. As 3 perguntas sao para pensar em edge cases do resume.\n💡 Analogia: 'o que acontece se trocar o motor do carro no meio de uma viagem?' Provavelmente funciona, mas precisa verificar.\n➡️ Respostas: (1) provavelmente funciona, comportamento pode diferir; (2) sem versionamento quebra, com migracao lazy funciona; (3) rodar com e sem interrupcao, comparar outputs.\n➡️ Gestao de contexto.", 29, T)

s = section_slide(3, "Gerenciamento de Contexto")
add_notes(s, "Inicio da gestao de contexto. Checkpointer persiste entre sessoes, mas DENTRO de uma sessao, a context window e finita. Custo, estrategias: janela deslizante, sumarizacao, eviction, entity-centric.\n➡️ O problema.")

s = content_slide("Context Window e Finita", [
    "Mesmo com 200k tokens, a context window ENCHE",
    "Conversas longas + tools + retrieval → estouro de contexto",
    "",
    "Sintomas: agente 'esquece' inicio, respostas degradam, erro 400",
    "Tensao: mais contexto = mais informacao, mas tambem mais custo, latencia e ruido",
    "",
    "Solucao: estrategias ATIVAS de gerenciamento (nao apenas 'encher ate estourar')",
], "📖 A context window e finita. Mesmo 200k tokens enchem em conversa longa com tools e retrieval.\n💡 Analogia: mala de viagem. Pode encher ate fechar a forca, mas fica pesada e desorganizada. Melhor curar o que coloca.\n⚠️ Context window grande nao elimina a necessidade de gestao.\n➡️ Mito do custo quadratico.", 31, T)

s = content_slide("Custo — Quadratico ou Linear?", [
    "MITO: 'attention e O(n²) → custo quadraticamente'",
    "",
    "REALIDADE: attention e O(n²) em COMPUTACAO, mas:",
    "  KV cache: tokens ja processados nao sao recomputados",
    "  Sliding window attention: so atende janela local",
    "  Custo de API e LINEAR em tokens (preco por token)",
    "",
    "MAS: latencia sim cresce, e qualidade degrada com contexto longo",
    "",
    "Conclusao: o problema nao e so custo — e QUALIDADE e LATENCIA",
], "📖 Desmistificando. Computacao attention e O(n²), mas custo de API que voce paga e LINEAR. O que cresce de verdade e latencia e degradacao de qualidade (Lost in the Middle).\n💡 Analogia: dirigir estrada longa. Combustivel cresce linear, mas fadiga e chance de acidente crescem mais.\n❓ 'Se o custo e linear, por que preocupar?' (Latencia e qualidade — 200k tokens nao dao 200x melhor resposta)\n⚠️ Nao repita 'custo quadratico' sem entender.\n➡️ Estrategia 1.", 32, T)

s = content_slide("Estrategia 1 — Janela Deslizante", [
    "Manter apenas as ultimas N mensagens no contexto",
    "Simples, previsivel, barato",
    "",
    "Mas: perde contexto antigo ('do que estavamos falando?')",
    "Variacao: system prompt + ultimas N mensagens (preserva instrucao)",
    "",
    "Quando usar: chatbots simples, conversas curtas, nao criticas",
], "📖 Janela deslizante: mantem so as ultimas N mensagens. Simples e barato, mas perde contexto antigo.\n💡 Analogia: conversa em bar barulhento. So lembra dos ultimos 5 minutos.\n⚠️ Nao use para casos que precisam de memoria de longo prazo. Use checkpointer + memoria persistente.\n➡️ Sumarizacao.", 33, T)

s = content_slide("Estrategia 2 — Sumarizacao em Cascata", [
    "Quando contexto enche: sumarizar mensagens antigas → substituir por sumario",
    "Sumarizacao em cascata: sumario de sumarios (multiplos niveis)",
    "  L1: sumario de 50 mensagens em 2k tokens",
    "  L2: sumario de 5 sumarios L1 em 500 tokens",
    "  L3: sumario de tudo em 100 tokens",
    "",
    "Cada nivel: menos detalhe, mais abstracao",
    "Pegadinha: sumario pode perder informacao critica",
    "Quando usar: conversas longas, agentes de suporte, reunioes",
], "📖 Sumarizacao em cascata comprime mensagens antigas em sumarios multi-nivel. Cada nivel perde detalhe mas ganha abstracao.\n💡 Analogia: ler livro resumido. Original 500p, resumo 50p, resumo do resumo 5p, sinopse 1 paragrafo.\n⚠️ Sumario e lossy. Mitigue: mantenha eventos originais no vector DB com fallback para nivel inferior.\n➡️ Eviction por relevancia.", 34, T)

s = content_slide("Estrategia 3 — Eviction por Relevancia", [
    "Nem tudo no contexto e igualmente relevante",
    "Eviction por relevancia: score = f(relevancia, idade, frequencia de acesso)",
    "",
    "Manter: system prompt, ultimas mensagens, entidades ativas, fatos criticos",
    "Evitar: mensagens antigas de baixa importancia, tool outputs verbosos, redundancias",
    "",
    "Analogos a: LRU cache, mas com score semantico",
], "📖 Em vez de descartar por tempo ou comprimir tudo, calcule um SCORE por item. Score alto fica; baixo e evictado.\n💡 Analogia: limpar email. Nao apaga por data nem arquiva tudo. Decide: importante (manter), antigo irrelevante (arquivar), spam (apagar).\n⚠️ Recencia e um fator, mas importancia e frequencia tambem. Fato critico de 6 meses pode ser mais importante que mensagem de 5 min.\n➡️ Fluxo de eviction.", 35, T)

s = content_slide("Eviction Flow — Decisao de Retencao", [
    "Diagrama: 12-Diagrams/ETHAGT05/eviction-flow.mmd",
    "",
    "Novo evento → score = relevance × decay(age)",
    "score > threshold? → manter na memoria",
    "Entidade critica? (ex.: pagamento) → decay lento, manter mais tempo",
    "Recall frequente? → aumenta score (reforco)",
    "Caso contrario → arquivar / apagar",
], "📖 Cada evento recebe score. Relevance × decay(age). Threshold decide. Entidades criticas tem decay lento. Recall frequente aumenta score (reforco).\n💡 Analogia: algoritmo do Facebook. Relevancia + recencia + frequencia decidem o que fica no feed.\n➡️ Entity-centric memory.", 36, T, "Diagrama")

s = content_slide("Estrategia 4 — Entity-Centric Memory", [
    "Em vez de lista de mensagens, organizar memoria por ENTIDADE",
    "Cada entidade (usuario, projeto, tarefa) tem seu proprio 'perfil' de memoria",
    "O modelo decide quando page-in (carregar) e page-out (descarregar) entidades",
    "",
    "MemGPT: self-editing memory — o modelo atualiza suas proprias notas",
    "Zep: memoria de longo prazo com graph + vector",
    "",
    "Vantagem: contexto sempre relevante a entidade ativa",
], "📖 Entity-centric organiza memoria por entidade. O modelo faz page-in/out. Como MemGPT e Zep.\n💡 Analogia: organizar arquivos por cliente em vez de por data. Abre pasta do cliente X (page-in); fecha do Y (page-out).\n⚠️ Entity-centric nao e vector DB — e uma forma de ORGANIZAR. Vector DB pode ser o backend.\n➡️ Vector DB nem sempre e a resposta.", 37, T)

s = comparison_slide("Quando Vector DB e Pior que Relacional",
    "Vector DB e BOM para", [
        "Recall por significado",
        "Fuzzy matching ('mais ou menos')",
        "Busca semantica",
        "'O que se parece com X?'",
    ],
    "Vector DB e PESSIMO para", [
        "Exact match ('saldo do usuario X')",
        "Ordenacao temporal ('eventos do mes Y')",
        "Agregacao ('total de X por categoria')",
        "Joins ('usuario X do projeto Y')",
    ],
    "📖 Vector DB e ferramenta, nao solucao universal. Bom para similaridade, pessimo para exact match, ordenacao, agregacao, joins.\n💡 Analogia: martelo e otimo para pregar, pessimo para aparafusar. Nao use martelo para parafuso.\n❓ 'Citem um caso onde vector DB seria a escolha ERRADA.' (Saldo bancario, lista de usuarios, agregacoes)\n⚠️ Regra: vector para recall semantico, relacional para fatos estruturados, KG para relacoes.\n➡️ Exercicio de eviction.", 38, T, "Comparacao")

s = exercise_slide("Exercicio — Politica de Eviction", [
    "Cenario: assistente pessoal com 1 ano de interacoes (~365 conversas)",
    "",
    "Em trios, escrevam uma politica de eviction combinando relevancia e idade.",
    "Template:",
    "  score = ___ × relevancia + ___ × recencia + ___ × frequencia",
    "  threshold = ___",
    "  entidades_criticas = [___, ___, ___]",
    "  TTL_PII = ___ dias",
    "",
    "Considerar: pagamentos/prazos, frequencia de recall, PII",
    "3 min discussao + 1 min compartilhar",
], "📖 Hora de praticar. Em trios, 3 min. Politica de eviction combinando relevancia e idade.\n💡 Analogia: geladeira. Conservas duram semanas, verduras dias, comida quente horas. Politica por tipo.\n❓ 'Quais entidades classificariam como criticas?' (Pagamentos, contratos, prazos, saude)\n⚠️ Nao esquecam PII — TTL maximo por compliance.\n➡️ Transicao.", 39, T, "Exercicio")

s = section_slide("→", "De Contexto para Recall")
add_notes(s, "Gestao de contexto resolve o CURTO PRAZO. E o longo prazo? Memoria vetorial.\n➡️ Seccao E.")

s = section_slide(4, "Memoria Vetorial para Recall")
add_notes(s, "Inicio do bloco de memoria vetorial. Como converter eventos em embeddings? Como recuperar por similaridade? Metadata filtering? Re-ranking? Pipeline completo. Lab 2 no final.\n➡️ Embedding de eventos.")

s = content_slide("Embedding de Eventos", [
    "Cada evento (mensagem, acao, observacao) → texto → embedding",
    "Embedding = vetor de N dimensoes (384, 768, 1536, 3072)",
    "",
    "Modelos: OpenAI text-embedding-3-small/large, Cohere, BGE, E5",
    "Armazenar: vector DB (Qdrant, Chroma, Pinecone, pgvector)",
    "",
    "Metadata: timestamp, user_id, session_id, type, entities",
    "Pipeline: embed_event(event) → (vector, metadata)",
], "📖 Primeiro passo: converter cada evento em embedding (vetor de N dimensoes). Modelos: OpenAI, BGE, E5. Armazene vetor + metadata.\n💡 Analogia: catalogar livro. Embedding = codigo de barras semantico; metadata = ficha catalografica (autor, data, genero).\n⚠️ Nao esqueca metadata. Sem metadata, so da para vector search puro — nao filtra por usuario/sessao/tempo.\n➡️ Recall por similaridade.", 42, T)

s = content_slide("Recall por Similaridade Semantica", [
    "Query: converter pergunta atual em embedding",
    "Buscar: top-K eventos mais similares (cosine similarity)",
    "Reinserir: colocar eventos recuperados no contexto do agente",
    "",
    "Exemplo: 'como resolveu aquele bug de autenticacao?'",
    "  → recall de evento passado similar",
    "",
    "Pos-recuperacao: nem sempre o mais similar e o mais util → re-ranking",
], "📖 Para recuperar, converta a query em embedding e busque top-K por cosine similarity. Insira no contexto.\n💡 Analogia: buscar no Google. Digita query, recupera milhoes 'similares', ranqueia top-10.\n❓ 'Se usuario pergunta lembra do deploy que falhou ontem, o que recuperamos?' (Eventos com embedding similar a 'deploy falhou', filtrados por timestamp=ontem)\n➡️ Metadata filtering.", 43, T)

s = content_slide("Metadata Filtering", [
    "Vector search sozinho e insuficiente: 'aquela reuniao de terca'",
    "Metadata filter: session_id = X AND timestamp > Y AND type = 'meeting'",
    "Combinacao: filter por metadata + rank por similaridade",
    "",
    "Padrao: pre-filter (filtra antes da busca) vs post-filter (filtra depois)",
    "Vector DBs modernas suportam hybrid search (vector + structured)",
], "📖 Vector search puro nao resolve tudo. Combine metadata filter + vector search. Pre-filter e mais eficiente que post-filter.\n💡 Analogia: biblioteca. Voce nao busca em milhoes de livros por similaridade. Vai a secao 'Ficcao Cientifica' (filter) depois busca por similaridade.\n⚠️ Nao faca vector search puro e filtre no codigo. Ineficiente. Deixe o vector DB fazer o filter.\n➡️ Re-ranking.", 44, T)

s = content_slide("Pos-Recuperacao — Re-ranking", [
    "Vector search e rapido mas IMPRECISO (aproximacao ANN)",
    "Re-ranking: modelo separado re-avuala os top-K resultados",
    "",
    "Modelos: Cohere Rerank, BGE Reranker, cross-encoder",
    "Pipeline: vector search (top 20) → re-ranker (top 5) → contexto",
    "",
    "Trade-off: re-ranking adiciona latencia mas melhora precisao",
], "📖 Vector search usa approximate nearest neighbors — rapido mas impreciso. Re-ranker (cross-encoder) refina top-K para top-N mais preciso.\n💡 Analogia: contratar. Triagem rapida (vector: top 20 curriculos), depois entrevistas (re-ranker: top 5).\n⚠️ Para producao, re-ranking melhora muito a qualidade. Latencia extra (~200ms) e justificavel.\n➡️ Pipeline completo.", 45, T)

s = content_slide("Pipeline Completo de Recall Episodico", [
    "1. Query atual → embedding",
    "2. Metadata filter (user, session, time range)",
    "3. Vector search (top-K candidatos)",
    "4. Re-ranking (top-N final)",
    "5. Insercao no contexto do agente",
    "6. Agente responde com memoria recuperada",
    "",
    "Este pipeline e a base do Lab 2 (Memoria Episodica)",
], "📖 Slide-chave da Seccao E. Pipeline completo em 6 etapas. Memorizem — e o que voces vao implementar no Lab 2.\n💡 Analogia: investigacao policial. Pista (query) → filtra suspeitos (metadata) → busca por perfil (vector) → re-avalia (re-rank) → apresenta evidencias (contexto) → detetive conclui (responde).\n⚠️ Cada etapa tem proposito. Nao pule nenhuma.\n➡️ Comparacao de vector DBs.", 46, T, "Pipeline")

s = content_slide("Comparacao de Vector DBs", [
    "Qdrant: Rust, open-source, self-hosted ou cloud, alta performance",
    "Chroma: Python, easy setup, ideal para prototipos",
    "Pinecone: managed, serverless, escala automatica",
    "pgvector: extensao PostgreSQL, unifica relacional + vetorial",
    "",
    "Criterio: managed vs self-hosted, escala, custo, latencia, hybrid search",
], "📖 Quatro opcoes. Qdrant (controle/performance), Chroma (prototipo), Pinecone (managed/zero-ops), pgvector (unifica com Postgres).\n💡 Analogia: Chroma = hostal (simples), Qdrant = Airbnb (controle), Pinecone = hotel 5 estrelas (managed), pgvector = morar com os pais (reutiliza).\n❓ 'Se ja tem Postgres, qual escolheria?' (pgvector — elimina componente)\n➡️ Lab 2.", 47, T)

s = content_slide("Lab 2 — Memoria Episodica", [
    "Lab 2 (5h): 'Memoria episodica'",
    "Agente recorda interacoes anteriores relevantes via recall vetorial",
    "Comparar agente COM e SEM memoria em sessoes espacadas",
    "",
    "Stack: LangGraph + Qdrant + OpenAI embeddings",
    "Referencia: 05-Labs/ETHAGT05/Lab2-Memoria-Episodica",
    "",
    "Metricas: hit rate do recall, latencia, custo extra por interacao",
], "📖 No Lab 2, voces implementam o pipeline completo de recall. Stack: LangGraph + Qdrant + OpenAI. Comparacao A/B: com vs sem memoria.\n💡 Analogia: medico com prontuario (com memoria) vs sem prontuario (sem memoria). Qualidade drasticamente diferente.\n➡️ Exercicio de escolha.", 48, T, "Lab 2")

s = exercise_slide("Embedding ou Metadata?", [
    "1. 'Aquela reuniao de terca' — embedding ou metadata?",
    "   Resposta: AMBOS! Metadata (terca) + vector (reuniao)",
    "",
    "2. 'Como configurar deploy?' — embedding ou metadata?",
    "   Resposta: Embedding (semantico) + filter por type='procedural'",
    "",
    "3. 'Qual o saldo da conta?' — embedding ou metadata?",
    "   Resposta: NENHUM! KB relacional, nao vector DB",
    "",
    "Discusao aberta (2 min)",
], "📖 Practicar a escolha. (1) Ambos — metadata temporal + vector semantico. (2) Embedding + filter procedural. (3) Nenhum — relacional.\n💡 Analogia: escolher ferramenta de busca. Semantica para 'encontre paragrafo sobre X', SQL para 'liste itens de ontem', lookup para 'saldo atual'.\n⚠️ Vector search nao e para tudo. Saiba quando usar filter, vector, ambos, ou nenhum.\n➡️ Transicao para semantica.", 49, T, "Exercicio")

s = section_slide("→", "De Episodica para Semantica")
add_notes(s, "Eventos acumulados viram conhecimento? Consolidação.\n➡️ Seccao F.")

s = section_slide(5, "Memoria Semantica e Grafos")
add_notes(s, "Inicio de semantica e grafos. Consolidacao (episodica → semantica), criterios de promocao, knowledge graph, Generative Agents.\n➡️ Consolidacao.")

s = content_slide("Consolidacao — Episodica → Semantica", [
    "Consolidacao: processo de extrair fatos de eventos acumulados",
    "Analogia biologica: consolidacao da memoria durante o sono",
    "",
    "Em agentes: processo offline ou em batch que:",
    "  1. Agrupa eventos relacionados",
    "  2. Extrai fatos consolidados",
    "  3. Remove redundancia",
    "  4. Promove para KB semantica",
    "",
    "Exemplo: 5 eventos 'Joao perguntou sobre Python'",
    "  → fato 'Joao e desenvolvedor Python'",
], "📖 Consolidacao transforma eventos em fatos. Processo offline/batch: agrupa, extrai, remove redundancia, promove.\n💡 Analogia: editor le 100 reportagens sobre tema, escreve 1 artigo de sintese. 100 reportagens = eventos; artigo = fato.\n⚠️ Nao consolide imediatamente. Espere massa critica (3-5 ocorrencias).\n➡️ Quando promover?", 52, T)

s = content_slide("Quando Promover uma Memoria?", [
    "Nem todo evento vira fato — criterios de promocao:",
    "",
    "  Frequencia: evento recorrente (multiplas ocorrencias)",
    "  Confianca: fonte confiavel ou confirmado multiplas vezes",
    "  Estabilidade: fato nao muda frequentemente",
    "  Utilidade: fato sera relevante em sessoes futuras",
    "",
    "Quando NAO promover:",
    "  Eventos unicos sem valor futuro",
    "  Fatos volateis (saldo, status de tarefa)",
    "  Informacao sensivel (PII desnecessaria)",
], "📖 Criterios de promocao: frequencia, confianca, estabilidade, utilidade. Nao promova volateis (saldo) ou PII desnecessaria.\n💡 Analogia: criterios da Wikipedia. Notoriedade, fontes confiaveis, estabilidade.\n⚠️ Nao promova fatos volateis. Saldo muda a cada transacao — mantenha como lookup relacional.\n➡️ Knowledge graph.", 53, T)

s = content_slide("Knowledge Graph como Memoria", [
    "KG = grafo de entidades + relacoes + atributos",
    "Exemplo: (Joao) --[trabalha_em]--> (Projeto X) --[usa]--> (Python)",
    "",
    "Vantagens:",
    "  Consulta estruturada",
    "  Raciocinio multi-hop",
    "  Explicitacao de relacoes",
    "",
    "Integracao com LLM: GraphRAG, Cypher + LLM, text-to-graph",
    "Profundidade em ETHAGT07 — Knowledge Graphs para Agentes",
], "📖 KG e poderoso para memoria semantica quando fatos tem RELACOES. Triplas (entidade)--[relacao]-->(entidade).\n💡 Analogia: organograma de empresa. Mostra quem reporta para quem, quem esta em qual time. Permite 'quem pode cobrir ferias do Joao?' percorrendo o grafo.\n⚠️ KG e para fatos com relacoes estruturadas. Se fatos sao isolados, tabela basta.\n➡️ Exemplo concreto.", 54, T)

s = content_slide("Exemplo — De Evento para Fato", [
    "Eventos episodicos (vector DB):",
    "  'Joao perguntou sobre decorators em Python' (12/jan)",
    "  'Joao pediu ajuda com asyncio' (25/jan)",
    "  'Joao compartilhou repositorio de FastAPI' (03/fev)",
    "",
    "Consolidacao → Fato semantico (KB):",
    "  'Joao e desenvolvedor Python (intermediario-avancado)'",
    "  Interesses: decorators, asyncio, FastAPI",
    "",
    "Acao: agente agora SABE quem e Joao sem recall de cada evento",
], "📖 3 eventos consolidados em 1 perfil semantico. Agente agora 'sabe' quem e Joao sem recall.\n💡 Analogia: RH le 3 avaliacoes, escreve 1 perfil sintetico. Para decisoes futuras, consulta perfil, nao as 3 avaliacoes.\n➡️ Generative Agents.", 55, T)

s = content_slide("Generative Agents — Memoria de Smallville", [
    "Park et al. (arXiv:2304.03442): 25 agentes em Smallville",
    "",
    "Cada agente tem:",
    "  Memory stream: log de observacoes com timestamp",
    "  Retrieval: por recencia + importancia + relevancia",
    "  Reflection: consolida eventos em insights de alto nivel",
    "",
    "Resultado: agentes 'lembram', planejam, socializam de forma emergente",
    "Licao: a arquitetura de memoria E a identidade do agente",
], "📖 Generative Agents e o caso canonico. Memory stream + retrieval scoreado (recencia + importancia + relevancia) + reflection. Agentes lembram e socializam de forma emergente.\n💡 Analogia: gemios com experiencias diferentes se tornam pessoas diferentes. Memoria = identidade.\n❓ 'A arquitetura de memoria e a identidade do agente? Concordam?' (Sim — memoria molda decisoes, preferencias, comportamento)\n⚠️ Generative Agents nao e so vector DB — e arquitetura completa (stream + score + reflection).\n➡️ Preview ETHAGT07.", 56, T)

s = content_slide("Integracao com Knowledge Graph (Preview ETHAGT07)", [
    "ETHAGT07 aprofunda: GraphRAG, ontologias, raciocinio multi-hop",
    "Memoria semantica + KG = raciocinio estruturado de longo prazo",
    "",
    "Exemplo: 'Quem trabalha com Python e esta livre na semana 20?'",
    "  → query multi-hop no KG",
    "",
    "Padrao: vector DB para recall + KG para raciocinio + relacional para transacional",
], "📖 ETHAGT07 aprofunda KG. Memoria semantica + KG = raciocinio estruturado. Padrao em producao: vector (recall) + KG (raciocinio) + relacional (transacional). Tres camadas, tres propositos.\n➡️ Producao.", 57, T)

s = section_slide(6, "Producao: Consistencia, Privacidade, Custo")
add_notes(s, "Inicio de producao. Consistencia multi-agente, PII, direito ao esquecimento, custo vs beneficio, observabilidade.\n➡️ Consistencia.")

s = content_slide("Consistencia em Multi-Agente", [
    "Multiplos agentes compartilham memoria → problemas de concorrencia",
    "Race condition: A le fato, B atualiza, A usa fato desatualizado",
    "",
    "Estrategias:",
    "  Eventual consistency: aceitar defasagem (mais simples)",
    "  Strong consistency: locks / transacoes (mais caro)",
    "  Event sourcing: log de mudancas (auditavel)",
    "",
    "Padrao: memoria local + memoria compartilhada via message bus",
], "📖 Em multi-agente, memoria compartilhada traz race conditions. Estrategias: eventual (simples), strong (caro), event sourcing (auditavel).\n💡 Analogia: 3 editores no mesmo Google Doc. Sem sync, um sobrescreve o outro.\n⚠️ Concorrencia SEMPRE traz race conditions. Planeje consistencia desde o design.\n➡️ PII.", 59, T)

s = content_slide("PII em Memoria — Redacao e Retencao", [
    "Memoria acumula PII: nomes, emails, preferencias, historico",
    "Riscos: vazamento, profiling, inferencia indevida",
    "",
    "Redacao: mascarar PII antes de armazenar (NER + replacement)",
    "  Ex.: 'Joao Silva (CPF 123...)' → 'USUARIO_42 (CPF_REDACTED)'",
    "",
    "Retencao: TTL por tipo de dado",
    "  Conversas: 90 dias | Perfil: ate consentimento",
    "",
    "Compliance: LGPD/GDPR exigem minimizacao de dados",
], "📖 Memoria acumula PII. Redija com NER antes de armazenar. Defina TTL por tipo. LGPD/GDPR exigem minimizacao.\n💡 Analogia: medico guardando prontuarios. Anota so relevante, regras de retencao, protege com senha.\n⚠️ Redija PII ANTES de armazenar. Se nao precisa do CPF explicito, armazene hash ou 'USER_X'.\n➡️ Direito ao esquecimento.", 60, T)

s = content_slide("Direito ao Esquecimento em Vector DB", [
    "Desafio: 'esquecer tudo sobre o usuario X' em vector DB",
    "Problema: embeddings nao sao reverasiveis",
    "",
    "Estrategias:",
    "  Delete por metadata: WHERE user_id = X (se metadata armazenada)",
    "  Re-embed: remover evento e re-gerar sumarios consolidados",
    "  Cryptographic erasure: encriptar por usuario, destruir chave",
    "",
    "Pergunta: Como implementar direito ao esquecimento em memoria vetorial?",
], "📖 Direito ao esquecimento e desafio em vector DB. 3 estrategias: delete por metadata, re-embed, cryptographic erasure. Nenhuma perfeita; combine.\n💡 Analogia: 'des-fazer' um bolo. Nao da para tirar so o acucar. Opcoes: jogar bolo fora (delete), fazer novo sem acucar (re-embed), ter assado em formas criptografadas separadas (cryptographic erasure).\n⚠️ Apagar do vector DB nao resolve se info foi consolidada em sumarios. Re-embed ou cryptographic erasure necessarios.\n➡️ Custo vs beneficio.", 61, T)

s = content_slide("Custo de Memoria vs Beneficio", [
    "Memoria nao e de gratis: storage, embedding, recall, contexto adicional",
    "",
    "Pergunta-chave: 'esta memoria vai melhorar uma decisao futura?'",
    "",
    "Quando NAO memorizar:",
    "  Eventos sem valor futuro ('qual a temperatura agora?')",
    "  Informacao facilmente recalculavel ('saldo atual')",
    "  Dados que mudam rapido demais para consolidar",
    "  Conteudo sem consentimento do usuario",
    "",
    "Regra: melhor memoria MINIMA e util que maxima e ruidosa",
], "📖 Memoria tem custo. Pergunta-chave: 'vai melhorar decisao futura?' Se nao, nao memorize.\n💡 Analogia: guardar moveis. Nao guarda tudo — so o que vale. Guardar inutil custa espaco e dificulta achar o util.\n❓ 'Citem 3 casos onde NAO vale a pena memorizar.' (Temperatura, saldo, status volatil, conversa trivial)\n⚠️ Memoria maxima e anti-pattern. Cura — qualidade sobre quantidade.\n➡️ Observabilidade.", 62, T)

s = content_slide("Observabilidade de Memoria", [
    "Memoria e um sistema — precisa de monitoring:",
    "",
    "  Quem acessou qual memoria, quando? (auditabilidade)",
    "  Hit rate do recall (quantos retornam algo util?)",
    "  Latencia do recall",
    "  Custo de storage vs custo de recall",
    "  Drift: a memoria esta desatualizada?",
    "",
    "Ferramentas: LangSmith (memory traces), dashboards custom, audit logs",
    "Profundidade em ETHAGT12 — AgentOps",
], "📖 Memoria precisa de observabilidade. Hit rate, latencia, custo, drift. Sem isso, e caixa preta.\n💡 Analogia: painel do carro. Temperatura, combustivel, velocidade. Sem isso, voce dirige no escuro.\n⚠️ Sem medir hit rate e latencia, nao sabe se sua memoria funciona. Observabilidade desde o dia 1.\n➡️ Fechamento.", 63, T)

s = content_slide("Boas Praticas de Memoria (DO)", [
    "✅ Comece com working memory + checkpointer (antes de vector DB)",
    "✅ Modele estado serializavel desde o dia 1",
    "✅ Use thread_id consistente (conversa = thread)",
    "✅ Versione o schema de estado (schema_version)",
    "✅ Combine vector + metadata (nao use so similaridade)",
    "✅ Defina politica de eviction desde o inicio",
    "✅ Redacao de PII antes de armazenar",
    "✅ Observe a memoria (hit rate, latencia, drift)",
], "📖 8 boas praticas. Comece simples (working + checkpointer). Versione schema. Hybrid search. Eviction desde o inicio. PII redacted. Observabilidade.\n💡 Analogia: os 10 mandamentos da memoria de agentes.\n❓ 'Qual dessas praticas voces ja aplicam? Qual e mais negligenciada?'\n➡️ Anti-patterns.", 64, T, "DO")

s = content_slide("Anti-Patterns de Memoria (DON'T)", [
    "❌ Encher context window ate estourar (sem gestao ativa)",
    "❌ Usar vector DB para tudo (quando relacional bastava)",
    "❌ Nao versionar schema de estado (quebra no resume)",
    "❌ Armazenar PII sem redacao",
    "❌ Sem politica de eviction (memoria cresce infinitamente)",
    "❌ Confiar cegamente no recall (sem re-ranking)",
    "❌ Misturar memoria de usuarios diferentes (cross-contamination)",
    "❌ Sem observabilidade de memoria ('caixa preta')",
    "❌ Promover tudo para semantica (sem criterio de consolidacao)",
], "📖 9 anti-patterns. Os mais comuns: encher context sem gestao, vector DB para tudo, sem eviction, sem re-ranking.\n💡 Analogia: 'nao faca' na cozinha. Parecem obvios, mas todo mundo comete.\n❓ 'Qual destes anti-patterns voces ja cometeram?' (humor ajuda a fixar)\n➡️ Resumo.", 65, T, "DON'T")

s = content_slide("Resumo da Aula", [
    "1. 4 camadas: working (context), episodica (vector), semantica (KB/KG), procedural (skills)",
    "2. Checkpointer = persistencia de estado → resume, replay, branching",
    "3. Gestao de contexto: janela deslizante, sumarizacao, eviction, entity-centric",
    "4. Memoria vetorial: embedding + metadata filter + re-ranking",
    "5. Consolidacao: episodica → semantica (criterios de promocao)",
    "6. Producao: consistencia, PII, esquecimento, custo vs beneficio, observabilidade",
    "7. MemGPT e Generative Agents como referencias canonicas",
], "📖 Resumo em 7 pontos. Memorizem estes — e a fundacao de memoria de agentes.\n💡 Analogia: resumo de curso de arquitetura. Tipos de fundacao, preservacao, gestao de espaco, catalogacao, consolidacao, seguranca.\n➡️ Quiz.", 66, T, "Resumo")

s = exercise_slide("Quiz — Pergunta 1", [
    "Qual tipo de memoria e mais apropriado para",
    "'o usuario prefere respostas em portugues'?",
    "",
    "A) Working memory (context window)",
    "B) Memoria episodica (vector DB)",
    "C) Memoria semantica (KB / perfil)",
    "D) Memoria procedural (skills)",
    "",
    "Resposta: C — e um fato estavel sobre o usuario, nao um evento",
], "📖 Resposta: C — semantica. Fato estavel, nao evento (episodica), nao skill (procedural), nao re-promptavel (working).\n➡️ Pergunta 2.", 67, T, "Quiz")

s = exercise_slide("Quiz — Pergunta 2", [
    "O que acontece se o schema de estado mudar",
    "e um checkpoint antigo for carregado?",
    "",
    "A) Erro 500 e crash do agente",
    "B) O agente ignora campos novos e continua",
    "C) Depende — sem versionamento quebra; com migracao lazy funciona",
    "D) O checkpointer reescreve o estado automaticamente",
    "",
    "Resposta: C — sem versionamento, campos novos causam erro; com migracao lazy, converte no load",
], "📖 Resposta: C. Sem versionamento quebra; com schema_version + migracao lazy, funciona. Licão: versione o schema desde o dia 1.\n➡️ Pergunta 3.", 68, T, "Quiz")

s = exercise_slide("Quiz — Pergunta 3", [
    "Usuario pede: 'lembra da conversa sobre deploy na semana passada?'",
    "Qual estrategia de recall usar?",
    "",
    "A) Apenas vector search (similaridade com 'deploy')",
    "B) Apenas metadata filter (timestamp da semana passada)",
    "C) Metadata filter (tempo) + vector search (topico) + re-ranking",
    "D) Sumarizacao em cascata",
    "",
    "Resposta: C — combinar filtro temporal + busca semantica + re-ranking",
], "📖 Resposta: C. Query tem componente temporal E semantico. Combine metadata filter + vector + re-ranking. Hybrid search e o padrao.\n➡️ Encerramento.", 69, T, "Quiz")

s = title_slide(
    "Conexao, Projeto e Q&A",
    "ETHAGT05 — Memoria de Agentes\nProximos: ETHAGT06 (RAG Agentico) · ETHAGT07 (KG)",
    "Universidade Etho"
)
add_notes(s, "Encerramento.\n📖 Proximos: ETHAGT06 (RAG Agentico), ETHAGT07 (KG), ETHAGT14 (Escalabilidade).\nProjeto: memoria de agente pessoal de longo prazo (4 camadas + ADR + politica de privacidade).\nLabs: Lab 1 (Checkpointer Postgres, 4h, 1 semana), Lab 2 (Memoria Episodica, 5h, 2 semanas).\nLeitura: MemGPT (arXiv:2310.08560), Generative Agents (arXiv:2304.03442).\n❓ 'Perguntas?' Se silencio: 'Qual camada de memoria achou menos clara?'\n⚠️ Confirme prazos e criterios antes de sair.")

output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT05-Apresentacao.pptx")
prs.save(output)
print(f"PPTX gerado: {output}")
print(f"Total de slides: {len(prs.slides)}")
