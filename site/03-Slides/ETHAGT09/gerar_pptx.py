#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT09: Comunicação e Coordenação Multi-Agente
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Paleta Etho ───
PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ───

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj=""):
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT09"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

T = 72  # Total

# ═══════════════════════════════════════
# SEÇÃO A — Abertura (1-7)
# ═══════════════════════════════════════

s = title_slide(
    "Comunicação e Coordenação Multi-Agente",
    "Universidade Etho · Especialização em Programação Agêntica\nFase C — Sistemas Multi-Agente · A2A · Blackboard · Actor Model",
    "ETHAGT09"
)
add_notes(s, "📖 Bem-vindos. Esta aula transforma '1 agente' em 'N agentes'. Hoje dominamos comunicação e coordenação — de troca direta a blackboard e actor model.\n💡 Analogia: música solo vs orquestra. Solo precisa tocar bem; orquestra precisa partitura (protocolo), regente (orquestrador), escuta mútua (coordenação).\n❓ 'Quantos já trabalharam com mais de 1 agente?'\n⚠️ Multi-agente não é 'LangGraph com vários nodes'. É um campo com padrões próprios.\n➡️ Vamos aos objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: dominar modelos de comunicação e coordenação entre agentes",
    "",
    "Objetivos específicos:",
    "1. Distinguir padrões de comunicação A2A (agent-to-agent)",
    "2. Implementar troca estruturada de mensagens (com schemas)",
    "3. Aplicar blackboard para coordenação com espaço compartilhado",
    "4. Aplicar actor model para concorrência e isolamento",
    "5. Lidar com negociação, conflito e consenso",
    "6. Escolher o padrão correto: 1 agente vs N agentes",
], "📖 Cada objetivo é mensurável: 'distinguir', 'implementar', 'aplicar', 'escolher'. O #6 é o mais importante.\n💡 Analogia: caixa de ferramentas. Você pode saber usar martelo e chave de fenda, mas se usa martelo no parafuso, falha.\n❓ 'Qual objetivo é mais desafiador?' (geralmente #4 ou #5)\n⚠️ A2A não é chamada de API — é entre agentes autônomos.\n➡️ Onde estamos no Framework Etho.", num=2, total=T)

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → A (Avançado)",
    "C2 Multi-Agent Systems → I (Intermediário) → aprofunda em ETHAGT10, 11, 15",
    "C3 MCP & Tool Use → B (Básico) → ETHAGT08",
    "C4 Agent Memory → B (Básico) → ETHAGT05",
    "C5 AgentOps & Avaliação → B (Básico) → ETHAGT12",
    "",
    "Pré-requisito: ETHAGT04 (Reasoning & Planning)",
], "📖 Este módulo leva C1 ao Avançado e C2 ao Intermediário. Você consolida programação agêntica e sobe em multi-agent.\n💡 Analogia: faixa preta em um martial art (C1) e faixa azul em outro (C2).\n⚠️ 'Avançado' não é 'sabe tudo' — é 'decide arquitetura e lida com edge cases'.\n➡️ Estrutura da aula.", num=3, total=T)

s = content_slide("Agenda da Aula", [
    "Bloco 1 (45 min):",
    "  Abertura (8 min) · Espectro A2A (10 min) · Padrões de Conversação (13 min) · Blackboard (10 min)",
    "  Intervalo (5 min)",
    "Bloco 2 (45 min):",
    "  Actor Model (13 min) · Negociação (10 min) · Protocolos (8 min) · Fechamento (12 min)",
    "",
    "DEMO ao vivo no Slide 43: duas arquiteturas, um problema",
], "📖 Dois blocos. Bloco 1: fundamentos de comunicação, padrões, blackboard. Bloco 2: actor model, negociação, protocolos.\n💡 Analogia: entrada, prato principal, sobremesa, café.\n⚠️ Avise que o Slide 5 define o tom da aula.\n➡️ Por que coordenação é tão difícil?", num=4, total=T)

s = content_slide("O Problema da Coordenação", [
    "Cenário: agente de pesquisa + agente de escrita precisam colaborar",
    "Sem padrão: mensagens se perdem, estado diverge, ninguém sabe quem faz o quê",
    "'Dois agentes inteligentes não produzem um sistema inteligente automaticamente'",
    "",
    "❓ Quantos agentes colaborando antes de virar caos?",
], "📖 Sem protocolo, mensagens se perdem, estado diverge. Inteligência individual não garante coordenação.\n💡 Analogia: duas pessoas inteligentes cozinhando sem combinar quem faz o quê. Ambas sabem cozinhar, mas sem divisão, nada funciona.\n❓ 'Quantos agentes antes do caos?' (3+ começa a ficar difícil)\n⚠️ 'LLMs são inteligentes, se coordenam naturalmente' — FALSO. Sem protocolo, loop, ignoram mensagens, agem sobre estado desatualizado.\n➡️ Por que só agora temos soluções viáveis?", num=5, total=T, obj="Motivação")

s = content_slide("Por Que Agora", [
    "Linha do tempo:",
    "  2022: ReAct (reasoning em loop)",
    "  2023: AutoGen (arXiv:2308.08155), CAMEL (arXiv:2303.17760)",
    "  2024: OpenAI Swarm (handoffs), Google A2A Protocol",
    "  2025: MetaGPT (ICLR 2024), padronização emergente",
    "",
    "Confluência: reasoning maduro + frameworks + custo baixo + padronização (MCP + A2A)",
], "📖 Multi-agent existe desde os anos 80 (HEARSAY-II, FIPA-ACL). O que mudou: LLMs bons em reasoning, frameworks reproduzíveis, custo viável.\n💡 Analogia: TCP/IP padronizou redes e explodiu a internet. MCP + A2A são o TCP/IP dos agentes.\n❓ 'Qual fator foi o gatilho mais recente?' (Padronização — A2A é 2024)\n⚠️ Multi-agent não é 'moda nova' — conceito é antigo; LLM tornou viável.\n➡️ Vamos ao espectro A2A.", num=6, total=T)

s = section_slide(1, "O Espectro da Comunicação A2A")
add_notes(s, "📖 Início do bloco de fundamentos. Como agentes se comunicam? Topologias? Como estruturar mensagens? O que pode dar errado?\n➡️ Síncrono vs assíncrono.")

# ═══════════════════════════════════════
# SEÇÃO B — Espectro A2A (8-15)
# ═══════════════════════════════════════

s = comparison_slide("Comunicação Direta vs Assíncrona",
    "Síncrona (request/response)",
    ["A pergunta → B responde → A continua", "Latência baixa", "Alto acoplamento", "A fica bloqueado esperando", "Analogia: telefone"],
    "Assíncrona (eventos)",
    ["A publica → B reage quando puder", "Latência variável", "Baixo acoplamento", "A não bloqueia", "Analogia: e-mail"],
    "📖 Síncrono = chamada telefônica (A pergunta, B responde, A bloqueado). Assíncrono = e-mail (A publica, B reage quando puder).\n❓ 'Pesquisa + escrita: qual modelo?' (Assíncrono)\n⚠️ Síncrono para tudo cria gargalos com N agentes.\n➡️ 3 topologias canônicas.", num=8, total=T, obj="Trade-off: latência vs desacoplamento")

s = content_slide("Topologias: Broadcast vs P2P vs Pub/Sub", [
    "Broadcast: 1 → todos (anúncio global, ex.: 'tarefa concluída')",
    "P2P (point-to-point): 1 → 1 (conversa direta, ex.: negociação)",
    "Pub/Sub: produtores → tópicos → consumidores (desacoplado)",
    "",
    "Escalabilidade: Broadcast O(N), P2P O(N²), Pub/Sub escala melhor",
    "",
    "❓ Qual topologia para 10 agentes debatendo um diagnóstico?",
], "📖 Broadcast = 1-para-todos (todos precisam saber). P2P = 1-para-1 (negociação). Pub/Sub = N-para-M via tópicos.\n💡 Analogia: Broadcast = outdoors; P2P = telefonema; Pub/Sub = revista por assinatura.\n❓ '10 agentes debatendo?' (Pub/Sub por tópicos ou blackboard. Broadcast satura; P2P fica O(100).)\n⚠️ Broadcast 'simples' satura contexto com N grande.\n➡️ Mensagens precisam de estrutura.", num=9, total=T, obj="Topologias")

s = code_slide("Schemas de Mensagem A2A",
"""{
  "message_id": "msg_abc123",
  "sender": "agent://researcher-01",
  "receiver": "agent://writer-02",
  "message_type": "task_result",
  "version": "1.2.0",
  "timestamp": "2026-07-07T14:30:00Z",
  "payload": {
    "task_id": "task_xyz",
    "result": "Resumo encontrado: ...",
    "confidence": 0.87
  },
  "in_reply_to": "msg_def456"
}""",
"📖 Mensagens A2A não são texto livre — são estruturadas. Campos obrigatórios: sender, receiver, message_type, payload, timestamp, version.\n💡 Analogia: envelope de correio. Sem destinatário/remetente/CEP, a carta se perde.\n❓ 'Dois agentes sem concordar no schema?' (Parsing error ou interpretação errada silenciosa)\n⚠️ 'LLM entende' — sim, mas o CÓDIGO que processa não. Schema é para o código.\n➡️ Schemas evoluem — versionamento.", num=10, total=T)

s = content_slide("Versionamento de Mensagens", [
    "Agentes evoluem independentemente (deploy separado)",
    "Estratégia: semver no campo version (major.minor.patch)",
    "",
    "Compatibilidade:",
    "  Backward: v2 lê mensagens v1 ✅",
    "  Forward: v1 lê mensagens v2 (ignora campos novos) ⚠️",
    "  Full: ambos (ideal)",
    "",
    "Exemplo: v1 sem 'priority', v2 adiciona — v1 ignora, v2 usa",
], "📖 Sem versionamento, mudança quebra o outro silenciosamente. Adicionar campos é backward-compatible; remover não é.\n💡 Analogia: formulário v1 (nome, email); v2 adiciona telefone. v1 ignora telefone — não quebra. Mas se v2 remove email, v1 quebra.\n⚠️ 'São só 2 agentes' — mas agentes evoluem. Sempre versione.\n➡️ Mesmo com schema, mensagens falham.", num=11, total=T)

s = content_slide("Erros: Mensagens Perdidas e Duplicadas", [
    "Perdida: agente envia, receptor nunca recebe (rede, crash)",
    "  Impacto: decisão com dado faltante",
    "Duplicada: agente envia 1x, receptor processa 2x (retry mal feito)",
    "  Impacto: ação executada 2x (ex.: cobrança duplicada)",
    "",
    "Mitigação:",
    "  Acknowledgments (ACK) — receptor confirma",
    "  Idempotência — mesma mensagem processa só 1x",
    "  Deduplication keys — message_id único",
], "📖 Comunicação não é confiável por padrão. Perdida = carta extraviada. Duplicada = convite enviado 2x, convidado aparece 2x.\n💡 ACK = confirmação de leitura. Idempotência = 'já recebi, ignorando'.\n⚠️ Sem idempotência, cada retry duplica a ação. Em produção, retries são inevitáveis.\n➡️ Outro erro: fora de ordem.", num=12, total=T, obj="Erros de comunicação")

s = content_slide("Erros: Fora de Ordem", [
    "Mensagem 2 chega antes da 1 → estado inconsistente",
    "Causa: rotas diferentes, latência variável",
    "",
    "Mitigação:",
    "  Sequence numbers — receptor ordena",
    "  Timestamps Lamport — ordem causal",
    "  Fila ordenada — broker garante ordem",
    "",
    "Em LLM agents: ordem do contexto importa — desordem = alucinação",
], "📖 Mensagens podem chegar fora de ordem (rotas, latência). Em LLM agents, ver resposta antes da pergunta = alucinação.\n💡 Analogia: receber páginas do livro fora de ordem — lê o final antes do início.\n⚠️ TCP garante ordem dentro de uma conexão, mas A2A pode usar múltiplas. Não assuma ordem global.\n➡️ Garantias de entrega.", num=13, total=T)

s = content_slide("Garantias de Entrega", [
    "At-most-once: pode perder, nunca duplica (fire-and-forget)",
    "  Uso: métricas, logs não-críticos",
    "At-least-once: nunca perde, pode duplicar (com retry)",
    "  Uso: tarefas com idempotência",
    "Exactly-once: nunca perde, nunca duplica (ideal, caro)",
    "  Uso: transações financeiras",
    "",
    "❓ Qual garantia para transferência bancária entre agentes?",
], "📖 At-most-once = jogar panfleto (pode cair). At-least-once = carta registrada (chega, pode chegar 2x). Exactly-once = transferência bancária (chega 1x com confirmação dupla).\n❓ 'Transferência bancária?' (Exactly-once, ou at-least-once + idempotência)\n⚠️ Exactly-once requer consenso (2PC, Paxos/Raft). Na prática, at-least-once + idempotência é mais viável.\n➡️ Exercício rápido.", num=14, total=T, obj="Semânticas de entrega")

s = exercise_slide("Exercício — Topologia e Garantia Ideal", [
    "1. 'Triager despacha para especialista' → ?",
    "2. '3 agentes votam em um diagnóstico' → ?",
    "3. 'Agente publica evento de conclusão' → ?",
    "4. 'Negociação comprador/vendedor' → ?",
    "",
    "Votação rápida (mãos levantadas)",
    "",
    "Gabarito: 1) P2P+handoff | 2) Broadcast+at-least-once | 3) Pub/Sub+at-most-once | 4) P2P+exactly-once",
], "📖 Votação rápida para internalizar.\n💡 Analogia: bicicleta (at-most-once), ônibus (at-least-once), Uber (exactly-once).\n⚠️ 'Exactly-once para tudo' é caro. Avalie custo real de perder/duplicar.\n➡️ Padrões de conversação.", num=15, total=T)

# ═══════════════════════════════════════
# SEÇÃO C — Padrões de Conversação (16-26)
# ═══════════════════════════════════════

s = section_slide(2, "Padrões de Conversação")
add_notes(s, "📖 Como agentes conversam? 3 padrões canônicos: CAMEL, AutoGen, Swarm.\n➡️ Visão geral dos 3.")

s = content_slide("Três Padrões Canônicos de Conversação", [
    "1. Two-agent dialogue (CAMEL) — 2 agentes em turnos",
    "   Origem: arXiv:2303.17760",
    "2. Group chat (AutoGen) — N agentes orquestrados por manager",
    "   Origem: arXiv:2308.08155",
    "3. Handoff / transfer (Swarm) — agente transfere controle",
    "   Origem: OpenAI Swarm (2024)",
    "",
    "Cada padrão resolve um problema de coordenação diferente",
], "📖 CAMEL = entrevista (2 pessoas alternando). AutoGen = reunião de equipe (N, alguém modera). Swarm = transferência de chamada.\n⚠️ 'Group chat = todos falam ao mesmo tempo' — FALSO. Manager sequencia quem fala.\n➡️ Aprofundar CAMEL.", num=17, total=T)

s = content_slide("Two-Agent Dialogue — CAMEL", [
    "CAMEL: Communicative Agents for Mind Exploration",
    "Role-playing: agente A como 'assistente', B como 'usuário simulado'",
    "Estrutura: task → inception prompting → turnos de diálogo",
    "",
    "Inception prompting: system prompt mantém os papéis",
    "",
    "Origem: arXiv:2303.17760",
    "",
    "Quando brilha: refinamento iterativo entre 2 perspectivas",
], "📖 CAMEL = 2 agentes em turnos. Inovação: inception prompting mantém papéis. Sem isso, agentes 'saem do personagem'.\n💡 Analogia: brainstorm entre copywriter e designer. Cada um traz uma perspectiva, alternam e refinam.\n⚠️ 'CAMEL é só 2 nodes no LangGraph' — a chave é o inception prompting.\n➡️ Trace real de CAMEL.", num=18, total=T)

s = code_slide("CAMEL em Ação — Role-Playing",
"""┌─────────────────────────────────────────────┐
│  Tarefa: Escrever artigo sobre energias     │
│  renováveis                                  │
├─────────────────────────────────────────────┤
│  [Usuário]    Turno 1: "Qual estrutura      │
│               você sugere?"                  │
│  [Assistente] Turno 2: "Introdução, 3        │
│               seções técnicas, conclusão."   │
│  [Usuário]    Turno 3: "Aprofunde a seção   │
│               solar."                        │
│  [Assistente] Turno 4: "Painéis fotovoltaic │
│               os convertem..."               │
└─────────────────────────────────────────────┘""",
"📖 Vejam o trace. Usuário guia com perguntas; assistente executa. Vai-e-vem = refinamento iterativo.\n💡 Analogia: editor e escritor. Editor guia, escritor executa.\n❓ 'Se o usuário começar a escrever o artigo?' (Quebra o role-playing. Inception prompt impede.)\n⚠️ Sem inception prompting, agentes colapsam em 'ambos fazem tudo'.\n➡️ Escalar para N: group chat.", num=19, total=T)

s = content_slide("Group Chat — AutoGen", [
    "AutoGen GroupChat: N agentes em um grupo orquestrado",
    "Componentes: agentes + group chat manager",
    "Manager decide quem fala a seguir (hub-and-spoke)",
    "",
    "Origem: arXiv:2308.08155",
    "",
    "Quando brilha: múltiplas especialidades colaborando",
    "Ex.: Researcher + Coder + Reviewer em um projeto",
], "📖 GroupChat escala para N. Manager no centro decide quem fala — evita caos de 'todos falando'.\n💡 Analogia: reunião com facilitador. Facilitador (manager) decide quem fala, não contribui com conteúdo.\n⚠️ Manager decide QUEM fala, não O QUÊ. Conteúdo vem dos especialistas.\n➡️ Como o manager decide? Estratégias.", num=20, total=T)

s = comparison_slide("AutoGen: Selector vs Round-Robin vs Dynamic",
    "Round-robin",
    ["Ordem fixa: A → B → C → A", "Previsível", "Rígido (ignora contexto)", "Gratuito (sem custo LLM)"],
    "Selector / Dynamic",
    ["Selector: LLM decide quem fala", "Adaptável ao contexto", "Custa 1 chamada LLM por turno", "Dynamic: agentes se inscrevem"],
    "📖 Round-robin = rodízio de fala. Selector = professor aponta quem responde. Dynamic = assembleia (quem quer, fala).\n❓ '5 especialistas debatendo diagnóstico?' (Selector — ordem imprevisível)\n⚠️ Round-robin ignora contexto — pode forçar agente irrelevante a falar.\n➡️ Terceiro padrão: handoff (Swarm).", num=21, total=T)

s = content_slide("Handoff / Transfer — OpenAI Swarm", [
    "Swarm: agente transfere controle para outro agente",
    "Handoff = 'passa a bola' — agente atual sai, novo entra",
    "",
    "Exemplo: Triager → handoff para Sales / Billing / Support",
    "",
    "Leve, sem orquestrador central (diferente de GroupChat)",
    "Origem: OpenAI Swarm (repo + paper técnico, 2024)",
    "",
    "Diagrama: 12-Diagrams/ETHAGT09/handoff.mmd",
], "📖 Handoff é TRANSFERÊNCIA, não conversa. Triager transfere para especialista e SAI. Sem manager.\n💡 Analogia: plantão médico. Generalista atende, transfere para cardiologista, e sai.\n⚠️ Confusão #1: handoff ≠ 'chamar para ajudar'. Handoff = agente original SAI. Se continua, é delegação.\n➡️ Vamos esclarecer handoff vs delegação.", num=22, total=T)

s = comparison_slide("Handoff vs Delegação",
    "Handoff (Swarm)",
    ["Controle TRANSFERIDO", "Agente original SAI", "'Transfere a chamada'", "Simples, sem retorno", "Ex.: triager → especialista"],
    "Delegação (Supervisor)",
    ["Controle DELEGADO", "Supervisor espera retorno", "'Coloca em espera e consulta'", "Supervisor continua envolvido", "Ex.: supervisor → sub-agente"],
    "📖 Handoff = 'transfere para vendas' (você sai). Delegação = 'consulto meu supervisor' (você fica, espera, volta).\n❓ 'Suporte ao cliente?' (Depende do SLA. Handoff = simples; delegação = monitora e escala.)\n⚠️ Usar handoff quando precisa de delegação perde contexto.\n➡️ MetaGPT: comunicação estruturada via SOPs.", num=23, total=T)

s = content_slide("MetaGPT — SOPs para Multi-Agente", [
    "MetaGPT: agentes seguem SOPs como uma equipe de software",
    "Papéis: Product Manager → Architect → Engineer → QA",
    "Cada papel tem entradas, processo e saídas definidos",
    "",
    "Comunicação via artefatos estruturados (não chat livre):",
    "  PRD → Design Doc → Code → Test Report",
    "",
    "Aprofundamento no caso de estudo (Slide 63)",
], "📖 MetaGPT usa SOPs (Standard Operating Procedures). Comunicação NÃO é chat — é passagem de artefatos estruturados.\n💡 Analogia: fábrica. Cada estação recebe insumo, processa, entrega produto. Sem improviso.\n⚠️ 'Chat livre é mais flexível' — MetaGPT prova o oposto. Estrutura > improvisação.\n➡️ Quando cada padrão brilha.", num=24, total=T)

s = content_slide("Quando Cada Padrão Brilha", [
    "Two-agent (CAMEL): refinamento iterativo entre 2 perspectivas",
    "Group chat (AutoGen): múltiplas especialidades colaborando",
    "Handoff (Swarm): roteamento baseado em especialização",
    "Selector group chat: ordem imprevisível",
    "Round-robin: todos contribuem igualmente",
    "MetaGPT (SOPs): workflow estruturado com papéis",
    "",
    "Regra: comece simples (2 agentes); escale só quando necessário",
], "📖 Cada padrão tem ponto forte. Two-agent para 2 perspectivas. Group chat para N especialistas. Handoff para roteamento. MetaGPT para workflow estruturado.\n💡 Analogia: escolher formato de reunião. 1:1 (CAMEL), brainstorm (group chat), plantão (handoff), montagem (MetaGPT).\n⚠️ Group chat para tudo explode em contexto e custo. Considere handoff/blackboard.\n➡️ Exercício.", num=25, total=T)

s = exercise_slide("Exercício — Padrão de Conversação Ideal", [
    "1. 'Writer e reviewer refinam um texto' → ?",
    "2. '3 especialistas debatem diagnóstico' → ?",
    "3. 'Triager encaminha para dept correto' → ?",
    "4. 'Equipe de dev constrói um feature' → ?",
    "",
    "Votação rápida",
    "",
    "Gabarito: 1) Two-agent/CAMEL | 2) Group chat/selector | 3) Handoff/Swarm | 4) MetaGPT/SOPs",
], "📖 Votação para internalizar critérios.\n⚠️ Usar group chat para triager (cenário 3) é overkill. Handoff é mais simples.\n➡️ Blackboard: alternativa à comunicação direta.", num=26, total=T)

# ═══════════════════════════════════════
# SEÇÃO D — Blackboard (27-33)
# ═══════════════════════════════════════

s = section_slide(3, "Blackboard: Espaço Compartilhado")
add_notes(s, "📖 Comunicação indireta via espaço compartilhado.\n➡️ O que é blackboard?")

s = content_slide("O Que É Blackboard", [
    "Espaço compartilhado de estado: agentes escrevem e leem",
    "Inspirado no blackboard físico de salas de aula",
    "Desacoplamento: agentes não precisam se conhecer",
    "",
    "Padrão clássico de IA (HEARSAY-II, 1970s)",
    "Renascido com LLM agents: estado compartilhado = 'memória coletiva'",
], "📖 Blackboard = espaço compartilhado. Agentes escrevem e leem, não falam diretamente. Desacoplamento total.\n💡 Analogia: lousa em sala de pesquisa. Cada pesquisador passa, lê, adiciona contribuição. Ninguém fala diretamente — a lousa coordena.\n⚠️ Blackboard ≠ banco de dados. BD armazena dados; blackboard coordena agentes.\n➡️ Diagrama do blackboard.", num=28, total=T)

s = content_slide("Arquitetura do Blackboard", [
    "3 agentes conectados a um blackboard central",
    "Blackboard armazena: facts, hypotheses, partial results",
    "Fluxo: agentes leem → processam → escrevem de volta",
    "Sem comunicação direta entre agentes",
    "",
    "Diagrama: 12-Diagrams/ETHAGT09/blackboard.mmd",
], "📖 Pesquisador escreve facts; Analista lê, escreve hypotheses; Redator lê facts+hypotheses, escreve partial result. Tudo via blackboard.\n💡 Analogia: Google Doc compartilhado. Múltiplas pessoas editam, cada uma na sua seção.\n❓ 'Analista sabe quando algo mudou?' (Lê periodicamente ou assina eventos de escrita.)\n⚠️ Sem mecanismo de notificação, agentes não sabem quando ler.\n➡️ Quando blackboard é a escolha certa.", num=29, total=T)

s = content_slide("Quando Blackboard Brilha", [
    "Problema dinâmico: estado muda conforme agentes contribuem",
    "Múltiplos especialistas: cada um contribui com seu conhecimento",
    "Contribuição incremental: ninguém tem a resposta completa",
    "",
    "Exemplos: diagnóstico multi-especialidade, planejamento estratégico, relatório por partes",
    "",
    "❓ Quando blackboard é preferível a mensagens diretas?",
], "📖 Blackboard brilha com problema dinâmico, múltiplos especialistas, contribuição incremental.\n💡 Analogia: investigação policial no quadro. Detetive adiciona pista, perito adiciona análise, testemunha relata. Solução emerge no quadro.\n❓ 'Quando preferível a mensagens diretas?' (N grande, contribuições independentes, estado dinâmico.)\n⚠️ Blackboard para 2 agentes é overkill. Mensagens diretas são mais simples.\n➡️ Trade-off sistemático.", num=30, total=T)

s = comparison_slide("Blackboard vs Mensagens Diretas",
    "Blackboard",
    ["Baixo acoplamento", "Estado compartilhado", "Sem ordem garantida", "Escala O(N) conexões", "N grande + contribuições independentes"],
    "Mensagens Diretas",
    ["Alto acoplamento", "Explícito", "Ordem controlada", "Escala O(N²) conexões", "2 agentes + ordem crítica"],
    "📖 Blackboard O(N) conexões; mensagens diretas O(N²). Com N=10, blackboard=10, diretas=45.\n💡 Analogia: blackboard = fórum online; mensagens = email individual.\n⚠️ Blackboard adiciona complexidade (locking, consistência). Para 2 agentes com ordem, diretas são melhores.\n➡️ Implementação.", num=31, total=T)

s = code_slide("Implementação do Blackboard",
"""class Blackboard:
    def __init__(self):
        self.entries = []

    def write(self, agent_id, entry_type, content):
        self.entries.append({
            "agent_id": agent_id,
            "type": entry_type,
            "content": content,
            "timestamp": time.time()
        })

    def read(self, entry_type=None):
        return [e for e in self.entries
                if entry_type is None
                or e["type"] == entry_type]""",
"📖 Em memória: dict/list compartilhado (rápido, volátil). Persistente: Postgres/Redis (resiliente, distribuído). Locking para concorrência.\n💡 Analogia: em memória = rascunho em papel (perde se cair); persistente = Google Drive (sobrevive).\n⚠️ Sem locking, 2 agentes escrevendo simultaneamente = race condition.\n➡️ Exercício em duplas.", num=32, total=T)

s = exercise_slide("Exercício — Blackboard ou Mensagens Diretas?", [
    "1. '5 agentes contribuem com partes de um relatório' → ?",
    "2. 'Agente pergunta preço a outro' → ?",
    "3. '3 agentes debatem com estado compartilhado' → ?",
    "",
    "Discutir em duplas (2 min)",
    "",
    "Gabarito: 1) Blackboard | 2) Mensagem direta (P2P) | 3) Blackboard",
], "📖 Duplas discutem 2 min.\n⚠️ Mensagens diretas para 5 agentes = O(10) conexões. Blackboard é mais simples.\n➡️ Actor model: paradigma diferente.", num=33, total=T)

# ═══════════════════════════════════════
# SEÇÃO E — Actor Model (34-44)
# ═══════════════════════════════════════

s = section_slide(4, "Actor Model: Concorrência sem Locks")
add_notes(s, "📖 Paradigma de concorrência inventado em 1973 por Carl Hewitt. Relevante para LLM agents.\n➡️ O que é?")

s = content_slide("O Que É Actor Model", [
    "Atores encapsulam estado; só se comunicam por mensagens",
    "Cada ator tem um mailbox (fila de mensagens)",
    "Processa uma mensagem por vez — sem locks, sem race conditions",
    "",
    "Origem: Hewitt, 1973 (clássico, MIT)",
    "Renascido com LLM agents: cada agente = um ator isolado",
], "📖 Ator = entidade com estado privado. Comunicação só por mensagem. Mailbox = fila. Uma mensagem por vez = sem locks.\n💡 Analogia: funcionários com mesa e caixa de entrada. Não mexem na mesa dos outros — só trocam memorandos.\n❓ 'Sem locks, como evitar conflito?' (Não há estado compartilhado. Cada ator isolado.)\n⚠️ Actor ≠ thread. Threads compartilham (locks). Atores não (sem locks).\n➡️ Diagrama.", num=35, total=T)

s = content_slide("Arquitetura do Actor Model", [
    "3 atores com mailboxes individuais",
    "Comunicação assíncrona via mensagens",
    "Estado privado encapsulado em cada ator",
    "Sem estado compartilhado = sem locks",
    "",
    "Diagrama: 12-Diagrams/ETHAGT09/actor-model.mmd",
], "📖 A envia para mailbox do B. B processa quando puder (assíncrono). Estado do B é inacessível ao A.\n💡 Analogia: correio. Você não entra na casa do destinatário — manda carta.\n❓ 'A precisa de info do estado do B?' (Envia mensagem pedindo. B responde. Não acessa diretamente.)\n⚠️ Expor estado como atributo público quebra encapsulamento.\n➡️ Princípios (continua no bloco 2).", num=36, total=T)

s = content_slide("Princípios: Encapsulamento e Isolamento", [
    "Encapsulamento: estado do ator é inacessível externamente",
    "Isolamento: falha em um ator não derruba outros",
    "'Share nothing' — único caminho de interação é mensagem",
    "Implicação para agentes: estado protegido",
    "",
    "❓ Actor model vs shared state — qual é mais seguro para dados críticos?",
], "📖 Encapsulamento = estado privado. Isolamento = crash não propaga. 'Share nothing'.\n💡 Analogia: apartamentos. Cada morador tem espaço privado. Incêndio no vizinho não afeta você.\n❓ 'Qual mais seguro para dados críticos?' (Actor — estado encapsulado, sem race conditions.)\n⚠️ Expor estado 'para facilitar' reintroduz race conditions.\n➡️ Concorrência sem locks.", num=37, total=T)

s = comparison_slide("Concorrência sem Locks",
    "Shared State (com locks)",
    ["Múltiplas threads + locks", "Deadlocks, race conditions", "Contenção de lock em alta concorrência", "Escala mal"],
    "Actor Model (sem locks)",
    ["1 mensagem por vez por ator", "Sem locks, sem deadlocks", "Escala horizontalmente", "Overhead: serialização de msgs"],
    "📖 Shared state: threads brigam por lock (deadlocks). Actor: cada um processa em paz.\n💡 Analogia: mesa compartilhada (brigam pela caneta) vs cada um com sua caneta (trocam recados).\nRegra: 'Don't communicate by sharing memory; share memory by communicating.'\n❓ 'Custo do actor model?' (Overhead de serialização, mas compensa em alta concorrência.)\n➡️ Localização transparente.", num=38, total=T)

s = content_slide("Localização Transparente", [
    "Ator local vs remoto: mesma interface (mensagem)",
    "Roteamento transparente: framework decide para onde enviar",
    "Escala: mova atores para outros nós sem mudar código",
    "Aplicação: agentes distribuídos em múltiplos servidores",
    "",
    "❓ Como isso muda a arquitetura de deployment?",
], "📖 Local e remoto têm mesma interface. Framework roteia. Pode mover atores entre servidores sem mudar código.\n💡 Analogia: ligar para colega — não sabe (nem precisa) se está no escritório ou em casa.\n❓ 'Muda deployment?' (Escala horizontal sem mudar código.)\n⚠️ 'Local e remoto precisam de código diferente' — não. Interface é a mesma.\n➡️ Implementações reais.", num=39, total=T)

s = comparison_slide("Implementações do Actor Model",
    "Erlang / Akka",
    ["Erlang/OTP: pioneiro, telecom", "Tolerância a falha (supervisores)", "Akka: JVM (Scala/Java)", "Maduro, enterprise"],
    "Python asyncio",
    ["Coroutines como atores leves", "Sem supervisão nativa", "Para LLM agents: asyncio + framework", "LangGraph, AutoGen"],
    "📖 Erlang = blind-car (tolerante a falha). Akka = SUV enterprise. asyncio = econômico (você adiciona acessórios).\n⚠️ Não reinvente: use asyncio + framework (LangGraph/AutoGen). Supervisão, mailbox, roteamento já prontos.\n➡️ Aplicação a agentes distribuídos.", num=40, total=T)

s = content_slide("Actor Model para Agentes Distribuídos", [
    "Cada agente = um ator com mailbox",
    "Tool execution = mensagem para ator especializado",
    "Memória = estado privado do ator",
    "Escala: agentes em nós diferentes, coordenados por mensagens",
    "",
    "Caso de uso: 10 agentes de pesquisa processando em paralelo",
], "📖 Agente = ator: mailbox, estado privado (memória), 1 msg por vez. Tools viram atores especializados.\n💡 Analogia: equipe remota. Cada um em casa (nó), computador próprio (estado), coordena por Slack (mensagens).\n⚠️ Todos em um processo perde vantagem de distribuição. Use localização transparente.\n➡️ Comparação sistemática.", num=41, total=T)

s = comparison_slide("Actor Model vs Shared State",
    "Shared State (blackboard)",
    ["Leitura/escrita compartilhada", "Precisa de locks", "Simples de entender", "Race conditions em concorrência"],
    "Actor Model",
    ["Estado isolado", "Mensagens assíncronas", "Sem locks", "Overhead de serialização"],
    "📖 Híbrido: ator gerencia blackboard. Combina isolamento do actor com compartilhamento do blackboard.\n💡 Analogia: cozinha compartilhada (locks) vs cozinha própria (actor) vs chef central (blackboard-ator).\n❓ V/F 'Actor model é mais lento que shared-state.' (FALSO — alta concorrência escala melhor sem locks.)\n➡️ DEMO.", num=42, total=T)

s = code_slide("DEMO — Duas Arquiteturas, um Problema",
"""Tarefa: pesquisar + resumir + formatar um relatório

═══════════════════════════════════════════
VERSÃO 1: AutoGen-style Group Chat
═══════════════════════════════════════════
Manager ↔ Researcher ↔ Summarizer ↔ Formatter
  (agentes conversam, manager coordena)

═══════════════════════════════════════════
VERSÃO 2: Blackboard
═══════════════════════════════════════════
Researcher → write(facts)     ↘
Summarizer → read/write(summary) → Blackboard
Formatter  → read → output    ↗
  (agentes compartilham espaço)

Comparar: nº mensagens, tempo, acoplamento
Ref: 05-Labs/ETHAGT09/Lab1-Duas-Arquiteturas""",
"📖 DEMO ao vivo. Mesma tarefa, 2 arquiteturas. Group chat gera mais mensagens; blackboard pode paralelizar.\n⚠️ Se API falhar, screenshots dos traces. A lição (comparar arquiteturas) > código rodando.\n➡️ Discussão da demo.", num=43, total=T, obj="DEMO")

s = exercise_slide("Discussão da DEMO", [
    "1. Qual arquitetura foi mais fácil de debugar? Por quê?",
    "2. Em qual delas adicionar um 4º agente seria mais simples?",
    "",
    "Discussão em duplas (2 min)",
    "",
    "Gabarito: 1) Blackboard (estado centralizado fácil de inspecionar) | 2) Blackboard (novo agente só lê/escreve)",
], "📖 Blackboard tipicamente mais fácil de debugar (estado central) e estender (novo agente só lê/escreve). Group chat exige rastrear sequência e registrar no manager.\n➡️ Negociação e conflito.", num=44, total=T)

# ═══════════════════════════════════════
# SEÇÃO F — Negociação e Conflito (45-52)
# ═══════════════════════════════════════

s = section_slide(5, "Negociação e Conflito")
add_notes(s, "📖 Agentes têm objetivos próprios que podem conflitar. Negociação é o tópico do projeto.\n➡️ Por que negociar?")

s = content_slide("Por Que Agentes Negociam", [
    "Agentes têm objetivos próprios — podem conflitar",
    "Exemplo: comprador quer preço baixo, vendedor quer alto",
    "Sem negociação: impasse ou força bruta",
    "Com negociação: convergência para acordo mutuamente aceitável",
    "",
    "❓ O que acontece quando objetivos são parcialmente conflitantes?",
], "📖 Sem negociação, impasse ou força bruta. Com negociação, convergência.\n💡 Analogia: negociação de salário. Candidato quer mais; empresa quer menos. Sem negociação, não há contratação.\n❓ 'Objetivos parcialmente conflitantes?' (Tensão. Sem protocolo, impasse. Com negociação, convergência.)\n⚠️ 'Agentes colaboram naturalmente' — não quando representam interesses diferentes.\n➡️ Padrão bargaining.", num=46, total=T)

s = content_slide("Bargaining — Barganha", [
    "Fluxo: proposta → contraproposta → ... → acordo ou impasse",
    "",
    "Estratégias:",
    "  Conceder gradativamente ('eu baixo 5, você baixa 5')",
    "  Fixar reserva (não abaixo de X)",
    "  BATNA (Best Alternative to a Negotiated Agreement)",
    "",
    "Exemplo: comprador 100 → vendedor 150 → convergem em 120",
    "Risco: loop infinito sem estratégia de cessão",
], "📖 Bargaining = proposta/contraproposta até acordo. Sem estratégia de cessão, loop infinito.\n💡 Analogia: regatear no mercado. Vendedor pede 200, você oferece 100, convergem em 150.\n⚠️ Sem estratégia de cessão, deadlock garantido. Defina como cada parte cede.\n➡️ Padrão auction.", num=47, total=T)

s = content_slide("Auction — Leilão", [
    "Auction: 1 vendedor, N compradores competem",
    "",
    "Tipos:",
    "  English (ascendente): lances sobem até ninguém oferecer mais",
    "  Dutch (descendente): preço desce até alguém aceitar",
    "  Sealed-bid: lances secretos, maior ganha",
    "  Vickrey: sealed-bid, vencedor paga o 2º maior lance",
    "",
    "Exemplo: agentes competindo por tempo de GPU",
], "📖 English = leilão de arte. Dutch = florada holandesa. Vickrey = leilão selado, paga 2º preço (incentiva lance honesto).\n💡 Vickrey é poderoso: você sempre dá valor verdadeiro, porque paga o segundo preço.\n⚠️ Vickrey pouco conhecido, mas valioso para agentes (incentiva honestidade).\n➡️ Diagrama de negociação.", num=48, total=T)

s = content_slide("Fluxo de Negociação", [
    "Comprador propõe 100 → Vendedor contrapropõe 150",
    "Comprador propõe 120 → Vendedor aceita → acordo",
    "",
    "Caminho alternativo: rejeita após N rounds → timeout / mediator",
    "",
    "Diagrama: 12-Diagrams/ETHAGT09/negotiation.mmd",
    "",
    "Critério do projeto: convergência ≥ 80% dos casos",
], "📖 Fallback (timeout/mediator) é crítico. Sem ele, loop infinito. Os 20% que não convergem devem cair no fallback.\n💡 Analogia: negociação com prazo. Se não fechar até sexta, contrato expira. Prazo força convergência.\n❓ 'Vendedor nunca aceita <140, comprador nunca oferece >130?' (Sem zona de acordo. Timeout/mediator.)\n⚠️ Sem max_rounds, agentes negociam indefinidamente.\n➡️ Resolução de conflito.", num=49, total=T)

s = comparison_slide("Resolução de Conflito",
    "Voting",
    ["Cada agente vota, maioria vence", "Simples", "Pode errar (2-1 onde 1 estava certo)", "Weighted: pondera por expertise", "Para decisões de baixo custo"],
    "Mediator",
    ["Agente neutro facilita consenso", "Mais caro (agente extra)", "Mais robusto", "Ouve todos, propõe solução", "Para decisões críticas"],
    "📖 Voting = eleição (maioria vence, pode errar). Mediator = terapia de casal (profissional facilita).\n❓ '3 especialistas com diagnóstico divergente?' (Mediator — stakes altos. Voting pode errar 2-1.)\n⚠️ Voting simples para decisões críticas pode errar. Mediator ou humano no loop.\n➡️ Deadlock.", num=50, total=T)

s = content_slide("Convergência e Deadlock", [
    "Convergência: acordo em N rounds",
    "Deadlock: nenhum agente cede, loop infinito",
    "",
    "Causas: estratégia rígida, sem BATNA, 'orgulho' do agente",
    "",
    "Mitigação:",
    "  max_rounds (força parada após N)",
    "  timeout (para após T segundos)",
    "  Concessão forçada (quem menos cedeu, cede)",
    "  Escalar para mediator",
    "",
    "Critério do projeto: convergência ≥ 80%",
], "📖 Deadlock = 2 carros em rua estreita, nenhum recua. Sem deadline (max_rounds) ou policial (mediator), ficam para sempre.\n⚠️ Sem max_rounds no projeto, deadlock garantido em divergência total.\n➡️ Exemplo concreto.", num=51, total=T)

s = content_slide("Exemplo: Objetivos Parcialmente Conflitantes", [
    "Cenário: agente de velocidade (rápido) vs agente de qualidade (completo)",
    "Negociação: trade-off entre tempo e profundidade",
    "Solução: orçamento de tempo + threshold de qualidade mínima",
    "Convergência: ambos aceitam meio-termo documentado",
    "",
    "❓ Como evitar que o agente de qualidade nunca aceite?",
], "📖 Sempre dá para melhorar qualidade — agente nunca aceita sem limite. Solução: threshold mínimo + tempo máximo. Interseção = zona de acordo.\n💡 Analogia: escolher restaurante. Você quer barato/rápido; amigo quer chique/elaborado. Negociam meio-termo.\n❓ 'Como evitar que qualidade nunca aceite?' (Threshold mínimo + tempo máximo. Sem limites, qualidade sempre melhora.)\n➡️ Protocolos emergentes.", num=52, total=T)

# ═══════════════════════════════════════
# SEÇÃO G — Protocolos Emergentes (53-59)
# ═══════════════════════════════════════

s = section_slide(6, "Protocolos e Padrões Emergentes")
add_notes(s, "📖 Padrões que estão padronizando A2A: A2A Protocol do Google e relação com MCP.\n➡️ A2A Protocol.")

s = content_slide("A2A Protocol (Google, 2024)", [
    "A2A (Agent-to-Agent): padrão aberto para comunicação entre agentes",
    "Objetivo: interoperabilidade entre frameworks diferentes",
    "",
    "Conceitos centrais:",
    "  Agent Card: manifesto de capacidades (URL /.well-known/agent.json)",
    "  Tasks: unidade de trabalho delegável",
    "",
    "Status: emerging spec (ainda em evolução)",
], "📖 A2A = interoperabilidade. Agent Card diz capacidades. Tasks são trabalho delegável.\n💡 Analogia: Agent Card = cartão de visitas digital. Tasks = ordem de serviço.\n⚠️ Confundir A2A com MCP. A2A = agent↔agent. MCP = agent↔system.\n➡️ Como funciona.", num=54, total=T)

s = content_slide("A2A Protocol — Como Funciona", [
    "1. Agente A descobre Agent Card de B (GET /.well-known/agent.json)",
    "2. A envia Task para B (POST /tasks)",
    "3. B processa e retorna status + resultado",
    "4. Streaming opcional para tarefas longas (SSE — Server-Sent Events)",
    "",
    "Protocolo HTTP-based — não reinventa a roda (REST + JSON)",
], "📖 Fluxo: descobre Agent Card → envia Task → B processa → status + resultado. SSE para updates em tempo real.\n💡 Analogia: contratar freelancer. Vê perfil (Agent Card), envia brief (Task), ele entrega com updates (SSE).\n⚠️ 'A2A é novo e complicado' — é HTTP + JSON. Inovação é padronização, não tecnologia.\n➡️ A pergunta: A2A vs MCP?", num=55, total=T)

s = comparison_slide("MCP vs A2A",
    "MCP (Model Context Protocol)",
    ["Agent ↔ System", "Agent ↔ Tools (banco, API)", "Agent ↔ Data Sources", "Da Anthropic", "Mais maduro"],
    "A2A (Agent-to-Agent)",
    ["Agent ↔ Agent", "Delegação de tarefas", "Colaboração", "Do Google", "Emergente"],
    "📖 MCP = agent↔system (ferramentas). A2A = agent↔agent (agentes). COMPLEMENTARES, não competidores.\n💡 Analogia: MCP = como você usa ferramentas (martelo). A2A = como fala com colega. Precisa dos dois.\n❓ 'Complementares ou competidores?' (Complementares.)\n⚠️ 'Preciso escolher' — não, use ambos.\n➡️ Padrões de orquestração.", num=56, total=T)

s = content_slide("Padrões de Orquestração Multi-Agente", [
    "Centralizada (supervisor): um coordena os demais — controle, gargalo",
    "Descentralizada (P2P): sem líder — flexível, sem garantia",
    "Hierárquica: supervisor → sub-supervisores → agentes — escala",
    "Market-based: agentes 'licitam' tarefas — eficiente",
    "",
    "Trade-off: controle vs flexibilidade vs escalabilidade",
], "📖 Centralizada = ditador (controle total, gargalo). Descentralizada = assembleia (flexível). Hierárquica = corporação (escala). Market = leilão (eficiente).\n⚠️ Centralizada não escala — supervisor vira gargalo.\n➡️ Contexto histórico: FIPA.", num=57, total=T)

s = content_slide("FIPA e Padrões Históricos", [
    "FIPA (Foundation for Intelligent Physical Agents): padrões ACL dos anos 90",
    "FIPA-ACL: performatives (request, inform, propose, accept, reject)",
    "KQML: predecessor de FIPA",
    "",
    "Aprendizado: LLM agents redescobrem conceitos de FIPA",
    "Lição: muita coisa 'nova' já foi pensada — aprender com o passado",
], "📖 FIPA já tinha performatives nos anos 90. LLM agents fazem o mesmo, mas sem saber que já tem nome.\n💡 Analogia: redescobrir a roda. FIPA já tinha 'request' e 'inform'. Aprender com o passado economiza tempo.\n⚠️ 'LLM agents inventaram multi-agent' — não. FIPA, HEARSAY-II, KQML já existiam. LLM tornou viável.\n➡️ Estado da padronização.", num=58, total=T)

s = content_slide("Estado da Padronização", [
    "MCP: mais maduro (Anthropic, amplamente adotado)",
    "A2A: emergente (Google, adoção crescente)",
    "FIPA: maduro mas sub-adotado em LLM agents",
    "",
    "Problema atual: fragmentação (cada framework tem seu protocolo)",
    "Tendência: convergência para MCP + A2A como camadas complementares",
    "",
    "Estamos em 2026 — ainda cedo para dizer que venceu",
], "📖 MCP maduro. A2A emergente. FIPA sub-adotado. Fragmentação atual. Tendência: MCP + A2A complementares.\n💡 Analogia: formatos de vídeo nos anos 2000 (VHS vs Betamax). Fragmentação até convergir.\n⚠️ Adotar protocolo cegamente = risco. Avalie maturidade, adoção, direção.\n➡️ Fechamento.", num=59, total=T)

# ═══════════════════════════════════════
# SEÇÃO H — Fechamento (60-72)
# ═══════════════════════════════════════

s = section_slide(7, "Boas Práticas e Anti-Patterns")
add_notes(s, "📖 Consolidação: o que fazer e o que não fazer.\n➡️ Boas práticas.")

s = content_slide("Boas Práticas (DO)", [
    "✓ Comece com 2 agentes antes de escalar para N",
    "✓ Defina schemas de mensagem com versionamento desde o dia 1",
    "✓ Use blackboard quando N é grande e contribuições independentes",
    "✓ Defina max_rounds em negociações (evitar deadlock)",
    "✓ Logue TODAS as mensagens A2A (traces = debugging)",
    "✓ Isole agentes como atores (estado protegido)",
    "✓ Documente o protocolo de handoff/delegação explicitamente",
], "📖 7 boas práticas. 2 agentes antes de N. Schemas desde dia 1. Blackboard para N grande. max_rounds. Logs sempre. Atores isolados. Handoff documentado.\n💡 Analogia: boas práticas de trânsito. Comece devagar, use cinto, siga padrões, tenha freio.\n⚠️ Pular versionamento 'para ir mais rápido' dói quando schema muda.\n➡️ Anti-patterns.", num=61, total=T, obj="Boas práticas")

s = content_slide("Anti-Patterns (DON'T)", [
    "✗ Começar com group chat de 5 agentes sem necessidade",
    "✗ Sem schema de mensagem ('telefone sem fio')",
    "✗ Negociação sem max_rounds (deadlock infinito)",
    "✗ Shared state sem locks (race condition)",
    "✗ Actor model para tudo (overhead em tarefas simples)",
    "✗ Não logar mensagens (debug cego em multi-agente)",
    "✗ Misturar MCP e A2A sem entender a diferença",
    "✗ Confiar que agentes 'vão se entender' sem protocolo",
], "📖 8 anti-patterns. Group chat de 5 no início. Sem schema. Sem max_rounds. Shared state sem locks. Actor para tudo. Sem logs. Confundir MCP/A2A. Confiar em 'se entendem'.\n💡 Analogia: anti-patterns de dieta. Cada um parece boa ideia, mas faz mal.\n⚠️ 'Agentes inteligentes se coordenam sozinhos' — não. Sem protocolo, caos.\n➡️ Caso MetaGPT.", num=62, total=T, obj="Anti-patterns")

s = content_slide("Caso de Estudo — MetaGPT", [
    "MetaGPT: framework multi-agente para desenvolvimento de software",
    "Papéis: Product Manager → Architect → Engineer → QA",
    "Comunicação via artefatos estruturados (não chat livre)",
    "SOPs codificam o workflow: cada papel sabe o que esperar",
    "Resultado: gera código funcional de ponta a ponta",
    "",
    "Lição: estrutura > improvisação em sistemas multi-agente",
], "📖 MetaGPT junta tudo. Papéis claros. Artefatos estruturados (PRD → Design → Code → Test Report). SOPs reduzem ambiguidade.\n💡 Analogia: linha de montagem de fábrica. Cada estação recebe insumo, processa, entrega produto.\n❓ 'Por que funciona melhor que group chat?' (Estrutura. SOPs reduzem ambiguidade.)\n⚠️ 'Chat livre é mais flexível' — MetaGPT prova o oposto.\n➡️ Exercício de schema.", num=63, total=T)

s = exercise_slide("Exercício — Schema de Mensagem A2A", [
    "Cenário: agentes de compra e venda negociando preço",
    "Em duplas: escrever schema JSON com versionamento",
    "",
    "Incluir: sender, receiver, message_type, payload, timestamp, version",
    "Considerar: proposta, contraproposta, aceitação (campo message_type)",
    "",
    "3 min discussão, compartilhar no quadro",
    "",
    "Dica: campo in_reply_to referencia a mensagem original",
], "📖 Duplas escrevem schema. Prepara para o projeto (sistema de negociação).\n💡 Analogia: desenhar formulário de proposta comercial. Campos: remetente, destinatário, tipo, valor, prazo.\n❓ 'Como representar contraproposta?' (message_type: counter_propose + in_reply_to.)\n⚠️ Esquecer version — quando schema evolui, quebra.\n➡️ Resumo.", num=64, total=T)

s = content_slide("Resumo da Aula", [
    "Espectro A2A: síncrono/assíncrono, topologias, schemas, garantias",
    "Padrões de conversação: CAMEL (2-agent), AutoGen (group), Swarm (handoff)",
    "Blackboard: espaço compartilhado, baixo acoplamento, escala com N",
    "Actor model: isolamento, concorrência sem locks, distribuição transparente",
    "Negociação: bargaining, auction, voting, mediator, evitar deadlock",
    "Protocolos: MCP (agent↔system) + A2A (agent↔agent) = complementares",
    "MetaGPT: SOPs estruturados > chat livre",
], "📖 Recapitulando 7 blocos: espectro A2A, padrões, blackboard, actor, negociação, protocolos, MetaGPT.\n💡 Analogia: resumo de curso de orquestra. Notas, partitura, seções, regente, afinação, padrão, apresentação.\n➡️ Checklist.", num=65, total=T)

s = content_slide("Checklist da Aula", [
    "☐ Definiu o espectro A2A (síncrono/assíncrono, topologias)",
    "☐ Apresentou schemas de mensagem com versionamento",
    "☐ Explicou 3 padrões de conversação (CAMEL, AutoGen, Swarm)",
    "☐ Demontrou blackboard vs mensagens diretas",
    "☐ Implementou actor model ou mostrou a demo",
    "☐ Discutiu negociação, conflito e deadlock",
    "☐ Comparou MCP vs A2A Protocol",
], "📖 Checklist de cobertura. Se algo faltou, mencione ou indique leitura.\n➡️ Quiz.", num=66, total=T)

# Quiz (67-70)
s = exercise_slide("Quiz — Pergunta 1", [
    "Quando blackboard é preferível a mensagens diretas?",
    "",
    "A) Quando há apenas 2 agentes",
    "B) Quando N é grande e contribuições são independentes",
    "C) Quando a ordem das mensagens é crítica",
    "D) Quando o acoplamento entre agentes deve ser máximo",
    "",
    "Resposta: B",
], "📖 Resposta B. Blackboard brilha com N grande. 2 agentes (A) = diretas. Ordem crítica (C) = diretas. Acoplamento máximo (D) = oposto do blackboard.\n➡️ Pergunta 2.", num=67, total=T, obj="Quiz")

s = exercise_slide("Quiz — Pergunta 2", [
    "Qual é a diferença entre handoff (Swarm) e delegação (supervisor)?",
    "",
    "A) Handoff é síncrono, delegação é assíncrona",
    "B) Handoff transfere controle (agente sai), delegação mantém supervisor",
    "C) Handoff é para 2 agentes, delegação para N",
    "D) Não há diferença",
    "",
    "Resposta: B",
], "📖 Resposta B. Handoff = transfere (sai). Delegação = delega (supervisor espera). Não é sincronicidade (A) nem número (C).\n➡️ Pergunta 3.", num=68, total=T, obj="Quiz")

s = exercise_slide("Quiz — Pergunta 3", [
    "MCP e A2A Protocol são:",
    "",
    "A) Competidores — um substitui o outro",
    "B) Complementares — MCP é agent↔system, A2A é agent↔agent",
    "C) A mesma coisa com nomes diferentes",
    "D) Ambos obsoletos",
    "",
    "Resposta: B",
], "📖 Resposta B. Complementares. MCP = agent↔system (tools). A2A = agent↔agent (agentes). Não competem (A), não são iguais (C), não obsoletos (D).\n➡️ Última pergunta.", num=69, total=T, obj="Quiz")

s = exercise_slide("Quiz — Pergunta 4", [
    "V/F: 'Actor model é mais lento que shared-state.'",
    "",
    "A) Verdadeiro — locks são mais rápidos",
    "B) Falso — em alta concorrência, actor model escala melhor sem locks",
    "C) Depende — sempre (contexto importa)",
    "D) Impossível determinar sem benchmark",
    "",
    "Resposta: B (afirmação geral é falsa)",
], "📖 Resposta B. Afirmação geral é falsa. Alta concorrência: actor model escala melhor. Overhead de serialização existe, mas compensa pela ausência de contenção de lock. C tem ponta de verdade (contexto importa), mas afirmação é falsa.\n➡️ Referências.", num=70, total=T, obj="Quiz")

s = content_slide("Conexão com Próximos Módulos + Referências", [
    "Próximos módulos:",
    "  ETHAGT10 — Padrões de Arquitetura Multi-Agente",
    "  ETHAGT11 — Event-Driven Agents",
    "  ETHAGT15 — Sociedades de Agentes",
    "",
    "Leitura obrigatória:",
    "  AutoGen (arXiv:2308.08155) · CAMEL (arXiv:2303.17760)",
    "  Hewitt, Actor Model (1973)",
    "",
    "Recomendado: Google A2A Protocol spec · OpenAI Swarm (repo)",
    "Survey: arXiv:2308.00352",
], "📖 ETHAGT09 é fundação para ETHAGT10 (topologias), 11 (event-driven), 15 (sociedades). Leiam os papers — detalhes que a aula não cobre.\n➡️ Encerramento.", num=71, total=T)

s = title_slide(
    "Projeto · Labs · Q&A",
    "Projeto: sistema de negociação entre agentes (convergência ≥ 80%)\nLab 1: Duas Arquiteturas (4h) · Lab 2: Actor Model com Handoffs (4h)\nPróxima aula: ETHAGT10 — Padrões de Arquitetura Multi-Agente",
    "ETHAGT09 · Obrigado!"
)
add_notes(s, "📖 Encerramento. Projeto = sistema de negociação. Labs complementam. Próxima aula: ETHAGT10.\n💡 Analogia: hoje aprenderam a reger a orquestra. Projeto = primeira apresentação. Labs = ensaios. ETHAGT10 aprofunda regência.\n❓ 'Perguntas?' (Se não houver: 'Qual parte foi menos clara?')\n⚠️ Se tempo: 'Um padrão de comunicação que você não conhecia.'\n➡️ FIM. Nos vemos na ETHAGT10.")

# ─── Salvar ───
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT09-slides.pptx")
prs.save(output_path)
print(f"✅ PPTX gerado: {output_path}")
print(f"   Total de slides: {len(prs.slides)}")
