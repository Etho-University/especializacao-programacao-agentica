#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT07: Knowledge Graphs & Vector Databases para Agentes
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj="", acronyms=""):
    if acronyms:
        add_textbox(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.35), acronyms, size=9, color=INFO)
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT07"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=18, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    if acronyms:
        add_textbox(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4), acronyms, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=30, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

T = 80

s = title_slide(
    "Knowledge Graphs & Vector Databases para Agentes",
    "Universidade Etho · Especialização em Programação Agêntica\nFase B — Infraestrutura Cognitiva · 30 h",
    "ETHAGT07"
)
add_notes(s, "📖 Bem-vindos. Esta aula é sobre a infraestrutura cognitiva dos agentes — vector databases e knowledge graphs. Vector DB é memória semântica; knowledge graph é memória relacional. Um agente maduro usa os dois.\n💡 Analogia: vector DB = memória de similaridade; grafo = memória de relacionamento.\n❓ 'Quem já usou vector DB? E knowledge graph?'\n⚠️ Alunos chegam achando que vector DB resolve tudo. Plantar a dúvida: similaridade ≠ raciocínio.\n➡️ Vamos aos objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: escolher, modelar e operar vector DBs e knowledge graphs",
    "",
    "Objetivos específicos:",
    "1. Comparar vector databases por cenário (Qdrant, Milvus, Weaviate, Chroma, pgvector)",
    "2. Modelar knowledge graphs para raciocínio estruturado",
    "3. Implementar GraphRAG (local + global)",
    "4. Construir pipelines híbridos (vector + grafo)",
    "5. Avaliar custo/latência/escala em volumes realistas"
], "📖 Cada objetivo é mensurável: comparar, modelar, implementar, construir, avaliar.\n💡 Analogia: serra e martelo não competem — vector e grafo são ferramentas cognitivas diferentes.\n❓ 'Qual objetivo parece mais desafiador?' (costuma ser GraphRAG ou híbrido)\n⚠️ GraphRAG NÃO é 'RAG com grafo' — é técnica específica (comunidades + sumarização).\n➡️ Onde estamos no mapa de competências.", 2, T, obj="Objetivos")

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → A (Avançado)",
    "C3 MCP & Tool Use → B (Básico)",
    "C4 Agent Memory → A (Avançado)",
    "C5 AgentOps & Avaliação → I (Intermediário)",
    "",
    "Vector DB = memória semântica do agente",
    "Knowledge Graph = memória relacional do agente"
], "📖 C1 e C4 chegam ao Avançado. Vector DB é memória semântica; grafo é memória relacional.\n💡 Analogia: um agente sem memória relacional é amnésico em causa-e-efeito.\n⚠️ 'Avançado em C1' ≠ saber muitos frameworks. É justificar arquitetura com evidência.\n➡️ Agenda da aula.", 3, T, obj="Competências", acronyms="AgentOps = Agent Operations  ·  MCP = Model Context Protocol")

s = content_slide("Agenda da Aula", [
    "Bloco 1 (45 min):",
    "  Abertura (8 min) — objetivos, motivação, contexto",
    "  Vector DBs (12 min) — ANN, HNSW, IVF, métricas, filtering",
    "  Comparativo (10 min) — 5 DBs, quando NÃO usar",
    "  Knowledge Graphs (15 min) — triplas, Neo4j, Cypher",
    "  Intervalo (5 min)",
    "",
    "Bloco 2 (45 min):",
    "  GraphRAG (15 min) — Microsoft, local vs global, DEMO",
    "  Híbridos (12 min) — vector + grafo, router, benchmark",
    "  Escala (8 min) — sharding, reindexação, observabilidade",
    "  Fechamento (10 min) — boas práticas, caso, quiz, Q&A"
], "📖 Dois blocos. Bloco 1: fundamentos separados. Bloco 2: GraphRAG, híbridos, escala.\n💡 Analogia: bloco 1 = ingredientes; bloco 2 = receita; fechamento = avaliação.\n⚠️ DEMO do GraphRAG (Slide 46) é o clímax — vector RAG falha, GraphRAG acerta.\n➡️ Por que estamos aqui? Quando similaridade não basta.", 4, T, obj="Agenda", acronyms="ANN = Approximate Nearest Neighbor")

s = comparison_slide(
    "Quando Similaridade Não Basta",
    "Vector Search", [
        "Encontra documentos SIMILARES por semântica",
        "Mas não sabe que 'Dipirona' e 'Warfarin' interagem",
        "Retorna textos SOBRE interações, não RACIOCINA sobre a cadeia",
        "Falha em perguntas multi-hop"
    ],
    "Knowledge Graph", [
        "Codifica RELACIONAMENTOS estruturados",
        "Dipirona → interage_com → Warfarin",
        "Encadeia: Paciente X → Dr. Silva → Dipirona → Warfarin",
        "Raciocínio multi-hop nativo"
    ],
    "📖 Vector search é ótimo para 'encontre documentos parecidos'. Mas 'quais medicações interagem com a do paciente X?' exige encadear 3 fatos. Vector não encadeia.\n💡 Analogia: vector = estante por tema; grafo = GPS com rotas.\n❓ 'No sistema de vocês: existe pergunta que exige encadear 2+ fatos?'\n⚠️ 'Prompt maior' ou 'mais chunks' não resolve multi-hop confiavelmente.\n➡️ A evolução do RAG até GraphRAG.", 5, T, obj="Motivação"
, acronyms="RAG = Retrieval-Augmented Generation")

s = content_slide("A Evolução do RAG", [
    "2020: RAG original (Lewis et al., arXiv:2005.11401)",
    "2022: Dense retrieval amadurece",
    "2023: RAG agêntico (ETHAGT06)",
    "2024: GraphRAG (Microsoft, arXiv:2404.16130)",
    "",
    "Confluência de 3 fatores:",
    "  Embeddings melhores e baratos",
    "  Graph databases acessíveis (Neo4j cloud)",
    "  LLMs bons em extração de entidades/relações",
    "",
    "Survey: Pan et al. (arXiv:2306.08302) — LLMs + KGs"
], "📖 RAG existe desde 2020, mas era chunk + generate. GraphRAG (2024) usa grafo + comunidades + sumarização para perguntas globais.\n💡 Analogia: vector RAG = atlas (página certa). GraphRAG = GPS (navega entre pontos).\n❓ 'Qual fator foi o gatilho mais recente?' (LLMs para extração automática)\n⚠️ GraphRAG não é 'novo' — KGs existem há décadas. A novidade é a automação via LLM.\n➡️ Bloco 1: como funciona um vector DB.", 6, T, obj="Contexto")

s = section_slide(1, "Vector DBs: Modelo Mental")
add_notes(s, "📖 Início do bloco de fundamentos. Vamos ver o que é um vector DB por dentro, como busca em milhões de vetores em ms, e as métricas que importam.\n➡️ O que é um vector DB?")

s = content_slide("O Que É um Vector DB?", [
    "Armazena embeddings (vetores densos) + metadata",
    "Index para busca ANN (não brute force)",
    "Query: vetor → top-k mais similares",
    "Componentes: collection → points (id, vector, payload) → index",
    "",
    "Diferença de DB relacional:",
    "  Relacional: busca por IGUALDADE (WHERE id = 42)",
    "  Vector DB: busca por SIMILARIDADE (top-k próximos)"
], "📖 Vector DB armazena vetores + metadata. A diferença: você busca por similaridade, não igualdade. Para isso em escala, usa ANN (não brute force).\n💡 Analogia: relacional = buscar livro pelo ISBN; vector DB = 'livros parecidos com este'.\n❓ 'Por que não brute force em tudo?' (O(N) não escala)\n⚠️ Vector DB ≠ Postgres com coluna de array. O index ANN é o que faz ser vector DB.\n➡️ O que armazenamos? Embeddings.", 8, T, obj="Vector DB")

s = content_slide("Embeddings: A Base", [
    "Texto → modelo de embedding → vetor denso (ex: 1536 dims)",
    "Semântica capturada na geometria do espaço",
    "Modelos: OpenAI text-embedding-3, Cohere, BGE, GTE",
    "Dimensões: trade-off riqueza vs custo de armazenamento/busca",
    "",
    "Aprofundamento em ETHAGT06",
    "Hoje: foco em COMO armazenar e buscar eficientemente"
], "📖 Embeddings = representação vetorial da semântica. Proximidade = semelhança. Dimensão é trade-off (riqueza vs custo). Aprofundado em ETHAGT06.\n💡 Analogia: embedding = CEP da semântica. Vector DB = agência dos correios que acha CEPs próximos.\n⚠️ Trocar de modelo sem reindexar = vetores incompatíveis (cosine vira lixo). Voltamos no Slide 66.\n➡️ Como buscar rápido em milhões? ANN.", 9, T, obj="Embeddings")

s = content_slide("ANN: Approximate Nearest Neighbor", [
    "Brute force: O(N) — inviável para milhões de vetores",
    "ANN: sacrifica exatidão por velocidade (recall vs latência)",
    "Duas famílias principais:",
    "  HNSW (graph-based) — dominante",
    "  IVF (cluster-based)",
    "Trade-off: recall ↑ → latência ↑",
    "Controle: ef_search (HNSW), nprobe (IVF)"
], "📖 Brute force compara query com TODOS os vetores — inviável em escala. ANN sacrifica recall por velocidade. 90% recall em ms > 100% recall em segundos.\n💡 Analogia: brute force = ler todas as capas; ANN = ir direto na seção do gênero.\n❓ 'Quando aceita 90% recall?' (Quase sempre em RAG — LLM tolera ruído)\n⚠️ Recall baixo ≠ ruim. Recall@10 de 90% é excelente em RAG.\n➡️ Aprofundar HNSW.", 10, T, obj="ANN")

s = content_slide("HNSW: Hierarchical Navigable Small World", [
    "Grafo hierárquico de camadas:",
    "  Top (sparse, fast) → Bottom (dense, precise)",
    "Busca: começa no topo, desce camada por camada",
    "Vantagem: recall alto, latência baixa",
    "Desvantagem: construção lenta, memória",
    "Usado por: Qdrant, Milvus, Weaviate, pgvector (opcional)",
    "Fonte: Malkov & Yashunin, arXiv:1603.09320"
], "📖 HNSW = grafo em camadas. Topo esparsa (pousar perto da região), base densa (refinar). Busca desce camada por camada.\n💡 Analogia: fly, drive, walk dos mapas. Topo = voar para a cidade; meio = dirigir ao bairro; base = caminhar à porta.\n❓ 'Por que construção é lenta?' (Cada inserção conecta arestas em múltiplas camadas)\n⚠️ HNSW é APROXIMADO, não exato. Mas com ef_search alto, recall ≈ 100%.\n➡️ Alternativa: IVF.", 11, T, obj="HNSW")

s = content_slide("IVF: Inverted File Index", [
    "K-means divide o espaço em clusters (Voronoi cells)",
    "Busca: encontra cluster mais próximo → busca dentro dele",
    "nprobe: quantos clusters vizinhos explorar",
    "Vantagem: construção mais rápida, bom para datasets grandes",
    "Desvantagem: recall menor que HNSW para mesmo custo",
    "Combinável com PQ (Product Quantization) para compressão"
], "📖 IVF divide o espaço em clusters (k-means). Busca vai ao cluster mais próximo. nprobe controla quantos vizinhos explorar.\n💡 Analogia: sistema de CEPs. Não busca no país todo — vai ao CEP do query + vizinhos (nprobe).\n⚠️ Deixar nprobe=1 (default) = recall baixo. nprobe é o parâmetro #1 de tuning de IVF.\n➡️ Independente do index, você precisa de uma métrica.", 12, T, obj="IVF")

s = content_slide("Métricas: Cosine, Dot, Euclidean", [
    "Cosine: ângulo entre vetores — ignora magnitude",
    "Dot product: cosine × magnitudes — sensível a norma",
    "Euclidean (L2): distância geométrica — magnitude + direção",
    "",
    "Insight chave:",
    "Se embeddings são normalizados → cosine = dot product",
    "Por isso muita gente normaliza e usa dot (mais rápido)"
], "📖 Métrica define 'similar'. Cosine = ângulo. Dot = ângulo × magnitude. Euclidean = distância. Se normalizados, cosine = dot.\n💡 Analogia: cosine = comparar DIREÇÃO de setas. Dot = direção + força. Euclidean = distância das pontas.\n⚠️ Misturar métricas entre index e query = recall destruído. Sempre consistente.\n➡️ Como escolher na prática?", 13, T, obj="Métricas")

s = content_slide("Quando Cada Métrica?", [
    "Cosine: default para texto (embeddings não normalizados)",
    "Dot: quando embeddings já são normalizados (mais rápido)",
    "Euclidean: quando magnitude importa (features multimodais)",
    "",
    "Regra prática: se não sabe, use cosine",
    "Normalize embeddings → mude para dot (mais barato)",
    "",
    "Armadilha: misturar métricas = recall destruído"
], "📖 Regra: cosine default para texto. Se normalizou, use dot (mais rápido). Euclidean só quando magnitude importa (raro em texto).\n💡 Analogia: cosine = escolha segura (dirigir devagar). Dot com normalização = otimização (velocidade ideal).\n❓ 'Cosine vs dot — quando a diferença importa?' (Quando NÃO normalizados com magnitudes variáveis)\n⚠️ Usar cosine em embeddings normalizados = desperdício de computação.\n➡️ Além da métrica: filtering de metadata.", 14, T, obj="Métricas")

s = content_slide("Metadata Filtering (Payloads)", [
    "Payload: JSON anexado a cada vetor ({source, page, date})",
    "Pre-filtering: filtra ANTES do ANN (preciso, pode ser lento)",
    "Post-filtering: filtra DEPOIS do ANN (rápido, recall pode cair)",
    "Qdrant: filtering otimizado DENTRO do index (vantagem chave)",
    "Padrão: vector + filter → top-k"
], "📖 Quase sempre você busca + filtra. Pre-filter filtra antes (preciso, lento se muito restritivo). Post-filter filtra depois (rápido, mas recall pode desabar). Qdrant otimiza o filter dentro do index.\n💡 Analogia: pre-filter = peneirar farinha antes de medir. Post-filter = medir tudo e descartar grumos. Qdrant = peneira integrada na medida.\n⚠️ Post-filter com filtro restritivo = top-10 vira top-2. Teste recall com filtros reais.\n➡️ Retrieval híbrido sparse + dense.", 15, T, obj="Filtering")

s = content_slide("Híbrido: Sparse + Dense (BM25 + Dense)", [
    "Dense: semântica (sinônimos, paráfrase) — falha em termos exatos",
    "Sparse (BM25/SPLADE): termos exatos — falha em semântica",
    "Híbrido: combina ambos com fusão de scores",
    "Reciprocal Rank Fusion (RRF): simples e eficaz",
    "Qdrant 1.10+: sparse vectors nativos"
], "📖 Dense acerta sinônimos, falha em termos exatos. Sparse acerta termos exatos, falha em paráfrase. Híbrido une os dois. RRF funde rankings sem calibrar pesos.\n💡 Analogia: dense = tradutor (entende sinonímia). Sparse = índice de livro (acerta a página). Híbrido usa os dois.\n⚠️ Tentar ponderar scores (0.7*dense + 0.3*sparse) = semanas tunando. RRF evita isso.\n➡️ Comparativo dos 5 vector DBs.", 16, T, obj="Híbrido Sparse+Dense", acronyms="BM25 = Best Matching 25")

s = section_slide(2, "Comparativo de Vector DBs")
add_notes(s, "📖 Vamos comparar os 5 vector DBs canônicos. Mensagem central: NÃO existe 'melhor' — existe 'melhor para seu cenário'.\n➡️ O landscape.")

s = content_slide("O Landscape: 5 Contenders", [
    "5 vector DBs canônicos:",
    "  Qdrant — Rust, filtering forte",
    "  Milvus — escala massiva, distribuído",
    "  Weaviate — modules integrados",
    "  Chroma — simplicidade, prototipagem",
    "  pgvector — extensão do PostgreSQL",
    "",
    "Eixos: linguagem, escala, filtering, ecosystem, maturidade",
    "Não existe 'melhor' — existe 'melhor para seu cenário'"
], "📖 5 DBs, cada um brilha em um cenário. Memorizem o TRADE-OFF, não features.\n💡 Analogia: Chroma = Uber (simples). pgvector = picape (já tem na garagem). Qdrant = esportivo (rápido). Milvus = caminhão (escala). Weaviate = van (integrado).\n➡️ Qdrant.", 18, T, obj="Landscape")

s = content_slide("Qdrant: Filtering Forte", [
    "Escrito em Rust: performance e safety",
    "Filtering otimizado no index (não é post-filter)",
    "Sparse vectors nativos (híbrido sem infra extra)",
    "Multi-tenancy via payloads",
    "Melhor quando: filtering complexo é requisito"
], "📖 Qdrant = Rust + filtering-first. O filter é otimizado DENTRO do index HNSW. Sparse vectors nativos desde 1.10.\n💡 Analogia: Qdrant é o DB 'filtering-first'. Se seu caso tem muitos filtros, brilha.\n⚠️ Escolher Qdrant 'porque é Rust' = escolher por hype. Escolha pelo filtering otimizado.\n➡️ Milvus.", 19, T, obj="Qdrant")

s = content_slide("Milvus: Escala", [
    "Arquitetura distribuída desde o início (cloud-native)",
    "Escala horizontal para bilhões de vetores",
    "Suporte a múltiplos index types (HNSW, IVF, DiskANN)",
    "Melhor quando: escala massiva é o requisito #1",
    "Trade-off: complexidade operacional maior"
], "📖 Milvus = distribuído nativo (QueryNode, DataNode, Proxy). Escala para bilhões. Custo: complexidade operacional.\n💡 Analogia: Milvus = big rig. Carga pesada (bilhões)? É ele. 1M de vetores? Overkill.\n⚠️ Escolher Milvus por 'escala futura' sem volume hoje = complexidade que não se paga.\n➡️ Weaviate.", 20, T, obj="Milvus")

s = content_slide("Weaviate: Modules", [
    "Modules integrados: generative, Q&A, multi-modal",
    "Schema-first: define classes e propriedades",
    "GraphQL API nativa",
    "Melhor quando: pipeline RAG completo sem orquestrar 5 serviços",
    "Trade-off: menos flexível para casos não previstos"
], "📖 Weaviate = batteries-included. Modules de geração, Q&A, multi-modal. Pipeline RAG sem orquestrar Qdrant + LangChain + Cohere.\n💡 Analogia: Weaviate = carro com GPS, ar, banco elétrico de fábrica. Prático. Customizar o motor = mais difícil.\n➡️ Chroma.", 21, T, obj="Weaviate")

s = content_slide("Chroma: Simplicidade", [
    "API minimalista: client.add(docs) → client.query(text)",
    "Embutido no processo Python (sem servidor separado)",
    "Melhor quando: prototipagem, MVP, dev local",
    "Trade-off: não é production-ready para escala (ainda)"
], "📖 Chroma = importe e use. add() e query() em uma linha. Embutido no Python, sem Docker. Perfeito para protótipo/MVP.\n💡 Analogia: Chroma = SQLite dos vector DBs. Embutido, simples. Para produção com escala, migre.\n⚠️ Começar com Chroma e levar para produção sem migrar = funciona até parar. Plano de migração desde o dia 1.\n➡️ pgvector.", 22, T, obj="Chroma", acronyms="MVP = Minimum Viable Product")

s = content_slide("pgvector: Operacional", [
    "Extensão do PostgreSQL: vector DB dentro do DB que você já tem",
    "HNSW e IVFFlat disponíveis",
    "ACID, transações, JOINs com vetores",
    "Melhor quando: já usa Postgres e não quer adicionar infra",
    "Trade-off: não escala como solução dedicada para bilhões"
], "📖 pgvector = extensão do Postgres. Vector DB + JOINs relacionais + ACID, tudo no mesmo DB. Para ~1M de vetores, resolve sem nova infra.\n💡 Analogia: pgvector = picape que já está na garagem. Não troque de veículo se a picada aguenta a carga.\n❓ 'Quando pgvector é melhor que Qdrant?' (Já tem Postgres + quer ACID + volume até ~1M)\n⚠️ Subestimar pgvector = erro. Para 90% das aplicações RAG, resolve.\n➡️ Tabela comparativa.", 23, T, obj="pgvector", acronyms="ACID = Atomicity Consistency Isolation Durability")

s = content_slide("Tabela Comparativa", [
    "              Qdrant    Milvus    Weaviate   Chroma    pgvector",
    "Linguagem     Rust      Go        Go         Python    C (ext PG)",
    "Escala        Alta      Muito alta Alta       Baixa     Média",
    "Filtering     ★★★★★     ★★★       ★★★        ★★        ★★★★",
    "Híbrido       Nativo    Add-on    Add-on      Não       Add-on",
    "Multi-tenant  ★★★★★     ★★★       ★★★        ★★        ★★★★",
    "Custo ops     Médio     Alto      Médio       Baixo      Baixo",
    "Curva         Média     Íngreme   Média       Suave      Suave",
    "Production    ★★★★★     ★★★★★     ★★★★★      ★★★        ★★★★★"
], "📖 Tabela dos trade-offs. Dedicados (Qdrant, Milvus) escalam mais, custam mais em complexidade. Embutidos (Chroma, pgvector) são simples com teto.\n❓ 'Para MVP, qual? E para 1M usuários?' (MVP: Chroma/pgvector. 1M: Qdrant/Milvus)\n⚠️ Escolher por popularidade no GitHub ≠ fit. Decida pelos trade-offs que importam.\n➡️ Mas antes: você precisa MESMO de vector DB?", 24, T, obj="Comparativo")

s = content_slide("Quando NÃO Usar Vector DB", [
    "Busca exata (igualdade): DB relacional é melhor",
    "Dados estruturados com relacionamentos: knowledge graph",
    "Volume pequeno (< 10k docs): brute force em memória",
    "Latência ultra-baixa: cache + keyword search",
    "",
    "Regra: não adicione vector DB se Postgres resolve",
    "Vector DB é para BUSCA POR SIMILARIDADE EM ESCALA"
], "📖 Vector DB é on-trend, mas nem tudo precisa. Igualdade? Relacional. Relacionamentos? Grafo. <10k docs? Brute force. Vector DB é para similaridade em escala.\n💡 Analogia: comprar trator para cuidar de jardim. Ótimo para fazenda; jardim = enxada.\n⚠️ Adicionar vector DB 'porque é moderno' quando JOIN resolve = dívida técnica. Comece simples.\n➡️ Knowledge Graphs — próximo bloco.", 25, T, obj="Quando NÃO usar")

s = section_slide(3, "Knowledge Graphs")
add_notes(s, "📖 Segundo paradigma de armazenamento cognitivo: knowledge graphs. Vector = similaridade; grafo = relacionamento. Vamos ver triplas, Neo4j, Cypher, e construção via LLM.\n➡️ O que é um KG?")

s = comparison_slide(
    "O Que É um Knowledge Graph?",
    "Vector DB", [
        "Similaridade semântica",
        "Proximidade no espaço vetorial",
        "'Documentos parecidos com X'",
        "Forte em lookup semântico"
    ],
    "Knowledge Graph", [
        "Relacionamentos estruturados",
        "Travessia de nós e arestas",
        "'Entidades conectadas a X via Y'",
        "Forte em raciocínio multi-hop"
    ],
    "📖 Vector responde 'o que é similar a X?'. Grafo responde 'o que está conectado a X, e como?'. São complementares.\n💡 Analogia: vector = biblioteca por tema (livros parecidos). Grafo = árvore genealógica (quem é parente de quem).\n⚠️ Não é 'escolha binária'. A escolha madura é pipeline híbrido (Seção F).\n➡️ O modelo de dados: a tripla.", 27, T, obj="Knowledge Graph"
)

s = content_slide("Triplas: Sujeito → Predicado → Objeto", [
    "Tripla: (sujeito, predicado, objeto)",
    "",
    "Exemplos:",
    "  ('Dipirona', 'interage_com', 'Warfarin')",
    "  ('Dr. Silva', 'prescreveu', 'Dipirona')",
    "  ('Paciente X', 'consultou', 'Dr. Silva')",
    "",
    "Encadeamento (multi-hop!):",
    "  Paciente X → Dr. Silva → Dipirona → Warfarin → RISCO!"
], "📖 KG = conjunto de triplas. A mágica está no ENCADEAMENTO: Paciente X → Dr. Silva → Dipirona → Warfarin = risco. Isso vector não faz.\n💡 Analogia: triplas = frases de 3 palavras. KG = texto feito dessas frases; significado emerge das conexões.\n❓ 'No domínio de vocês: quais 3 triplas seriam mais valiosas?'\n⚠️ Predicados vagos ('RELATED_TO') destroem raciocínio. Predicados específicos = queries poderosas.\n➡️ Triplas são esqueleto; propriedades são carne.", 28, T, obj="Triplas")

s = content_slide("Propriedades e Labels", [
    "Nó: label (tipo) + propriedades (key-value)",
    "Aresta: tipo + propriedades",
    "",
    "Exemplos:",
    "  Nó 'Dipirona' {tipo: Medicamento, classe: analgésico}",
    "  Aresta 'interage_com' {severidade: moderada, fonte: FDA}",
    "",
    "Labeled Property Graph (LPG) — modelo do Neo4j"
], "📖 No LPG, nó tem LABEL (tipo) + PROPRIEDADES. Aresta tem TIPO + PROPriedades. Permite queries ricas: 'interações de severidade alta pela FDA em 2024'.\n💡 Analogia: label = espécie (Medicamento?). Propriedade = atributo (classe?). Aresta com propriedade = verbo com advérbio.\n⚠️ Sobrecarregar nós com propriedades que deveriam ser arestas = confusão. Se é conceito com relações próprias, vire nó.\n➡️ Neo4j e Cypher.", 29, T, obj="Propriedades")

s = content_slide("Neo4j e Cypher", [
    "Neo4j: graph database nativa (Labeled Property Graph)",
    "Cypher: linguagem declarativa para grafos (como SQL para relacional)",
    "Alternativas: Memgraph (in-memory), ArangoDB (multi-modelo)",
    "RDF/SPARQL: paradigma W3C (Linked Open Data, ontologias)",
    "Foco do curso: Neo4j + Cypher (ecossistema, maturidade)"
], "📖 Neo4j é a graph DB dominante. Cypher é declarativa (descreve o padrão, não o algoritmo). Alternativas: Memgraph, ArangoDB. RDF/SPARQL é outro paradigma (W3C).\n💡 Analogia: Cypher é para grafos o que SQL é para tabelas. Você descreve o PADRÃO; o engine encontra.\n➡️ Cypher na prática.", 30, T, obj="Neo4j/Cypher")

s = code_slide(
    "Query Cypher: 'Co-autores de X'",
    'MATCH (a:Author {name: "X"})-[:WROTE]->(p:Paper)\n<-[:WROTE]-(b:Author)\nWHERE b <> a\nRETURN DISTINCT b.name AS coauthor,\n       count(p) AS papers\nORDER BY papers DESC\nLIMIT 10',
    "📖 Encontra co-autores de X: Author → WROTE → Paper ← WROTE ← Author. Padrão visual: nós entre (), arestas entre -[]->.\n💡 Analogia: MATCH = ctrl+F no grafo. Descreve o padrão; o engine encontra ocorrências.\n❓ 'Como adaptar para co-autores de co-autores (2-hop)?' (Mais um nível de travessia)\n⚠️ Esquecer WHERE b <> a = o próprio autor aparece como co-autor.\n➡️ Cypher é linguagem. A MODELAGEM é o que importa.", 31, T, obj="Cypher"
)

s = comparison_slide(
    "Modelagem para Raciocínio",
    "Anti-pattern: tudo RELATED_TO", [
        "Nós genéricos, arestas vagas",
        "Não distingue interage de substitui",
        "Queries fracas, raciocínio pobre",
        "Sem contexto na aresta"
    ],
    "Bom: relações tipadas", [
        "Labels específicos (Author, Paper)",
        "Relações tipadas (WROTE, CITES, DEPENDE_DE)",
        "Propriedades em nós e arestas",
        "Raciocínio multi-hop poderoso"
    ],
    "📖 Modelagem determina o que você raciocina. Tudo 'RELATED_TO' = não distingue relações. Modele orientado a QUERIES.\n💡 Analogia: como modelar DB relacional. Tabela genérica 'Coisa' = queries fracas. Tabelas específicas = queries poderosas.\n⚠️ Modelar 'para parecer com a realidade' sem pensar nas queries = armadilha. Pergunte 'quais perguntas?' primeiro.\n➡️ Como construir grafo a partir de texto? NER.", 32, T, obj="Modelagem"
, acronyms="NER = Named Entity Recognition")

s = content_slide("Extração de Entidades com LLM (NER)", [
    "Texto → LLM com prompt de NER → entidades estruturadas",
    "Prompt: 'Extraia todas as entidades com tipo e propriedades'",
    "Output: [{name: Dipirona, type: Medicamento}, ...]",
    "Desafios: entidades ambíguas, resolução de coreferência",
    "Ferramentas: LangChain LLMGraphTransformer, GLiNER, spaCy"
], "📖 NER = extração de entidades via LLM com prompt estruturado. Desafios: ambiguidade (Dipirona = Dipirona sódica?), coreferência.\n💡 Analogia: NER = estudante sublinhando termos médicos no prontuário. LLM escala para milhões de docs.\n⚠️ Confiar cegamente no NER do LLM = alucinação. Schema restrito + validação + amostragem humana.\n➡️ Entidades são nós. Relações são arestas.", 33, T, obj="NER", acronyms="LLM = Large Language Model")

s = content_slide("Extração de Relações com LLM", [
    "Dadas 2 entidades → LLM classifica a relação",
    "Prompt: 'Qual a relação entre A e B? Escolha: [interage_com, prescreveu, ...]'",
    "Output: triplas (sujeito, predicado, objeto)",
    "Desafio crítico: ALUCINAÇÃO de relações inexistentes",
    "Mitigação: schema restrito + validação + confiança na aresta"
], "📖 Dadas 2 entidades, LLM classifica a relação. Output = tripla. Desafio crítico: ALUCINAÇÃO. LLM inventa relações.\n💡 Analogia: jornalista reportando (LLM) + editor verificando (validação). Sem editor = notícia falsa no grafo.\n⚠️ Extração sem validação = grafo cheio de alucinações. Agente raciocina sobre falso com confiança. Lineage permite auditar.\n➡️ Com entidades e relações, podemos inferir.", 34, T, obj="Relações")

s = content_slide("Inferência no Grafo", [
    "Explícito: 'A interage com B' e 'B interage com C'",
    "Inferido: 'A pode interagir com C' (transitividade — com cautela!)",
    "Cypher: pathfinding (shortestPath, allShortestPaths)",
    "Multi-hop: 'Todos os medicamentos que X toma indiretamente'",
    "Padrão: LLM formula Cypher → grafo executa → LLM interpreta"
], "📖 Grafos permitem raciocínio além do explícito. Cypher tem pathfinding. Padrão moderno: LLM gera query Cypher, grafo executa, LLM interpreta.\n💡 Analogia: detetive conectindo pistas. LLM = detetive (formula investigação). Grafo = quadro de investigação (executa busca).\n⚠️ Inferir por transitividade sem cautela = armadilha. 'A interage B' + 'B interage C' NÃO garante 'A interage C'. Use como HIPÓTESE.\n➡️ Exercício de modelagem.", 35, T, obj="Inferência")

s = exercise_slide(
    "Cardinalidade de um KG Empresarial",
    [
        "Pergunta: Qual a cardinalidade típica de um KG empresarial?",
        "",
        "Quantas entidades (nós)? Quantas relações (arestas)?",
        "Qual a densidade (relações por entidade)?",
        "",
        "Pistas: depende do domínio (farmacêutica vs jurídico vs RH)",
        "Ordem de magnitude: 10k-100k entidades, 50k-500k relações",
        "Densidade típica: 5-15 relações por entidade",
        "",
        "Discutam em duplas (2 min). Depois compartilhamos."
    ],
    "📖 Cardinalidade = tamanho do 'cérebro' do agente. Mais relações = mais raciocínio, mais custo de manutenção.\n❓ 'No domínio de vocês: estimem entidades e relações.' (2 min em duplas)\n⚠️ Subestimar cardinalidade = surpresa com volume. Para 50k docs, espere 100k+ entidades.\n➡️ Intervalo. Depois: GraphRAG.", 36, T
)

s = section_slide(4, "GraphRAG")
add_notes(s, "📖 Início do bloco GraphRAG. Técnica da Microsoft (2024): knowledge graph + comunidades + sumarização hierárquica. Resolve perguntas GLOBAIS que vector RAG não consegue.\n➡️ Por que GraphRAG?")

s = comparison_slide(
    "Por Que GraphRAG?",
    "Vector RAG", [
        "Recupera chunks isolados",
        "Perde visão global do corpus",
        "Falha em 'quais os temas principais?'",
        "Falha em multi-hop (3+ hops)"
    ],
    "GraphRAG", [
        "Constrói grafo de entidades",
        "Detecta comunidades temáticas",
        "Sumariza hierarquicamente",
        "Responde perguntas globais"
    ],
    "📖 Vector RAG recupera chunks isolados. Falha em perguntas globais ('quais os temas principais?'). GraphRAG constrói grafo + comunidades + sumários para responder.\n💡 Analogia: vector RAG = índice remissivo (página certa). GraphRAG = sumário de capítulos + índice analítico (visão global).\n⚠️ GraphRAG não SUBSTITUI vector RAG. São complementares. Híbrido usa os dois.\n➡️ Arquitetura do Microsoft GraphRAG.", 38, T, obj="GraphRAG"
)

s = content_slide("Microsoft GraphRAG: Visão Geral", [
    "Pipeline de CONSTRUÇÃO (offline):",
    "  1. Text chunking",
    "  2. Extração de entidades + relações (LLM)",
    "  3. Construção do knowledge graph",
    "  4. Detecção de comunidades (Leiden algorithm)",
    "  5. Sumarização hierárquica de comunidades (LLM)",
    "",
    "Pipeline de QUERY (online):",
    "  Local search: entidades + vizinhança (perguntas específicas)",
    "  Global search: map-reduce sobre sumários (perguntas amplas)",
    "",
    "Fonte: Edge et al., arXiv:2404.16130"
], "📖 Duas fases. CONSTRUÇÃO: chunking → NER+relações → grafo → comunidades (Leiden) → sumários hierárquicos. QUERY: Local (subgrafo) ou Global (map-reduce sobre sumários).\n💡 Analogia: construção = escrever enciclopédia (ler tudo, organizar por tema, escrever verbetes). Query = consultar enciclopédia.\n⚠️ GraphRAG não é 'build once'. Grafo envelhece. Manutenção é recorrente e cara.\n➡️ A etapa mais mágica: detecção de comunidades.", 39, T, obj="GraphRAG Pipeline")

s = content_slide("Construção do Grafo de Comunidades", [
    "Grafo de entidades → algoritmo Leiden → comunidades",
    "Comunidade: grupo de entidades densamente conectadas entre si",
    "Hierarquia: comunidades grandes contêm subcomunidades",
    "Análogo a 'capítulos de um livro' — estrutura temática emergente",
    "Resultado: grafo particionado em níveis hierárquicos"
], "📖 Leiden detecta comunidades: entidades densamente conectadas entre si. A estrutura EMERGE dos dados — você não define os temas. É hierárquica (comunidades contêm subcomunidades).\n💡 Analogia: Amazon categorizando produtos automaticamente. Produtos comprados juntos formam clusters. Leiden faz isso com entidades.\n⚠️ Comunidades não são 'perfeitamente temáticas'. Leiden é heurístico. Sumário hierárquico suaviza, mas valide com amostragem.\n➡️ Comunidades viram sumários.", 40, T, obj="Comunidades")

s = content_slide("Sumarização Hierárquica", [
    "Para cada comunidade (nível mais baixo): LLM gera sumário",
    "Sumários de subcomunidades → LLM gera sumário da comunidade pai",
    "Resultado: árvore de sumários do específico ao geral",
    "Vantagem: responde perguntas globais sem ler todos os chunks",
    "Custo: uma chamada de LLM por comunidade (caro na construção)"
], "📖 Cada comunidade na folha → LLM gera sumário. Subcomunidades agregadas → sumário do pai. Resultado: árvore do específico ao geral. Consulta a raiz para visão global.\n💡 Analogia: relatório executivo hierárquico. Analista (folha) → gerente (intermediário) → diretor (raiz). CEO lê só a raiz.\n⚠️ Subestimar custo: 1000 comunidades = 1000 chamadas só na sumarização. Voltamos no Slide 48.\n➡️ Como consultar? Dois modos.", 41, T, obj="Sumarização")

s = content_slide("Local Search", [
    "Input: pergunta do usuário",
    "Passo 1: identificar entidades relevantes na pergunta (LLM)",
    "Passo 2: expandir vizinhança no grafo (entidades + relações + chunks)",
    "Passo 3: LLM gera resposta com o contexto local",
    "Melhor para: perguntas específicas sobre entidades conhecidas",
    "Exemplo: 'Quem são os co-autores de X?'"
], "📖 Local search = zoom in no grafo ao redor das entidades da pergunta. (1) identificar entidades; (2) expandir vizinhança; (3) LLM responde com contexto local.\n💡 Analogia: perguntar a especialista sobre tópico específico. Ela vai direto à seção relevante, pega os livros, responde.\n⚠️ Usar local para pergunta global = resposta parcial. Local é para entidades específicas.\n➡️ Global search.", 42, T, obj="Local Search")

s = content_slide("Global Search", [
    "Input: pergunta ampla ('Quais temas principais do corpus?')",
    "Map: para cada comunidade, LLM gera resposta parcial",
    "Reduce: LLM agrega respostas parciais em resposta final",
    "Melhor para: perguntas que exigem visão holística",
    "Exemplo: 'Quais as tendências de pesquisa nesta área?'"
], "📖 Global search = map-reduce. MAP: LLM por comunidade gera resposta parcial. REDUCE: LLM agrega tudo em resposta final. Mais lento, mas responde o que vector não consegue.\n💡 Analogia: censo. Não pergunta a uma pessoa (local). Pergunta a cada região (map), agrega em relatório nacional (reduce).\n⚠️ Global é caro (N chamadas). Use para perguntas genuinamente globais; para específicas, local.\n➡️ Como escolher?", 43, T, obj="Global Search")

s = content_slide("Local vs Global: Quando Cada", [
    "Local: pergunta menciona entidades específicas → foco no subgrafo",
    "Global: pergunta é temática/ampla → síntese entre comunidades",
    "Híbrido: começa local, escala para global se insuficiente",
    "",
    "Heurística: 'Esta pergunta precisa de 1 entidade ou de todo o corpus?'",
    "",
    "Na dúvida: comece local (mais barato), escale se insuficiente"
], "📖 Critério: entidades específicas → Local. Temática/ampla → Global. Heurística: '1 entidade ou corpus todo?'\n💡 Analogia: Local = 'o que o João faz?' (vá ao João). Global = 'profissões mais comuns na cidade?' (precisa de todos).\n➡️ Quando GraphRAG supera vector RAG claramente?", 44, T, obj="Local vs Global")

s = comparison_slide(
    "GraphRAG vs Vector RAG",
    "Vector RAG", [
        "Melhor: lookup factual, 1 hop",
        "Latência baixa",
        "Custo baixo",
        "Construção trivial (só embeddings)"
    ],
    "GraphRAG", [
        "Melhor: multi-hop, visão global",
        "Latência maior (global = map-reduce)",
        "Custo maior",
        "Construção 10-100x mais cara"
    ],
    "📖 Vector ganha em lookup factual. GraphRAG ganha em multi-hop/global. Mas GraphRAG é 10-100x mais caro na construção. Regra: GraphRAG quando VALOR do raciocínio justifica CUSTO.\n💡 Analogia: vector = sedã (eficiente, 90% dos trajetos). GraphRAG = helicóptero (caro, mas vê o que sedã não vê).\n⚠️ Implementar GraphRAG 'por ser moderno' sem perguntas que exigem multi-hop = custo não se paga.\n➡️ DEMO ao vivo!", 45, T, obj="GraphRAG vs Vector RAG"
)

s = content_slide("DEMO: GraphRAG em Neo4j", [
    "Referência: 05-Labs/ETHAGT07/Lab2-GraphRAG-Neo4j",
    "",
    "Passo 1: carregar corpus técnico",
    "Passo 2: extrair entidades/relações com LLM",
    "Passo 3: construir grafo no Neo4j",
    "Passo 4: query Cypher multi-hop",
    "",
    "Momento-chave: vector RAG falha na mesma pergunta, GraphRAG acerta",
    "Se API falhar: screenshots pré-gravados do Neo4j Browser"
], "📖 DEMO clímax. Mesma pergunta com vector RAG (falha) e GraphRAG (acerta, com path explicável). Se API falhar, screenshots.\n💡 Analogia: antes e depois. Vector = informações soltas. GraphRAG = cadeia encadeada.\n⚠️ Demo pode demorar (extração é lenta). Use corpus pequeno pré-carregado. Objetivo: ver a query multi-hop, não a construção.\n➡️ Vamos discutir o que vimos.", 46, T, obj="DEMO"
)

s = exercise_slide(
    "Pergunta da DEMO",
    [
        "Discutam em duplas (2 min):",
        "",
        "1. O GraphRAG errou alguma coisa na demo? Onde?",
        "2. Qual o custo aproximado da construção do grafo (em tokens)?",
        "3. Se o corpus mudar (novos docs), quanto custa reindexar?",
        "",
        "Dicas:",
        "Erros típicos: relações alucinadas, entidades ambíguas",
        "Custo: estimar chamadas LLM (extração + relações + sumários)",
        "Reindexação: incremental é mais barato, mas entity resolution custa"
    ],
    "📖 Reflexão em duplas. GraphRAG é poderoso mas FALÍVEL — alucinações, comunidades ruidosas, sumários que perdem nuance.\n❓ 'Em duplas: onde o GraphRAG pode ter errado?' (2 min)\n⚠️ Sair da demo achando que GraphRAG é 'perfeito' = armadilha. Auditoria e lineage são essenciais.\n➡️ O custo é real. Vamos quantificar.", 47, T
)

s = content_slide("Custos de Construção", [
    "Extração de entidades: 1 chamada LLM por chunk",
    "Extração de relações: 1 chamada LLM por par de entidades",
    "Sumarização de comunidades: 1 chamada LLM por comunidade",
    "",
    "Ordem de magnitude (1000 docs):",
    "  ~5k-15k chamadas LLM no total",
    "  Custo: $50-$500 dependendo do modelo"
], "📖 Para 1000 docs: ~1000 extrações + ~3000 relações + ~1500 sumários = ~5500 chamadas. Com GPT-4o: $50-$500. Upfront, não recorrente (até manutenção).\n💡 Analogia: construir biblioteca física. Fundação + prateleiras + catalogação = upfront. Manutenção = contínua.\n❓ 'Por que GraphRAG é caro? Justifique o custo.' (Cada etapa = chamadas LLM; valor = raciocínio multi-hop/global)\n➡️ E a manutenção? Também custa.", 48, T, obj="Custos")

s = content_slide("Custos de Manutenção", [
    "Novos documentos: extrair + atualizar comunidades",
    "Atualização incremental: mais barato que reconstrução, mas complexo",
    "Drift de entidades: 'Dipirona' = 'Dipirona sódica'? → merge",
    "Comunidades mudam: re-sumarização necessária",
    "Lineage: rastrear qual doc originou qual tripla"
], "📖 GraphRAG não é 'build once'. Novos docs → extração → comunidades mudam → re-sumarização. Atualização incremental é complexa (entity resolution).\n💡 Analogia: jardim. Não planta uma vez e esquece — rega, poda, remove ervas. GraphRAG precisa manutenção contínua.\n⚠️ Subestimar manutenção = grafo degrada em 3 meses. Planeje custo de manutenção desde o início.\n➡️ Vale a pena? Depende do domínio.", 49, T, obj="Manutenção")

s = exercise_slide(
    "Vale a Pena?",
    [
        "Para qual domínio GraphRAG vale o custo de construção?",
        "E para qual domínio é overkill?",
        "",
        "Critério: valor do raciocínio multi-hop vs custo de manutenção",
        "",
        "Exemplos:",
        "  Farmacêutica (interações medicamentosas) → VALE",
        "  Jurídico (precedências, referências) → VALE",
        "  Blog pessoal → NÃO VALE (overkill)",
        "",
        "Discutam: no domínio de vocês, vale?"
    ],
    "📖 GraphRAG vale quando valor do raciocínio é ALTO e custo é SUSTENTÁVEL. Farmacêutica/jurídico = vale. Blog = não vale.\n💡 Analogia: comprar vs alugar. GraphRAG = comprar (investimento alto, custo de manutenção, retorno a longo prazo).\n❓ 'No domínio de vocês: vale ou é overkill?' (2-3 respostas)\n⚠️ Aplicar GraphRAG a tudo = over-engineering. Vector RAG resolve 95% em muitos casos.\n➡️ Recap das 3 abordagens.", 50, T
)

s = content_slide("Recap: Vector vs Graph vs GraphRAG", [
    "Vector DB: similaridade semântica, barato, rápido",
    "Knowledge Graph: relacionamentos estruturados, multi-hop, modelagem cara",
    "GraphRAG: KG + comunidades + sumários, perguntas globais, caro na construção",
    "",
    "Próximo: combinar vector + grafo em pipeline híbrido"
], "📖 Vector = similaridade (barato, rápido, sem estrutura). Grafo = relacionamentos (multi-hop, modelagem cara). GraphRAG = grafo + comunidades (perguntas globais, caríssimo). Próximo: híbrido.\n💡 Analogia: vector = atleta de velocidade. Grafo = estrategista. GraphRAG = general. Exército maduro usa os três.\n➡️ Pipelines híbridos.", 51, T, obj="Recap")

s = section_slide(5, "Pipelines Híbridos")
add_notes(s, "📖 Combinação de vector + grafo. Um agente que escolhe a estratégia certa para cada pergunta. Coração do projeto do módulo.\n➡️ Por que combinar?")

s = content_slide("Vector + Grafo: Por Que Combinar?", [
    "Vector: bom em 'encontrar documentos sobre X'",
    "Grafo: bom em 'encontrar entidades relacionadas a X'",
    "Híbrido: 'documentos sobre entidades relacionadas a X'",
    "",
    "Exemplo: 'Artigos sobre medicamentos que interagem com o que paciente X toma'",
    "  Vector sozinho: não sabe interações",
    "  Grafo sozinho: não tem texto dos artigos",
    "  Híbrido: une os dois"
], "📖 Vector e grafo têm forças complementares. A pergunta 'artigos sobre medicamentos que interagem com X' precisa dos DOIS: grafo (interações) + vector (artigos).\n💡 Analogia: vector = bibliotecário (encontra por tema). Grafo = detetive (conecta suspeitos). 'Livros sobre suspeitos do caso X' precisa dos dois.\n⚠️ Híbrido é vector E grafo combinados, não vector OU grafo. A fusão é a mágica.\n➡️ Arquitetura híbrida canônica.", 53, T, obj="Por que híbrido")

s = content_slide("Arquitetura Híbrida", [
    "Componentes:",
    "  1. Ingestão: texto → embeddings (vector DB) + entidades/relações (grafo)",
    "  2. Query: pergunta → vector search + graph traversal",
    "  3. Fusão: combinar resultados (chunks + subgrafo)",
    "  4. Geração: LLM com chunks + subgrafo como contexto",
    "",
    "Sync: vector DB e grafo referenciam os mesmos doc_ids"
], "📖 INGESTÃO: cada doc vira embedding E entidades/relações. QUERY: vector + graph traversal. FUSÃO: chunks + subgrafo como contexto. GERAÇÃO: LLM cita texto e relações.\n💡 Analogia: jornalista consultando arquivo (vector) + rede de contatos (grafo). Matéria usa os dois, com fontes citadas.\n⚠️ Esquecer SYNC entre vector e grafo = inconsistência. doc_id compartilhado + reconciliation job.\n➡️ Nem toda pergunta precisa de híbrido. O agente escolhe.", 54, T, obj="Arquitetura Híbrida")

s = content_slide("RetrievalAgent: Escolhendo Estratégia", [
    "Nem toda pergunta precisa de grafo",
    "RetrievalAgent classifica a pergunta:",
    "  'O que é X?' → vector search (simples)",
    "  'Quem está conectado a X?' → graph traversal",
    "  'Documentos sobre entidades relacionadas a X?' → híbrido",
    "Implementação: LLM como router (classifica intenção)",
    "Alternativa: sempre híbrido (mais caro, mais simples)"
], "📖 RetrievalAgent CLASSIFICA a pergunta e despacha para vector/graph/híbrido. Nem toda pergunta precisa de híbrido — router economiza retrieval desnecessário.\n💡 Analogia: médico que tria. Dor de cabeça? clínico (vector). Interação medicamentosa? especialista (grafo). Caso complexo? equipe (híbrido).\n⚠️ 'Sempre híbrido' = latência e custo altos. Router adiciona 1 chamada LLM mas economiza muito.\n➡️ Código do router.", 55, T, obj="RetrievalAgent")

s = code_slide(
    "Implementação do Router",
    'def retrieve(question: str) -> list[Document]:\n    strategy = classify_strategy(question)\n    if strategy == "vector":\n        return vector_search(question)\n    elif strategy == "graph":\n        return graph_traverse(question)\n    else:  # híbrido\n        entities = extract_entities(question)\n        related = graph.find_related(entities)\n        return vector_search(question, filter=related)',
    "📖 retrieve() é o coração do RetrievalAgent. classify_strategy usa LLM para classificar. Híbrido: extrai entidades, acha relacionadas no grafo, vector search filtrado.\n💡 Analogia: gerente de suporte que encaminha ticket. Lê título (1 chamada LLM), decide: FAQ (vector), especialista (grafo), equipe (híbrido).\n⚠️ Pular router e fazer híbrido sempre = latência e custo altos. Router se paga rapidamente.\n➡️ Manutenção do pipeline híbrido.", 56, T, obj="Código Router"
)

s = content_slide("Manutenção: Incremental Update", [
    "Novo documento → 2 pipelines disparam:",
    "  1. Embedding → insert no vector DB",
    "  2. Extração de entidades/relações → insert/update no grafo",
    "Sync: ambos referenciam o mesmo doc_id",
    "Desafio: se extração falha, grafo fica incompleto",
    "Solução: retry queue + reconciliation job"
], "📖 Novo doc dispara 2 pipelines: embedding (vector) + extração (grafo). Ambos com mesmo doc_id. Se extração falha, inconsistência. Retry queue + reconciliation.\n💡 Analogia: casal dividindo tarefas. Um lava (vector), outro seca (grafo). Se um falha, louça molhada na prateleira. Reconciliation = checagem periódica.\n⚠️ Não implementar reconciliation = inconsistência cresce silenciosamente. 'Agente não sabe sobre doc X' meses depois.\n➡️ Como rastrear respostas? Lineage.", 57, T, obj="Manutenção")

s = content_slide("Lineage e Provenance", [
    "Lineage: qual documento originou qual tripla do grafo",
    "Provenance: qual tripla originou qual parte da resposta",
    "Implementação: doc_id como propriedade em nós e arestas",
    "Valor: debugging, auditoria, correção",
    "Sem lineage: grafo é caixa preta — impossível explicar respostas"
], "📖 Lineage = rastreabilidade. Cada tripla tem doc_id. Resposta ← tripla ← documento. Permite debugging, auditoria, correção.\n💡 Analogia: lineage = citação em paper. Sem citação = 'dizem que...' (caixa preta). Com citação = verificável.\n⚠️ Sem lineage, 'agente disse X, mas não sei de onde' = impossível debugar. doc_id em cada nó e aresta, SEMPRE.\n➡️ Caso do projeto.", 58, T, obj="Lineage")

s = content_slide("Caso: Base de Conhecimento Técnica", [
    "Corpus: documentação técnica + issues + changelogs",
    "Vector DB: busca por 'como fazer X'",
    "Grafo: 'X depende de Y', 'X substituiu Z', 'X introduzido na v2.3'",
    "Híbrido: 'documentação sobre features que dependem de X'",
    "Agente: escolhe estratégia e JUSTIFICA a escolha"
], "📖 Caso do projeto. Vector = 'como fazer X'. Grafo = 'X depende de Y'. Híbrido = 'docs sobre features que dependem de X'. Agente justifica a escolha.\n💡 Analogia: dev sênior consultando wiki (vector) + código-fonte (grafo de dependências). 'Quais features quebram se mudar X?' precisa dos dois.\n➡️ Vamos praticar a decisão.", 59, T, obj="Caso Técnico")

s = exercise_slide(
    "Vector, Graph, ou Híbrido?",
    [
        "Para cada cenário: V, G, ou H? Justifiquem.",
        "",
        "1. Catálogo de produtos com busca por similaridade",
        "2. Base médica: doenças, sintomas, tratamentos",
        "3. Documentos jurídicos com referências cruzadas",
        "4. Chatbot de RH: perguntas sobre políticas",
        "5. Sistema de recomendação de filmes",
        "",
        "Em grupos (2 min). Depois compartilhamos."
    ],
    "📖 Exercício em grupos. Praticar o CRITÉRIO: similaridade (V), relacionamentos (G), ambos (H). Não há resposta única.\n❓ 'Em grupos: V, G, ou H para cada? Por quê?' (2 min)\n⚠️ Responder 'híbrido para tudo' = armadilha. Catálogo de produtos = vector puro. Híbrido só com evidência de gap.\n➡️ Justamente: sempre híbrido?", 60, T
)

s = content_slide("Sempre Híbrido?", [
    "Verdadeiro/falso: 'Sempre preferir híbrido.'",
    "",
    "Resposta: FALSO",
    "Híbrido adiciona complexidade (2 sistemas, router, reconciliation) e custo",
    "Híbrido só quando: vector falha E o valor justifica",
    "",
    "Anti-pattern: 'híbrido por precaução' → over-engineering",
    "Regra: comece com vector, adicione grafo com EVIDÊNCIA de gap"
], "📖 'Sempre híbrido' = FALSO. Híbrido adiciona complexidade e custo. Justificado quando vector falha em perguntas importantes E valor justifica.\n💡 Analogia: canivete suíço quando só precisa de faca. Canivete faz tudo, mas cada ferramenta é menos boa que a dedicada.\n⚠️ Implementar híbrido desde o dia 1 'para estar pronto' = 2x complexidade, 2x bugs, sem evidência. Comece simples.\n➡️ Mas quando híbrido vale, os números falam.", 61, T, obj="Sempre Híbrido?")

s = content_slide("Benchmark: Vector vs Grafo vs Híbrido", [
    "Dataset: corpus técnico com perguntas multi-hop",
    "",
    "Resultado típico:",
    "  Vector:  60% success · 200ms · $0.01/query",
    "  Graph:   75% success · 400ms · $0.02/query",
    "  Híbrido: 85% success · 500ms · $0.03/query",
    "",
    "Híbrido melhora success rate em multi-hop sem explodir custo",
    "Critério de sucesso do projeto"
], "📖 Vector: 60% success (falha em multi-hop). Graph: 75%. Híbrido: 85%. Híbrido é 3x mais caro que vector, mas success rate salta 25 pontos.\n💡 Analogia: vector = bicicleta (barata, rápida no plano). Graph = jipe (sobe morro, lento no plano). Híbrido = 4x4 (faz os dois).\n⚠️ Olhar só success rate e ignorar latência/custo = erro. 85% a $0.03 pode ser pior que 75% a $0.02 dependendo do orçamento.\n➡️ Operação em escala.", 62, T, obj="Benchmark")

s = section_slide(6, "Operação em Escala")
add_notes(s, "📖 Último bloco técnico. Operar pipeline híbrido em produção: sharding, consistência, reindexação, quantização, observabilidade.\n➡️ Sharding e replicação.")

s = content_slide("Sharding e Replicação", [
    "Sharding: dividir vetores em partições (hash ou range)",
    "Replicação: cópias para disponibilidade e leitura paralela",
    "Qdrant: sharding por collection, replicação configurável",
    "Milvus: distribuído nativo (QueryNode, DataNode, Proxy)",
    "Trade-off: mais shards = mais paralelismo, mas mais coordenação"
], "📖 Sharding divide vetores em partições (nós diferentes = paralelismo). Replicação faz cópias (disponibilidade + leitura paralela). Mais shards = mais coordenação.\n💡 Analogia: sharding = biblioteca em filiais (cada uma com parte do acervo). Replicação = cópias do best-seller em cada filial.\n⚠️ Shards demais 'para escalar' = latência de coordenação. Comece com poucos, adicione com evidência.\n➡️ Consistência.", 64, T, obj="Sharding")

s = content_slide("Consistência", [
    "Strong: toda leitura vê o último write (mais caro)",
    "Eventual: leituras podem ver dados stale (mais rápido)",
    "Bounded staleness: stale por no máximo T segundos",
    "Qdrant: configurável por collection",
    "Para RAG: eventual costuma bastar (embeddings não mudam a cada segundo)"
], "📖 Strong = quorum síncrono (caro). Eventual = propagação eventual (rápido). Para RAG, eventual basta — embeddings não mudam a cada segundo.\n💡 Analogia: strong = banco (saldo exato agora). Eventual = extrato do cartão (alguns minutos defasado). RAG tolera 'extrato'.\n⚠️ Exigir strong 'por segurança' = latência desnecessária em RAG.\n➡️ Reindexação e drift.", 65, T, obj="Consistência")

s = content_slide("Reindexação e Drift de Embeddings", [
    "Modelo de embedding novo → vetores antigos INCOMPATÍVEIS",
    "Reindexação total: re-embeddar todos os docs (caro, demorado)",
    "Estratégias:",
    "  Dual index: manter novo e velho durante transição",
    "  Zero-downtime: construir novo em background, swap quando pronto",
    "Drift: documentos mudam → embeddings antigos perdem relevância"
], "📖 Trocar de modelo = vetores incompatíveis. Reindexação total. Zero-downtime: dual index (velho ativo + novo em background) → swap → desligar velho.\n💡 Analogia: trocar linguagem de catalogação da biblioteca. Dewey → LOC. Não dá para misturar. Zero-downtime = manter as duas durante transição.\n⚠️ Trocar modelo sem reindexar = recall despenca. Vetores em 'espaços' diferentes. SEMPRE reindex.\n➡️ Custo de armazenamento: quantização.", 66, T, obj="Reindexação")

s = content_slide("Custo de Armazenamento (Quantização)", [
    "Vetor 1536 dims float32 = 6 KB · 100M vetores = 600 GB",
    "",
    "Quantização:",
    "  float32 → float16: 50% redução, perda mínima",
    "  Product Quantization (PQ): 90%+ redução, perda moderada",
    "  Binary quantization: 97% redução, perda maior",
    "",
    "Qdrant: quantização com re-scoring (rápido + preciso)"
], "📖 100M vetores de 1536 dims = 600 GB. Quantização reduz: float16 (-50%), PQ (-90%), binary (-97%). Qdrant faz re-scoring: busca no quantizado (rápido) + re-rank com originais (preciso).\n💡 Analogia: quantização = comprimir foto em JPEG. RAW (float32) → PNG (float16) → JPEG alta (PQ) → JPEG baixa (binary). Re-scoring = tirar em JPEG, abrir RAW dos top.\n⚠️ Quantização agressiva sem re-scoring = recall baixo. Teste recall com filtros reais.\n➡️ Observabilidade.", 67, T, obj="Quantização")

s = content_slide("Observabilidade", [
    "Latência de query: p50, p95, p99",
    "Recall@k: verdadeiros vizinhos vs retornados (amostragem)",
    "Cache hit rate: queries repetidas",
    "Index size vs collection size",
    "Throughput: queries por segundo",
    "Alertas: p99 > threshold, recall < threshold, index desatualizado"
], "📖 Métricas mínimas. Latência (p50/p95/p99). Recall@k (amostral vs brute force). Cache hit. Index vs collection size. Throughput. Alertas.\n💡 Analogia: painel do carro. Não dirige só olhando velocímetro (latência) — também combustível (index size), temperatura (recall), conta-quilômetros (throughput).\n⚠️ Monitorar só latência = respostas rápidas mas ERRADAS (recall ruim). Monitore recall desde o dia 1.\n➡️ Fechamento.", 68, T, obj="Observabilidade")

s = section_slide(7, "Boas Práticas e Anti-Patterns")
add_notes(s, "📖 Última seção. Consolidar em boas práticas (DO), anti-patterns (DON'T), caso, resumo, quiz, Q&A.\n➡️ O que FAZER.")

s = content_slide("Boas Práticas (DO)", [
    "Comece com vector DB, adicione grafo com EVIDÊNCIA de gap",
    "Normalize embeddings se usar cosine (vira dot, mais rápido)",
    "Defina schema de entidades/relações ANTES de extrair",
    "Valide triplas extraídas por LLM (confiança + amostragem)",
    "Mantenha lineage: doc_id em cada nó e aresta",
    "Monitore recall@k em produção (não só latência)",
    "Reindexação planejada ao trocar modelo de embedding",
    "Híbrido só quando success rate do vector é insuficiente"
], "📖 Checklist DO. Comece simples, adicione complexidade com evidência. Normalize embeddings. Defina schema antes. Valide triplas. Mantenha lineage. Monitore recall. Planeje reindexação.\n💡 Analogia: código de boas práticas de engenharia. Não constrói sem fundação. Não pinta sem projeto. Audita a obra.\n➡️ E o que NÃO fazer.", 70, T, obj="Boas Práticas")

s = content_slide("Anti-Patterns (DON'T)", [
    "Começar com GraphRAG quando vector RAG basta (overkill)",
    "Misturar métricas entre index e query",
    "Modelar grafo com 'RELATED_TO' genérico em tudo",
    "Confiar cegamente em triplas extraídas por LLM (alucinação)",
    "Sem lineage — impossível debugar respostas erradas",
    "Não monitorar recall (latência boa ≠ qualidade boa)",
    "Reindexar em pico sem dual index",
    "Adicionar vector DB quando Postgres já resolve",
    "Híbrido 'por precaução' sem evidência de gap"
], "📖 Checklist DON'T. Overkill (GraphRAG desnecessário). Métricas misturadas. RELATED_TO genérico. Triplas sem validação. Sem lineage. Sem recall. Reindex em pico. Vector DB quando Postgres resolve. Híbrido sem evidência.\n💡 Analogia: cada anti-pattern é um pecado capital. Overkill = orgulho. Caixa preta = preguiça. Over-engineering = ganância por complexidade.\n➡️ Caso real.", 71, T, obj="Anti-Patterns")

s = content_slide("Caso: GraphRAG em Farmacêutica", [
    "Domínio: interações medicamentosas",
    "Corpus: 50k documentos (bulas, papers, guidelines)",
    "Vector DB: busca por 'documentos sobre efeito colateral X'",
    "Grafo: med → interage_com → med → contraindicado_para → condição",
    "GraphRAG: comunidades de medicamentos → classes terapêuticas",
    "Resultado: multi-hop 'paciente com X + medicação Y → risco?'",
    "Custo: ~$200 de construção, justificado por valor clínico"
], "📖 Caso real que une tudo. Vector (bulas), grafo (interações), GraphRAG (classes terapêuticas). Pergunta-chave: 'risco do paciente?' — multi-hop que vector não responde. Custo $200, justificado por vida.\n💡 Analogia: exército completo. Vector = médico de plantão. Grafo = farmacêutico. GraphRAG = comitê de especialistas. Para segurança do paciente, precisa dos três.\n⚠️ TODO domínio NÃO precisa de GraphRAG. Farmacêutica = vida ou morte. Blog de receitas = vector basta. Contexto define.\n➡️ Resumo.", 72, T, obj="Caso Farmacêutica")

s = content_slide("Resumo da Aula", [
    "Vector DB: ANN (HNSW/IVF), métricas, filtering, sparse+dense",
    "5 vector DBs: Qdrant, Milvus, Weaviate, Chroma, pgvector",
    "Knowledge Graph: triplas, propriedades, Cypher, multi-hop",
    "GraphRAG: comunidades + sumarização → local vs global",
    "Híbrido: vector + grafo, RetrievalAgent como router",
    "Escala: sharding, reindexação, quantização, observabilidade",
    "",
    "Regra de ouro: comece simples, adicione complexidade com EVIDÊNCIA"
], "📖 7 pontos-chave. (1) Vector DB: ANN, métricas, filtering. (2) 5 DBs com trade-offs. (3) KG: triplas, Cypher, multi-hop. (4) GraphRAG: comunidades, local/global. (5) Híbrido: router. (6) Escala. (7) Regra: comece simples, evidência antes de complexidade.\n💡 Analogia: a 'mochila' que levam. 7 ferramentas + critério de decisão. Usem com critério.\n➡️ Checklist da aula.", 73, T, obj="Resumo")

s = content_slide("Checklist da Aula", [
    "[ ] Explicou ANN (HNSW, IVF) e métricas",
    "[ ] Comparou 5 vector DBs por cenário",
    "[ ] Modelou knowledge graph com triplas e Cypher",
    "[ ] Apresentou GraphRAG (local + global)",
    "[ ] Construiu pipeline híbrido com router",
    "[ ] Discutiu escala, reindexação e observabilidade"
], "📖 Checklist para confirmar cobertura. Se algo não ficou claro, é hora de perguntar no Q&A.\n➡️ Quiz.", 74, T, obj="Checklist")

s = content_slide("Quiz — Pergunta 1", [
    "Qual métrica é equivalente a dot product",
    "quando embeddings estão normalizados?",
    "",
    "A) Euclidean (L2)",
    "B) Cosine",
    "C) Manhattan",
    "D) Jaccard"
], "📖 Resposta: B) Cosine. Normalizados (norma=1) → cosine = dot. Por isso muita gente normaliza e usa dot (mais rápido).\n➡️ Pergunta 2.", 75, T, obj="Quiz 1")

s = content_slide("Quiz — Pergunta 2", [
    "Quando GraphRAG supera vector RAG mais claramente?",
    "",
    "A) Fato específico em um documento",
    "B) Raciocínio multi-hop entre entidades",
    "C) Latência é o requisito mais importante",
    "D) Corpus pequeno (< 100 documentos)"
], "📖 Resposta: B) multi-hop. GraphRAG brilha em encadear relacionamentos. Fato específico = vector. Latência = vector. Corpus pequeno = não justifica custo.\n➡️ Pergunta 3.", 76, T, obj="Quiz 2")

s = content_slide("Quiz — Pergunta 3", [
    "Em um pipeline híbrido, o que o RetrievalAgent faz?",
    "",
    "A) Sempre busca vetorial e grafo em paralelo",
    "B) Classifica a pergunta: vector, graph ou híbrido",
    "C) Substitui o vector DB por um grafo",
    "D) Gera embeddings para o grafo"
], "📖 Resposta: B). RetrievalAgent CLASSIFICA a pergunta e despacha. Nem toda pergunta precisa de híbrido — router economiza retrieval desnecessário.\n➡️ Conexão com próximos módulos.", 77, T, obj="Quiz 3")

s = content_slide("Conexão com Próximos Módulos", [
    "ETHAGT14 — Escalabilidade: vector DBs em produção multi-tenant",
    "ETHAGT90 — Capstone: integra vector + grafo no projeto final",
    "ETHAGT05 — Memória: vector DB como memória semântica",
    "ETHAGT12 — AgentOps: observabilidade de pipelines híbridos",
    "",
    "O que aprenderam hoje é a INFRAESTRUTURA que sustenta esses módulos"
], "📖 ETHAGT07 conecta com 4 módulos. ETHAGT14 (escala), ETHAGT90 (capstone), ETHAGT05 (memória), ETHAGT12 (AgentOps). Hoje = infraestrutura.\n➡️ Leitura recomendada.", 78, T, obj="Conexão")

s = content_slide("Leitura Recomendada e Referências", [
    "Obrigatório: Edge et al. GraphRAG (arXiv:2404.16130)",
    "Obrigatório: Lewis et al. RAG original (arXiv:2005.11401)",
    "Recomendado: Pan et al. LLMs + KGs (arXiv:2306.08302)",
    "",
    "Documentação: Qdrant, Milvus, Weaviate, Neo4j, pgvector",
    "",
    "Próxima aula: ETHAGT14 — Escalabilidade"
], "📖 Leitura obrigatória: GraphRAG (Edge et al.) + RAG (Lewis et al.). Recomendado: survey Pan et al. Docs dos DBs que vão usar.\n➡️ Q&A.", 79, T, obj="Referências")

s = title_slide(
    "Perguntas?",
    "Universidade Etho · ETHAGT07\nPróxima aula: ETHAGT14 — Escalabilidade\nLembrete: iniciar Lab 1 (Vector DB Bake-off)",
    "Obrigado!"
)
add_notes(s, "📖 Hora do Q&A. Se não houver perguntas: 'Qual parte foi menos clara — vector, grafo ou híbrido?'\n❓ 'Perguntas? Ou: qual parte foi menos clara?'\n⚠️ Nenhuma pergunta = alunos sobrecarregados. Sugira revisar slides + trazer dúvidas no fórum.\n➡️ Fim da aula. Vejo vocês em ETHAGT14!")

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT07-slides.pptx")
prs.save(out)
print(f"PPTX gerado: {out}")
print(f"Total de slides: {len(prs.slides)}")
