#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT04: Reasoning & Planning
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj=""):
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT04"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

T = 90

# ═══════════════════════════════════════
# SEÇÃO A — Abertura (1-7)
# ═══════════════════════════════════════

s = title_slide(
    "Reasoning & Planning",
    "Universidade Etho · Especialização em Programação Agêntica\nFase B — Razão, Memória e Conhecimento · 30 h",
    "ETHAGT04"
)
add_notes(s, "📖 Bem-vindos ao ETHAGT04. Hoje vamos além do ReAct: planejamento, exploração em árvore, auto-crítica e reasoning nativo.\n💡 ETHAGT01 foi aprender a andar. Hoje vocês vão correr, pular e escalar.\n❓ 'Quem já viu um agente preso em loop?'\n⚠️ Alunos acham que raciocínio é só CoT. É um espectro.\n➡️ Objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: dominar o espectro de estratégias de raciocínio",
    "",
    "Objetivos específicos:",
    "1. Caracterizar planejamento antes vs durante a ação",
    "2. Implementar: ReAct, Plan-and-Execute, ReWOO, ToT, LATS, Reflexion, Self-Discover",
    "3. Compreender reasoning model nativo vs prompting",
    "4. Avaliar trade-offs: qualidade × custo × latência × robustez",
    "5. Lidar com falhas: plano rígido, loops, re-planejamento ausente",
], "📖 Cada objetivo é mensurável. O mais importante é #4 — trade-offs.\n❓ 'Qual objetivo é mais desafiador?' (#2 ou #4)\n➡️ Competências.", 2, T, "5 objetivos mensuráveis")

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → I (Intermediário)",
    "C2 Multi-Agent Systems → B (Básico)",
    "C4 Agent Memory → B (Básico)",
    "C5 AgentOps & Avaliação → B (Básico)",
], "📖 C1 atinge Intermediário. C4 é tocado por Reflexion (memória de erros).\n➡️ Agenda.", 3, T, "Framework Etho")

s = content_slide("Agenda da Aula", [
    "Bloco 1 (60 min):",
    "  Abertura (8 min) — motivação, espectro",
    "  Tipologia (14 min) — CoT, Self-Consistency, estruturas",
    "  Plan-and-Execute (16 min) — Planner, ReWOO, trade-offs",
    "  ToT e LATS (18 min) — busca em árvore, MCTS",
    "  Intervalo (5 min)",
    "",
    "Bloco 2 (60 min):",
    "  Reflexion (12 min) — auto-crítica com memória",
    "  Self-Discover (8 min) — composição de estratégia",
    "  Reasoning Nativo (14 min) — o1/o3, Claude thinking",
    "  Falhas e Orçamento (12 min) — loops, budget, benchmarks",
    "  Fechamento (8 min) — comparação, quiz, Q&A",
], "📖 Dois blocos. Primeiro: técnicas de planejamento e busca. Segundo: auto-crítica, meta-raciocínio e reasoning nativo.\n➡️ Por que estamos aqui?", 4, T, "Roteiro da aula")

s = content_slide("Quando ReAct Não Basta", [
    "Cenário: 'Planeje viagem de 7 dias com R$5000'",
    "",
    "ReAct puro:",
    "  Passo 1: buscar voos → R$3000 (50% do orçamento)",
    "  Passo 2: buscar hotel → R$2000 (80%)",
    "  Passo 3: atividades → não cabe mais. Loop.",
    "  Sem backtracking: não pode voltar e escolher voos mais baratos",
    "",
    "O problema: um passo errado no meio inviabiliza todo o plano",
    "Solução: raciocínio estruturado — planejar, explorar, refletir",
], "📖 ReAct é excelente para passos independentes. Mas quando dependem uns dos outros, um erro inicial compromete tudo.\n💡 Analogia: xadrez. ReAct = jogar move-by-move. Mestre = pensar adiante (ToT), aprender de jogos (Reflexion).\n❓ 'Quando um passo errado inviabiliza todo o plano?'\n⚠️ 'Mais max_steps' não resolve — continua na direção errada.\n➡️ O espectro de soluções.", 5, T, "Criar tensão — ReAct tem limites")

s = content_slide("O Espectro do Raciocínio", [
    "Linha do tempo:",
    "  2022: Chain-of-Thought (Wei et al.)",
    "  Mar/2023: ReAct (Yao et al.) — raciocínio durante a ação",
    "  Mai/2023: Tree of Thoughts — busca em árvore",
    "  Mar/2023: Reflexion (Shinn et al.) — auto-crítica com memória",
    "  2024: LATS (Zhou et al.) — MCTS + LLM",
    "  Fev/2024: Self-Discover — composição de estratégia",
    "  Set/2024: OpenAI o1 — reasoning nativo via RL",
    "",
    "Cada técnica resolve uma limitação da anterior",
    "Diagrama: 12-Diagrams/ETHAGT04/reasoning-spectrum.mmd",
], "📖 Cada técnica adiciona algo: CoT adiciona raciocínio. ReAct adiciona ação. ToT adiciona exploração. LATS adiciona busca sistemática. Reflexion adiciona aprendizado. Self-Discover adiciona meta-raciocínio. o1 traz raciocínio para dentro do modelo.\n➡️ Começando pela fundação: CoT.", 6, T, "Panorama das técnicas")

s = section_slide(1, "Tipologia do Raciocínio")
add_notes(s, "Vamos classificar as estratégias de raciocínio por estrutura e timing.")

# ═══════════════════════════════════════
# SEÇÃO B — Tipologia (8-18)
# ═══════════════════════════════════════

s = content_slide("Chain-of-Thought (CoT)", [
    "A fundação de todo raciocínio em LLMs",
    "",
    "CoT zero-shot: 'Let's think step by step'",
    "CoT few-shot: exemplos de raciocínio no prompt",
    "",
    "Por que funciona: o LLM pensa ATRAVÉS da geração",
    "  Cada token gerado é parte do raciocínio",
    "",
    "Fonte: Wei et al., NeurIPS 2022 (arXiv:2201.11903)",
    "Efeito: até +20 pts em benchmarks matemáticos",
], "📖 CoT é a técnica mais simples e poderosa. A ideia: não pedir a resposta direta — pedir que o modelo MOSTRE o raciocínio.\n💡 Analogia: pedir para o aluno mostrar o cálculo, não só a resposta.\n⚠️ CoT não é 'só adicionar prompt' — muda a distribuição de geração.\n➡️ E se amostrarmos múltiplas vezes?", 8, T, "Fundação do raciocínio")

s = content_slide("Self-Consistency", [
    "Ideia: gerar N raciocínios independentes (temperatura > 0)",
    "Votação: escolher a resposta mais frequente (majority vote)",
    "",
    "Exemplo: 10 amostras → 7 dão '42', 3 dão '38' → escolher '42'",
    "",
    "Custo: N× mais caro que CoT simples",
    "Quantas amostras? Tipicamente 5-20 (mais que 20 = retorno marginal)",
    "",
    "Fonte: Wang et al., 2022 (arXiv:2203.11171)",
], "📖 LLMs podem chegar à mesma resposta por caminhos diferentes. Self-consistency explora isso.\n💡 Analogia: perguntar a 10 pessoas qual a capital da Austrália. Maioria vence.\n➡️ Antes ou durante a ação?", 9, T, "Múltiplas amostras + votação")

s = comparison_slide("Antes vs Durante a Ação",
    "Durante (ReAct)", [
        "pensa → age → observa → pensa → age...",
        "Adaptativo, flexível",
        "Mas: pode se perder",
        "Re-escolhe a cada passo",
        "Mais caro (N raciocínios)",
    ],
    "Antes (Plan-and-Execute)", [
        "planeja tudo → executa",
        "Eficiente, estruturado",
        "Mas: rígido",
        "Não adapta se inesperado",
        "Mais barato (1 raciocínio)",
    ],
    "📖 A distinção mais importante da aula. ReAct raciocina DURANTE — reconsidera a cada passo. Plan-and-Execute raciocina ANTES — gera plano completo.\n💡 Analogia: ReAct = dirigir olhando pelo espelho. Plan-and-Execute = programar GPS no início.\n➡️ Estruturalmente?", 10, T, "Duas famílias de planejamento")

s = content_slide("Linear vs Árvore vs Grafo", [
    "Linear (ReAct):",
    "  Thought → Action → Observation → Thought → ...",
    "  Um caminho, sem bifurcação",
    "",
    "Árvore (ToT):",
    "  Gerar múltiplos candidatos → avaliar → explorar melhor",
    "  Backtracking possível",
    "",
    "Grafo (LATS):",
    "  MCTS — seleção, expansão, avaliação, backpropagation",
    "  Busca sistemática com memória de estados",
], "📖 A evolução é estrutural. Linear → Árvore → Grafo. Cada nível adiciona poder mas também custo.\n➡️ Reasoning + tools.", 11, T, "3 estruturas de raciocínio")

s = content_slide("Reasoning + Tools", [
    "Raciocínio e tool use não são separados — são integrados",
    "",
    "ReAct: raciocínio guia QUANDO chamar tools",
    "ToT: pode ramificar a árvore chamando diferentes tools",
    "Reflexion: reflete sobre POR QUE uma tool falhou",
    "Reasoning nativo (o1): chama tools entre 'pensamentos'",
    "",
    "A qualidade do raciocínio determina a qualidade do uso de tools",
], "📖 Uma confusão comum: achar que raciocínio e tools são separados. Não são. Em todas as técnicas, raciocínio guia o uso de tools.\n➡️ Plan-and-Execute.", 12, T, "Raciocínio guia tool use")

s = section_slide(2, "Plan-and-Execute e ReWOO")
add_notes(s, "Vamos ver como planejar ANTES de agir torna agentes mais eficientes.")

# ═══════════════════════════════════════
# SEÇÃO C — Plan-and-Execute (19-30)
# ═══════════════════════════════════════

s = content_slide("Plan-and-Execute", [
    "Arquitetura: Planner → Executor → Replanner (opcional)",
    "",
    "Planner: LLM gera lista de passos (sem executar)",
    "  Ex: ['1. Buscar voos', '2. Buscar hotéis', '3. Calcular total', '4. Verificar orçamento']",
    "",
    "Executor: executa cada passo (pode usar ReAct internamente)",
    "",
    "Replanner: se um passo falha, re-planeja o restante",
    "",
    "Diagrama: 12-Diagrams/ETHAGT04/plan-execute.mmd",
], "📖 Plan-and-Execute separa planejamento e execução. Planner só decomponhe. Executor só executa. Replanner torna o sistema adaptativo.\n💡 Analogia: arquiteto (Planner) + pedreiro (Executor). Se encontra tubulação, arquiteto redesenha.\n❓ 'Para que tarefa é melhor que ReAct?'\n⚠️ Sem Replanner = plano rígido que falha ao primeiro obstáculo.\n➡️ ReWOO.", 14, T, "Planner → Executor → Replanner")

s = content_slide("ReWOO: Reasoning WithOut Observation", [
    "Ideia: Planner gera plano 'cego' (sem ver observações)",
    "",
    "1. Planner: 'Preciso saber X, Y, Z' (sem esperar respostas)",
    "2. Todas as tools chamadas EM PARALELO",
    "3. Solver: combina as evidências para a resposta final",
    "",
    "Vantagem: 1 raciocínio + N tools paralelas + 1 solver",
    "Latência reduzida (paralelismo)",
    "",
    "Fonte: Xu et al., 2023 (arXiv:2305.18323)",
], "📖 ReWOO separa totalmente raciocínio e execução. Planner diz 'preciso X, Y, Z' sem esperar. Tools rodam em paralelo. Solver combina.\n💡 Analogia: 3 pesquisadores simultâneos em vez de um de cada vez.\n⚠️ ReWOO não funciona se tool B depende do resultado de tool A.\n➡️ Comparação.", 15, T, "Paralelização de evidências")

s = comparison_slide("Comparação: 3 Abordagens",
    "ReAct", [
        "Flexibilidade: ✅ Alta",
        "Tokens: ❌ Muitos",
        "Latência: ❌ Serial",
        "Re-planejamento: ✅ Nativo",
        "Complexidade: ✅ Simples",
    ],
    "Plan-Execute / ReWOO", [
        "Flexibilidade: ⚠️ Média / ❌ Baixa",
        "Tokens: ⚠️ Médios / ✅ Poucos",
        "Latência: ⚠️ Serial / ✅ Paralelo",
        "Re-planejamento: ⚠️ Opcional / ❌ Não",
        "Complexidade: ⚠️ Média",
    ],
    "📖 Não há vencedor. ReAct = flexibilidade. ReWOO = eficiência. Plan-and-Execute = meio-termo.\nSequencial dependente → ReAct. Paralelizável independente → ReWOO. Estruturado com falha → Plan-and-Execute + Replanner.\n➡️ Quando re-planejar.", 16, T, "Tabela comparativa")

s = content_slide("Quando Re-planejar", [
    "Sinais de falha do plano:",
    "  • Tool retorna erro ou resultado inesperado",
    "  • Resultado contradiz premissa do plano",
    "  • Orçamento estourando (custo, tempo)",
    "  • Loop detectado (mesma ação repetida)",
    "",
    "Estratégias:",
    "  • Re-planejamento automático (LLM re-gera plano restante)",
    "  • HITL: humano aprova ou ajusta o plano",
    "  • Fallback: plano simplificado pré-definido",
], "📖 Re-planejamento torna Plan-and-Execute adaptativo. Sem ele, o plano é rígido. A questão é QUANDO disparar.\nSinais objetivos: erro, timeout, custo. Sinais subjetivos: resultado 'estranho'.\nEm produção, HITL no re-planejamento é essencial para alto risco.\n➡️ ToT e LATS.", 17, T, "Detecção de falha de plano")

s = content_slide("Intervalo", [
    "5 minutos",
    "",
    "Pense em: qual técnica de raciocínio você testaria primeiro no seu projeto?",
], "📖 Intervalo. Aproveitem para pensar: qual técnica se aplica ao problema de vocês?\n➡️ Voltamos com Tree of Thoughts.", 18, T, "Intervalo")

# ═══════════════════════════════════════
# SEÇÃO D — ToT e LATS (19-24)
# ═══════════════════════════════════════

s = section_slide(3, "Tree of Thoughts e LATS")
add_notes(s, "Agora vamos além do linear: busca em árvore e MCTS.")

s = content_slide("Tree of Thoughts (ToT)", [
    "Generalização de CoT para uma árvore",
    "",
    "4 passos:",
    "  1. Decompor — quebrar problema em sub-problemas",
    "  2. Gerar candidatos — para cada estado, múltiplos 'pensamentos'",
    "  3. Avaliar — LLM pontua cada candidato (viable/not viable, ou 1-10)",
    "  4. Buscar — BFS ou DFS através dos estados avaliados",
    "",
    "Backtracking: se ramo leva a beco sem saída, voltar e tentar outro",
    "",
    "Fonte: Yao et al., NeurIPS 2023 (arXiv:2305.10601)",
    "Diagrama: 12-Diagrams/ETHAGT04/tot-search-tree.mmd",
], "📖 ToT explora múltiplos caminhos em vez de um só. Em cada nó, LLM gera candidatos. Avaliador pontua. Explora os melhores. Backtracking se beco sem saída.\n💡 Analogia: xadrez. Não joga uma sequência — considera várias, avalia, escolhe.\n❓ 'Em que problema ToT vale o custo?' (espaço de busca grande + função de avaliação clara)\n⚠️ ToT para problema simples = overkill. 5-10x mais caro.\n➡️ LATS.", 20, T, "Busca em árvore")

s = content_slide("LATS: Language Agent Tree Search", [
    "Base: Monte Carlo Tree Search (MCTS) — algoritmo de AlphaGo",
    "",
    "4 fases do MCTS adaptadas para LLM:",
    "  1. Seleção — escolher nó via UCT (explora vs explora)",
    "  2. Expansão — gerar nova ação (LLM)",
    "  3. Simulação — executar e avaliar (LLM + environment)",
    "  4. Backpropagation — propagar valor para nós ancestrais",
    "",
    "Vantagem: busca sistemática com memória de estados",
    "Fonte: Zhou et al., 2024 (arXiv:2310.01757)",
], "📖 LATS é o estado da arte em busca agêntica. MCTS + LLM. A inovação é a backpropagation: valor propagado para ancestrais, permitindo APRENDER durante a busca.\n💡 Analogia: explorador que marca no mapa quais trilhas testou e resultado.\n⚠️ LATS ≠ ToT. ToT é BFS/DFS. LATS é MCTS com UCT e backpropagation.\n➡️ Custo.", 21, T, "MCTS + LLM")

s = content_slide("Custo vs Qualidade", [
    "ReAct: ~N chamadas (N passos)",
    "ToT: ~B^D chamadas (branching^profundidade)",
    "  Ex: D=3, B=3 → 27 nós",
    "LATS: ~(M+s)×k chamadas",
    "  Ex: 5 iterações × 3 simulações = 15+ nós",
    "",
    "Em produção: $0.50-$1.00 por execução para problemas difíceis",
    "",
    "Quando vale: problema difícil onde acertar vale mais que economizar",
    "Não vale para: tarefas simples onde CoT basta",
], "📖 ToT e LATS são caros. 27 chamadas para ToT com D=3, B=3. Vale quando o problema é difícil o suficiente.\n💡 Analogia: cirurgia vs consulta. Cirurgia custa 100x mais, mas você não faz para resfriado.\n➡️ Demo.", 22, T, "Trade-off de custo")

s = code_slide("DEMO: ToT em Ação", """from langgraph.tot import ToTState

# Problema: jogo dos 24
# Objetivo: combinar 4 números com +,-,*,/ para obter 24

def generate_thoughts(state):
    # LLM gera múltiplas próximas operações candidatas
    return llm.generate_candidates(state.numbers, n=3)

def evaluate_thought(state, thought):
    # LLM avalia: o thought leva a 24?
    return llm.evaluate(state, thought)  # 1-10

def search(initial_state, max_depth=4):
    frontier = [initial_state]
    for depth in range(max_depth):
        candidates = []
        for state in frontier:
            for thought in generate_thoughts(state):
                score = evaluate_thought(state, thought)
                if score >= 7:
                    new_state = apply(state, thought)
                    if is_solution(new_state):
                        return new_state
                    candidates.append(new_state)
        frontier = candidates
    return None  # Backtrack ou falha""",
    "📖 Demo ao vivo. Prestem atenção em: (1) geração de candidatos; (2) avaliação; (3) backtracking.\n⚠️ Se API falhar, tenho screenshot do trace.\n➡️ Intervalo.", 23, T, "Demo ToT")

# ═══════════════════════════════════════
# SEÇÃO E — Reflexion (24-27)
# ═══════════════════════════════════════

s = section_slide(4, "Reflexion")
add_notes(s, "E se o agente pudesse aprender com os próprios erros?")

s = content_slide("Reflexion: Auto-Crítica com Memória", [
    "Ciclo Reflexion:",
    "  1. Actor: agente tenta resolver (ReAct)",
    "  2. Evaluator: LLM avalia (sucesso/falha + motivo)",
    "  3. Self-Reflection: LLM gera reflexão verbal",
    "     'Errei porque usei fórmula errada'",
    "  4. Memory: reflexão armazenada (texto natural)",
    "  5. Actor: tenta novamente, com reflexão no contexto",
    "",
    "Memória é VERBAL — legível por humanos",
    "Diagrama: 12-Diagrams/ETHAGT04/reflexion-loop.mmd",
    "Fonte: Shinn et al., NeurIPS 2023 (arXiv:2303.11366)",
], "📖 Reflexion: quando falha, reflete SOBRE O PORQUÊ e guarda. Próxima tentativa tem a reflexão no contexto.\n💡 Analogia: chef que erra receita, anota 'excesso de sal', consulta antes da próxima.\n❓ 'Quantas tentativas até convergir?' (3-5 tipicamente)\n⚠️ Reflexion ≠ retry. Retry = tentar de novo. Reflexion = refletir PRIMEIRO.\n➡️ Reflexion vs reflection.", 25, T, "Auto-crítica com memória")

s = comparison_slide("Reflexion vs Reflection Pattern",
    "Reflection Pattern", [
        "Auto-revisão em UMA passada",
        "gerar → criticar → revisar",
        "Sem memória entre tentativas",
        "Simples, barato",
        "Bom para tarefas de uma tentativa",
    ],
    "Reflexion (paper)", [
        "Múltiplas tentativas com MEMÓRIA",
        "tentar → refletir → memorizar → tentar",
        "Erros anteriores consultados",
        "Iterativo, mais caro",
        "Bom para benchmarks, debugging iterativo",
    ],
    "📖 A diferença é sutil mas crítica. Reflection = uma passada. Reflexion = iterativo com memória.\n➡️ Convergência.", 26, T, "Técnica vs padrão")

s = content_slide("Convergência e Limites", [
    "Convergência: tipicamente 3-5 tentativas melhoram",
    "  Após isso, retorno marginal",
    "",
    "Divergência: se reflexão é ruim, piora próximas tentativas",
    "  Ex: 'errei porque fórmula é X' (mas é Y)",
    "",
    "Mitigação:",
    "  • Limite de tentativas (max_reflections)",
    "  • Resetar memória se performance piora",
    "  • HITL na reflexão (humano valida a lição)",
], "📖 Reflexion nem sempre converge. Reflexão incorreta guia para o erro. Em produção: limite + reset + HITL para críticas.\n➡️ Self-Discover.", 27, T, "Quando Reflexion para de ajudar")

# ═══════════════════════════════════════
# SEÇÃO F — Self-Discover (28-30)
# ═══════════════════════════════════════

s = section_slide(5, "Self-Discover")
add_notes(s, "E se o agente pudesse criar sua própria estratégia de raciocínio?")

s = content_slide("Self-Discover: Composição de Estratégia", [
    "Problema: nem toda tarefa se encaixa em CoT, ToT ou Reflexion",
    "",
    "Solução: o agente DESCOBRE qual estratégia usar",
    "",
    "Processo:",
    "  1. Selecionar — das primitivas, quais são relevantes?",
    "  2. Adaptar — como combinar para ESTE problema?",
    "  3. Implementar — gerar o raciocínio composto",
    "",
    "Fonte: Major et al., 2024 (arXiv:2402.03620)",
], "📖 Self-Discover é meta-raciocínio. O agente raciocina sobre QUAL estratégia usar para o problema específico.\n💡 Analogia: médico que não receita o mesmo remédio — diagnostica e escolhe.\n⚠️ Mais caro que CoT. Não vale onde estratégia ótima é conhecida.\n➡️ Primitivas.", 29, T, "Meta-raciocínio")

s = content_slide("Primitivas como Building Blocks", [
    "CoT: pensar passo-a-passo",
    "Decomposition: quebrar em sub-problemas",
    "ReAct: pensar + agir",
    "ToT: explorar múltiplos caminhos",
    "Reflexion: aprender com erros",
    "",
    "O agente combina 2-3 primitivas relevantes:",
    "  Matemática → decomposition + CoT",
    "  Debugging → ReAct + Reflexion",
    "  Viagem → ToT + decomposition",
], "📖 As primitivas são os LEGO blocks. Self-Discover monta a estrutura certa para o problema.\n➡️ Reasoning nativo.", 30, T, "Primitivas de raciocínio")

# ═══════════════════════════════════════
# SEÇÃO G — Reasoning Nativo (31-34)
# ═══════════════════════════════════════

s = section_slide(6, "Inference-Time Reasoning Nativo")
add_notes(s, "A fronteira: modelos treinados para pensar.")

s = content_slide("Modelos com Reasoning Nativo", [
    "OpenAI o1/o3: treinados com RL para raciocínio interno",
    "Claude extended thinking: raciocínio extendido antes de responder",
    "",
    "Como funcionam:",
    "  • Geram tokens de raciocínio internos (não visíveis)",
    "  • Otimizados via RL para pensar antes de agir",
    "  • 'Reasoning tokens' têm custo mas melhoram qualidade",
    "",
    "A maior mudança desde o CoT",
], "📖 Modelos como o1 e o3 são TREINADOS para raciocinar — não precisam de 'let's think step by step'. Raciocínio interno automático.\n💡 Analogia: estagiário que precisa ser lembrado (CoT) vs sênior que já sabe (modelo treinado).\n➡️ Impacto no design.", 32, T, "Modelos treinados para pensar")

s = content_slide("Impacto no Design do Agente", [
    "Sem CoT promptado: o modelo já raciocina",
    "  'Let's think step by step' pode ATRAPALHAR",
    "",
    "Mais tools, menos prompting:",
    "  Invista em tools e contexto, não em engenharia de prompt",
    "",
    "Orçamento de reasoning tokens:",
    "  Definir max_completion_tokens para o raciocínio",
    "",
    "Menos steps no loop:",
    "  Cada chamada é mais 'pensada', precisa de menos iterações",
], "📖 Com reasoning models, o design muda. Sem CoT promptado. Mais tools. Orçamento de reasoning tokens. Loop mais curto.\n⚠️ Continuar usando 'let's think step by step' com o1 pode PIORAR performance.\n➡️ Quando escolher cada um.", 33, T, "Design muda com reasoning models")

s = comparison_slide("Reasoning Model vs Prompting",
    "Reasoning Model (o1/o3)", [
        "Problema difícil: ✅ Ideal",
        "Custo: ❌ Mais caro/chamada",
        "Latência: ❌ Alta",
        "Controle: ❌ Caixa preta",
        "Simplicidade: ✅ Sem prompting",
        "Debugging: ❌ Raciocínio oculto",
    ],
    "Prompting (CoT/ToT/Reflexion)", [
        "Problema difícil: ⚠️ Funciona mas frágil",
        "Custo: ✅ Mais barato/chamada",
        "Latência: ⚠️ Variável",
        "Controle: ✅ Transparente",
        "Simplicidade: ❌ Receita manual",
        "Debugging: ✅ Vê cada passo",
    ],
    "📖 Reasoning models: melhor para problemas difíceis onde opacidade é aceitável. Prompting: melhor para transparência, controle, debugging. Em produção com auditoria, prompting vence.\n➡️ Custo.", 34, T, "Quando escolher")

s = content_slide("Gerenciando Custo do Reasoning Nativo", [
    "Estratégias de mitigação:",
    "",
    "1. Roteamento — classificar dificuldade, usar reasoning model só para difícil",
    "2. Caching de raciocínio — reusar resultado para mesma query",
    "3. Orçamento — max_completion_tokens limita quanto pensa",
    "4. Modelo cascata — barato primeiro, escalar para reasoning se falhar",
    "",
    "o1 pode custar 5-10x mais que GPT-4o",
    "Roteamento + caching + cascata reduzem custo mantendo qualidade",
], "📖 Reasoning models são caros (5-10x). Para gerenciar: roteamento (difícil só), caching (reusar), cascata (barato primeiro, escalar se falhar).\n➡️ Falhas.", 35, T, "Estratégias de custo")

# ═══════════════════════════════════════
# SEÇÃO H — Falhas (36-38)
# ═══════════════════════════════════════

s = section_slide(7, "Falhas, Loops e Orçamento")
add_notes(s, "Toda técnica tem modos de falha. Vamos cobri-los.")

s = content_slide("Detecção e Quebra de Loops", [
    "Sintoma: agente repete A→B→A→B ou A→A→A",
    "",
    "Detecção:",
    "  • Hash das últimas N ações — se repetidas, alerta",
    "  • Embedding similarity — se semanticamente iguais",
    "  • Contador — se mesma action > 3x, intervir",
    "",
    "Quebra:",
    "  • Forçar resposta final",
    "  • Disparar Reflexion (refletir sobre o loop)",
    "  • HITL: humano decide",
], "📖 Loops são o bug mais comum. Detecção: se mesma ação 3+ vezes, algo errado. Quebra: forçar resposta, refletir, ou HITL.\nEm produção, detecção de loops é guardrail obrigatório.\n➡️ Orçamento.", 36, T, "Detecção de loops")

s = content_slide("Orçamento como Guardrail", [
    "max_steps: limite de iterações do loop (obrigatório)",
    "max_tokens: limite de tokens totais por execução",
    "max_cost: limite em dólares por execução",
    "max_time: timeout em segundos",
    "",
    "Orçamento por componente: raciocínio vs tools vs retrieval",
    "",
    "Sem orçamento: loop infinito, custo ilimitado, usuário esperando",
    "Defina TODOS desde o dia 1",
], "📖 Orçamento é segurança, não otimização. Sem max_steps = loop infinito. Sem max_cost = surpresa no fim do mês. Defina TODOS desde o dia 1.\n➡️ Re-planejamento supervisionado.", 37, T, "Guardrails obrigatórios")

s = content_slide("Re-planejamento Supervisionado (HITL)", [
    "Em tarefas de alto risco: humano aprova plano antes da execução",
    "",
    "Checkpoints de validação:",
    "  Após cada passo crítico, humano valida",
    "  Se rejeita: agente re-planeja com feedback",
    "",
    "Frameworks:",
    "  LangGraph interrupt_before",
    "  OpenAI handoffs",
    "",
    "Essencial em compliance: finanças, saúde, operações irreversíveis",
], "📖 Em alto risco, o plano não executa sem aprovação humana. Checkpoints adicionais validam passos críticos. Rejeição → re-planejamento com feedback.\n➡️ Fechamento.", 38, T, "HITL no planejamento")

# ═══════════════════════════════════════
# SEÇÃO I — Fechamento (39-44)
# ═══════════════════════════════════════

s = section_slide(8, "Fechamento")
add_notes(s, "Vamos consolidar tudo.")

s = content_slide("Comparação Geral das Estratégias", [
    "ReAct — Linear — Baixo custo — Tarefas adaptativas",
    "Plan-and-Execute — Linear+plano — Médio — Estruturadas",
    "ReWOO — Paralelo — Baixo — Independentes",
    "ToT — Árvore — Alto — Problemas de busca",
    "LATS — Grafo (MCTS) — Muito alto — Muito difíceis",
    "Reflexion — Linear+memória — Médio — Iterativas",
    "Self-Discover — Composição — Médio-Alto — Sem padrão",
    "Reasoning nativo — Interno — Alto — Difíceis",
    "",
    "Não existe 'melhor' — existe a certa para o problema",
], "📖 Este quadro é o resumo da aula. Cada técnica tem seu lugar. A arte é saber escolher.\n➡️ Anti-patterns.", 39, T, "Quadro comparativo")

s = content_slide("Anti-Patterns (DON'T)", [
    "❌ ToT/LATS para problemas simples (overkill)",
    "❌ Reflexion sem limite de tentativas (loop de reflexão)",
    "❌ Plan-and-Execute sem Replanner (plano rígido)",
    "❌ ReWOO com tools dependentes (não paralelizável)",
    "❌ 'Let's think step by step' com reasoning model (atrapalha)",
    "❌ Sem orçamento de tokens/custo (surpresa no fim do mês)",
    "❌ Sem detecção de loops (gasto infinito)",
], "📖 Cada anti-pattern é um bug visto em produção. O mais comum: usar ToT onde CoT basta.\n➡️ Quiz.", 40, T, "Checklist de anti-patterns")

s = exercise_slide("Quiz Final", [
    "5 perguntas individuais, sem consulta:",
    "",
    "1. Qual técnica adiciona backtracking ao raciocínio?",
    "2. ReWOO é melhor que ReAct quando...?",
    "3. Reflexion difere de retry porque...?",
    "4. Com o1, 'let's think step by step' deve ser...?",
    "5. Qual guardrail é OBRIGATÓRIO em todo agente?",
    "",
    "≥3 acertos = compreensão básica",
], "📖 Respostas: 1) ToT. 2) Tools são independentes e paralelizáveis. 3) Reflexion reflete antes de tentar de novo. 4) Removido (modelo já raciocina). 5) max_steps + orçamento.\n➡️ Conexão.", 41, T)

s = content_slide("Conexão com Próximos Módulos", [
    "ETHAGT05 — Memória de Agentes",
    "  Estende Reflexion com memória persistente",
    "",
    "ETHAGT06 — RAG Agêntico",
    "  Raciocínio + recuperação combinados",
    "",
    "ETHAGT09 — Multi-Agente",
    "  Múltiplos agentes raciocinando juntos",
    "",
    "ETHAGT15 — Meta-Agentes",
    "  Self-Discover como bloco para meta-raciocínio",
], "📖 ETHAGT04 é a base de raciocínio. Memória (ETHAGT05) estende Reflexion. RAG (ETHAGT06) combina raciocínio + recuperação. Multi-agente (ETHAGT09) distribui raciocínio.\n➡️ Leitura.", 42, T, "Mapa da especialização")

s = content_slide("Leitura Recomendada", [
    "Obrigatório:",
    "  Yao et al. Tree of Thoughts (arXiv:2305.10601)",
    "  Shinn et al. Reflexion (arXiv:2303.11366)",
    "",
    "Recomendado:",
    "  Zhou et al. LATS (arXiv:2310.01757)",
    "  Major et al. Self-Discover (arXiv:2402.03620)",
    "  OpenAI. Learning to Reason with LLMs (o1 paper)",
    "",
    "Complementar:",
    "  Wei et al. Chain-of-Thought (arXiv:2201.11903)",
    "  Xu et al. ReWOO (arXiv:2305.18323)",
], "📖 Leiam ToT e Reflexion antes da próxima aula. LATS e Self-Discover são recomendados.\n➡️ Q&A.", 43, T, "Leitura obrigatória e recomendada")

s = title_slide(
    "Perguntas?",
    "Na próxima aula: ETHAGT05 — Memória de Agentes",
    "ETHAGT04 · Obrigado!"
)
add_notes(s, "📖 Q&A aberto. 'Qual técnica vocês vão testar primeiro no projeto de vocês?'\n➡️ Encerramento.")

# ═══════════════════════════════════════
# Salvar
# ═══════════════════════════════════════

output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT04-Apresentacao.pptx")
prs.save(output_path)
print(f"PPTX gerado: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
