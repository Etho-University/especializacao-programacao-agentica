#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT12: AgentOps, Observabilidade & Avaliação
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Paleta Etho ───
PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ───

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj="", acronyms=""):
    if acronyms:
        add_textbox(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.35), acronyms, size=9, color=INFO)
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT12"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    if acronyms:
        add_textbox(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4), acronyms, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

T = 78  # Total

# ═══════════════════════════════════════
# SEÇÃO A — Abertura (1-6)
# ═══════════════════════════════════════

s = title_slide(
    "AgentOps, Observabilidade & Avaliação",
    "Universidade Etho · Especialização em Programação Agêntica\nFase D — Produção, Governança e Fronteira · 30 h",
    "ETHAGT12"
)
add_notes(s, "📖 Bem-vindos. Esta aula separa quem brinca de agentes de quem opera agentes em produção. Cobrimos observar (traces), medir (métricas), avaliar (evals) e iterar (CI).\n💡 Analogia: ter um painel com velocímetro, combustível e temperatura.\n❓ 'Quem já deployou agente sem saber se estava melhor ou pior?'\n⚠️ AgentOps é mais que logs — é observar, medir, avaliar, iterar.\n➡️ Objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: estabelecer rigor experimental em sistemas de agentes",
    "",
    "Objetivos específicos:",
    "1. Implementar observabilidade end-to-end (traces, spans, métricas)",
    "2. Construir pipelines de avaliação automatizada (LLM-as-judge, golden cases)",
    "3. Aplicar benchmarks canônicos (SWE-bench, GAIA, τ-bench, AgentBench, WebArena)",
    "4. Operar ciclos de melhoria contínua com dados",
    "5. Reportar resultados com rigor (eval report)",
], "📖 Objetivos mensuráveis: implementar, construir, aplicar, operar, reportar. Vamos revisar no Slide 70.\n❓ 'Qual objetivo é mais desafiador no trabalho de vocês?' (geralmente #2 ou #4)\n⚠️ LLM-as-judge não é plug-and-play — tem vieses. Cobrimos na Seção D.\n➡️ Competências.", 2, T, "5 objetivos mensuráveis", acronyms="LLM = Large Language Model — Modelo de Linguagem de Grande Escala  ·  SWE-bench = Software Engineering Benchmark — benchmark de issues reais do GitHub")

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → A (Avançado)",
    "C2 Multi-Agent Systems → B (Básico)",
    "C4 Agent Memory → B (Básico)",
    "C5 AgentOps & Avaliação → A (Avançado)",
    "",
    "Módulo atinge nível A em C1 e C5: consolida maturidade técnica",
    "C2 e C4 em B: fundação para módulos de multi-agente e memória",
], "📖 Framework Etho: 6 competências em 3 níveis. ETHAGT12 atinge A em duas competências críticas (C1 e C5).\n💡 Analogia: sair de motorista amador para piloto com telemetria.\n⚠️ 'Avaliação' não é quiz — é rigor experimental aplicado.\n➡️ Agenda.", 3, T, "Conectar ao Framework Etho")

s = content_slide("Agenda da Aula", [
    "Bloco 1 (45 min):",
    "  Abertura (8 min) — objetivos, motivação, contexto",
    "  Por que difícil (10 min) — não-determinismo, falácias",
    "  Observabilidade (16 min) — traces, OTel, DEMO",
    "  Intervalo (5 min)",
    "",
    "Bloco 2 (45 min):",
    "  Avaliação (18 min) — LLM-as-judge, golden cases, CI, DEMO",
    "  Benchmarks (14 min) — SWE-bench, GAIA, τ-bench",
    "  Melhoria & Report (12 min) — dataset, eval report",
    "  Fechamento — boas práticas, caso, quiz, projeto, Q&A",
], "📖 Dois blocos. Primeiro: observabilidade (base). Segundo: avaliação e melhoria.\n⚠️ Alunos atrasados perdem o Slide 5 (motivação) — define o tom.\n➡️ Por que estamos aqui?", 4, T, "Roteiro da aula")

s = content_slide("5/5 no Teste, 30% em Produção", [
    "Cenário real: agente funcionou 5/5 no teste, 30% em produção",
    "",
    "Sem métrica, você não sabia que estava falhando",
    "'Parece que funciona' não é medida",
    "Usuários desistem silenciosamente — não reclamam, somem",
    "",
    "Pergunta: Como você sabe se seu agente está melhor ou pior?",
], "📖 Slide mais importante da abertura. Aconteceu com todo mundo. Você testa 5 casos, faz deploy, retenção cai. Sem observabilidade e eval, você está cego.\n💡 Analogia: aprovar remédio testando em 5 pessoas. Funcionou — mas e nos outros 9995?\n❓ Silêncio da turma é o gancho da aula.\n⚠️ 'Testar mais 10' não resolve. Precisa de observabilidade + eval sistemática.\n➡️ Por que AgentOps agora?", 5, T, "Criar tensão — sem métrica, você está cego")

s = content_slide("Por Que AgentOps Agora", [
    "Linha do tempo:",
    "  2022: ChatGPT (LLM como oráculo)",
    "  2023: Primeiros agentes (AutoGPT, ReAct em produção)",
    "  2024: SWE-bench, GAIA, Phoenix/Langfuse surgem",
    "  2025: AgentOps consolida como disciplina",
    "  2026: AgentOps é pré-requisito para produção séria",
    "",
    "Confluência: não-determinismo + custo de runs + iterar com confiança",
    "LLMOps tradicional não cobre: tools, loops, estado, ambiente",
    "Hamel Husain: 'Evals for LLMs' como marco conceitual",
], "📖 LLMOps cobre treinar, deployar, monitorar modelo. AgentOps cobre isso + comportamento emergente: tools, loops, estado, ambiente. É uma camada a mais.\n💡 Analogia: monitorar servidor web (LLMOps) vs sistema distribuído com microserviços (AgentOps).\n➡️ Por que agentes são difíceis de avaliar?", 6, T, "Confluência histórica", acronyms="LLMOps = LLM Operations — praticas de MLOps adaptadas a LLMs  ·  ReAct = Reasoning and Acting — padrao de loop Thought / Action / Observation")

# ═══════════════════════════════════════
# SEÇÃO B — Por que difícil (7-14)
# ═══════════════════════════════════════

s = section_slide(1, "Por que Agentes São Difíceis de Avaliar")
add_notes(s, "Transição: entender o problema antes das soluções.")

s = content_slide("Não-Determinismo em Agentes", [
    "LLMs são probabilísticos: mesma pergunta, respostas diferentes",
    "Agentes amplificam: cada step tem não-determinismo que se acumula",
    "Temperatura, sampling, tool choice podem variar",
    "",
    "Implicação: avaliação única não é suficiente — precisamos de N runs",
    "Mínimo 3 runs por caso; idealmente 5-10",
    "",
    "Pergunta: Você já rodou o mesmo prompt duas vezes com resultados diferentes?",
], "📖 Em código, f(2)=4 sempre. Em agente, f('reserve voo') pode dar 10 caminhos. Variância se acumula em N steps.\n💡 Analogia: jogar moeda 10x em sequência — mesmo começando igual, caminhos diferentes.\n❓ Todos concordam.\n⚠️ Eval com 1 run se surpreende com variância. Sempre média + desvio.\n➡️ Ambiente.", 8, T, "Variância se acumula")

s = content_slide("Dependência de Ambiente", [
    "Tools chamam APIs que mudam de comportamento",
    "RAG depende de base de conhecimento que evolui",
    "Ambiente de execução (FS, web) muda entre runs",
    "",
    "Exemplo: agente que usa API de cotação — depende do momento",
    "",
    "Dilema: mock (previsível mas irreal) vs real (variante mas autêntico)",
    "Solução: ambiente controlado (sandbox) com fixtures conhecidas",
], "📖 Em código você mocka dependências. Em agente, mockar tools perde realismo. Use sandbox com fixtures.\n💡 Analogia: piloto de F1. Simulador (previsível, irreal) vs pista molhada (real, caótico). Precisa de ambos.\n⚠️ Eval em produção direto gera custo e variância alta.\n➡️ Custo.", 9, T, "Ambiente mutável", acronyms="RAG = Retrieval-Augmented Generation — Geracao Aumentada por Recuperacao")

s = content_slide("Custo de Runs", [
    "Cada run = tokens de entrada + tokens de saída × N steps",
    "Agente complexo: 10-50 chamadas de LLM por tarefa",
    "",
    "Cálculo: 1000 casos × 50 chamadas × $0.01/chamada = $500",
    "Benchmarks completos (SWE-bench): centenas de dólares por modelo",
    "",
    "Estratégias: amostragem, subconjuntos, eval incremental",
    "Comece com 10-50 casos; expanda quando confiante",
], "📖 Diferente de testes unitários (gratuitos), cada execução de agente custa tokens. Estratégia importa.\n💡 Analogia: pesquisa de mercado. Comece com 30, valide questionário, depois expanda.\n❓ 'Qual o orçamento de eval de vocês hoje?' (geralmente zero)\n⚠️ Eval completo a cada commit custa caro e demora. Subconjunto no CI; completo noturno.\n➡️ Falácias.", 10, T, "Eval custa dinheiro real")

s = content_slide("Falácias Comuns de Avaliação", [
    "'Funcionou uma vez' → sobrevivência de amostra size=1",
    "'Parece bom' → vibes-based eval, não escala",
    "'O usuário não reclamou' → usuários desistem silenciosamente",
    "'Passou no benchmark' → benchmark ≠ produção",
    "'LLM disse que está correto' → LLM-as-judge sem calibração",
    "",
    "Pergunta: Qual dessas falácias você já cometeu?",
], "📖 5 falácias = causa #1 de agentes que parecem funcionar mas falham. Cada uma é armadilha.\n💡 Analogia: dirigir bêbado e chegar em casa. 'Funcionou!' Desta vez. Na próxima você bate.\n❓ Todos já cometemos.\n⚠️ 'Funcionou uma vez' é sorte, não validação.\n➡️ Pontual vs contínua.", 11, T, "5 armadilhas comuns")

s = content_slide("Avaliação Contínua vs Pontual", [
    "Pontual: 'rodamos eval antes do deploy' → estátua (congela)",
    "Contínua: eval roda a cada mudança, em produção, com drift detection",
    "",
    "Agentes mudam: prompt, tool, modelo, ambiente evoluem",
    "Sem eval contínua: regressão passa despercebida por semanas",
    "",
    "Pontual é foto. Contínua é vídeo. Em agentes, você precisa do vídeo.",
], "📖 Tudo muda em agentes: modelo do provedor atualiza, API muda formato, RAG cresce. Cada mudada pode causar regressão silenciosa.\n💡 Analogia: checkup anual (pontual) vs smartwatch (contínuo). Agentes precisam do smartwatch.\n➡️ A falácia mais comum tem nome.", 12, T, "Processo, não evento")

s = content_slide("Vibes-Based Eval: A Armadilha", [
    "'Se você não tem conjunto de avaliação, você está fazendo",
    " vibes-based development.' — Hamel Husain",
    "",
    "Vibes = intuição, não medida",
    "Intuição é útil para explorar, não para decidir deploy",
    "",
    "Regra: todo deploy precisa de número, não de feeling",
], "📖 Hamel Husain cunhou o termo. Seu cérebro é péssimo em avaliar sistemas estocásticos sem dados.\n💡 Analogia: decidir se dieta funciona 'no feeling'. Sem balança, você mente para si mesmo.\n➡️ Vamos praticar reconhecendo falácias.", 13, T, "Intuição ≠ medida")

s = exercise_slide("Exercício — Identificando Falácias", [
    "Para cada afirmação, identifique a falácia:",
    "",
    "1. 'Testei com 3 exemplos e funcionou, pode ir para produção.'",
    "2. 'O benchmark mostra 80%, então o agente está bom.'",
    "3. 'O GPT-4 disse que a resposta está correta.'",
    "4. 'Ninguém reclamou desde o deploy.'",
    "",
    "Votação rápida (mãos levantadas)",
    "Discussão: 1 min por afirmação",
], "📖 Respostas: 1) sample size=1; 2) benchmark ≠ produção; 3) self-preference sem calibração; 4) sobrevivência silenciosa.\n❓ Mostrar uma por vez, votação. Anotar no quadro.\n⚠️ #4 não é OK — usuários desistem sem reclamar.\n➡️ Observabilidade.", 14, T)

# ═══════════════════════════════════════
# SEÇÃO C — Observabilidade (15-27)
# ═══════════════════════════════════════

s = section_slide(2, "Observabilidade: Traces, Spans, Métricas")
add_notes(s, "Observabilidade é a base. Sem ela, você não sabe o que medir.")

s = content_slide("Traces: O Que São", [
    "Trace = registro completo de uma execução de agente",
    "Estrutura de árvore: root span → child spans → sub-spans",
    "",
    "Cada span: nome, duração, input, output, atributos",
    "Diferença de log: trace mostra relações e hierarquia",
    "",
    "'Sem traces, você está debugando no escuro'",
], "📖 Trace é a unidade fundamental. Árvore hierárquica de spans. Cada span é uma operação. Mostra quem chamou quem, timing, o que entrou/saiu.\n💡 Analogia: árvore genealógica. Logs são lista telefônica; traces mostram hierarquia.\n⚠️ Trace ≠ log. Log é linha; trace é árvore com timing.\n➡️ Anatomia.", 16, T, "Árvore de execução")

s = content_slide("Anatomia de um Trace", [
    "Root span: tarefa completa (ex: 'reservar voo')",
    "Child spans: chamada LLM 1, tool call A, chamada LLM 2, tool call B",
    "Sub-spans: API request dentro de tool call",
    "",
    "Atributos por span: model, tokens, latency, tool name, args, result",
    "Bags: contexto propagado entre spans (user_id, session_id)",
    "",
    "Diagrama: 12-Diagrams/ETHAGT12/trace-anatomy.mmd",
], "📖 Root = tarefa. Child = operações (LLM, tool). Sub-spans = dentro de tools. Atributos = metadados. Bags = contexto propagado.\n💡 Analogia: receita de bolo. Root = 'fazer bolo'. Child = 'misturar', 'assar'. Sub = 'ralar chocolate'.\n⚠️ Sem propagar bags, traces isolados não conectam a usuários.\n➡️ OpenTelemetry.", 17, T, "Hierarquia de spans")

s = content_slide("OpenTelemetry para LLMs", [
    "OpenTelemetry (OTel): padrão CNCF para observabilidade distribuída",
    "",
    "GenAI semantic conventions — atributos padronizados:",
    "  • gen_ai.system (provedor)",
    "  • gen_ai.request.model (modelo)",
    "  • gen_ai.usage.prompt_tokens / completion_tokens",
    "",
    "Vantagem: vendor-neutral — funciona com qualquer backend",
    "Para agentes: traces cruzam LLM calls, tool calls, retrieval",
], "📖 OTel é padrão CNCF. Para LLMs, define GenAI semantic conventions. Vendor-neutral: uma vez, qualquer backend.\n💡 Analogia: USB-C. Antes cada marca tinha conector. Agora padrão. OTel é USB-C da observabilidade.\n⚠️ Sem OTel, você fica preso a uma ferramenta.\n➡️ Tooling.", 18, T, "Padrão vendor-neutral")

s = content_slide("Tooling: LangSmith, Phoenix, Langfuse", [
    "LangSmith — integrado ao LangChain/LangGraph; comercial SaaS",
    "Phoenix (Arize) — open source; foco em LLM observability + eval",
    "Langfuse — open source self-hostable; traces + eval + prompts",
    "OpenLLMetry — instrumentação OTel automática",
    "",
    "Critério de escolha: stack existente, self-host vs SaaS, custo",
    "",
    "Pergunta: Quem já usa alguma dessas em produção?",
], "📖 4 ferramentas dominam. LangSmith se já usa LangChain. Phoenix se quer OSS agnóstico. Langfuse se quer self-host. OpenLLMetry para instrumentação automática.\n💡 Analogia: AWS vs GCP vs Azure. Mesma função. Escolha pelo que você já usa.\n❓ Maioria: nenhuma. Esse é o ponto.\n⚠️ Não escolha pelo hype. Escolha pelo fit.\n➡️ Logs vs traces.", 19, T, "4 ferramentas principais")

s = content_slide("Logs Estruturados vs Traces", [
    "Logs estruturados: JSON com timestamp, step, thought, action, observation",
    "Traces: árvore hierárquica com timing e relações",
    "",
    "Logs: simples, barato, bom para debug rápido",
    "Traces: rico, relacional, bom para análise e produção",
    "",
    "Prática: comece com logs estruturados, evolua para traces",
    "Complementares, não mutuamente exclusivos",
], "📖 Logs e traces não competem — complementam. Logs: simples e baratos, debug rápido. Traces: ricos mas custam mais storage.\n💡 Analogia: diário (logs) vs organograma (traces).\n⚠️ Não pule direto para traces sem logs primeiro.\n➡️ Custo de observabilidade.", 20, T, "Complementares")

s = content_slide("Custo de Observabilidade (Amostragem)", [
    "Cada trace gera storage + processamento",
    "Em alto volume: 100% de traces é caro",
    "",
    "Amostragem: registrar apenas 1 em N traces",
    "Estratégias: head-based (decide no início), tail-based (decide no fim)",
    "Tail sampling: sempre logar traces com erro, amostrar sucesso",
    "",
    "Regra: 100% em dev, amostrado em produção, 100% em erro",
], "📖 Observabilidade tem custo. Em alto volume, 100% é caro. Amostragem é a solução. Tail sampling prioriza erros.\n💡 Analogia: câmera de segurança. Não grava tudo; grava quando detecta movimento.\n⚠️ Amostrar produção e perder erros é erro grave. Erros = 100% capturados.\n➡️ Dashboard.", 21, T, "Trade-off de custo")

s = content_slide("Dashboard Mínimo de Observabilidade", [
    "Painel 1: Success rate por tarefa (últimas 24h)",
    "Painel 2: Latência P50/P95/P99 por step",
    "Painel 3: Custo por execução (tokens × preço)",
    "Painel 4: Tool usage (qual tool, quantas vezes, sucesso/falha)",
    "Painel 5: Erros agrupados por tipo",
    "Painel 6: Distribuição de steps (loops?)",
    "",
    "Sem esses 6 painéis, você está operando cego",
], "📖 Dashboard mínimo: 6 painéis. Success rate, latência P50/P95/P99, custo, tool usage, erros, steps. Com isso você opera em produção.\n💡 Analogia: painel do carro. Velocímetro, combustível, temperatura — cada um essencial.\n➡️ Custo e latência são métricas de primeira classe.", 22, T, "6 painéis essenciais", acronyms="P95 = Percentil 95 — latencia abaixo da qual 95% das requisicoes ficam  ·  P99 = Percentil 99 — latencia abaixo da qual 99% das requisicoes ficam")

s = content_slide("Métricas de Primeira Classe: Custo e Latência", [
    "Custo por execução: tokens in/out × preço por modelo",
    "Latência cumulativa: serial = soma; paralelo = max",
    "",
    "P95 e P99 importam mais que média (cauda longa)",
    "Se P99 é 30s, 1 em 100 usuários espera meio minuto — inaceitável",
    "",
    "Orçamento por execução: abortar se custo > limite",
    "Em produção: alerta se custo médio sobe 20%",
], "📖 Custo e latência são métricas de produto, não detalhes. Média mente — P95 e P99 importam mais.\n💡 Analogia: salário. Média R$ 5k mas se P95 é R$ 500, tem gente sofrendo. Em latência, sofrimento está na cauda.\n⚠️ Reportar média de latência é incompleto. Sempre P50/P95/P99.\n➡️ DEMO!", 23, T, "Custo e latência importam")

s = code_slide("DEMO — Traces Everywhere", "Lab 1: adicionar LangSmith a agente ReAct existente\n\nAntes:\n  agent = ReActAgent(tools=[...])\n  result = agent.run('Qual o preço do iPhone?')\n\nDepois (3 linhas):\n  import os\n  os.environ['LANGSMITH_TRACING'] = 'true'\n  os.environ['LANGSMITH_PROJECT'] = 'ethagt12-demo'\n  # + @traceable decorator ou tracing_callback\n\n  @traceable\n  def run_agent(query):\n      return agent.run(query)\n\nMostrar no trace viewer:\n  - Root span: tarefa completa\n  - Child spans: LLM calls, tool calls\n  - Latência por span\n  - Custo por chamada (tokens × preço)\n  - Identificar gargalo (geralmente tool call)", "📖 Vamos ao vivo. Agente ReAct sem observabilidade. Adicionar LangSmith — 3 linhas. Ver trace completo: cada step, latência, custo.\n⚠️ Se API falhar, tenho screenshot do trace gravado.\n💡 Analogia: ligar modo desenvolvedor do navegador. De repente você vê tudo.\n➡️ Pergunta da demo.", 24, T, "DEMO ao vivo")

s = exercise_slide("Pergunta da DEMO", [
    "Discussão em duplas (2 min):",
    "",
    "1. O que o trace revelou que logs simples não mostrariam?",
    "2. Qual step do agente é o gargalo de latência?",
    "3. Como você decidiria se vale otimizar esse step?",
    "",
    "Respostas esperadas:",
    "  • Hierarquia, timing cumulativo, custo por step",
    "  • Geralmente a tool call (latência de API)",
    "  • ROI: tempo economizado vs esforço",
], "📖 Resposta chave: trace revela hierarquia (quem chamou quem), timing cumulativo, custo por step, atributos estruturados. Esses são os superpoderes.\n❓ Pedir 2 duplas compartilharem. Anotar no quadro.\n➡️ Exercício de leitura de trace.", 25, T)

s = exercise_slide("Exercício — Lendo um Trace", [
    "Trace com problema (agente em loop de 8 steps):",
    "",
    "  [STEP 0] Action: search_price('Produto X')",
    "  [STEP 1] Action: search_price('Produto X')",
    "  [STEP 2] Action: search_price('Produto X')",
    "  [STEP 3] Action: search_price('Produto X')  ← LOOP",
    "  ...",
    "  [STEP 7] max_steps reached. Abortando.",
    "",
    "Em duplas: identificar onde e por quê. Propor 2 correções.",
    "2 min discussão, 1 min compartilhar",
], "📖 Loop: agente repete mesma action sem processar Observation. Causas: prompt não instrui usar observations; contexto não inclui anteriores; modelo fraco.\n⚠️ 'Aumentar max_steps' não resolve — só mascara.\n➡️ Pergunta custo vs investimento.", 26, T)

s = content_slide("Observabilidade — Custo ou Investimento?", [
    "Argumento custo: storage, tooling, overhead de instrumentação (1-5% latency)",
    "Argumento investimento: debugging 10x mais rápido, regressão detectada,",
    "  otimização guiada por dados",
    "",
    "ROI quase sempre positivo",
    "'Sem observabilidade, você não tem como melhorar o que não mede'",
    "",
    "A pergunta não é 'se' — é 'quanto'",
], "📖 Custo direto: storage, tooling, overhead. Investimento: debugging rápido, regressão detectada, otimização com dados. ROI positivo.\n➡️ Avaliação automatizada.", 27, T, "ROI positivo")

# ═══════════════════════════════════════
# SEÇÃO D — Avaliação (28-43)
# ═══════════════════════════════════════

s = section_slide(3, "Avaliação Automatizada: LLM-as-Judge, Golden Cases, Regressão")
add_notes(s, "Observabilidade é olhar. Avaliação é julgar. Vamos automatizar o julgamento.")

s = content_slide("Avaliação Manual Não Escala", [
    "Manual: humano lê output e julga → lento, caro, inconsistente",
    "",
    "Não escala: 1000 casos × 5 min/caso = 83 horas",
    "Humano fica cansado: qualidade decai ao longo do dia",
    "Não é reproduzível: mesmo humano julga diferente em outro dia",
    "",
    "Solução: LLM-as-judge + golden cases + métricas programáticas",
], "📖 Manual funciona para 10 casos. Para 1000, não. Humano cansa, é inconsistente, custa caro. Combine três técnicas.\n➡️ LLM-as-judge.", 29, T, "Manual não escala")

s = content_slide("LLM-as-Judge: Conceito", [
    "LLM-as-judge: um LLM avalia a saída de outro LLM",
    "",
    "Padrão: input + output + critério (rubrica) → judge → score + justificativa",
    "",
    "Vantagem: escala, barato, reproduzível (com temperatura 0)",
    "Quando usar: tarefas subjetivas (qualidade de resposta, completude)",
    "Quando NÃO usar: ground truth exato (use string match)",
], "📖 Elegante: LLM avalia LLM. Escala, barato, reproduzível. Mas não é mágica — tem vieses.\n💡 Analogia: corretor de prova. Treina com rubrica, corrige 1000 redações. Mas pode ter vieses.\n⚠️ Para tarefas objetivas, métricas programáticas são melhores e mais baratas.\n➡️ Vieses.", 30, T, "LLM avalia LLM", acronyms="LLM-as-Judge = uso de um LLM como avaliador automatico de saidas de outro LLM")

s = content_slide("LLM-as-Judge: Vieses", [
    "Positional bias: prefere primeira ou última opção",
    "Sycophancy: concorda com o que o humano parece querer",
    "Verbosity bias: prefere respostas mais longas",
    "Self-preference: modelo prefere saídas do mesmo modelo",
    "Knowledge bias: judge não sabe o que o agente deveria saber",
    "",
    "Pergunta: Qual viés é mais perigoso para seu caso de uso?",
], "📖 5 vieses. Cada um tem mitigação. Importante: acima de tudo, calibre com humano.\n❓ Resposta varia por domínio. Comparações A/B: positional. Assistentes: sycophancy.\n⚠️ Ignorar vieses e confiar cegamente = fé cega.\n➡️ Mitigações.", 31, T, "5 vieses conhecidos")

s = content_slide("LLM-as-Judge: Mitigações", [
    "Rubrica clara: critérios explícitos, não 'avalie a qualidade'",
    "Exemplos: few-shot com casos calibrados por humano",
    "Múltiplos judges: 3+ modelos, média ou votação",
    "Calibração: comparar judge com humano em subconjunto (≥80%)",
    "Swap de posição: rodar A vs B e B vs A, média dos resultados",
    "Cross-model: usar modelo diferente do agente como judge",
], "📖 Cada viés tem mitigação. Positional → swap. Sycophancy → rubrica explícita. Verbosity → penalizar. Self → cross-model. Knowledge → contexto. Acima de tudo: calibre com humano.\n💡 Analogia: treinar corretor. Não entrega rubrica e vai embora — calibra, mostra exemplos.\n⚠️ Sem calibração, você não sabe se judge está certo.\n➡️ Golden cases.", 32, T, "Mitigação por viés")

s = content_slide("Golden Cases", [
    "Golden case = par (input, critério de sucesso)",
    "",
    "Input: pergunta/tarefa para o agente",
    "Critério: como saber se a resposta está correta",
    "  • Exato: string match, regex, JSON schema",
    "  • Subjetivo: rubrica para LLM-as-judge",
    "  • Funcional: tool foi chamada? ação foi executada?",
    "",
    "Conjunto crescente: cada bug encontrado vira golden case",
], "📖 Golden cases são testes unitários do mundo de agentes. Todo bug vira caso novo — nunca mais regredirá.\n💡 Analogia: teste de regressão em software tradicional.\n⚠️ 'Resposta deve ser boa' não é critério. Use critério operacionalizável.\n➡️ Exemplo concreto.", 33, T, "Testes unitários de agentes")

s = code_slide("Escrevendo um Golden Case", "# Caso objetivo (factual)\ngolden_case = {\n    'id': 'GC-042',\n    'input': 'Qual a capital da França?',\n    'expected_behavior': 'Resposta deve mencionar Paris',\n    'eval_fn': lambda response: 'Paris' in response,\n    'category': 'factual'\n}\n\n# Caso subjetivo (LLM-as-judge)\ngolden_case = {\n    'id': 'GC-043',\n    'input': 'Resuma este artigo em 3 pontos.',\n    'expected_behavior': '3 pontos fiéis ao original',\n    'eval_fn': llm_as_judge(\n        rubrica='3 pontos, fiel ao original, claro'\n    ),\n    'category': 'subjective'\n}\n\n# Estrutura: id, input, expected_behavior, eval_fn, category\n# Versione no Git — casos são código", "📖 Estrutura: id, input, expected_behavior, eval_fn, category. Objetivo: assertion. Subjetivo: LLM-as-judge com rubrica.\n⚠️ Sempre versione golden cases no Git.\n➡️ Conjuntos de regressão.", 34, T)

s = content_slide("Conjuntos de Regressão", [
    "Conjunto de regressão: N golden cases que rodam a cada mudança",
    "Toda mudança de prompt, tool, ou modelo roda o conjunto",
    "Se score cai: deploy bloqueado (CI gate)",
    "",
    "Crescimento: 10 → 100 → 1000+ casos ao longo de meses",
    "Categorização: factual, multi-step, tool-use, edge-case, error-handling",
    "Manutenção: remover casos obsoletos, adicionar casos de bug",
], "📖 Conjunto de regressão = ativo mais valioso do projeto. Comece com 10, cresça para centenas. Categorize para diagnosticar.\n💡 Analogia: testes de aplicação madura. Começa com 10, cresce para milhares.\n⚠️ Conjunto sem manutenção fica obsoleto.\n➡️ Métricas de tarefa.", 35, T, "Ativo do projeto")

s = content_slide("Métricas de Tarefa", [
    "Success rate: % de casos que passaram no critério",
    "Partial credit: não é tudo ou nada (0.7 se 7 de 10 sub-tarefas)",
    "Custo por execução: tokens × preço",
    "Latência P50/P95: tempo total da tarefa",
    "Score oficial (benchmarks): ex: SWE-bench usa % resolved",
    "",
    "Reporte TODAS — success rate sozinha é incompleta",
], "📖 Métricas de tarefa medem resultado. Success rate, partial credit, custo, latência. Reporte todas — sem custo e latência, métrica é incompleta.\n⚠️ Só success rate é incompleto. 90% com 60s de latência é inutilizável.\n➡️ Métricas de processo.", 36, T, "Resultado final")

s = content_slide("Métricas de Processo", [
    "Número de steps: quantas iterações do loop ReAct",
    "Loops detectados: agente repetiu a mesma action?",
    "Tool misuse rate: % de tool calls que falharam ou foram inúteis",
    "Token efficiency: tokens usados / tokens necessários",
    "Recovery rate: % de vezes que agente se recuperou de erro",
    "",
    "Processo explica POR QUE o agente falhou — direciona correção",
], "📖 Métricas de processo dizem COMO o agente chegou ao resultado. Explicam sucesso/fracasso e direcionam onde corrigir.\n💡 Analogia: tempo final da maratona (tarefa) vs ritmo, cadência, onde tropeçou (processo).\n➡️ A/B testing.", 37, T, "Como chegou ao resultado")

s = content_slide("A/B Testing de Prompts e Tools", [
    "A/B test: versão A (atual) vs versão B (nova) no mesmo conjunto",
    "Variáveis: prompt, tool description, modelo, temperatura, max_steps",
    "",
    "Uma variável por vez (ou design de experimentos)",
    "Métricas: success rate, custo, latência, satisfação",
    "Significância: mínimo 30 casos por versão, idealmente 100+",
    "",
    "Pergunta: Você mudaria o prompt sem A/B test?",
], "📖 Igual A/B test em produto. Versão A vs B no mesmo conjunto. Mude uma variável por vez. Significância: mínimo 30 por versão.\n💡 Analogia: teste A/B de produto. Não muda botão, cor e texto ao mesmo tempo.\n⚠️ Mudar 3 coisas e atribuir a uma = erro metodológico.\n➡️ Pipeline de CI.", 38, T, "Experimentação controlada")

s = content_slide("Pipeline de Eval com CI", [
    "Mudança (prompt/tool/modelo) → CI roda eval",
    "CI executa: golden cases + subconjunto de benchmark",
    "LLM-as-judge avalia resultados",
    "Compara com baseline: houve regressão?",
    "Sim → bloquear deploy; Não → permitir deploy",
    "",
    "Pipeline automatizado, não manual",
    "Diagrama: 12-Diagrams/ETHAGT12/eval-pipeline.mmd",
], "📖 Diagrama mais importante da aula. Toda mudança dispara eval. CI roda golden + benchmark subset. Gate: regressão > threshold bloqueia. Transforma 'vibes' em 'evidence'.\n💡 Analogia: CI tradicional bloqueia merge sem testes passarem. Em agentes, 'testes' são evals.\n⚠️ Eval manual antes do deploy = esquecem de rodar. Automatize sempre.\n➡️ DEMO!", 39, T, "CI gate")

s = code_slide("DEMO — Eval Automatizado", "Lab 2: pipeline de eval com LLM-as-judge\n\n# Agente com 10 golden cases\nagent = build_agent()\ngolden_cases = load_golden_cases('cases.json')  # 10 casos\n\n# Versão baseline do prompt\nscore_baseline = run_eval(agent, golden_cases)\n# Score: 85%\n\n# Mudar prompt (piorar de propósito)\nagent.set_system_prompt('Novo prompt, menos instrutivo...')\n\n# Rodar eval novamente\nscore_novo = run_eval(agent, golden_cases)\n# Score: 72%  ← REGRESSÃO DETECTADA\n\n# Eval report\nreport = generate_eval_report({\n    'baseline': score_baseline,\n    'novo': score_novo,\n})\n# Mostra: quais casos regrediram? Por quê?\n# LLM-as-judge justifica cada falha", "📖 Vamos ao vivo. Agente com 10 golden cases. Score 85%. Mudo prompt. Score cai para 72%. Eval report mostra caso por caso, qual passou, qual falhou, e por quê.\n⚠️ Se API falhar, tenho screenshot do eval report.\n➡️ Pergunta da demo.", 40, T)

s = exercise_slide("Pergunta da DEMO", [
    "Discussão aberta (2 min):",
    "",
    "1. A regressão é real ou o judge está errado?",
    "2. Como distinguir bug do agente de viés do judge?",
    "3. O que fazer: reverter o prompt ou corrigir os casos?",
    "",
    "Como distinguir: olhar caso por caso.",
    "Se falhas fazem sentido → regressão real.",
    "Se não fazem sentido → viés do judge.",
], "📖 Pergunta difícil. Pode ser: (1) regressão real; (2) viés do judge; (3) casos errados. Como distinguir? Olhe caso por caso. Se faz sentido, é regressão. Se não, é viés.\n❓ Pedir 2-3 opiniões. Anotar no quadro.\n➡️ Exercício de golden cases.", 41, T)

s = exercise_slide("Exercício — Escrevendo um Golden Case", [
    "Cenário: agente de reserva de voo",
    "",
    "Em duplas, escrevam 2 golden cases:",
    "",
    "1. Caso factual (ex: 'reserve voo direto SP-Rio amanhã')",
    "   Critério: assertion que resposta menciona voo direto",
    "",
    "2. Caso subjetivo (ex: 'quais opções de viagem para o RJ?')",
    "   Critério: rubrica para LLM-as-judge (clareza, utilidade)",
    "",
    "3 min escrita, 2 min compartilhar",
], "📖 Em duplas. Critério deve ser operacionalizável — não 'resposta boa'. 'Menciona preço e horário' é critério.\n❓ Pedir 2 duplas compartilharem. Validar critérios.\n⚠️ 'Resposta útil' não é critério. 'Menciona preço e horário' é.\n➡️ V/F.", 42, T)

s = exercise_slide("V/F — LLM-as-Judge É Sempre Confiável", [
    "Verdadeiro ou Falso:",
    "",
    "'LLM-as-judge é sempre confiável.'",
    "",
    "Resposta: FALSO",
    "",
    "Tem 5 vieses (positional, sycophancy, verbosity, self-preference, knowledge)",
    "Mitigações necessárias: rubrica, calibração, múltiplos judges",
    "Sem calibração com humano, pode estar sistematicamente errado",
], "📖 Falso. Não é always-right. Sem mitigação, pode estar sistematicamente errado — e você nem sabe. Regra: calibre sempre com humano em subconjunto. <80% concordância = ajuste rubrica.\n➡️ Benchmarks.", 43, T)

# ═══════════════════════════════════════
# SEÇÃO E — Benchmarks (44-54)
# ═══════════════════════════════════════

s = section_slide(4, "Benchmarks Canônicos: SWE-bench, GAIA, τ-bench")
add_notes(s, "Depois de avaliação custom, vamos para avaliação comparável: benchmarks.")

s = content_slide("Por que Benchmarks?", [
    "Benchmark = conjunto padronizado de tarefas com critério objetivo",
    "Permite comparar modelos e arquiteturas de forma justa",
    "",
    "Reprodutível: mesma tarefa, mesmo critério, resultado comparável",
    "Referência da comunidade: 'meu agente faz X% no SWE-bench'",
    "",
    "Limitação: benchmark ≠ produção",
    "Pergunta: Você confia em score de benchmark sem ter rodado?",
], "📖 Benchmark é padronizado. Permite comparar. Mas é ambiente controlado — não representa seu caso de uso real. Combine: benchmark (comparável) + eval custom (representativo).\n💡 Analogia: prova do ENEM. Padronizada, comparável, mas não diz tudo.\n❓ Resposta certa: não, não confie sem rodar.\n⚠️ Score de benchmark só é marketing sem eval custom.\n➡️ Landscape.", 45, T, "Comparável mas ≠ produção")

s = content_slide("Landscape de Benchmarks", [
    "SWE-bench — código (resolver issues de GitHub)",
    "GAIA — raciocínio geral multi-step",
    "τ-bench — tool use em domínios (airline, retail)",
    "WebArena — navegação web autônoma",
    "AgentDojo — segurança (injeção em agentes)",
    "τ²-bench — multi-agent tool use",
    "",
    "Diagrama: 12-Diagrams/ETHAGT12/benchmark-landscape.mmd",
], "📖 6 benchmarks canônicos. Cada um testa capacidade diferente. Escolha pelo fit com seu domínio.\n⚠️ Não escolha pelo hype.\n➡️ SWE-bench.", 46, T, "6 benchmarks canônicos")

s = content_slide("SWE-bench / SWE-bench Verified", [
    "SWE-bench: 2.294 issues reais de 12 repositórios Python open source",
    "Tarefa: resolver issue → gerar patch → testes passam",
    "",
    "SWE-bench Verified: subconjunto de 500 com validação humana",
    "Claude 3.5 Sonnet: ~49% no Verified (dez/2024)",
    "",
    "Por que importa: código é verificável (testes = ground truth)",
    "Fonte: Jimenez et al., arXiv:2310.06770",
], "📖 Benchmark de código mais influente. Issues reais de GitHub, agente gera patch, testes validam. Ground truth objetivo. Claude 3.5: 49% no Verified.\n💡 Analogia: prova prática de programação.\n⚠️ 49% é altíssimo para engenharia real.\n➡️ GAIA.", 47, T, "Código verificável")

s = content_slide("GAIA", [
    "GAIA: 466 questões de raciocínio multi-step com tools",
    "Níveis: Level 1 (simples) → Level 3 (complexo, dezenas de steps)",
    "",
    "Requer: web search, file processing, reasoning, tool use",
    "Human: ~92% no Level 1",
    "GPT-4 + plugins: ~15%",
    "",
    "Por que importa: testa capacidade geral de agente assistente",
    "Fonte: Mialon et al., arXiv:2311.12983",
], "📖 Testa raciocínio geral. Gap humano (92%) vs GPT-4 (15%) mostra quanto agentes ainda são fracos em multi-step.\n💡 Analogia: caça ao tesouro. Cada passo te leva ao próximo.\n⚠️ Level 1 não é fácil — exige múltiplos steps e tools.\n➡️ τ-bench.", 48, T, "Raciocínio geral")

s = content_slide("τ-bench", [
    "τ-bench: tool-agent-user interaction em domínios",
    "Agente simula atendente com acesso a APIs (policy, booking, etc.)",
    "",
    "Avalia: success rate em tarefas com usuários simulados",
    "Domínios: airline (policy-driven), retail (catalog-driven)",
    "",
    "Por que importa: testa tool use realista com políticas",
    "Fonte: Yao et al., arXiv:2404.44529",
    "Usado no projeto do módulo",
], "📖 Elegante: simula atendente com usuário simulado e APIs com policy. Sucesso medido pela tarefa completada. Testa tool use realista.\n💡 Analogia: simulação de atendimento. Atendente, cliente, sistemas.\n➡️ AgentBench e WebArena.", 49, T, "Tool use com policy")

s = content_slide("AgentBench e WebArena", [
    "AgentBench: 8 ambientes (OS, DB, KG, web, card game, etc.)",
    "Avalia capacidades diversas: reasoning, planning, tool use",
    "",
    "WebArena: navegação web autônoma em sites simulados",
    "Tarefas: 'encontre o produto mais barato', 'agende reunião'",
    "VisualWebArena: adiciona compreensão visual",
    "",
    "Fontes: Liu et al. (arXiv:2308.03688), Zhou et al. (arXiv:2307.13854)",
], "📖 Dois de panorama amplo. AgentBench: 8 ambientes (decatlo). WebArena: navegação web específica. VisualWebArena adiciona visão.\n💡 Analogia: AgentBench = decatlo. WebArena = prova específica de natação.\n➡️ Como rodar.", 50, T, "Panorama amplo")

s = content_slide("Como Rodar Localmente", [
    "SWE-bench: Docker + eval harness (repositório GitHub)",
    "GAIA: download do dataset HuggingFace + tools (web search, file tools)",
    "τ-bench: pip install + simulate user",
    "",
    "Custo: SWE-bench completo = $$ (muitas runs); subconjuntos = controlável",
    "Dica: comece com subconjunto (10-50 casos) para iterar rápido",
    "Tempo: completo = horas; subconjunto = minutos",
], "📖 Prático. Cada benchmark tem setup próprio. Comece com subconjunto (10-50 casos) para iterar. Só rode completo quando otimizar.\n⚠️ Rodar completo de primeira: demora, custa, não itera.\n➡️ Limites.", 51, T, "Setup prático")

s = content_slide("Limites e Contaminação", [
    "Contaminação: dados de treino podem incluir o benchmark",
    "Overfitting: otimizar para benchmark ≠ melhorar para produção",
    "Coverage: benchmark cobre domínio específico, não seu caso",
    "Saturação: modelos melhores → benchmark fica fácil",
    "",
    "Detecção de contaminação: overlaps de n-gramas, memorização",
    "Pergunta: Score alto em SWE-bench garante produção?",
], "📖 4 limites. Contaminação, overfitting, coverage, saturação. Score de benchmark é sinal, não prova.\n❓ Resposta: não.\n⚠️ Score de benchmark só é marketing sem eval custom.\n➡️ Exercício.", 52, T, "4 limites críticos")

s = exercise_slide("Exercício — Escolhendo um Benchmark", [
    "Para cada cenário, escolha o benchmark:",
    "",
    "1. 'Agente de coding que resolve issues de GitHub' → ?",
    "2. 'Assistente de pesquisa geral multi-step' → ?",
    "3. 'Agente de atendimento com APIs e políticas' → ?",
    "4. 'Bot que navega e-commerce e completa compras' → ?",
    "",
    "Votação rápida (mãos levantadas)",
], "📖 Respostas: 1) SWE-bench; 2) GAIA; 3) τ-bench; 4) WebArena. Escolha pelo domínio.\n❓ Votação, um por vez. Anotar.\n➡️ Benchmark vs produção.", 53, T)

s = content_slide("Benchmark vs Produção", [
    "Pergunta: Score alto em benchmark garante produção?",
    "Resposta: NÃO necessariamente",
    "",
    "Benchmark é ambiente controlado; produção é mundo real",
    "Combinar: benchmark (comparável) + eval custom (representativo)",
    "",
    "'Benchmark diz o que é possível; eval custom diz o que é real'",
], "📖 Fechamento da seção. Benchmark é laboratório. Produção é mundo real. Score alto é NECESSÁRIO mas não SUFICIENTE.\n➡️ Melhoria contínua.", 54, T, "Lab vs mundo real")

# ═══════════════════════════════════════
# SEÇÃO F — Melhoria & Report (55-64)
# ═══════════════════════════════════════

s = section_slide(5, "Ciclo de Melhoria Contínua e Reportando Resultados")
add_notes(s, "Avaliar é uma coisa. Melhorar com base nela é outra.")

s = content_slide("Dataset Crescente", [
    "'Você não mediu se não mediu' — cada execução é dado",
    "",
    "Produção gera casos: logs de conversas, feedback de usuário",
    "Triagem: humano classifica (bom/ruim/edge case)",
    "Bug encontrado → novo golden case (nunca mais regredir)",
    "",
    "Crescimento: 10 → 50 → 200 → 1000+ casos ao longo de meses",
    "Dataset é VANTAGEM COMPETITIVA: concorrente não tem seus casos",
], "📖 Ativo mais valioso do projeto. Produção gera dados. Tria. Cada bug vira caso. Cresce de 10 para 1000. Esse dataset é vantagem competitiva.\n💡 Analogia: biblioteca de casos médicos. Cada paciente gera conhecimento. Em anos, milhares de casos.\n⚠️ Não triar dados de produção = jogar ouro fora.\n➡️ CI.", 56, T, "Ativo competitivo")

s = content_slide("CI para Agentes (Testes de Regressão)", [
    "Toda mudança (prompt, tool, modelo, config) dispara eval",
    "CI roda: golden cases + subconjunto de benchmark",
    "",
    "Gate: se regressão > threshold → bloquear merge",
    "Threshold: 0% para casos críticos, 5% para não-críticos",
    "Velocidade: CI eval precisa rodar em < 10 min (subconjunto)",
    "Full eval: noturno ou pré-deploy",
], "📖 CI para agentes = CI tradicional mas 'testes' são evals. Gate bloqueia regressão. Threshold varia por criticidade. Velocidade importa — subconjunto no CI, completo noturno.\n💡 Analogia: CI tradicional. Não faz merge sem testes passarem.\n⚠️ Full eval no CI demora 1h e bloqueia dev. Use subconjunto.\n➡️ Shadow e canary.", 57, T, "CI gate")

s = content_slide("Shadow Runs e Canary", [
    "Shadow run: nova versão roda em paralelo sem afetar usuário",
    "Compara outputs: versão atual vs versão nova",
    "",
    "Canary: 5% → 25% → 50% → 100% do tráfego",
    "Rollback automático: se métricas degradam, volta para anterior",
    "",
    "Para agentes: shadow é ideal — sem risco para usuário",
    "Agente não-determinístico precisa de mais amostras",
], "📖 Duas estratégias de deploy seguro. Shadow: paralelo sem risco. Canary: gradual com rollback. Para agentes, shadow é ideal.\n💡 Analogia: Shadow = ensaio geral. Canary = estreiar em poucas cidades antes do lançamento nacional.\n➡️ Feedback humano.", 58, T, "Deploy seguro")

s = content_slide("Feedback Humano Estruturado", [
    "Thumbs up/down: simples, mas sem contexto",
    "Feedback estruturado: categoria (errado/incompleto/lento/off-topic) + texto",
    "Implicit feedback: usuário refez a pergunta? Abandonou?",
    "",
    "Loop: feedback → triagem → golden case → regressão",
    "Cuidado: feedback enviesado (só frustrados respondem)",
], "📖 Feedback humano vira dado de eval. Estruturado é melhor que thumbs. Implicit é poderoso (refez pergunta? abandonou?).\n💡 Analogia: pesquisa de satisfação. Quem responde é quem amou ou odiou.\n⚠️ Confiança cega em feedback é erro. Cruze com métricas de uso.\n➡️ Eval report.", 59, T, "Loop de feedback")

s = section_slide(6, "Reportando Resultados com Rigor")
add_notes(s, "Medir é metade. Reportar com rigor é a outra metade.")

s = content_slide("Eval Report — Template", [
    "Template: 24-Templates/template-eval-report.md",
    "",
    "Seções:",
    "1. Sumário executivo (1 parágrafo)",
    "2. Metodologia (dataset, métricas, N runs)",
    "3. Resultados (tabela: métrica × versão)",
    "4. Análise de falhas (categorização)",
    "5. Comparações (vs baseline, vs humano)",
    "6. Recomendações (deploy? corrigir? investigar?)",
    "",
    "Princípio: REPRODUTÍVEL — outro engenheiro consegue rerodar",
], "📖 Eval report é entrega profissional. 6 seções. Princípio fundamental: REPRODUTÍVEL.\n💡 Analogia: paper científico. Metodologia, resultados, discussão, conclusão. Reprodutível por outros.\n⚠️ Sem metodologia, ninguém reproduz.\n➡️ Análise de falhas.", 61, T, "Template de 6 seções")

s = content_slide("Análise de Falhas (Categorização)", [
    "Categorias canônicas:",
    "  1. Não entendeu a tarefa (interpretação errada)",
    "  2. Tool errada (escolheu tool inadequada)",
    "  3. Alucinação (inventou fato/tool/result)",
    "  4. Loop infinito (não convergiu)",
    "  5. Erro de tool (API falhou, formato errado)",
    "  6. Incompleto (parou antes de terminar)",
    "",
    "Por que categorizar: direciona onde investir esforço",
    "Exemplo: 60% 'tool errada' → melhorar ACI",
], "📖 'Falhou em 30%' é inútil. '60% tool errada, 20% alucinação...' é actionable. Categoria direciona esforço.\n💡 Analogia: triagem médica. Não basta 'está doente'. Tem que dizer 'gripe', 'dengue', 'fratura'. Cada categoria tem tratamento.\n➡️ Comparações honestas.", 62, T, "Direciona esforço", acronyms="ACI = Agent-Computer Interface — Interface Agente-Computador")

s = content_slide("Comparações Honestas", [
    "Sempre comparar vs baseline (versão anterior)",
    "Comparar vs humano quando possível (ground truth)",
    "Mesmas condições: mesmo dataset, ambiente, N",
    "",
    "Reportar intervalo de confiança, não só média",
    "Honestidade: reportar onde PERDEU, não só onde ganhou",
    "",
    "Pergunta: O que reportar para o CEO — accuracy ou custo?",
], "📖 Comparação honesta. Sempre vs baseline. Mesmas condições. Intervalo de confiança. Reporte onde perdeu.\n❓ Depende do KPI do CEO. Geralmente: ambos.\n⚠️ Esconder regressão é antiético e prejudica decisão.\n➡️ Exercício.", 63, T, "Rigor e honestidade")

s = exercise_slide("Exercício — Detectando Regressão", [
    "Cenário: time mudou prompt, eval score caiu de 85% para 72%",
    "",
    "Em grupos, analisar:",
    "  • Viés do judge?",
    "  • Mudança real no comportamento?",
    "  • Casos que regrediram — padrão?",
    "",
    "Propor próximos passos:",
    "  A/B test? Reverter prompt? Revisar casos?",
    "",
    "3 min discussão, 2 min compartilhar",
], "📖 Diagnóstico: olhe caso por caso. Se faz sentido, regressão real. Se não, viés. Padrão? Agrupe por categoria. Próximos passos dependem do diagnóstico.\n❓ Pedir 2 grupos compartilharem.\n➡️ Fechamento.", 64, T)

# ═══════════════════════════════════════
# SEÇÃO G — Fechamento (65-78)
# ═══════════════════════════════════════

s = section_slide(7, "Fechamento: Boas Práticas, Caso de Estudo, Quiz")
add_notes(s, "Vamos consolidar com boas práticas, anti-patterns e caso real.")

s = content_slide("Boas Práticas (DO)", [
    "✓ Comece com logs estruturados desde a primeira linha de código",
    "✓ Adicione traces antes de adicionar features",
    "✓ Construa golden cases desde o dia 1 (comece com 10)",
    "✓ Rode eval a cada mudança (CI gate)",
    "✓ Use LLM-as-judge com calibração humana",
    "✓ Categorize falhas para direcionar esforço",
    "✓ Cresça o dataset continuamente (produção → casos)",
    "✓ Reporte com honestidade: inclua onde perdeu",
], "📖 8 boas práticas. Mais importante: COMECE com observabilidade desde a linha 1. Não deixe para depois.\n➡️ Anti-patterns.", 66, T, "8 DOs")

s = content_slide("Anti-Patterns (DON'T)", [
    "✗ Vibes-based eval ('parece bom')",
    "✗ Sem observabilidade em produção",
    "✗ Eval manual que não escala",
    "✗ LLM-as-judge sem calibração",
    "✗ Benchmark como única medida de qualidade",
    "✗ Deploy sem gate de regressão",
    "✗ Não categorizar falhas",
    "✗ Dataset estagnado (mesmos 10 casos para sempre)",
    "✗ Overfitting para benchmark",
], "📖 9 anti-patterns. Pior: vibes-based eval. Deploy sem gate é jogar dados. Dataset estagnado é estagnação técnica.\n❓ 'Quantos desses vocês já cometeram?' Todos vão rir.\n➡️ Caso de estudo.", 67, T, "9 DON'Ts")

s = content_slide("Caso de Estudo — Anthropic × Claude × SWE-bench", [
    "Anthropic avaliando Claude 3.5 Sonnet em SWE-bench",
    "Metodologia: SWE-bench Verified (500 casos, validação humana)",
    "Resultado: ~49% resolved (dez/2024)",
    "",
    "Processo de melhoria: iterar com traces + eval contínua",
    "Análise de falhas: categorizou onde Claude falhava",
    "",
    "Lição: melhoria guiada por dados, não por intuição",
    "'Não há magia — é rigor experimental aplicado'",
], "📖 Anthropic não chegou em 49% por mágica. Rigor: rodar SWE-bench, analisar traces, categorizar falhas, corrigir, re-avaliar. Ciclo iterativo guiado por dados.\n💡 Analogia: treinar atleta de elite. Mede tudo, identifica gargalos, corrige, mede de novo.\n➡️ Resumo.", 68, T, "Sem magia — rigor")

s = content_slide("Resumo da Aula", [
    "1. Agentes são difíceis: não-determinismo, ambiente, custo",
    "2. Observabilidade desde o dia 1: traces, spans, métricas",
    "3. LLM-as-judge com calibração: escala, mas tem vieses",
    "4. Golden cases + regressão: o CI gate de agentes",
    "5. Benchmarks canônicos: comparáveis mas ≠ produção",
    "6. Dataset crescente: sua vantagem competitiva",
    "7. Eval report: rigor e honestidade",
], "📖 7 pontos para levar. Fundação de AgentOps.\n➡️ Checklist.", 69, T, "7 pontos-chave")

s = content_slide("Checklist da Aula", [
    "[x] Explicou por que agentes são difíceis de avaliar",
    "[x] Implementou observabilidade com traces",
    "[x] Construiu pipeline de eval com LLM-as-judge",
    "[x] Descreveu benchmarks canônicos",
    "[x] Explicou CI para agentes e regressão",
    "[x] Apresentou estrutura de eval report",
], "📖 Revisão dos objetivos. Todos cobertos? Se algum não ficou claro, perguntem no Q&A.\n➡️ Quiz.", 70, T, "Objetivos cobertos")

# Quiz (71-75)
s = exercise_slide("Quiz — Pergunta 1", [
    "Qual é a principal limitação do LLM-as-judge?",
    "",
    "A) É muito caro para usar",
    "B) Não consegue avaliar texto",
    "C) Tem vieses que precisam de mitigação",
    "   (positional, sycophancy, verbosity, self-preference, knowledge)",
    "D) Só funciona com GPT-4",
    "",
    "Resposta: C",
], "📖 Resposta: C. Barato, avalia texto, funciona com vários modelos. Limitação: vieses com mitigação.\n➡️ Próxima.", 71, T)

s = exercise_slide("Quiz — Pergunta 2", [
    "O que é um golden case?",
    "",
    "A) Caso de sucesso para marketing",
    "B) Par (input, critério de sucesso) usado como teste de regressão",
    "C) Caso que sempre passa no eval",
    "D) Benchmark padronizado",
    "",
    "Resposta: B",
], "📖 Resposta: B. Teste unitário do mundo de agentes. Não é marketing (A), não é caso que sempre passa (C), não é benchmark (D).\n➡️ Próxima.", 72, T)

s = exercise_slide("Quiz — Pergunta 3", [
    "O que o SWE-bench avalia?",
    "",
    "A) Navegação web autônoma",
    "B) Atendimento ao cliente com tools",
    "C) Resolução de issues reais de GitHub (código)",
    "D) Raciocínio geral multi-step",
    "",
    "Resposta: C",
], "📖 Resposta: C. A é WebArena. B é τ-bench. D é GAIA.\n➡️ Próxima.", 73, T)

s = exercise_slide("Quiz — Pergunta 4", [
    "Em CI para agentes, o que PRINCIPALMENTE bloqueia o deploy?",
    "",
    "A) Latência acima de 1 segundo",
    "B) Regressão no eval score acima do threshold",
    "C) Número de steps diferente da versão anterior",
    "D) Custo acima de $0.01 por run",
    "",
    "Resposta: B",
], "📖 Resposta: B. CI gate é regressão. Os outros geram alertas mas não bloqueiam sozinhos.\n➡️ Última.", 74, T)

s = exercise_slide("Quiz — Pergunta 5", [
    "V/F: 'Bom score em benchmark garante bom desempenho em produção.'",
    "",
    "A) Verdadeiro — benchmarks são representativos",
    "B) Falso — benchmark é ambiente controlado, produção é mundo real",
    "C) Depende — só se o benchmark for do mesmo domínio",
    "D) Verdadeiro — se passou em SWE-bench, está pronto",
    "",
    "Resposta: B",
], "📖 Resposta: B. Benchmark é controlado. Produção é mundo real. Combine benchmark (comparável) + eval custom (representativo).\n➡️ Conexão.", 75, T)

s = content_slide("Conexão e Projeto", [
    "ETHAGT13 — Segurança & Governança: observabilidade como defesa",
    "ETHAGT90 — Capstone: eval report completo",
    "",
    "Projeto do módulo: avaliar agente em subconjunto de τ-bench ou GAIA",
    "  Entrega: eval report + dataset + código de rerun + análise de falhas",
    "  Critério: eval reproduzível; ≥3 categorias de falha documentadas",
    "",
    "Lab 1 (5h): 'Traces Everywhere' — adicionar observabilidade",
    "Lab 2 (5h): 'Eval automatizado' — pipeline com LLM-as-judge + golden cases",
], "📖 ETHAGT12 conecta em duas direções. ETHAGT13 (observabilidade = defesa). ETHAGT90 (eval report). Projeto: avaliem agente em subconjunto de τ-bench ou GAIA. Entreguem eval report completo.\n➡️ Leitura.", 76, T, "Conexões e projeto")

s = content_slide("Leitura Recomendada", [
    "Obrigatório: Jimenez et al., SWE-bench (arXiv:2310.06770)",
    "Obrigatório: Mialon et al., GAIA (arXiv:2311.12983)",
    "Obrigatório: Yao et al., τ-bench (arXiv:2404.44529)",
    "",
    "Recomendado: Hamel Husain, Evals for LLMs",
    "Recomendado: Liu et al., AgentBench (arXiv:2308.03688)",
    "Recomendado: Zhou et al., WebArena (arXiv:2307.13854)",
    "",
    "Docs: LangSmith, Phoenix, Langfuse, OpenLLMetry, OpenTelemetry GenAI",
], "📖 Obrigatórios: 3 papers canônicos. Recomendados: Hamel Husain (conceitual), AgentBench, WebArena. Docs das ferramentas.\n➡️ Q&A.", 77, T, "Obrigatórios e recomendados")

s = title_slide(
    "Perguntas?",
    "ETHAGT12 — AgentOps, Observabilidade & Avaliação\nPróxima aula: ETHAGT13 — Segurança & Governança de Agentes",
    "Q&A"
)
add_notes(s, "📖 Abrir para Q&A. Se não houver perguntas: 'Qual parte foi menos clara?' Lembrar prazo dos labs e projeto.\n💡 MENSAGEM FINAL: 'Vocês saem com o toolkit completo para operar agentes com confiança. Observabilidade, avaliação, melhoria contínua. Usem. A diferença entre um agente que parece funcionar e um que funciona é rigor experimental.'")

# ─── Salvar ───
out = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out, "ETHAGT12-slides.pptx")
prs.save(out_path)
print(f"PPTX gerado: {out_path}")
print(f"Total de slides: {len(prs.slides)}")
