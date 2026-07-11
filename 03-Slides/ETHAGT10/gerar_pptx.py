#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT10: Padroes de Arquitetura Multi-Agente (topologias)
Universidade Etho · Especializacao em Programacao Agentica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj="", acronyms=""):
    if acronyms:
        add_textbox(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.35), acronyms, size=9, color=INFO)
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT10"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    if acronyms:
        add_textbox(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4), acronyms, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercicio", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

T = 62

s = title_slide(
    "Padroes de Arquitetura Multi-Agente",
    "Universidade Etho · Especializacao em Programacao Agentica\nFase C — Sistemas Multi-Agente · 30 h",
    "ETHAGT10"
)
add_notes(s, "L Bem-vindos. Hoje o foco nao e 'como construir agentes' — e 'como organiza-los'. A pergunta central: dada uma tarefa complexa, qual topologia voce escolhe e por que? Se voce responder 'supervisor' sem justificativa, voce falhou esta aula.\nL Analogia: organizar uma equipe — chefe que delega (supervisor), hierarquia (hierarchical), time que se repassa (swarm), esteira (pipeline).\n? 'Quantos de voces ja usaram multi-agente em producao?'\n! Alunos chegam querendo 'a melhor topologia'. Nao existe melhor — existe a mais adequada.\n> Vamos aos objetivos.")

s = content_slide("Objetivos do Modulo", [
    "Objetivo geral: dominar as topologias multi-agente — quando usar, trade-offs, e como escolher",
    "",
    "Objetivos especificos:",
    "1. Caracterizar 12 topologias (when-to-use, when-to-avoid)",
    "2. Justificar a escolha via ADR",
    "3. Implementar ao menos 4 topologias em um mesmo dominio",
    "4. Medir trade-offs (custo, latencia, qualidade)",
    "5. Identificar sinais de que a topologia precisa evoluir"
], "L Cada objetivo e mensuravel. O foco de hoje e decisao arquitetural justificada, nao 'vamos de supervisor'.\nL Analogia: arquiteto de software escolhendo entre monolito, microsservicos e serverless.\n? 'Qual desses objetivos voces acham mais desafiador?' (costuma ser #4 ou #5)\n! Conhecer = saber when-to-use e when-to-avoid, nao decorar.\n> Vamos as competencias.", 2, T, obj="Objetivo: caracterizar, justificar, implementar, medir")

s = content_slide("Competencias Desenvolvidas", [
    "C1 Programacao Agentica — A (Avancado)",
    "C2 Multi-Agent Systems — A (Avancado)",
    "C3 MCP & Tool Use — B (Basico)",
    "C4 Agent Memory — B (Basico)",
    "C6 Agent Security — B (Basico)",
    "",
    "C1 e C2 atingem nivel A: voce decide arquitetura multi-agente com justificativa"
], "L Este modulo atinge nivel Avancado em duas competencias: C1 e C2.\nL Analogia: aprender a reger uma orquestra. Voce ja toca um instrumento (single agent). Hoje aprende a conduzir.\n! 'Avancado' nao significa 'dominio completo' — significa decidir arquitetura com justificativa.\n> Vamos ver a agenda.", 3, T, obj="Competencias: C1 (A), C2 (A), C3-C6 (B)")

s = content_slide("Agenda da Aula", [
    "Bloco 1 (45 min):",
    "  Abertura (8 min) — objetivos, motivacao, contexto",
    "  Catalogo das 12 Topologias (12 min)",
    "  Supervisor e Hierarchical (15 min) — DEMO 1",
    "  Swarm e Handoffs (10 min)",
    "  Intervalo (5 min)",
    "Bloco 2 (45 min):",
    "  Pipeline e Orchestrator-Workers (10 min)",
    "  Event-Driven, Actor Model, Mesh (10 min)",
    "  Tree, Recursive + DEMO 2 (10 min)",
    "  Fechamento (15 min) — matriz, ADR, quiz, Q&A"
], "L Aula em 2 blocos. Bloco 1: panorama + topologias mais comuns. Bloco 2: topologias avancadas + decisao arquitetural.\nL Analogia: tour por uma cidade — mapa geral, bairros principais, e no final aprender a escolher onde morar.\n! Alunos podem pular o catalogo? Nao — sem conhecer as opcoes, nao ha decisao.\n> Vamos ao por que.", 4, T)

s = content_slide("O Problema da Topologia Padrao", [
    "'Vamos de supervisor' — a decisao default sem justificativa",
    "",
    "Cenario: revisao de PR com 3 especialistas (code, security, docs)",
    "  Supervisor funciona, mas swarm talvez seja mais simples e barato",
    "",
    "Custo de escolher errado:",
    "  - Gargalo de latencia (supervisor serializa)",
    "  - Custo excessivo (LLM call a cada roteamento)",
    "  - Complexidade desnecessaria (hierarchical onde flat basta)",
    "",
    "? Qual foi a topologia mais comum que voces ja usaram em projetos?"
], "L O problema mais comum em producao nao e tecnico — e de decisao. Times escolhem 'supervisor' por default de tutorial.\nL Analogia: hospital onde todo paciente passa por clinico geral. Se ja e cardiologico, por que o clinico? Swarm e triagem direta.\n? 'Qual topologia mais comum que voces usaram?' (maioria diz supervisor ou nao sabe)\n! Supervisor e o mais comum, raramente o mais correto sem analise.\n> Multi-agente se tornou viavel recentemente. Vamos ver por que.", 5, T, acronyms="LLM = Large Language Model — Modelo de Linguagem de Grande Escala  ·  PR = Pull Request — requisicao de pull (GitHub)")

s = content_slide("Do Single Agent ao Multi-Agente", [
    "Linha do tempo:",
    "  2023: ReAct -> AutoGen (arXiv:2308.08155)",
    "  2023: MetaGPT (arXiv:2308.00352) — SOPs para multi-agente",
    "  2023: AgentVerse (arXiv:2308.10848) — assembling agents",
    "  2024: OpenAI Swarm, CrewAI, LangGraph multi-agent",
    "  2026: arXiv:2601.12560 — survey de taxonomia",
    "",
    "Confluencia: reasoning + tool calling + frameworks + custo reduzido",
    "MetaGPT como marco: SOPs aplicadas a agentes",
    "A pergunta deixou de ser 'se' multi-agente, mas 'qual topologia'"
], "L Multi-agente existe desde os anos 90 (Weiss). Tres coisas mudaram: LLMs bons em role-playing, frameworks viaveis, custo baixo.\nL MetaGPT (2023) e o marco: SOPs de software house codificadas em agentes.\n? 'Qual destes marcos voces conheciam?' (maioria conhece Swarm/CrewAI)\n! Multi-agente nao e 'tecnologia nova'. O conceito e dos anos 90; o que e novo e a viabilidade com LLMs.\n> Vamos ao panorama das 12 topologias.", 6, T, acronyms="ReAct = Reasoning and Acting — padrao de loop Thought / Action / Observation")

s = section_slide(1, "Catalogo das 12 Topologias")
add_notes(s, "L Inicio do catalogo. Vamos ver 12 topologias em grid, posiciona-las em um eixo, e fazer matching.\n> Primeiro: quais sao as 12?")

s = content_slide("As 12 Topologias", [
    "1. Single Agent · 1 agente + N tools (baseline)",
    "2. Supervisor · roteador central + workers",
    "3. Hierarchical · arvore de supervisores",
    "4. Blackboard · espaco compartilhado incremental",
    "5. Actor Model · atores com mailbox, async",
    "6. Pipeline · sequencia fixa de agentes",
    "7. Event-Driven · agentes reagem a eventos",
    "8. Swarm · handoffs sem central",
    "9. Tree of Agents · arvore de exploracao (LATS)",
    "10. Recursive · meta-agentes criam sub-agentes",
    "11. Agent Mesh · P2P flat",
    "12. Hybrid · composicao das anteriores",
    "",
    "Fonte: 10-Architecture/architectures/catalog.md + arXiv:2601.12560"
], "L 12 topologias agrupadas: centralizadas, fluxo, descentralizadas, estruturadas, hibridas.\nL Analogia: catalogo de plantas de construtora — voce escolhe a que cabe no terreno e orcamento.\n? 'Quantas destas voces ja viram em producao?' (maioria viu supervisor e pipeline)\n! Nao decorar 12 — lembrar dos 5 grupos.\n> Para comparar, precisamos de um eixo.", 8, T, obj="Catalogo: 12 topologias em 5 grupos")

s = content_slide("O Espectro: Centralizado <-> Descentralizado", [
    "CENTRALIZADO: Supervisor, Hierarchical, Pipeline, Orchestrator-Workers",
    "  Controle forte, ordem previsivel, gargalo central",
    "",
    "INTERMEDIARIO: Blackboard, Event-Driven",
    "  Desacoplamento parcial, coordenacao via barramento",
    "",
    "DESCENTRALIZADO: Swarm, Mesh, Actor Model",
    "  Sem autoridade central, flexibilidade maxima",
    "",
    "ESTRUTURADO (ortogonal): Tree, Recursive",
    "  Estrutura recursiva",
    "",
    "Trade-off fundamental: controle x flexibilidade"
], "L Eixo mais util da aula. Esquerda = controle; direita = flexibilidade.\nL Analogia: empresa vertical (CEO no topo) vs horizontal (squads autonomos).\n? 'Em qual ponto do eixo fica o sistema de voces hoje?'\n! Descentralizado nao e 'sempre melhor'. Troca controle por flexibilidade.\n> Vamos detalhar por grupo.", 9, T, obj="Eixo: centralizado <-> descentralizado")

s = comparison_slide("Topologias Centralizadas",
    "Supervisor", [
        "1 roteador (LLM) + N workers (tools)",
        "Supervisor decide: qual worker, ordem, sintese",
        "LangGraph: cada worker = tool",
        "When-to-use: tarefas decomponiveis, <=10 workers",
        "When-to-avoid: baixa latencia, muitos workers"
    ],
    "Hierarchical", [
        "Arvore de supervisores -> workers -> sub-workers",
        "Top supervisor -> sub-supervisores -> workers",
        "LangGraph: hierarchical_agent_teams",
        "When-to-use: muitos workers, sub-dominios",
        "When-to-avoid: sistemas pequenos"
    ],
    "L Supervisor = tech lead com devs. Hierarchical = VP -> diretores -> gerentes -> devs.\n? 'Ja usou supervisor em producao? Funcionou?' (anotar experiencias de gargalo)\n! Hierarchical nao e 'sempre melhor'. Adiciona latencia e complexidade.\n> Agora as de fluxo.", 10, T)

s = content_slide("Topologias de Fluxo: Pipeline e Orchestrator-Workers", [
    "Pipeline: sequencia fixa de agentes (A -> B -> C)",
    "  Cada step e um agente especializado; ordem predefinida",
    "Pipeline dinamico: proximo step decidido em runtime",
    "",
    "Orchestrator-Workers: orchestrator particiona UMA tarefa em sub-tarefas",
    "  Diferenca do supervisor: orchestrator divide; supervisor roteia",
    "",
    "When-to-use: passos predefinidos, previsibilidade, ordem importa",
    "When-to-avoid: flexibilidade alta necessaria",
    "Aprofundamento: ETHAGT03 (5 workflows canonicos)"
], "L Pipeline = ordem fixa no codigo. Orchestrator = divide tarefa e sintetiza.\nL Analogia: Pipeline = linha de montagem. Orchestrator = arquiteto que divide projeto em modulos.\n! Pipeline com if para proximo step = pipeline dinamico, nao supervisor.\n> Agora as descentralizadas.", 11, T)

s = content_slide("Topologias Descentralizadas: Swarm, Mesh, Actor Model", [
    "Swarm: agentes leves com transfer() de controle (OpenAI)",
    "  A recebe tarefa -> se nao for dele, transfere para B",
    "  Sem orquestrador — controle passa de agente para agente",
    "",
    "Agent Mesh: topologia flat peer-to-peer",
    "  Cada agente fala com qualquer outro diretamente",
    "",
    "Actor Model: cada agente = ator com mailbox, async",
    "  Estado privado, sem estado compartilhado",
    "  Frameworks: Akka (JVM), Ray (Python), AutoGen",
    "",
    "When-to-use: flexibilidade, escalabilidade",
    "When-to-avoid: consistencia forte, coordenacao complexa"
], "L Swarm transfere controle (1 ativo por vez). Mesh = todos com todos. Actor = mailbox async.\nL Analogia: Swarm = jogo de queimada (bola passa). Mesh = grupo de WhatsApp. Actor = correio.\n! Swarm nao e 'mesh leve'. Swarm transfere; mesh permite multiplos ativos.\n> Agora as estruturadas.", 12, T)

s = content_slide("Topologias Estruturadas: Tree e Recursive", [
    "Tree of Agents: arvore de exploracao (LATS)",
    "  Agente expande possibilidades -> avalia -> seleciona melhor",
    "  Como MCTS aplicado a agentes LLM",
    "",
    "Recursive: meta-agentes que instanciam sub-agentes",
    "  'Preciso de especialista em X' -> cria agente X -> delega",
    "  Sub-agentes podem criar sub-sub-agentes",
    "",
    "When-to-use: exploracao de solucoes, decomposicao recursiva",
    "When-to-avoid: custo explode, profundidade incontrolavel",
    "Preview: ETHAGT15 (meta-agentes em profundidade)"
], "L Tree = LATS, arvore de exploracao. Recursive = meta-programacao, agentes criando agentes.\nL Analogia: Tree = jogador de xadrez analisando movimentos. Recursive = gerente que contrata consultores.\n! Recursive nao e 'sempre adaptativo e bom'. Sem max_depth e anti-pattern.\n> Agora as reativas e hibridas.", 13, T, acronyms="MCTS = Monte Carlo Tree Search — Busca em Arvore de Monte Carlo")

s = content_slide("Reativas e Hibridas: Event-Driven, Blackboard, Hybrid", [
    "Event-Driven: agentes reagem a eventos via mensageria (Kafka, NATS)",
    "  Preview: ETHAGT11 (Event-Driven Agents em profundidade)",
    "",
    "Blackboard: espaco compartilhado onde agentes leem/escrevem",
    "  Origem: Hearsay-II (1971-1976), IA classica",
    "  Precursor do event-driven",
    "",
    "Hybrid: composicao das topologias anteriores",
    "  A regra em producao real, nao a excecao",
    "",
    "When-to-use: assincronia, desacoplamento, mundo real (hybrid)",
    "When-to-avoid: simplicidade (over-engineering em prototipos)"
], "L Event-driven = pub/sub com broker. Blackboard = quadro compartilhado classico. Hybrid = realidade.\nL Analogia: Event-driven = rede social. Blackboard = mural de equipe. Hybrid = empresa real.\n! Nao comece prototipo com hybrid. Comece simples (supervisor ou pipeline).\n> Antes de aprofundar, lembremos do baseline.", 14, T)

s = content_slide("Single Agent — O Baseline Esquecido", [
    "Single Agent = 1 agente com N tools (nao e multi-agente, mas e o baseline)",
    "",
    "Quando NAO ir para multi-agente:",
    "  - Tarefa nao se decompoe naturalmente",
    "  - Latencia nao e problema",
    "  - Custo precisa ser minimo",
    "  - Um agente com boas tools resolve",
    "",
    "Regra: comece com single agent, so suba para multi-agente com evidencia",
    "Conexao: ETHAGT01 (escalada de complexidade)"
], "L Antes de qualquer topologia, pergunte: single agent resolve? Um agente ReAct com boas tools faz muita coisa.\nL Analogia: se um funcionario bom resolve, nao contrate 3 medianos.\n! Multi-agente por 'moda' e overhead. Se single agent resolve, single agent e a topologia certa.\n> Vamos praticar o matching.", 15, T)

s = exercise_slide("Exercicio: Matching Cenario -> Topologia", [
    "Para cada cenario, vote na topologia:",
    "",
    "1. 'Chatbot de suporte com escalonamento' -> ?",
    "2. 'Revisao de PR com 3 especialistas' -> ?",
    "3. 'Relatorio financeiro em etapas' -> ?",
    "4. 'Simulacao de mercado' -> ?",
    "5. 'Assistente pessoal multi-funcao' -> ?",
    "6. 'Pesquisa autonomo com exploracao' -> ?",
    "",
    "Votacao rapida (maos levantadas)"
], "L Respostas: 1=Supervisor/Swarm, 2=Supervisor/Swarm, 3=Pipeline, 4=Actor/Mesh, 5=Swarm, 6=Tree.\n! Alunos votam 'supervisor' para tudo. Supervisor e default, nao resposta universal.\n> Vamos aprofundar nas duas mais comuns.", 16, T)

s = section_slide(2, "Supervisor e Hierarchical: A Topologia Default")
add_notes(s, "L Entramos na topologia mais comum em producao. Supervisor pattern, hierarchical, quando escalar, casos de falha, DEMO e MetaGPT.\n> Primeiro: o que e o supervisor pattern?")

s = content_slide("Supervisor Pattern: O Roteador", [
    "Supervisor = LLM que roteia tarefas para workers especializados",
    "Workers sao especializados: pesquisa, escrita, analise, etc.",
    "",
    "Supervisor decide:",
    "  - Qual worker chamar",
    "  - Em que ordem",
    "  - Quando sintetizar",
    "",
    "Analogia: supervisor = tech lead, workers = desenvolvedores",
    "Fonte: LangGraph multi-agent-collaboration",
    "",
    "Diagrama: 12-Diagrams/ETHAGT10/supervisor-topology.mmd"
], "L Supervisor = agente ReAct onde as 'tools' sao outros agentes. Recebe tarefa, raciocina sobre qual worker, chama, recebe, decide proximo.\nL Analogia: tech lead que nao escreve codigo — ele orquestra.\n! Supervisor 'pensa diferente' dos workers — raciocina sobre roteamento, nao execucao.\n> Como implementar em LangGraph? Via tool calls.", 18, T, obj="Diagrama: supervisor-topology.mmd")

s = code_slide("Supervisor como Tool Calls (LangGraph)", """from langgraph.prebuilt import create_react_agent

researcher = create_react_agent(model, tools=[search])
writer = create_react_agent(model, tools=[write])

supervisor = create_react_agent(
    model,
    tools=[researcher, writer, analyst]
)

# Supervisor e um ReAct onde as "tools" sao agentes
# Fluxo: supervisor -> tool_call -> worker -> return -> supervisor""",
    "L Em LangGraph, cada worker e uma tool. Supervisor = ReAct com bind_tools.\nL Elegante: reutiliza todo o mecanismo de tool calling.\n? 'O que impede o supervisor de chamar a mesma tool em loop?' (nada, sem prompt e observabilidade)\n! Nao implemente supervisor com if/else. Volta para pipeline. Supervisor = LLM decide.\n> Quando supervisor nao basta? Hierarchical.", 19, T)

s = content_slide("Hierarchical — Arvore de Supervisores", [
    "Hierarchical = supervisor de supervisores",
    "",
    "Estrutura: top supervisor -> sub-supervisores -> workers -> (sub-workers)",
    "Cada nivel delega para o nivel abaixo",
    "",
    "Caso: revisao de PR -> top supervisor ->",
    "  (code reviewer sup, security sup, docs sup) -> workers",
    "",
    "Fonte: LangGraph hierarchical_agent_teams",
    "",
    "Diagrama: 12-Diagrams/ETHAGT10/hierarchical-topology.mmd"
], "L Hierarchical surge quando supervisor tem workers demais. Crie sub-dominios.\nL Analogia: empresa com departamentos. CEO -> diretores -> gerentes -> funcionarios.\n! 3 niveis raramente necessario. Maioria resolve com flat ou 2 niveis.\n> Quando vale escalar?", 20, T, obj="Diagrama: hierarchical-topology.mmd")

s = content_slide("Quando Escalar Hierarquia", [
    "Flat (1 nivel): supervisor + workers — simples, baixa latencia",
    "2 niveis: top supervisor + sub-supervisores + workers — dominios distintos",
    "3 niveis: raramente necessario — custo e latencia explodem",
    "",
    "Criterio: escalar quando workers de um supervisor > 5-7 (limite cognitivo)",
    "Regra: prefira flat; so adicione nivel com evidencia",
    "",
    "Cada nivel adiciona: latencia (hop de LLM), custo (LLM calls), complexidade (debug)"
], "L Criterio: limite cognitivo. >7 workers em um supervisor = roteamento degrada.\nL Analogia: gerente com 3-5 subordinados funciona. Com 10+, vira gargalo.\n! Hierarquia desnecessaria e over-engineering. Comece flat.\n> Supervisor tem modos de falha. Vamos ser honestos.", 21, T)

s = content_slide("Casos de Falha — Supervisor Gargalo", [
    "Supervisor vira bottleneck: todas as requisicoes passam por ele",
    "Latencia cumulativa: supervisor serializa todas as chamadas",
    "Custo: supervisor faz LLM call a cada roteamento",
    "Single point of failure: se supervisor alucina, sistema falha",
    "",
    "Sinal de alerta: latencia > 30s para tarefas simples",
    "",
    "Solucoes: paralelizar workers, migrar para swarm, hierarchical"
], "L Supervisor e o ponto mais critico. Tudo passa por ele. Se lento, tudo lento.\nL Analogia: portaria de predio onde TODOS passam. Em pico, fila cresce.\n! 'Supervisor mais inteligente' nao resolve gargalo. Gargalo e estrutural.\n> Outro anti-pattern: workers redundantes.", 22, T)

s = content_slide("Casos de Falha — Workers Redundantes", [
    "Workers com responsabilidades sobrepostas -> supervisor nao sabe qual escolher",
    "Resultado: roteamento errado, retrabalho, custo duplicado",
    "",
    "Solucao: responsabilidade exclusiva por worker (SRP)",
    "Sinal: supervisor chama 2+ workers para a mesma sub-tarefa"
], "L Workers com overlap = supervisor confuso. Resultado: roteamento errado ou retrabalho.\nL Analogia: dois funcionarios com mesma descricao de cargo — ninguem sabe quem faz.\n! Workers genericos sempre se sobrepoem. Seja especifico: 'pesquisador de web', 'analisador de codigo'.\n> Vamos ver tudo em acao na DEMO.", 23, T)

s = code_slide("DEMO — Hierarchical Teams (Lab 1)", """Topologia: supervisor -> (researcher, writer) -> writer -> formatter

Tarefa: "Escreva um resumo sobre topico X"

Fluxo:
1. Supervisor delega para researcher
2. Researcher pesquisa -> retorna
3. Supervisor delega para writer
4. Writer escreve -> delega formatacao para formatter
5. Formatter formata -> retorna para writer
6. Writer retorna para supervisor
7. Supervisor sintetiza resposta final

Codigo: 05-Labs/ETHAGT10/Lab1-Hierarchical-Teams""",
    "L DEMO 1 ao vivo. Observem no trace: cada nivel aparece como span.\n? 'Notem quantos LLM calls acontecem. Cada nivel = 1+ call. Isso e custo.'\n! Se API falhar: mostrar screenshot do trace salvo.\n> Vamos discutir.", 24, T, obj="DEMO: Hierarchical Teams")

s = exercise_slide("Pergunta da DEMO", [
    "Discutam em duplas (2 min):",
    "",
    "1. Em quantos niveis de hierarquia o supervisor se torna gargalo?",
    "2. O que acontece se o researcher retornar informacao errada?",
    "   Quem detecta?",
    "3. O formatter deveria ser sub-worker do writer",
    "   ou worker do supervisor?"
], "L Respostas: 1) Gargalo em 3+ niveis. 2) Ninguem detecta por padrao — precisa validacao explicita. 3) Depende do contexto.\n> Conexao academica: MetaGPT.", 25, T)

s = content_slide("MetaGPT — SOPs como Hierarquia", [
    "MetaGPT: framework multi-agente baseado em SOPs",
    "  (Standard Operating Procedures)",
    "",
    "Papeis: Product Manager -> Architect -> Engineer -> QA",
    "Cada papel = um agente com responsabilidade especifica",
    "",
    "SOPs codificam o fluxo: PRD -> design doc -> codigo -> testes",
    "Licao: a hierarquia reflete a organizacao humana de uma software house",
    "",
    "Fonte: arXiv:2308.00352"
], "L MetaGPT (Hong et al., 2023) e o paper canonico de hierarchical multi-agent. SOPs de software house codificadas em agentes.\nL Analogia: simular uma empresa de software inteira. Cada agente = funcionario.\n! MetaGPT nao e 'apenas pipeline'. E pipeline + hierarchical.\n> Conceito relacionado: crew formation.", 26, T, obj="Paper: arXiv:2308.00352")

s = content_slide("Crew Formation — Montando uma Equipe", [
    "Crew = conjunto de agentes com papeis, objetivos e backstory",
    "",
    "CrewAI: Agent(role, goal, backstory) + Task(description, agent)",
    "         + Crew(agents, tasks, process)",
    "Process: sequential (pipeline) ou hierarchical (supervisor)",
    "",
    "AgentVerse (arXiv:2308.10848): assembling agents como montar equipe",
    "",
    "? Como voce montaria uma crew para revisar um PR?"
], "L Crew formation (CrewAI/AgentVerse): defina papeis, tarefas e processo. Abstracao de orquestracao.\nL Analogia: montar um time de projeto — papeis, tarefas, processo.\n? 'Como montaria uma crew para revisar um PR?' (code reviewer, security, docs; processo hierarchical)\n! Backstory importa — da contexto ao agente.\n> Vamos praticar a decisao de profundidade.", 27, T)

s = exercise_slide("Exercicio — Hierarchical ou Flat?", [
    "Para cada cenario, vote: Flat, 2 niveis, ou 3 niveis?",
    "",
    "1. 'Sistema com 3 dominios distintos (code, security, docs)'",
    "2. 'Chatbot simples com 3 ferramentas'",
    "3. 'Software house simulado com 6 papeis'",
    "4. 'Assistente de viagem com 4 sub-tarefas'",
    "",
    "Votacao rapida"
], "L Respostas: 1) 2 niveis, 2) Flat, 3) 2-3 niveis (MetaGPT), 4) Flat.\n! Alunos votam hierarchical 'para ser seguro'. Hierarchical desnecessario adiciona latencia e custo.\n> Vimos a topologia default. Agora a alternativa descentralizada.", 28, T)

s = section_slide(3, "Swarm e Handoffs: Controle Descentralizado")
add_notes(s, "L Swarm e a topologia descentralizada mais popular. Handoffs, quando e melhor que supervisor, limites.\n> Primeiro: o que e swarm?")

s = content_slide("OpenAI Swarm — Transfer de Controle", [
    "OpenAI Swarm (2024): agentes leves com transfer() de controle",
    "",
    "Fluxo: Agente A recebe tarefa -> se nao for dele, transfer(to=agent_b)",
    "Nao ha orquestrador central — controle passa de agente para agente",
    "",
    "Analogia: triagem em hospital — paciente repassado ate o especialista",
    "Estado e passado no handoff (context transfer)",
    "",
    "Diagrama: 12-Diagrams/ETHAGT10/swarm-topology.mmd"
], "L Swarm (OpenAI, 2024): handoffs sem orquestrador. Agente A transfere para B com contexto.\nL Analogia: jogo de 'quem sabe, responde'. A pergunta viaja ate encontrar quem sabe.\n! Swarm nao e mesh. Swarm transfere controle (1 ativo). Mesh permite multiplos ativos.\n> Quando swarm e melhor?", 30, T, obj="Diagrama: swarm-topology.mmd")

s = comparison_slide("Quando Transfer e Melhor que Roteamento Central",
    "Swarm e melhor quando", [
        "Nao ha necessidade de sintese central",
        "Agentes especializados em dominios distintos",
        "Latencia importa (sem supervisor como hop)",
        "Tarefa e 'encontrar o especialista certo'"
    ],
    "Supervisor e melhor quando", [
        "Sintese de multiplos workers necessaria",
        "Controle de ordem importa",
        "Estado global precisa ser mantido"
    ],
    "L Criterio chave: precisa de sintese? Sim -> supervisor. Nao -> swarm.\nL Analogia: Swarm = balcao de atendimento (encaminhado de guiche). Supervisor = gerente que coordena 3 guiches.\n! Swarm nao e 'melhor porque e moderno'. E mais leve. Se precisa sintese, nao serve.\n> Mas swarm tem limites.", 31, T)

s = content_slide("Limites do Swarm", [
    "Coordenacao complexa: swarm nao orquestra, so transfere",
    "Estado compartilhado: cada handoff passa contexto, pode crescer",
    "Loops: A transfere para B que transfere de volta para A",
    "Sem visao global: nenhum agente 'sabe' o estado do sistema",
    "Debugging dificil: tracar o caminho do handoff e nao-trivial",
    "",
    "Em producao: adicione max_handoffs e deteccao de ciclo"
], "L 5 problemas: coordenacao, estado, loops, sem visao global, debugging.\nL Analogia: linha de producao sem supervisor. Se algo da errado, ninguem percebe.\n! Sem max_handoffs, loops infinitos gastam orcamento. Adicione max_handoffs=10.\n> Vamos sistematizar swarm vs supervisor.", 32, T)

s = content_slide("Swarm vs Supervisor — Comparacao Estrutural", [
    "Eixo               | Swarm          | Supervisor",
    "-------------------|----------------|------------------",
    "Controle central   | Nenhum         | Total",
    "Sintese            | Nao nativa     | Nativa",
    "Latencia           | Baixa          | Media (1 hop)",
    "Custo              | Baixo          | Medio",
    "Complexidade       | Baixa          | Media",
    "Escalabilidade     | Alta           | Limitada",
    "Debugging          | Dificil        | Mais facil",
    "Estado global      | Nao            | Sim",
    "",
    "Nao ha vencedor — ha trade-offs"
], "L Tabela mais util da aula. Swarm ganha em latencia/custo/escalabilidade. Supervisor em sintese/controle/debugging.\nL Analogia: bicicleta (swarm) vs carro (supervisor). Qual e melhor? Depende do trajeto.\n! Nao existe 'melhor'. Existe 'melhor PARA ESTE CONTEXTO'.\n> Vamos discutir.", 33, T)

s = exercise_slide("Pergunta — Swarm ou Supervisor?", [
    "Discussao aberta (3 min):",
    "",
    "1. 'Para revisao de PR com 3 especialistas: swarm ou supervisor?'",
    "2. 'E se precisarmos sintetizar as 3 revisoes em um relatorio?'",
    "3. 'E se cada especialista trabalha de forma independente?'"
], "L A resposta muda com o requisito. Precisa sintese -> supervisor. Independentes -> swarm.\nL NAO existe resposta universal. A topologia depende dos requisitos.\n> Fim do Bloco 1. Intervalo de 5 min.", 34, T)

s = section_slide(4, "Pipeline e Orchestrator-Workers: Fluxo Controlado")
add_notes(s, "L Pipeline fixo vs dinamico, orchestrator-workers, composicao. Conexao com ETHAGT03.\n> Primeiro: pipeline fixo vs dinamico.")

s = content_slide("Pipeline Fixo vs Dinamico", [
    "Pipeline fixo: A -> B -> C (sequencia predefinida, codigo orquestra)",
    "  Cada step e um agente especializado; ordem hardcoded",
    "",
    "Pipeline dinamico: A decide se vai para B ou C (runtime decision)",
    "  Proximo step decidido em runtime",
    "",
    "Fixo = workflow (ETHAGT03); dinamico = agente orquestrando agentes",
    "",
    "Exemplo multi-agente: research -> draft -> review -> publish"
], "L Pipeline fixo = ordem no codigo. Dinamico = step decide o proximo.\nL Analogia: linha de montagem (fixo) vs fluxograma com decisao (dinamico).\n! Se a ordem e sempre a mesma, fixo e mais barato. Dinamico so vale quando ordem varia.\n> Orchestrator-workers revisitado.", 36, T)

s = content_slide("Orchestrator-Workers Multi-Agente", [
    "Orchestrator = agente que particiona a tarefa e delega para workers",
    "",
    "Diferenca do supervisor:",
    "  Orchestrator PARTICIONA uma tarefa em sub-tarefas distintas",
    "  Supervisor ROTEIA a tarefa inteira para um worker",
    "",
    "Orchestrator: 'divida e conquiste' -> sintese no final",
    "Fonte: Anthropic, Building Effective Agents (2024)",
    "Aprofundamento: ETHAGT03"
], "L Orchestrator divide a tarefa; supervisor escolhe quem faz a tarefa inteira.\nL Analogia: Supervisor = recepcionista que encaminha. Orchestrator = arquiteto que divide projeto em modulos.\n! Orchestrator PARTICIONA; supervisor ROTEIA. Nao confunda.\n> Topologias se compoem.", 37, T)

s = content_slide("Composicao — Pipeline + Hierarchical", [
    "Pipeline onde um 'step' e um sub-sistema hierarchical",
    "",
    "Exemplo: pipeline (research -> [hierarchical: code, security, docs] -> publish)",
    "  O step de revisao e um supervisor com 3 workers",
    "",
    "Composicao e a regra, nao excecao — sistemas reais sao hybrid",
    "Padrao: pipeline no macro, hierarchical/swarm no micro"
], "L Topologias nao sao mutuamente exclusivas. Sistemas reais sao hybrid.\nL Analogia: empresa. Macro = processo (pipeline). Micro = hierarchical dentro de cada etapa.\n! Nao force uma unica topologia. Componha. Hybrid nao e bagunca — e realismo.\n> Quando pipeline > supervisor?", 38, T)

s = content_slide("Quando Pipeline > Supervisor", [
    "Pipeline e melhor quando:",
    "  - Passos sao conhecidos e fixos (previsibilidade)",
    "  - Ordem importa (A antes de B)",
    "  - Custo precisa ser controlado (sem roteamento dinamico)",
    "  - Nao ha necessidade de decisao em runtime",
    "",
    "Supervisor e melhor quando:",
    "  - Passos sao imprevisiveis",
    "  - Ordem depende do contexto",
    "  - Multiplos workers em ordem variavel"
], "L Pergunta-chave: os passos sao previsiveis? Sim -> pipeline. Nao -> supervisor.\nL Analogia: Pipeline = receita de bolo (sempre mesmos passos). Supervisor = cozinheiro criativo.\n! Pipeline dinamico quando fixo basta = custo imprevisivel.\n> Vamos praticar.", 39, T)

s = exercise_slide("Exercicio — Pipeline ou Agente?", [
    "Para cada cenario, vote Pipeline (P) ou Agente (A) e justifique:",
    "",
    "1. 'Geracao de relatorio financeiro trimestral' -> ?",
    "2. 'Investigacao de incidente de seguranca' -> ?",
    "3. 'Processamento de pedidos de e-commerce' -> ?"
], "L Respostas: 1) Pipeline (etapas fixas), 2) Supervisor (passos imprevisiveis), 3) Pipeline (processamento padronizado).\n! Alunos votam supervisor 'para ter flexibilidade'. Se passos sao fixos, pipeline e mais barato.\n> Agora as assincronas e escalaveis.", 40, T)

s = section_slide(5, "Event-Driven, Actor Model e Mesh")
add_notes(s, "L Topologias que escalam. Event-driven, actor model, mesh. Assincronia.\n> Primeiro: event-driven.")

s = content_slide("Event-Driven Multi-Agente", [
    "Agentes reagem a eventos (nao a chamadas diretas)",
    "Event broker (Kafka, Redis Streams, NATS) como barramento",
    "",
    "Fluxo: Agente A publica evento -> Agente B consome -> reage",
    "Desacoplamento total: A nao sabe quem consome",
    "",
    "Preview: ETHAGT11 (Event-Driven Agents em profundidade)"
], "L Event-driven = pub/sub. Total desacoplamento e resiliencia.\nL Analogia: rede social. Voce posta; seguidores reagem. Feed = broker.\n! Event-driven para tarefas sincronas simples = overkill. Chamada direta e melhor.\n> A fundacao teorica: actor model.", 42, T)

s = content_slide("Actor Model — Fundacao de Escalabilidade", [
    "Cada agente = ator com: estado privado + mailbox + comportamento",
    "Atores se comunicam apenas por mensagens assincronas",
    "Sem estado compartilhado -> sem race conditions",
    "",
    "Escalabilidade natural: milhares de atores em paralelo",
    "Frameworks: Akka (JVM), Ray (Python), Microsoft AutoGen",
    "Fonte: Hewitt, 1973 (origem); AutoGen (arXiv:2308.08155)"
], "L Actor model (1973): estado privado + mailbox + mensagens async. Escala naturalmente.\nL Analogia: funcionarios em cubiculos. Cada um com caixa de entrada. Comunicam por memorandos.\n! Actor model nao e 'academico'. E a base de Erlang (99.9999999% uptime) e WhatsApp.\n> A mais descentralizada: mesh.", 43, T)

s = content_slide("Agent Mesh — P2P sem Central", [
    "Agent Mesh: topologia flat peer-to-peer, sem orquestrador",
    "Cada agente pode falar com qualquer outro diretamente",
    "",
    "Vantagem: sem single point of failure, escalabilidade maxima",
    "Desafio: coordenacao sem central -> protocolos de consenso",
    "",
    "Quando usar: simulacoes, sistemas distribuidos",
    "Quando evitar: controle global, consistencia forte, muitos agentes (N^2)"
], "L Mesh = todos com todos, P2P. Sem ponto de falha. Mas conexoes crescem O(N^2).\nL Analogia: sala de aula onde todos conversam com todos. Funciona com 5-6; caos com 50.\n! Mesh nao e 'a topologia do futuro'. E cara e complexa.\n> Padrao classico: blackboard.", 44, T)

s = content_slide("Blackboard — O Padrao Esquecido", [
    "Blackboard: espaco compartilhado onde agentes leem e escrevem",
    "Origem: Hearsay-II (1971-1976), IA classica",
    "",
    "Fluxo: A escreve hipotese -> B le e refina -> C le e valida",
    "Nao ha orquestrador — agentes monitoram e reagem",
    "",
    "Quando usar: multiplas perspectivas, solucao incremental",
    "Conexao: precursor do event-driven"
], "L Blackboard = quadro compartilhado classico (anos 70). Precursor do event-driven.\nL Analogia: mural de equipe. Cada um passa, le, adiciona. Conhecimento emerge.\n! Blackboard nao e 'obsoleto'. Util para diagnostico incremental, brainstorming.\n> Vamos discutir mesh vs hierarchical.", 45, T)

s = exercise_slide("Pergunta — Mesh vs Hierarchical?", [
    "Discussao aberta (3 min):",
    "",
    "1. 'Mesh e sempre a topologia mais escalavel?' (V/F — justifique)",
    "2. 'Em que cenario mesh e melhor que hierarchical?'",
    "3. 'E se precisarmos de uma decisao global coordenada?'"
], "L 'Mesh e sempre mais escalavel?' FALSO. Conexoes crescem O(N^2). Hierarchical e O(N).\nL Mesh e melhor: poucos agentes altamente colaborativos. Hierarchical: muitos agentes com sub-dominios.\n? 'Se mesh fosse sempre melhor, por que empresas usam hierarquia?'\n> Agora as estruturadas avancadas e a segunda DEMO.", 46, T)

s = section_slide(6, "Tree of Agents, Recursive e DEMO")
add_notes(s, "L Tree of Agents (LATS), recursive (meta-agentes), DEMO 2 swarm vs supervisor.\n> Primeiro: tree of agents.")

s = content_slide("Tree of Agents (LATS)", [
    "LATS (Language Agent Tree Search): arvore de possiveis acoes",
    "Agente expande possibilidades -> avalia -> seleciona melhor caminho",
    "Como MCTS (Monte Carlo Tree Search) aplicado a agentes LLM",
    "",
    "Quando usar: problemas com multiplos caminhos, needing exploration",
    "Custo: exponencial com profundidade da arvore",
    "",
    "Fonte: arXiv:2310.01757; ETHAGT04"
], "L LATS = MCTS aplicado a agentes LLM. Expande, avalia, seleciona. Poderoso mas caro.\nL Analogia: jogador de xadrez pensando 3 movimentos a frente.\n! Tree para problemas simples e desperdicio. Single-path (ReAct) e mais barato.\n> Agora recursive — a mais poderosa e perigosa.", 48, T)

s = content_slide("Recursive — Meta-Agentes", [
    "Recursive: agente que instancia sub-agentes para resolver sub-tarefas",
    "",
    "Meta-agente: 'preciso de especialista em X' -> cria agente X -> delega",
    "Auto-referencia: sub-agentes podem criar sub-sub-agentes",
    "",
    "Preview: ETHAGT15 (meta-agentes e auto-aprendizado)",
    "Quando usar: problemas abertos onde especialistas sao desconhecidos"
], "L Recursive = meta-programacao. Agentes criando agentes. Adaptativo mas perigoso.\nL Analogia: detetive que 'contrata' especialistas conforme descobre necessidades.\n! Recursive precisa de guardrails: max_depth, validacao, orcamento.\n> Por isso recursive pode ser anti-pattern.", 49, T)

s = content_slide("Recursive — Anti-Pattern?", [
    "Recursive e anti-pattern quando:",
    "  - Custo explode (cada nivel = LLM calls adicionais)",
    "  - Profundidade incontrolavel (sem max_depth)",
    "  - Sub-agentes redundantes",
    "  - Latencia inaceitavel",
    "",
    "Recursive NAO e anti-pattern quando:",
    "  - Decomposicao natural do problema",
    "  - max_depth definido",
    "  - Sub-tarefas genuinamente distintas",
    "",
    "? Recursive e sempre anti-pattern? Quando nao?"
], "L Recursive polariza: poderosa bem usada, desastrosa mal usada. Precisa de guardrails.\nL Analogia: empresa que contrata consultores que contratam sub-consultores. Sem limite = torre de intermediarios.\n? 'Recursive e sempre anti-pattern?' Nao. E anti-pattern sem guardrails.\n! SEMPRE defina max_depth (ex.: 3) antes de qualquer outra coisa.\n> Segunda DEMO.", 50, T)

s = code_slide("DEMO — Swarm vs Supervisor (Lab 2)", """Tarefa: "Revisao de PR com 3 especialistas"

Versao 1: Supervisor
  - 3 workers: code, security, docs
  - Supervisor roteia via tool calls
  - Sintetiza as 3 revisoes

Versao 2: Swarm
  - 3 agentes com handoffs
  - Sem orquestrador central
  - Cada agente transfere para o especialista

Medir:
  - Latencia (tempo total)
  - Custo (tokens / LLM calls)
  - Qualidade (subjetiva)

Codigo: 05-Labs/ETHAGT10/Lab2-Swarm-vs-Supervisor""",
    "L DEMO 2 ao vivo. Mesma tarefa, duas topologias. Observem as metricas.\n? 'Supervisor: 1 (rotear) + 3 (workers) + 1 (sintetizar) = 5 LLM calls. Swarm: ~3-4.'\n! Se API falhar: traces comparativos salvos.\n> Vamos discutir os resultados.", 51, T, obj="DEMO: Swarm vs Supervisor")

s = exercise_slide("Pergunta da DEMO", [
    "Discutam em duplas (2 min):",
    "",
    "1. Qual topologia teve menor latencia? Por que?",
    "2. Qual teve menor custo? E melhor qualidade?",
    "3. Em que cenario voce inverteria a escolha?"
], "L Respostas: 1) Swarm (sem hop de supervisor). 2) Swarm menor custo; supervisor melhor qualidade (se sintese boa). 3) Precisa sintese -> supervisor.\n> Ultima secao: a decisao arquitetural.", 52, T)

s = section_slide(7, "Escolha de Topologia e ADR")
add_notes(s, "L Climax da aula: como decidir. Matriz, ADR, sinais de evolucao, exercicio.\n> Primeiro: a matriz de decisao.")

s = content_slide("Matriz de Decisao", [
    "4 eixos: consistencia x latencia x custo x flexibilidade",
    "",
    "Supervisor:    alta consistencia, latencia media, custo medio, flex. baixa",
    "Swarm:         consistencia baixa, latencia baixa, custo baixo, flex. alta",
    "Pipeline:      alta consistencia, latencia alta (serial), custo baixo, flex. baixa",
    "Mesh:          consistencia baixa, latencia baixa, custo medio, flex. maxima",
    "Hierarchical:  alta consistencia, latencia alta, custo alto, flex. media",
    "Hybrid:        balanceado (depende da composicao)",
    "",
    "A escolha e um trade-off: qual eixo e mais importante?",
    "",
    "Diagrama: 12-Diagrams/ETHAGT10/decision-matrix.mmd"
], "L Matriz = ferramenta mais pratica da aula. 4 eixos. Cada topologia avaliada.\nL Analogia: escolher carro. Potencia x economia x espaco x preco. Qual e melhor? Depende.\n! Nao da para otimizar todos os eixos. Escolha qual sacrificar.\n> A decisao precisa ser documentada: ADR.", 54, T, obj="Diagrama: decision-matrix.mmd")

s = content_slide("ADR de Topologia", [
    "ADR (Architecture Decision Record) para topologia multi-agente",
    "",
    "Estrutura:",
    "  1. Contexto: problema, requisitos, restricoes",
    "  2. Decision: topologia escolhida + justificativa",
    "  3. Consequencias: trade-offs aceitos, riscos, sinais de evolucao",
    "",
    "Exemplo: 'Revisao de PR -> Supervisor (sintese necessaria,",
    "  3 especialistas, baixa latencia nao critica)'",
    "",
    "Template: 08-ADRs/ETHAGT10/"
], "L ADR = documentar decisao arquitetural. 3 secoes: Contexto, Decision, Consequencias.\nL Analogia: prontuario medico. Sintomas, diagnostico, efeitos colaterais.\n! ADR durante a decisao, nao depois como justificativa.\n> Como saber quando mudar?", 55, T)

s = content_slide("Sinais de Evolucao", [
    "Supervisor precisa evoluir quando:",
    "  - Latencia crescente (supervisor gargalo)",
    "  - Workers > 7 (limite cognitivo)",
    "  - Necessidade de paralelismo",
    "",
    "Swarm precisa evoluir quando:",
    "  - Necessidade de sintese global",
    "  - Loops de handoff",
    "  - Estado compartilhado crescendo",
    "",
    "Pipeline precisa evoluir quando:",
    "  - Passos imprevisiveis",
    "  - Necessidade de decisao dinamica",
    "",
    "Regra: monitore metricas, nao intuicao"
], "L Nenhuma topologia e definitiva. Sinais sao metricas observaveis.\nL Analogia: painel do carro. Troca quando consumo sobre, manutencao frequente.\n! Trocar topologia 'por moda' = erro. Mude com base em metricas.\n> Caso de estudo: MetaGPT.", 56, T)

s = content_slide("Caso de Estudo — MetaGPT Software House", [
    "MetaGPT simula uma software house completa",
    "",
    "Papeis: Product Manager -> Architect -> Engineer -> QA -> Tester",
    "Topologia: hierarchical com SOPs codificando o fluxo",
    "",
    "Entrada: 'Build a 2048 game'",
    "Saida: codigo + testes + docs",
    "Resultado: codigo funcional em minutos, nao dias",
    "",
    "Licao: a topologia reflete a estrutura organizacional humana",
    "Fonte: arXiv:2308.00352"
], "L MetaGPT aplica tudo: hierarchical + pipeline + SOPs. Entrada NL, saida = codigo funcional.\nL Analogia: simular empresa inteira em software. Cada agente = funcionario. SOPs = processos.\n! MetaGPT nao e 'magica'. Sao SOPs bem definidas + papeis claros + topologia adequada.\n> Vamos praticar a decisao completa.", 57, T, obj="Paper: arXiv:2308.00352")

s = exercise_slide("Exercicio — 6 Cenarios + ADR", [
    "Em grupos (3 min discussao + 2 min compartilhar):",
    "",
    "Para cada cenario, propor topologia + esqueleto de ADR",
    "",
    "1. 'Chatbot de suporte com escalonamento humano'",
    "2. 'Sistema de revisao de codigo com 3 especialistas'",
    "3. 'Pipeline de geracao de relatorio financeiro'",
    "4. 'Simulacao de mercado com multiplos agentes'",
    "5. 'Assistente pessoal multi-funcao'",
    "6. 'Sistema de pesquisa autonomo com exploracao'",
    "",
    "Escrever: Contexto, Decision, Consequencias"
], "L Respostas: 1=Supervisor, 2=Supervisor/Swarm, 3=Pipeline, 4=Actor/Mesh, 5=Swarm, 6=Tree/Recursive.\n! ADR generico nao conta. ADR especifico: 'Supervisor porque sintese de 3 revisoes e requisito, latencia <30s, custo de 5 LLM calls within budget.'\n> Vamos resumir.", 58, T)

s = content_slide("Resumo da Aula + Checklist", [
    "12 topologias catalogadas (5 grupos)",
    "Supervisor = roteador com tool calls; Hierarchical = arvore de supervisores",
    "Swarm = handoffs sem central; Pipeline = sequencia controlada",
    "Actor Model = mensagens async; Mesh = P2P flat",
    "Tree = exploracao; Recursive = meta-agentes (cuidado com custo)",
    "Escolha via matriz + ADR; Sinais de evolucao por metricas",
    "",
    "Checklist:",
    "  [ ] Catalogou as 12 topologias",
    "  [ ] Implementou supervisor + hierarchical",
    "  [ ] Comparou swarm vs supervisor",
    "  [ ] Escreveu ADR de topologia",
    "  [ ] Identificou sinais de evolucao"
], "L Recap: 12 topologias, 5 grupos, matriz de decisao, ADR, sinais de evolucao.\n> Quiz rapido.", 59, T)

s = exercise_slide("Quiz — 3 Perguntas", [
    "P1: Qual topologia NAO tem orquestrador central?",
    "  A) Supervisor  B) Hierarchical  C) Swarm  D) Pipeline",
    "",
    "P2: Em LangGraph, o supervisor roteia para workers via:",
    "  A) Eventos  B) Tool calls  C) Handoffs  D) Mailbox",
    "",
    "P3: Recursive e anti-pattern quando:",
    "  A) Decomposicao natural  B) max_depth nao definido",
    "  C) Sub-tarefas distintas  D) Profundidade e 1"
], "L Respostas: P1=C (Swarm), P2=B (Tool calls), P3=B (sem max_depth).\nL >=2 acertos = compreensao basica.\n> Conexoes com proximos modulos.", 60, T)

s = content_slide("Conexao com Proximos Modulos + Referencias", [
    "Conexoes:",
    "  ETHAGT11 — Event-Driven Agents (event-driven em profundidade)",
    "  ETHAGT14 — Escalabilidade (topologias em producao)",
    "  ETHAGT15 — Meta-Agentes (recursive e auto-aprendizado)",
    "  ETHAGT90 — Capstone (projeto final)",
    "",
    "Leitura obrigatoria:",
    "  Hong et al. MetaGPT (arXiv:2308.00352)",
    "  LangGraph Multi-Agent (hierarchical_agent_teams)",
    "  OpenAI Swarm (repo, 2024)",
    "",
    "Leitura recomendada:",
    "  Chen et al. AgentVerse (arXiv:2308.10848)",
    "  AutoGen (arXiv:2308.08155)",
    "",
    "Proxima aula: ETHAGT11 — Event-Driven Agents"
], "L ETHAGT10 e o hub dos sistemas multi-agente. ETHAGT11 aprofunda event-driven. ETHAGT14 = escalabilidade. ETHAGT15 = meta-agentes.\nL MetaGPT e leitura obrigatoria — inspirou esta aula.\n> Perguntas?", 61, T)

s = title_slide("Perguntas?", "Universidade Etho · ETHAGT10\nProxima aula: ETHAGT11 — Event-Driven Agents\nLembrete: entregar Lab 1 antes da proxima aula", "Q&A")
add_notes(s, "L Se nao houver perguntas: 'Qual topologia voces usariam no projeto de voces hoje, e por que?'\nL Lab 1 (Hierarchical Teams) deve ser entregue antes da proxima aula.\n> Obrigado. Na proxima: event-driven em profundidade.")

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT10-slides.pptx")
prs.save(out)
print(f"PPTX gerado: {out}")
print(f"Total de slides: {len(prs.slides)}")
