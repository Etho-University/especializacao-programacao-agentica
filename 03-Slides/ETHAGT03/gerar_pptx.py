#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT03: Padrões de Workflow Agêntico
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj="", acronyms=""):
    if acronyms:
        add_textbox(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.35), acronyms, size=9, color=INFO)
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT03"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    if acronyms:
        add_textbox(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4), acronyms, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

T = 63

# ═══════════════════════════════════════
# SEÇÃO A — Abertura (1-5)
# ═══════════════════════════════════════

s = title_slide(
    "Padrões de Workflow Agêntico",
    "Universidade Etho · Especialização em Programação Agêntica\nFase A — Fundamentos Agênticos · 30 h",
    "ETHAGT03"
)
add_notes(s, "📖 Bem-vindos à terceira aula. Em ETHAGT01 estabelecemos o bloco fundamental e a distinção workflow vs agente. Hoje aprofundamos os 5 padrões canônicos de workflow.\n💡 Analogia: arquiteto que conhece 5 tipologias de edifício. Cada padrão resolve uma classe de problema.\n❓ 'Quantos já usaram workflow estruturado em produção?'\n⚠️ Alunos acham que workflow é 'menos avançado'. Redirecionar: workflow é escolha por previsibilidade.\n➡️ Vamos aos objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: dominar os 5 padrões canônicos de workflow e suas composições",
    "",
    "Objetivos específicos:",
    "1. Implementar os 5 workflows (chaining, routing, parallelization, orchestrator, evaluator)",
    "2. Identificar gates programáticos e onde inseri-los",
    "3. Combinar workflows em pipelines complexos",
    "4. Justificar workflow vs agente autônomo em cenário real",
    "5. Medir trade-offs de custo/latência/qualidade",
], "📖 Cada objetivo é mensurável: implementar, identificar, combinar, justificar, medir. O #4 (justificar) é o que diferencia sênior.\n❓ 'Qual objetivo vocês acham mais desafiador?' (geralmente #4 ou #5)\n⚠️ Alunos confundem 'conhecer os 5 padrões' com 'saber escolher entre eles'.\n➡️ Competências.", 2, T, "Objetivos mensuráveis")

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → I (Intermediário)",
    "C2 Multi-Agent Systems → B (Básico)",
    "C3 MCP & Tool Use → B (Básico)",
    "C5 AgentOps & Avaliação → B (Básico)",
    "",
    "C1 consolida Intermediário: projetar workflow que melhor se ajusta ao problema",
    "C2 começa: orchestrator-workers é proto-multi-agente",
    "C5 na prática: avaliar classificadores e convergência de loops",
], "📖 C1 consolida Intermediário. C2 começa (orchestrator-workers é proto-multi-agente). C5 na prática (avaliar classificadores e loops).\n💡 Analogia: ETHAGT01 foi sair do estacionamento. Hoje dirige na cidade com trânsito.\n⚠️ 'Básico' não significa 'superficial' — significa fundação.\n➡️ Agenda.", 3, T, "Conectar ao Framework Etho", acronyms="AgentOps = Agent Operations — operacao e monitoramento de agentes em producao  ·  MCP = Model Context Protocol — Protocolo de Contexto de Modelo")

s = content_slide("Agenda da Aula", [
    "Bloco 1 (45 min):",
    "  Abertura (6 min) — objetivos, motivação",
    "  Por Que Workflows (7 min) — princípio, 5 padrões",
    "  Prompt Chaining (8 min) — estrutura, gates, código",
    "  Routing (9 min) — classificação, modelos",
    "  Parallelization (15 min) — sectioning, voting, DEMO",
    "  Intervalo (5 min)",
    "",
    "Bloco 2 (45 min):",
    "  Orchestrator-Workers (10 min) — delegação dinâmica",
    "  Evaluator-Optimizer (10 min) — loop, critérios",
    "  Composições (8 min) — composição, trade-offs",
    "  Fechamento (17 min) — caso, quiz, projeto, Q&A",
], "📖 Dois blocos. Primeiro: fundamentação + 3 padrões. Segundo: 2 padrões complexos + composições + fechamento. DEMO de latência é o ponto alto.\n💡 Analogia: menu degustação. Cada padrão é um prato. A composição é o prato final.\n➡️ Por que estamos aqui?", 4, T, "Roteiro da aula")

s = content_slide("Comece Simples", [
    "Problema: time usa agente autônomo onde um 'if' bastava",
    "",
    "Exemplo: chatbot de suporte com 3 etapas fixas",
    "  (classificar → buscar → responder) — agente livre é overkill",
    "",
    "Princípio de Anthropic: comece simples, só aumente com evidência",
    "",
    "Custo de complexidade prematura:",
    "  • Debugging difícil (caminho não determinístico)",
    "  • Custo imprevisível (nº de chamadas varia)",
    "  • Latência alta (loops extras)",
    "",
    "Pergunta: Quantos já viram projeto que usou agente onde um if bastava?",
], "📖 O erro mais caro da indústria de IA hoje é complexidade prematura. Teams pulam para agente onde workflow de 3 passos bastava.\n💡 Analogia: construir casa. Não começa pela piscina. Começa pela fundação.\n❓ Mãos levantadas: 'Quem já viu agente onde if bastava?'\n⚠️ Workflow não é 'iniciante'. É escolha de engenharia por previsibilidade.\n➡️ Esse princípio é de Anthropic.", 5, T, "Criar tensão — complexidade prematura")

# ═══════════════════════════════════════
# SEÇÃO B — Por Que Workflows (6-10)
# ═══════════════════════════════════════

s = section_slide(1, "Por Que Workflows Antes de Agentes")
add_notes(s, "Início do bloco de fundamentos. Por que workflows primeiro? Qual a hierarquia de complexidade?")

s = content_slide("O Princípio de Anthropic: Comece Simples", [
    "'Comece simples, só aumente complexidade com evidência' — Anthropic",
    "",
    "Níveis de escalada:",
    "  Nível 0: Single LLM call + retrieval + examples  (90% dos casos)",
    "  Nível 1: Workflow simples (prompt chaining)",
    "  Nível 2: Workflow complexo (orchestrator-workers)",
    "  Nível 3: Agente autônomo",
    "  Nível 4: Multi-agente",
    "",
    "Regra: só suba um nível com EVIDÊNCIA de que o anterior é insuficiente",
], "📖 Hierarquia de complexidade. Base (Nível 0) resolve 90%. Só escale com evidência (métricas, testes, custos).\n💡 Analogia: escada de atendimento médico. Auto-cuidado → clínico → especialista → hospital → UTI. Não manda tudo para UTI.\n❓ 'Em que nível está o sistema de vocês?' (maioria em 0-1 achando que está em 3)\n⚠️ Pular para Nível 3 sem tentar Nível 1 é 'resume-driven development'.\n➡️ Recap da distinção central.", 7, T, "Hierarquia de complexidade", acronyms="LLM = Large Language Model — Modelo de Linguagem de Grande Escala")

s = content_slide("Workflows vs Agentes (Recap)", [
    "Workflows: LLMs e tools orquestrados via CÓDIGO PREDEFINIDO",
    "  → previsibilidade",
    "",
    "Agentes: LLM DIRECIONA seu próprio processo e tool usage",
    "  → flexibilidade",
    "",
    "Workflows = controle do DESENVOLVEDOR",
    "Agentes = controle do MODELO",
    "",
    "Este módulo: foco total em workflows",
    "ETHAGT04 aprofunda agentes",
], "📖 Recap da distinção canônica. Workflow = código orquestra LLMs. Agente = LLM orquestra a si mesmo. A diferença é CONTROLE.\n💡 Analogia: trem (workflow — trilhos predefinidos) vs táxi (agente — motorista decide).\n⚠️ Workflow não é 'menos poderoso'. É mais PREVISÍVEL. Em produção, previsibilidade é poder.\n➡️ Os 5 padrões.", 8, T, "Recap de ETHAGT01")

s = content_slide("Os 5 Workflows Canônicos", [
    "Fonte: Anthropic, Building Effective Agents (2024)",
    "",
    "1. Prompt Chaining — sequência de steps com gates",
    "2. Routing — classificação + dispatch",
    "3. Parallelization — sectioning / voting",
    "4. Orchestrator-Workers — delegação dinâmica",
    "5. Evaluator-Optimizer — loop de refinamento",
    "",
    "Cada padrão resolve uma CLASSE de problema",
    "Diagrama: 12-Diagrams/ETHAGT03/ (um .mmd por padrão)",
], "📖 Anthropic catalogou 5 padrões que cobrem a maioria dos casos de produção. Cada um resolve uma classe: linear (chaining), categorias (routing), independência (parallel), dinâmica (orchestrator), refinamento (evaluator).\n💡 Analogia: 5 ferramentas na caixa. Martelo, chave-inglesa, alicate, furadeira, lixadeira. Cada uma para um parafuso.\n❓ 'Qual desses 5 vocês já usaram em produção?' (geralmente routing e chaining)\n⚠️ Não use orchestrator para tudo porque 'parece avançado'.\n➡️ Quando cada um brilha.", 9, T, "5 padrões canônicos")

s = content_slide("Quando Cada Padrão Brilha", [
    "Prompt Chaining: tarefa linear com checkpoints de qualidade",
    "  → evitar: tarefas simples (uma chamada basta)",
    "",
    "Routing: categorias discretas com tratamento especializado",
    "  → evitar: continuum sem fronteiras claras",
    "",
    "Parallelization: subtarefas independentes ou votação",
    "  → evitar: subtarefas com dependências",
    "",
    "Orchestrator-Workers: subtarefas dinâmicas não conhecidas a priori",
    "  → evitar: subtarefas fixas e conhecidas",
    "",
    "Evaluator-Optimizer: feedback articulável e iterável melhora",
    "  → evitar: evaluator não é melhor que generator",
], "📖 O mapa mental da aula. A coluna 'evitar' é tão importante quanto 'usar'. Cada padrão tem armadilha.\n💡 Analogia: guia de restaurantes por ocasião. Encontro rápido → lanchonete. Jantar romântico → elegante.\n⚠️ Decorem 'usar' mas NÃO esqueçam 'evitar'.\n➡️ Primeiro padrão: prompt chaining.", 10, T, "Mapa de decisão")

# ═══════════════════════════════════════
# SEÇÃO C — Prompt Chaining (11-16)
# ═══════════════════════════════════════

s = section_slide(2, "Prompt Chaining: Sequência com Gates")
add_notes(s, "Primeiro padrão. Cadeia de chamadas LLM com gates programáticos entre elas.")

s = content_slide("Estrutura: LLM → Gate → LLM → …", [
    "Cadeia de chamadas LLM intercaladas com GATES PROGRAMÁTICOS",
    "",
    "Gate = código determinístico (NÃO LLM) que decide:",
    "  • continuar?  • rodar de novo?  • parar?",
    "",
    "Cada step tem prompt especializado e contexto focado",
    "Saída de um step = entrada do próximo",
    "",
    "Gates capturam erros cedo → menos retrabalho downstream",
    "",
    "Diagrama: 12-Diagrams/ETHAGT03/prompt-chaining.mmd",
], "📖 Estrutura: LLM → gate → LLM → gate → saída. Cada LLM foca em UMA coisa. Gates são código puro (determinísticos). A combinação de passos focados (qualidade) e gates (controle) é o poder.\n💡 Analogia: linha de montagem. Estações (LLM) + inspetores de qualidade (gates).\n❓ 'Por que gates são código e não LLM?' (determinismo)\n⚠️ Gate com LLM é evaluator-optimizer, não prompt chaining.\n➡️ Onde colocar gates.", 12, T, "Estrutura do prompt chaining")

s = content_slide("Onde Adicionar Gates", [
    "4 tipos principais:",
    "",
    "1. Validação estrutural — 'tem todos os campos obrigatórios?'",
    "2. Classificação — 'a resposta é do tipo esperado?'",
    "3. Formatação — 'está no formato JSON correto?'",
    "4. Filtro de qualidade — 'passou no check de segurança?'",
    "",
    "Gates são código puro (if/else, regex, schema) — NÃO LLM",
    "",
    "Gate falhou? → retry com feedback / fallback / escalada",
], "📖 4 tipos de gate. Validação estrutural (schema), classificação (tipo), formatação (JSON), filtro de qualidade (regras de negócio). Ao falhar: retry, fallback ou escalada.\n💡 Analogia: alfândega. Conferência documental, de tipo, de formato, de segurança.\n⚠️ Não coloque gates demais. Cada gate adiciona latência. Só se o custo de NÃO ter é maior.\n➡️ Vale a latência?", 13, T, "Tipos de gate")

s = content_slide("Trade-off: Latência por Accuracy", [
    "Latência = SOMA de todos os steps (serial)",
    "",
    "Mas: cada step focado = melhor qualidade que chamada gigante",
    "Gates reduzem retrabalho downstream (capturam erro cedo)",
    "",
    "Quando vale: tarefas com critérios de qualidade verificáveis",
    "Quando não vale: tarefas simples (uma chamada basta)",
    "",
    "Curva de ganho marginal: primeiros gates dão muito,",
    "os últimos dão pouco",
], "📖 Trade-off central: latência vs accuracy. Serial soma latência. Mas cada step focado = melhor qualidade. A questão: o ganho compensa a latência? Depende do domínio.\n💡 Analogia: revisar texto. Uma passada pega 80% dos erros. Duas pegam 95%. Três pegam 98%. Quarta mal agrega.\n⚠️ Sem medição, você não sabe se o gate ajuda ou só adiciona latência.\n➡️ Código concreto.", 14, T, "Trade-off latência/accuracy")

s = code_slide("Código: Gerar → Revisar → Reescrever", """def chain(ticket):
    # Step 1: gerar rascunho
    draft = llm("Gere rascunho", ticket)

    # Gate 1: validação estrutural
    if not validate_schema(draft,
            fields=["saudacao","corpo","fechamento"]):
        draft = llm("Inclua todos os campos", ticket)

    # Step 2: revisar contra critérios
    review = llm("Revise contra critérios", draft)

    # Gate 2: score >= threshold
    if review["score"] < 0.8:
        # Step 3: reescrever com feedback
        return llm("Reescreva com feedback",
                   draft, review["feedback"])
    return draft""", "📖 Cada LLM tem prompt focado. Gates de código: validate_schema e score < threshold. Gate falha? Retry com feedback específico.\n💡 Analogia: escritor revisando. Escreve → verifica estrutura → editor revisa → se nota baixa, reescreve.\n⚠️ Sem retry no gate, um gate falho vira erro fatal.\n➡️ Vamos praticar gates.", 15, T, "Implementação de prompt chaining")

s = exercise_slide("Exercício — Gate em Tradução", [
    "Cenário: prompt chaining para tradução técnica (EN → PT-BR)",
    "",
    "Pergunta: Descreva um gate programático útil entre o step de",
    "tradução e o step de revisão.",
    "",
    "Dica: pense em verificação ESTRUTURAL, não semântica",
    "(gates são código puro, não LLM)",
    "",
    "Formato: discussão em duplas (2 min)",
], "📖 Deixar turma pensar. Pergunta é sobre verificação ESTRUTURAL. Exemplos: nº de parágrafos, glossário preservado, variáveis {entre_chaves} não traduzidas, contagem de tokens.\n💡 Analogia: inspetor de obra verifica planta (estrutura), não se a casa é bonita (semântica).\n❓ 'Compartilhem 2 respostas.'\n⚠️ 'Verificar se a tradução é fiel' é evaluator, não gate.\n➡️ Segundo padrão: routing.", 16, T)

# ═══════════════════════════════════════
# SEÇÃO D — Routing (17-23)
# ═══════════════════════════════════════

s = section_slide(3, "Routing: Classificação e Dispatch")
add_notes(s, "Segundo padrão. Separação de concerns: classificar é mais barato que resolver.")

s = content_slide("Classificação como Etapa Separada", [
    "Router = LLM (ou código) que classifica a entrada em categorias",
    "",
    "Cada categoria tem seu próprio tratamento especializado",
    "",
    "Separação de concerns: CLASSIFICAR é mais barato que RESOLVER",
    "",
    "Router pode ser modelo menor (Haiku)",
    "  → classify é tarefa simples e barata",
    "",
    "Diagrama: 12-Diagrams/ETHAGT03/routing.mmd",
], "📖 Ideia central: separar classificação de resolução. Classificar (técnico/billing/geral) é mais barato que resolver. Modelo barato classifica, handler especializado resolve. Ganho duplo: custo + qualidade.\n💡 Analogia: recepção de hospital. Recepcionista (router) classifica, especialista resolve.\n⚠️ Um LLM para tudo sem classificar infla custo e reduz qualidade.\n➡️ Routing por modelo.", 18, T, "Princípio do routing")

s = content_slide("Routing por Modelo", [
    "Haiku para tarefas fáceis (classificação, resumo simples)",
    "  → barato e rápido (10x mais barato que Sonnet)",
    "",
    "Sonnet para tarefas difíceis (raciocínio, código complexo)",
    "  → caro mas capaz",
    "",
    "Economia: 80% dos tickets são fáceis → Haiku",
    "  → 80%×Haiku + 20%×Sonnet << 100%×Sonnet",
    "",
    "Risco: classificação errada envia difícil para Haiku → resposta ruim",
    "Mitigação: quando incerto, despachar para Sonnet",
], "📖 Routing por modelo. Router classifica dificuldade. Fácil → Haiku (10x mais barato). Difícil → Sonnet. 80% fáceis × Haiku é muito mais barato que 100% × Sonnet.\n💡 Analogia: hospital. Casos leves → clínico (Haiku). Complexos → especialista (Sonnet).\n⚠️ Confiar cegamente no router é perigoso. Incerteza → path robusto.\n➡️ Routing não é só por modelo.", 19, T, "Routing por modelo")

s = content_slide("Routing por Prompt/Tools Especializados", [
    "Mesmo modelo, diferentes SYSTEM PROMPTS especializados",
    "",
    "Exemplo:",
    "  ticket de billing → prompt + tools de fatura",
    "  ticket técnico → prompt + tools de docs",
    "  ticket geral → prompt genérico",
    "",
    "Vantagem: cada path tem CONTEXTO FOCADO",
    "  → menos tokens, melhor qualidade",
    "",
    "O router decide qual conjunto de tools/prompt carregar",
], "📖 Outra forma de routing: por prompt e tools. Mesmo modelo (Sonnet para todos), mas cada categoria carrega system prompt e tools diferentes. Vantagem: foco. Menos tools no contexto = menos confusão = melhor qualidade.\n💡 Analogia: médico com várias especialidades mas ativa só a relevante.\n⚠️ Carregar TODAS as tools sempre confunde o modelo. Menos tools bem selecionadas > muitas genéricas.\n➡️ Como avaliar o router?", 20, T, "Routing por prompt/tools")

s = content_slide("Avaliando a Qualidade do Classificador", [
    "O router é um CLASSIFICADOR → métricas de classificação",
    "",
    "Matriz de confusão: quais categorias confunde?",
    "Precision: dos classificados como X, quantos eram X?",
    "Recall: dos que eram X, quantos o router pegou?",
    "",
    "Erro mais caro: falso negativo (difícil → fácil → Haiku → ruim)",
    "",
    "Estratégia: quando incerto, enviar para path robusto (Sonnet)",
], "📖 Router é classificador. Tem matriz de confusão, precision, recall. O erro mais caro é falso negativo em routing por modelo: difícil classificado como fácil vai para Haiku e vira resposta ruim. Mitigação: incerteza → path robusto.\n💡 Analogia: triagem médica. Erro tipo I (mandar grave para casa) é pior que tipo II (manter leve no hospital).\n⚠️ Sem matriz de confusão, você não sabe quais categorias confundem.\n➡️ Código.", 21, T, "Avaliação do router")

s = code_slide("Código: Routing com Classifier", """def route(ticket):
    # Router classifica
    cat = router_llm(
        "Classifique: tecnico | billing | geral | unknown",
        ticket, response_format="json"
    )["categoria"]

    # Handlers especializados
    handlers = {
        "tecnico": lambda t: sonnet(
            PROMPT_TECNICO, t, tools=TOOLS_DOCS),
        "billing": lambda t: sonnet(
            PROMPT_BILLING, t, tools=TOOLS_FATURA),
        "geral":   lambda t: haiku(PROMPT_GERAL, t),
        "unknown": lambda t: sonnet(
            PROMPT_ROBUSTO, t, tools=ALL_TOOLS),
    }
    return handlers.get(cat, handlers["unknown"])(ticket)""", "📖 Padrão limpo: router classifica, switch/case despacha, cada handler tem prompt+modelo+tools. Fallback 'unknown' → Sonnet com todas as tools. response_format=json evita parsing frágil.\n💡 Analogia: PABX. Recepcionista virtual direciona para o ramal certo. Se não sabe, recepção humana.\n⚠️ Sem fallback, categoria não prevista quebra o sistema.\n➡️ Vamos discutir avaliação.", 22, T, "Implementação de routing")

s = exercise_slide("Exercício — Avaliando o Router", [
    "Discutam em duplas (2 min):",
    "",
    "1. Como saber se o roteador está errando? O que medir?",
    "",
    "2. Qual métrica é mais importante: precision ou recall?",
    "   Depende do quê?",
    "",
    "3. Se o router envia 10% dos fáceis para o path difícil,",
    "   isso é problema?",
    "",
    "1 min para compartilhar",
], "📖 Q1: logs de classificação vs ground truth → matriz de confusão. Q2: recall mais importante quando difícil→fácil é caro. Q3: 10% fáceis→difíceis é aceitável (custa mais mas resolve). O problema real é o oposto.\n💡 Analogia: triagem de aeroporto. Comum→detalhada só atrasa. Suspeito→rápida é perigoso.\n⚠️ Accuracy global esconde vieses por categoria. Sempre precision/recall por categoria.\n➡️ Intervalo! Voltamos com parallelization.", 23, T)

# ═══════════════════════════════════════
# SEÇÃO E — Parallelization (24-32)
# ═══════════════════════════════════════

s = section_slide(4, "Parallelization: Sectioning e Voting")
add_notes(s, "Terceiro padrão. Duas variantes: sectioning (subtarefas independentes) e voting (mesma tarefa N vezes para robustez).")

s = content_slide("Sectioning: Subtarefas Independentes", [
    "Tarefa dividida em N subtarefas INDEPENDENTES",
    "",
    "Cada subtarefa roda em PARALELO (chamadas simultâneas)",
    "",
    "Resultados agregados por um REDUCER/sintetizador",
    "",
    "Exemplo: revisar código → linter + testes + typecheck em paralelo",
    "",
    "Ganho de latência: max(subtasks) vs sum(subtasks)",
    "  3 tarefas de 3s: serial=9s, paralelo=3s",
    "",
    "Diagrama: 12-Diagrams/ETHAGT03/parallelization.mmd (sectioning)",
], "📖 Sectioning: dividir em partes independentes e rodar em paralelo. Ganho principal é latência: max em vez de sum. Mas NÃO reduz custo — paga todas as chamadas. Reducer agrega. Crítico: INDEPENDÊNCIA real.\n💡 Analogia: buffet. 3 pessoas servem pratos diferentes simultaneamente. Tempo = pessoa mais lenta.\n❓ 'Por que paralelismo não reduz custo?' (paga todas as chamadas)\n⚠️ Subtarefas com dependência oculta = erro. Precisa orchestrator.\n➡️ A outra variante: voting.", 25, T, "Sectioning")

s = content_slide("Voting: Mesma Tarefa N Vezes", [
    "Mesma tarefa executada N vezes (com temperatura/variação)",
    "",
    "Agregação por MAIORIA (classificação) ou SELEÇÃO (geração)",
    "",
    "Quando usar:",
    "  • classificação ambígua",
    "  • geração de código crítico",
    "  • decisões de segurança",
    "",
    "Custo: N× o custo base",
    "  mas 3×Haiku pode ser mais barato que 1×Sonnet",
    "",
    "Diagrama: 12-Diagrams/ETHAGT03/parallelization.mmd (voting)",
], "📖 Voting: mesma tarefa N vezes com variação. Agrega por maioria (classificação) ou seleção (geração). Ganho: robustez (reduz variância). Mas NÃO resolve viés — voting reforça erro sistemático.\n💡 Analogia: 3 médicos independentes. Se concordam, alta confiança. Se discordam, investiga. Mas mesmo viés = maioria reforça erro.\n⚠️ Voting não corrige viés. Corrige variância.\n➡️ Variação híbrida: guardrails.", 26, T, "Voting")

s = content_slide("Guardrails em Paralelo", [
    "Modelo A: gera a resposta principal",
    "Modelo B (paralelo): filtra/modera a resposta",
    "",
    "Se B sinaliza problema → bloqueia ou regenera",
    "",
    "Latência quase ZERO adicionada (roda em paralelo)",
    "Custo: 2× chamadas (mas segurança é crítica em produção)",
    "",
    "Exemplo: resposta de suporte + verificação de PII em paralelo",
    "         + moderação de toxicidade em paralelo",
], "📖 Guardrails em paralelo: modelo principal gera, segundo modelo avalia segurança. Se problema, bloqueia/regenera. Latência zero (paralelo). Custo 2× mas segurança crítica.\n💡 Analogia: segurança em evento. Garçom serve (modelo principal), segurança observa (guardrail). Trabalham em paralelo.\n⚠️ Guardrails em série adiciona latência. Paralelo é o ponto.\n➡️ Outra aplicação: LLM-as-judge.", 27, T, "Guardrails paralelos", acronyms="PII = Personally Identifiable Information — dados pessoais identificaveis")

s = content_slide("LLM-as-Judge em Paralelo", [
    "N julgamentos paralelos de uma saída",
    "",
    "Votação reduz VARIÂNCIA do juiz",
    "",
    "Usado em pipelines de qualidade antes de entregar ao usuário",
    "",
    "CUIDADO: viés do juiz é sistemático",
    "  → votação NÃO resolve viés, só variância",
    "  → para viés, calibre com rubric estruturada",
], "📖 LLM-as-judge: LLM avalia saída de outro LLM. Em paralelo, N juízes + agregação reduz variância. Mas viés é sistemático — voting reforça. Para viés, calibre o juiz (rubric).\n💡 Analogia: painel de jurados. Vários reduzem variância da nota. Mas mesmo viés = média não corrige.\n⚠️ N juízes idênticos (mesmo modelo, mesmo prompt) não reduz variância. Precisa variação.\n➡️ Armadilhas do parallelization.", 28, T, "LLM-as-judge paralelo", acronyms="LLM-as-Judge = uso de um LLM como avaliador automatico de saidas de outro LLM")

s = content_slide("Erros Comuns em Parallelization", [
    "1. Dependência oculta: 'independentes' que não são",
    "   (step 2 precisa de output de step 1)",
    "",
    "2. Custo explodindo: N=5 voting × 10 steps = 50 chamadas",
    "",
    "3. Reducer mal feito: síntese perde informação crítica",
    "",
    "4. Timeout não tratado: uma subtarefa lenta atrasa tudo",
    "",
    "Mitigação: medir custo/latência de cada subtarefa DESDE O INÍCIO",
], "📖 4 armadilhas. Dependência oculta (desenhe DAG antes). Custo explodindo (log de custo desde dia 1). Reducer mal feito (síntese perde insight). Timeout (use timeout + fallback).\n💡 Analogia: coordenar equipe. Dependência ociosa = 'espera seu relatório'. Custo = cada um contratou 5 freelancers. Reducer = gestor resume em 1 linha. Timeout = membro sumiu.\n⚠️ Sem log de custo desde o início, custo explode silenciosamente.\n➡️ Código de ambos os modos.", 29, T, "Armadilhas")

s = code_slide("Código: Parallelization", """import asyncio

# SECTIONING: subtarefas DIFERENTES em paralelo
async def sectioning(doc):
    parts = split(doc)
    results = await asyncio.gather(
        *[analyze(p) for p in parts]
    )
    return reducer(results)

# VOTING: MESMA tarefa N vezes com variação
async def voting(question, n=5):
    answers = await asyncio.gather(*[
        llm(question, temperature=0.7 + i*0.05)
        for i in range(n)
    ])
    return majority_vote(answers)""", "📖 asyncio.gather() para ambas. Sectioning: funções diferentes (analisa partes diferentes). Voting: mesma função com temperatura diferente. Reducer agrega. Única diferença é o que vai no gather.\n💡 Analogia: Sectioning = equipe onde cada um faz parte do relatório. Voting = mesma equipe escreve 5 vezes e escolhe o melhor.\n⚠️ gather sem return_exceptions=True lança exceção se uma chamada falha.\n➡️ DEMO ao vivo.", 30, T, "Implementação parallelization")

s = code_slide("DEMO — Latência Serial vs Paralelo", """# Mesma tarefa: 3 chamadas LLM

# VERSÃO SERIAL (sum)
import time
t0 = time.time()
r1 = await llm("análise A")
r2 = await llm("análise B")
r3 = await llm("análise C")
t_serial = time.time() - t0
# → ~9.0s (3 + 3 + 3)

# VERSÃO PARALELA (max)
t0 = time.time()
results = await asyncio.gather(
    llm("análise A"),
    llm("análise B"),
    llm("análise C"),
)
t_parallel = time.time() - t0
# → ~3.0s (max dos três)

print(f"Serial:   {t_serial:.1f}s")
print(f"Paralelo: {t_parallel:.1f}s")
print(f"Custo: IDÊNTICO nos dois casos")
print(f"Lição: paralelismo NÃO reduz custo, reduz LATÊNCIA")""", "📖 A demo mais importante do módulo. Mesma tarefa, duas formas. Serial: ~9s. Paralelo: ~3s. Mas custo IDÊNTICO — paguei as mesmas 3 chamadas. Paralelismo não reduz custo, reduz latência.\n💡 Analogia: lavar 3 pratos. Um de cada vez = 3 min. Três pessoas = 1 min. Mas pagou 3 pessoas nos dois casos.\n❓ 'Quando paralelismo NÃO vale a pena?' (quando custo é o gargalo, não latência)\n⚠️ Paralelismo NÃO reduz custo. Só reduz latência.\n➡️ Vamos praticar sectioning vs voting.", 31, T, "DEMO ao vivo")

s = exercise_slide("Exercício — Voting vs Sectioning", [
    "Para cada cenário, vote: Sectioning (S) ou Voting (V)?",
    "",
    "1. Traduzir 3 páginas independentes",
    "2. Decidir se um email é phishing",
    "3. Gerar 3 alternativas de copy para anúncio",
    "",
    "Votação rápida (mãos levantadas)",
    "",
    "Gabarito: 1=S, 2=V, 3=S (mas voting para escolher a melhor)",
], "📖 1=Sectioning (páginas diferentes). 2=Voting (mesma pergunta, robustez). 3=Sectioning para gerar (cópias diferentes), Voting para escolher a melhor. Regra: sectioning quando DIFERENTES; voting quando MESMA tarefa + robustez.\n💡 Analogia: Sectioning = 3 pessoas cozinhando pratos diferentes. Voting = 3 pessoas provando o mesmo vinho.\n⚠️ Voting para subtarefas independentes é desperdício.\n➡️ Intervalo! Voltamos com orchestrator-workers.", 32, T)

# ═══════════════════════════════════════
# SEÇÃO F — Orchestrator-Workers (33-39)
# ═══════════════════════════════════════

s = section_slide(5, "Orchestrator-Workers: Delegação Dinâmica")
add_notes(s, "Quarto padrão. O mais confundido com parallelization. Distinção: subtarefas DINÂMICAS vs FIXAS.")

s = comparison_slide(
    "Distinção Chave: Subtarefas Dinâmicas",
    "Parallelization (FIXAS)", [
        "Subtarefas CONHECIDAS a priori",
        "Você escreve N funções",
        "Roda em paralelo",
        "Simples, barato, previsível",
        "Ex: 3 fontes fixas de busca",
    ],
    "Orchestrator-Workers (DINÂMICAS)", [
        "Subtarefas NÃO conhecidas",
        "LLM orquestrador define em runtime",
        "Número e tipo variam por input",
        "Flexível, mas mais caro",
        "Ex: fontes emergem do problema",
    ],
    "📖 Distinção MAIS importante da aula. Parallelization: subtarefas fixas (você sabe). Orchestrator-workers: dinâmicas (orquestrador decide). Se fixas, use parallelization. Se dinâmicas, orchestrator.\n💡 Analogia: Parallelization = receita de bolo (passos conhecidos). Orchestrator = resolver mistério (descobre passos).\n❓ 'Para 3 fontes fixas, qual padrão?' (Parallelization)\n⚠️ Orchestrator para subtarefas fixas = custo e latência desnecessários.\n➡️ Papel do orquestrador.", 34, T, "Distinção crítica"
)

s = content_slide("Orquestrador: Planejar + Sintetizar", [
    "Fase 1 — PLANEJAR: LLM decompõe a tarefa em subtarefas",
    "Fase 2 — DELEGAR: despacha cada subtarefa para um worker",
    "Fase 3 — SINTETIZAR: agrega resultados em resposta final",
    "",
    "O orquestrador NÃO executa subtarefas — ele COORDENA",
    "",
    "Workers = LLMs com prompts/tools especializados",
    "",
    "Separação de planejamento e execução = qualidade (Plan-and-Solve)",
], "📖 Orquestrador tem 3 papéis. Planejar (decompõe). Delegar (despacha). Sintetizar (integra). NÃO executa — só coordena. Workers é que fazem. Separação de planejamento e execução melhora qualidade.\n💡 Analogia: maestro. Não toca instrumentos — coordena. Olha partitura (planejar), indica quem entra (delegar), integra sons (sintetizar).\n⚠️ Orquestrador executando subtarefas quebra o padrão.\n➡️ Arquitetura de implementação.", 35, T, "Papéis do orquestrador")

s = content_slide("Implementação: Decomposição + Workers + Reducer", [
    "Componente 1: DECOMPOSITOR (LLM com prompt de planejamento)",
    "Componente 2: WORKERS (pool de LLMs especializados)",
    "Componente 3: REDUCER (LLM que agrega)",
    "",
    "Estado: lista de subtarefas + resultados parciais",
    "",
    "Conexão com papers:",
    "  • Plan-and-Solve (arXiv:2305.17126) — separar planejar/executar",
    "  • ReWOO (arXiv:2305.04091) — plano cego + paralelismo total",
    "  • LLMCompiler (arXiv:2310.01757) — paralelização estruturada",
    "",
    "Diagrama: 12-Diagrams/ETHAGT03/orchestrator-workers.mmd",
], "📖 3 componentes: decompositor (LLM planeja), workers (pool executa), reducer (LLM agrega). Estado passa entre componentes. Plan-and-Solve: separar melhora. ReWOO: plano sem reagir permite paralelismo total. LLMCompiler: formaliza fan-out.\n💡 Analogia: empresa de consultoria. Sócio sênior (decompositor) divide. Consultores (workers) executam. Sócio (reducer) integra.\n⚠️ Sem estado persistente, reducer não acessa resultados dos workers.\n➡️ Casos de uso.", 36, T, "Arquitetura orchestrator-workers", acronyms="ReWOO = Reasoning WithOut Observation — raciocinio sem observacao intermediaria")

s = content_slide("Casos de Uso", [
    "1. Coding em múltiplos arquivos",
    "   orquestrador decide quais arquivos criar/modificar",
    "",
    "2. Search em múltiplas fontes",
    "   orquestrador decide quais fontes consultar",
    "",
    "3. Relatório multi-fonte",
    "   orquestrador divide por tópico, workers pesquisam, reducer sintetiza",
    "",
    "Sinal comum: 'não sei quantos steps vou precisar'",
    "As subtarefas EMERGEM do input",
], "📖 3 casos clássicos. Coding multi-arquivo (orquestrador decide arquivos). Search multi-fonte (orquestrador decide fontes). Relatório (orquestrador divide por tópico). Sinal comum: não sabe steps de antemão.\n💡 Analogia: general em campanha. Não sabe quantas batalhas — decide conforme terreno. Cada batalha diferente. Integra vitórias.\n❓ 'Pensem: subtarefas fixas ou dinâmicas no sistema de vocês?'\n⚠️ Orchestrator para passos fixos = overkill.\n➡️ Código.", 37, T, "Casos de uso reais")

s = code_slide("Código: Orchestrator-Workers", """async def orchestrate(task):
    # Fase 1: PLANEJAR (LLM decompõe)
    plan = await planner_llm(
        "Decomponha em subtarefas", task
    )
    subtasks = plan["subtasks"]

    # Fase 2: DELEGAR (workers em paralelo)
    results = await asyncio.gather(
        *[worker(sub) for sub in subtasks]
    )

    # Fase 3: SINTETIZAR (LLM integra)
    return await synthesizer_llm(
        "Integre os resultados", task, results
    )""", "📖 Comparem com parallelization (Slide 30). Diferença: planner LLM gera subtarefas dinamicamente (não hardcoded). Depois, mesmo padrão: gather para paralelizar. MAS há step extra: synthesizer integra. Sem ele, N resultados soltos.\n💡 Analogia: diretor de filme. Roteiro (planner) diz cenas. Câmeras (workers) filmam em paralelo. Editor (synthesizer) monta o filme.\n⚠️ Sem synthesizer: N respostas fragmentadas em vez de resposta integrada.\n➡️ Vamos quebrar um mito.", 38, T, "Implementação orchestrator-workers")

s = exercise_slide("V/F: Orchestrator-Workers É Sempre Melhor", [
    "Verdadeiro ou Falso:",
    "",
    "'Orchestrator-workers é sempre melhor que parallelization.'",
    "",
    "Resposta: FALSO",
    "",
    "Por quê: se subtarefas são FIXAS e conhecidas,",
    "parallelization é mais simples, barato e previsível.",
    "",
    "Orchestrator adiciona custo (LLM de planejamento)",
    "e latência (step extra serial)",
    "",
    "Regra: só use orchestrator quando subtarefas são",
    "genuinamente DINÂMICAS",
], "📖 Armadilha 'mais complexo = melhor'. FALSO. Orchestrator adiciona custo (chamada extra de planejamento) e latência (step serial). Se subtarefas fixas, parallelization vence.\n💡 Analogia: receita vs improvisação. Se tem receita (fixas), siga-a. Improvisar (orchestrator) só quando não tem.\n❓ 'Votem: V ou F?'\n⚠️ Complexidade sem necessidade = débito técnico.\n➡️ Quinto padrão: evaluator-optimizer.", 39, T)

# ═══════════════════════════════════════
# SEÇÃO G — Evaluator-Optimizer (40-46)
# ═══════════════════════════════════════

s = section_slide(6, "Evaluator-Optimizer: Loop de Refinamento")
add_notes(s, "Quinto e último padrão. Loop: gerar, avaliar, refinar até satisfazer critério.")

s = content_slide("O Loop: Gerar → Avaliar → Refinar", [
    "Step 1: LLM gera saída (GENERATOR)",
    "Step 2: LLM (ou código) avalia contra critérios (EVALUATOR)",
    "Step 3: Se abaixo do critério → LLM refina com feedback (OPTIMIZER)",
    "",
    "Loop até CRITÉRIO DE PARADA",
    "",
    "Diferença de ReAct:",
    "  ReAct = pensar-agir-observar (ambiente externo)",
    "  Evaluator = gerar-avaliar-refinar (critério interno)",
    "",
    "Diagrama: 12-Diagrams/ETHAGT03/evaluator-optimizer.mmd",
], "📖 Loop: generator → evaluator → optimizer. Evaluator avalia contra critério. Se abaixo, optimizer refina com feedback. Repete até parar. Diferença de ReAct: aqui é contra critério interno, não ambiente externo.\n💡 Analogia: escritor com editor. Escreve, editor avalia contra critérios, revisa com feedback. Repete até aprovar.\n⚠️ Confundir com ReAct é comum. ReAct age no ambiente; evaluator avalia contra critério.\n➡️ Quando vale?", 41, T, "Estrutura do loop", acronyms="ReAct = Reasoning and Acting — padrao de loop Thought / Action / Observation")

s = content_slide("Quando Tem Valor", [
    "3 condições devem ser VERDADEIRAS:",
    "",
    "1. Feedback é ARTICULÁVEL",
    "   (LLM explica O QUE está errado, não só QUE está)",
    "",
    "2. LLM consegue AVALIAR",
    "   (existe critério objetivo ou near-objetivo)",
    "",
    "3. Iteração MELHORA resultado",
    "   (nem sempre — às vezes o modelo trava)",
    "",
    "NÃO use se: evaluator não é melhor que generator → diverge",
    "Funciona: tradução literária (feedback sobre tom, ritmo)",
    "Não funciona: 'seja mais criativo' (não-articulável)",
], "📖 3 condições. Feedback articulável (explica o quê, não só que). LLM consegue avaliar (critério objetivo). Iteração melhora (modelo não trava). Se evaluator é fraco, loop diverge.\n💡 Analogia: professor de piano. 'Toque com mais sentimento' (não-articulável) = não melhora. 'Diminua tempo no 3º compasso' (articulável) = melhora.\n⚠️ Feedback vago ('melhore') não converge.\n➡️ Como estruturar critérios.", 42, T, "Quando usar evaluator-optimizer")

s = content_slide("Critérios Claros e Mensuráveis", [
    "Critério VAGO: 'a resposta é boa' → não converge",
    "",
    "Critério CLARO: 'cita ≥3 fontes, ≤500 palavras,",
    "  inclui recomendação'",
    "",
    "RUBRIC ESTRUTURADA: dimensão → escala → descrição por nível",
    "",
    "Exemplo — tradução:",
    "  Fidelidade: 1-5 (5 = preserva tom e ritmo)",
    "  Fluência: 1-5 (5 = natural em PT-BR)",
    "  Terminologia: 1-5 (5 = glossário correto)",
    "",
    "LLM-as-judge COM rubric >>> SEM rubric",
], "📖 Qualidade depende dos critérios. Vago não converge. Claro é verificável. Melhor: rubric (dimensões, escala, descrição por nível). LLM-as-judge com rubric é significativamente mais consistente.\n💡 Analogia: rubric de redação do ENEM. Sem rubric, cada corretor dá nota diferente. Com rubric (competências, níveis), mais consistente.\n⚠️ 'Bom/médio/ruim' é frágil. Sempre escala numérica (1-5) + descrição por nível.\n➡️ Quando parar o loop?", 43, T, "Rubric estruturada")

s = content_slide("Convergência: Critérios de Parada", [
    "Critério 1 — SCORE: parar quando ≥ threshold (ex.: 4.5/5)",
    "Critério 2 — MAX ITERS: parar após N (ex.: 5) — orçamento",
    "Critério 3 — DELTA ESTAGNADO: score não melhora por 2 iters",
    "",
    "SEMPRE usar ≥2 critérios (score OU max iters, no mínimo)",
    "",
    "Sem max iters → loop explode em custo",
    "Sem score → para por orçamento, talvez antes da qualidade",
    "",
    "Custo: cada iteração = generator + evaluator = 2× chamadas",
], "📖 3 critérios de parada. Score (qualidade). Max iters (orçamento). Delta estagnado (detecta platô). Sempre ≥2: score OU max iters. Sem max iters = custo explode. Cada iteração = 2 chamadas.\n💡 Analogia: treinar para maratona. Score = correr 42km (meta). Max iters = 12 semanas (tempo). Delta = se tempo não melhora por 2 semanas, pare.\n⚠️ Só score (sem max iters) = loop explode.\n➡️ Código.", 44, T, "Condições de parada")

s = code_slide("Código: Evaluator-Optimizer", """def evaluate_optimize(task, max_iters=5, threshold=4.5):
    output = generator(task)
    prev_score = 0

    for i in range(max_iters):
        # Evaluator avalia com rubric
        result = evaluator(output, rubric=RUBRIC)

        # Parada 1: score >= threshold
        if result["score"] >= threshold:
            return output

        # Parada 2: delta estagnado
        if abs(result["score"] - prev_score) < 0.1:
            return output

        # Optimizer refina com feedback
        output = optimizer(task, output, result["feedback"])
        prev_score = result["score"]

    # Parada 3: max iters
    return output""", "📖 3 condições no código: score ≥ threshold, delta < 0.1, iter ≥ max. Optimizer recebe output + feedback, produz versão melhorada. Em produção, loguem cada iteração (score, feedback) para debugar convergência.\n💡 Analogia: code review. Envia PR (generator). Revisor avalia (evaluator). Se rejeita com comentários, revisa (optimizer). Repete até aprovar ou limite.\n⚠️ Sem log de scores, não sabe se converge ou diverge.\n➡️ Vamos praticar.", 45, T, "Implementação evaluator-optimizer")

s = exercise_slide("Condição de Parada para Tradução Literária", [
    "Cenário: evaluator-optimizer para tradução de um poema",
    "(EN → PT-BR)",
    "",
    "Pergunta: Escreva a condição de parada.",
    "",
    "  • Quais critérios?",
    "  • Qual threshold?",
    "  • Quantas iterações máximas?",
    "",
    "Dica: tradução literária é subjetiva —",
    "como evitar loop infinito?",
    "",
    "Discussão em duplas (2 min), 1 min compartilhar",
], "📖 Threshold NÃO 5.0 (perfeição inatingível em literário). Use média de rubric (fidelidade, fluência, ritmo, tom) ≥ 4.0. Max iters 3-4 (literário não melhora indefinidamente; modelo trava). Delta estagnado.\n💡 Analogia: polir escultura. Primeiras passadas dão forma. Últimas mal mudam. Saber parar é arte.\n❓ 'Compartilhem 2 respostas.'\n⚠️ Threshold=5.0 = loop nunca converge = perde orçamento.\n➡️ Composições.", 46, T)

# ═══════════════════════════════════════
# SEÇÃO H — Composições e Limites (47-51)
# ═══════════════════════════════════════

s = section_slide(7, "Composições e Limites")
add_notes(s, "Na prática, raramente usamos um padrão isolado. Combinamos. Mas a fronteira entre workflow composto e agente é tênue.")

s = content_slide("Composição Típica: Routing → Parallelization → Evaluator", [
    "Camada 1 — ROUTING: classifica entrada (tipo de ticket)",
    "Camada 2 — PARALLELIZATION: subtarefas em paralelo",
    "  (buscar FAQ + docs + histórico simultaneamente)",
    "Camada 3 — EVALUATOR-OPTIMIZER: valida resposta antes de entregar",
    "",
    "Cada camada resolve um problema DIFERENTE:",
    "  routing = custo (modelo certo)",
    "  parallelization = latência (subtarefas em paralelo)",
    "  evaluator = qualidade (validação final)",
    "",
    "Outras: chaining + evaluator; routing + orchestrator",
    "",
    "Diagrama: 12-Diagrams/ETHAGT03/composition-routing-parallel-evaluator.mmd",
], "📖 Composição mais comum em produção. 3 camadas: routing (custo), parallelization (latência), evaluator (qualidade). Cada uma resolve problema diferente.\n💡 Analogia: restaurante. Hostess (routing) direciona. Garçom (parallelization) pede em paralelo. Chef (evaluator) prova antes de servir.\n⚠️ Não componha sem justificar. Cada camada adiciona custo+latência. Só se agrega valor mensurável.\n➡️ Quando vira agente?", 48, T, "Composição de 3 camadas")

s = content_slide("Quando Combinar Vira Agente", [
    "Workflow composto: camadas FIXAS, ordem PREDEFINIDA,",
    "  gates determinísticos",
    "",
    "Agente: LLM decide ordem, se itera, se busca mais info",
    "",
    "A linha é TÊNUE:",
    "  routing + loops + branches condicionais suficientes...",
    "  ...em algum ponto, vira agente",
    "",
    "Sinal de transição: workflow precisa de 'meta-step'",
    "  que decide QUAL caminho seguir dinamicamente",
    "",
    "ETHAGT04 (Reasoning & Planning) aprofunda essa ponte",
], "📖 Fronteira tênue. Workflow composto: você define camadas, ordem, gates. Agente: LLM decide. Mas routing dinâmico + loops + branches suficientes = semi-agêntico. Sinal de transição: meta-step decide caminho. ETHAGT04 aprofunda.\n💡 Analogia: receita (workflow) vs cozinhar 'no feeling' (agente). No meio: receitas com variações ('se molho grosso, adicione água') = semi-agêntico.\n❓ 'O que diferencia workflow composto de agente para vocês?' (debate)\n⚠️ Workflow composto ≠ agente. Enquanto caminho é predefinido (mesmo com condicionais), é workflow.\n➡️ Sinais de alerta.", 49, T, "Fronteira workflow/agente")

s = content_slide("Sinais de Que Você Está Forçando Workflow", [
    "⚠️ Gates cada vez mais complexos para cobrir casos especiais",
    "",
    "⚠️ Número de branches no routing cresce sem parar",
    "",
    "⚠️ Precisa de 'loops dentro de loops' para edge cases",
    "",
    "⚠️ Evaluator precisa de contexto que generator não teve",
    "",
    "⚠️ Descrição do fluxo: 'depende do que acontecer no step anterior'",
    "",
    "→ Esses sinais indicam: o problema pede AGENTE, não workflow",
], "📖 5 sinais de que problema pede agente. Gates viram LLM disfarçado. Routing tem 20 categorias crescendo. Loops aninhados. Evaluator sem contexto do generator. 'Depende do step anterior'.\n💡 Analogia: dirigir com marchas pré-programadas. Funciona em reta. Em trânsito caótico, precisa motorista que decide — agente.\n⚠️ Insistir em workflow por 'previsibilidade' em problema dinâmico = ilusão.\n➡️ Trade-offs consolidados.", 50, T, "Sinais de alerta")

s = content_slide("Trade-offs Consolidados", [
    "                   Previsib.  Flexib.  Custo      Latência",
    "Prompt Chaining:   Alta       Baixa    Médio      Alta (serial)",
    "Routing:           Alta       Média    Baixo      Baixa",
    "Parallelization:   Alta       Média    Alto (N×)  Baixa (max)",
    "Orchestrator-W:    Média      Alta     Alto       Média",
    "Evaluator-Opt:     Alta       Média    Alto(loop) Alta",
    "",
    "Escolha pelas RESTRIÇÕES do problema:",
    "  custo → routing | latência → parallelization",
    "  flexibilidade → orchestrator | qualidade → evaluator",
    "  linearidade → chaining",
], "📖 Tabela MAIS importante do módulo. 5 padrões × 4 eixos. Chaining: previsível mas lento. Routing: barato/rápido. Parallelization: rápido mas caro. Orchestrator: flexível mas menos previsível. Evaluator: previsível mas caro/lento. Escolha pelas restrições.\n💡 Analogia: escolher transporte. Bicicleta (routing). Carro (parallelization). Táxi (orchestrator). Avião (evaluator). Trem (chaining).\n⚠️ Escolher pelo padrão 'favorito' é erro. Escolha pelas restrições.\n➡️ Fechamento.", 51, T, "Matriz de trade-offs")

# ═══════════════════════════════════════
# SEÇÃO I — Fechamento (52-63)
# ═══════════════════════════════════════

s = section_slide(8, "Fechamento")
add_notes(s, "Fechamento. Caso real, resumo, quiz, exercício final, projeto, Q&A.")

s = content_slide("Caso de Estudo: Coinbase / Intercom", [
    "Coinbase / Intercom: workflows agênticos em suporte ao cliente",
    "",
    "Arquitetura típica:",
    "  routing (classificar ticket)",
    "  → parallelization (buscar fontes: KB, docs, histórico)",
    "  → evaluator-optimizer (validar resposta)",
    "",
    "Resultados: redução de tempo de resposta, escalabilidade,",
    "  custo controlado",
    "",
    "Lição: workflow composto > agente autônomo para suporte",
    "  (previsibilidade importa em escala)",
    "",
    "Referência: 09-CaseStudies/",
], "📖 Coinbase e Intercom usam composição que vimos (Slide 48). Por que workflow e não agente? Suporte exige previsibilidade. Cada resposta consistente, auditável, dentro de SLA. Agente autônomo = custo imprevisível, latência variável.\n💡 Analogia: call center estruturado (workflow) vs consultor freelancer (agente). Em escala, call center vence.\n⚠️ Maioria de produção usa workflow composto, não agente autônomo.\n➡️ Resumo.", 53, T, "Caso real de produção")

s = content_slide("Resumo da Aula", [
    "1. 5 padrões canônicos: chaining, routing, parallelization,",
    "   orchestrator-workers, evaluator-optimizer",
    "",
    "2. Gates programáticos = controle DETERMINÍSTICO em workflows",
    "",
    "3. Parallelization: sectioning (independentes) vs voting (robustez)",
    "",
    "4. Orchestrator-Workers: subtarefas DINÂMICAS (vs fixas)",
    "",
    "5. Evaluator-Optimizer: critérios claros + convergência mensurável",
    "",
    "6. Composição = realidade; mas quando vira agente, mude de módulo",
    "",
    "7. COMECE SIMPLES, só aumente com evidência",
], "📖 7 pontos-chave. 5 padrões = caixa de ferramentas. Gates = determinismo. Parallelization tem 2 modos. Orchestrator = dinâmico. Evaluator = critérios+parada. Componha mas saiba quando vira agente. Comece simples.\n💡 Analogia: decálogo do engenheiro de workflows. Sete mandamentos.\n➡️ Checklist.", 54, T, "Síntese da aula")

s = content_slide("Checklist da Aula", [
    "[ ] Explicou o princípio 'comece simples'",
    "[ ] Detalhou os 5 padrões canônicos com diagramas",
    "[ ] Mostrou código de cada padrão",
    "[ ] Discutiu gates programáticos",
    "[ ] Comparou parallelization vs orchestrator-workers",
    "[ ] Apresentou critérios de convergência do evaluator",
    "[ ] Discutiu composições e quando workflow vira agente",
], "📖 Checklist final. Se algum item sem check, vale revisitar. Os não-cobertos viram leitura recomendada.\n➡️ Quiz.", 55, T, "Confirmação de cobertura")

s = content_slide("Quiz — Pergunta 1", [
    "Qual é a diferença fundamental entre parallelization",
    "e orchestrator-workers?",
    "",
    "A) Orchestrator-workers é mais rápido",
    "B) Parallelization usa modelos menores",
    "C) Em orchestrator-workers, as subtarefas são DINÂMICAS",
    "   (definidas em runtime)",
    "D) Parallelization não tem reducer",
    "",
    "Resposta: C",
], "📖 C. Distinção: parallelization = fixas; orchestrator = dinâmicas. A é falso (orchestrator é mais lento por step de planejamento). B é falso (ambos usam qualquer modelo). D é falso (ambos têm reducer).\n➡️ Próxima.", 56, T, "Quiz 1/4")

s = content_slide("Quiz — Pergunta 2", [
    "O que é um gate em prompt chaining?",
    "",
    "A) Um modelo LLM que decide se continua",
    "B) Um checkpoint PROGRAMÁTICO (código) que valida",
    "   a saída antes do próximo step",
    "C) Um tipo de prompt especial",
    "D) Um framework de orquestração",
    "",
    "Resposta: B",
], "📖 B. Gate é CÓDIGO (determinístico), não LLM. Se fosse LLM, seria evaluator-optimizer. A é a armadilha clássica.\n➡️ Próxima.", 57, T, "Quiz 2/4")

s = content_slide("Quiz — Pergunta 3", [
    "Quando o evaluator-optimizer NÃO vale a pena?",
    "",
    "A) Quando o feedback é articulável",
    "B) Quando o evaluator não é melhor que o generator",
    "C) Quando há orçamento para múltiplas iterações",
    "D) Quando a tarefa tem critérios objetivos",
    "",
    "Resposta: B",
], "📖 B. Se evaluator é mais fraco que generator, feedback é pior que saída. Loop diverge. A, C, D são quando VALE a pena.\n➡️ Última pergunta.", 58, T, "Quiz 3/4")

s = content_slide("Quiz — Pergunta 4", [
    "Em routing por modelo, qual é o erro mais caro?",
    "",
    "A) Enviar tarefa fácil para modelo forte (Sonnet)",
    "B) Enviar tarefa difícil para modelo fraco (Haiku)",
    "C) Usar o mesmo modelo para todas as categorias",
    "D) Ter apenas 2 categorias",
    "",
    "Resposta: B",
], "📖 B. Falso negativo (difícil→fácil) gera resposta RUIM. A é apenas ineficiência (custo extra, resposta boa). Mitigação: incerteza → path robusto.\n➡️ Exercício final.", 59, T, "Quiz 4/4")

s = exercise_slide("Exercício — Escolha o Workflow (5 Cenários)", [
    "Para cada cenário, indique o workflow + justifique:",
    "",
    "1. Tradução com revisão de qualidade",
    "2. Análise de sentimentos de 10.000 tweets",
    "3. Geração de relatório com múltiplas fontes",
    "4. Chatbot de FAQ simples",
    "5. Correção de redação com feedback",
    "",
    "Grupos: 3 min discussão, 2 min compartilhar",
    "",
    "Gabarito: 1=Chaining+Evaluator | 2=Routing+Sectioning",
    "3=Orchestrator | 4=Routing | 5=Evaluator-Optimizer",
], "📖 Exercício de consolidação. A chave é a JUSTIFICATIVA — por que este padrão e não outro? Baseiem nas restrições do problema.\n💡 Analogia: diagnóstico médico. Cada sintoma pede tratamento diferente. O bom médico justifica.\n❓ 'Compartilhem 2 grupos com justificativas.'\n⚠️ Orchestrator para tudo ('parece avançado') = erro. Justifique nas restrições.\n➡️ Projeto.", 60, T)

s = content_slide("Projeto do Módulo + Labs", [
    "PROJETO: síntese de relatório a partir de múltiplas fontes",
    "  • Projetar e implementar workflow composto",
    "  • Comparar 2 abordagens (ex.: chaining vs orchestrator)",
    "  • Entrega: código + benchmark + ADR justificando",
    "  • Critério: ADR coerente; medições em ≥20 casos",
    "",
    "Lab 1 (5h): 'Os 5 em 1 dia'",
    "  Versões mínimas dos 5 workflows em domínio comum",
    "",
    "Lab 2 (5h): 'Composição'",
    "  routing → parallelization (3 workers) → evaluator-optimizer",
], "📖 Projeto é entrega principal. Código funciona é pré-requisito. O ADR é o que vale — justificativa com dados. Lab 1 implementa os 5 isoladamente. Lab 2 compõe.\n💡 Analogia: defesa de tese. Código é a tese (precisa existir). ADR é a defesa (precisa convencer).\n⚠️ Sem medições (≥20 casos), ADR é opinião, não engenharia.\n➡️ Referências.", 61, T, "Projeto e labs")

s = content_slide("Conexão com Próximos Módulos + Referências", [
    "Próximos módulos:",
    "  ETHAGT04 — Reasoning & Planning (quando workflow não basta)",
    "  ETHAGT09 — Multi-Agente (orchestrator expandido)",
    "  ETHAGT10 — Orquestração em escala",
    "",
    "Leitura obrigatória:",
    "  Anthropic, Building Effective Agents (2024)",
    "",
    "Papers:",
    "  Plan-and-Solve (arXiv:2305.17126)",
    "  ReWOO (arXiv:2305.04091)",
    "  LLMCompiler (arXiv:2310.01757)",
    "",
    "Vídeo: Schluntz & Albert (YouTube)",
], "📖 ETHAGT03 é fundação para ETHAGT04 (reasoning), 09 (multi-agente), 10 (orquestração). Leitura obrigatória: Anthropic (fonte dos 5 padrões). Papers aprofundam orchestrator (Plan-and-Solve, ReWOO) e parallelization (LLMCompiler).\n➡️ Q&A.", 62, T, "Referências e próximos passos")

s = title_slide(
    "Perguntas?",
    "ETHAGT03 — Padrões de Workflow Agêntico\nPróxima aula: ETHAGT04 — Reasoning & Planning",
    "Universidade Etho"
)
add_notes(s, "📖 Abrir para Q&A. Se não houver perguntas: 'Qual padrão foi menos claro?' Mensagem final: 'Os 5 padrões são a caixa de ferramentas. A escolha certa é o que diferencia engenharia de tentativa-e-erro. Pratiquem no Lab 1 — implementar os 5 em um dia é a melhor forma de fixar.'\n➡️ Fim da aula.")

output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT03-Apresentacao.pptx")
prs.save(output)
print(f"PPTX gerado: {output}")
print(f"Total de slides: {len(prs.slides)}")
