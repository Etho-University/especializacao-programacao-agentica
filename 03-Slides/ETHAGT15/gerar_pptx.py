#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT15: Meta-Agentes & Sistemas Autoaprendentes
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

CODE = "ETHAGT15"

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj="", acronyms=""):
    if acronyms:
        add_textbox(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.35), acronyms, size=9, color=INFO)
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code=CODE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    if acronyms:
        add_textbox(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4), acronyms, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s

T = 67

s = title_slide(
    "Meta-Agentes & Sistemas Autoaprendentes",
    "Universidade Etho · Fase D — Fronteira · 15 h",
    CODE,
)
add_notes(s, "Bem-vindos à fronteira. Hoje: agentes que criam agentes, com segurança.")

content_slide(
    "Objetivos do Módulo",
    [
        "Geral: explorar a fronteira dos meta-agentes com conscientização dos riscos",
        "1. Definir meta-agente, strategy evolver, meta-learning",
        "2. Implementar sistema onde agente gera agentes especializados",
        "3. Aplicar otimização automatizada (DSPy, Promptbreeder)",
        "4. Discutir auto-aprendizado contínuo com memória acumulada",
        "5. Identificar riscos (recursão, drift, segurança) e mitigações",
    ],
    "Cada objetivo é mensurável. O #5 é o mais importante: sair sem achar que é bala de prata.",
    2, T,
)

content_slide(
    "Competências Desenvolvidas",
    [
        "C1 Programação Agêntica → A (Avançado)",
        "C2 Multi-Agent Systems → A (Avançado)",
        "C3 MCP & Tool Use → B (Básico)",
        "C4 Agent Memory → A (Avançado)",
        "C6 Agent Security → I (Intermediário)",
    ],
    "Culminância técnica. C6 sobe para Intermediário por causa dos riscos novos.",
    3, T,
, acronyms="MCP = Model Context Protocol")

content_slide(
    "Agenda da Aula",
    [
        "Bloco 1: Abertura → Meta-Agência → Geração de Agentes",
        "Bloco 2: Otimização → Auto-Aprendizado → Riscos → Fechamento",
        "7 seções · 67 slides · 90 min",
    ],
    "Progressão clara: definir → gerar → otimizar → aprender → controlar riscos.",
    4, T,
)

content_slide(
    "A Fronteira dos Meta-Agentes",
    [
        "Otimização manual não escala: 5 engenheiros × 20 agentes = impraticável",
        "Cada novo domínio = novo agente manual = semanas",
        "Solução: meta-agente que otimiza automaticamente",
        "Fronteira: agente que cria outro agente — onde está o limite?",
        "Pergunta: Você deixaria um agente modificar o próprio prompt?",
    ],
    "Criar tensão. A pergunta final engaja — a maioria hesita.",
    5, T,
)

content_slide(
    "Por Que Meta-Agência Agora",
    [
        "2023: DSPy · Voyager · Promptbreeder",
        "2024: Meta-Prompting · AI Scientist",
        "2025: adoção industrial",
        "Confluência: modelos capazes + frameworks de eval + escala",
        "Canônico: DSPy (arXiv:2310.03714)",
    ],
    "Três fatores convergiram. 2023 foi o ano seminal.",
    6, T,
)

section_slide(1, "O Que É Meta-Agência")

content_slide(
    "Agentes que Operam sobre Agentes",
    [
        "Agente comum: opera sobre ambiente (tools, APIs, dados)",
        "Meta-agente: opera sobre agentes (criar, configurar, otimizar)",
        "O meta-agente é um agente cujo ambiente é outro agente",
        "Analogia: compilador é para código o que meta-agente é para agente",
        "Níveis: agente → meta-agente → meta-meta (cuidado com recursão)",
    ],
    "Definição central da aula. Meta-agente age sobre agentes, não sobre ambiente.",
    8, T,
)

content_slide(
    "Três Estratégias de Meta-Agência",
    [
        "Synthesis: criar agente do zero (risco médio, reward alto)",
        "Evolution: mutar config existente (risco baixo-médio)",
        "Optimization: ajustar parâmetros finos (risco baixo)",
        "Regra: tentar optimization antes de evolution, evolution antes de synthesis",
    ],
    "Espectro de drasticidade. Sempre começar pelo menos drástico.",
    9, T,
)

content_slide(
    "Arquitetura Conceitual do Meta-Agente",
    [
        "tarefa T → Meta-agente → compõe primitivas",
        "→ gera config JSON → validar → instanciar",
        "Primitivas: personas, tools, workflows",
        "Validação é crítica: eval antes de deploy",
        "Loop de feedback (falha → iterar) é o coração do sistema",
        "Diagrama: 12-Diagrams/ETHAGT15/meta-agent.mmd",
    ],
    "Diagrama central da aula. Destacar o loop de validação.",
    10, T,
)

content_slide(
    "Implementações Reais de Meta-Agência",
    [
        "DSPy (arXiv:2310.03714) — compilação declarativa",
        "Voyager (arXiv:2305.16291) — auto-skills no Minecraft",
        "Promptbreeder (arXiv:2309.16797) — evolução de prompts",
        "Meta-Prompting (arXiv:2311.11402) — decomposição com especialistas",
        "AI Scientist (arXiv:2408.06292) — pesquisa autônoma",
        "Padrão comum: gerar → avaliar → selecionar → iterar",
    ],
    "Cinco papers definiram o campo. Padrão comum é evolução aplicada.",
    11, T,
)

comparison_slide(
    "Risco vs Benefício da Meta-Agência",
    "Riscos",
    [
        "Perda de controle (goal drift)",
        "Recursão descontrolada",
        "Opacidade (por que a config é assim?)",
        "Drift para ambiente obsoleto",
        "Segurança: agente remove constraints",
    ],
    "Benefícios",
    [
        "Escala: otimizar 100 agentes",
        "Adaptação a mudanças de domínio",
        "Descoberta de configs contraintuitivas",
        "Velocidade: semanas → horas",
    ],
    "Encarar trade-offs de frente. A pergunta é 'usar com quais salvaguardas'.",
    12, T,
)

exercise_slide(
    "Você Deixaria um Agente Modificar o Próprio Prompt?",
    [
        "Votação: sim / não / depende",
        "Em quais condições você aceitaria?",
        "Sandbox? Eval rigoroso? Governor? HITL?",
        "Quem audita as mudanças?",
        "Discussão aberta (2 min)",
    ],
    "Resposta correta: depende — com 5 salvaguardas. Anotar votação.",
    13, T,
    obj="Votação",
, acronyms="HITL = Human-in-the-Loop")

exercise_slide(
    "É Meta-Agência ou Não?",
    [
        "1. Agente que escreve código Python → C",
        "2. Agente que gera prompt para outro → M",
        "3. Agente que ajusta própria temperatura → M",
        "4. Agente que busca na web → C",
        "Votação rápida por cenário",
    ],
    "Pegadinha: #3 é meta-agência (optimization). Ajustar config de agente é meta.",
    14, T,
    obj="Exercício rápido",
)

section_slide(2, "Geração de Agentes")

content_slide(
    "Meta-Agente que Produz Agentes",
    [
        "Input: descrição de tarefa",
        "Output: config completa (prompt + tools + modelo + params)",
        "Primitivas: personas, tools (MCP), workflows",
        "Composição: selecionar + instanciar + configurar",
        "Não inventa do zero — compõe a partir de biblioteca",
    ],
    "Segredo: composição, não criação ex nihilo. Torna auditável.",
    16, T,
)

content_slide(
    "Templates e Composição",
    [
        "Template = config parametrizada de agente",
        "Ex: SupportAgentTemplate(product, kb, escalation_policy)",
        "Composição: combinar templates (suporte + QA)",
        "Biblioteca de templates: catálogo versionado",
        "Vantagens: reuso, consistência, auditabilidade",
    ],
    "Templates são o mecanismo de reuso. Versionar tudo.",
    17, T,
)

content_slide(
    "Pipeline de Geração de Agentes",
    [
        "1. Interpretar tarefa (LLM analisa descrição)",
        "2. Selecionar primitivas (da biblioteca)",
        "3. Gerar config (JSON declarativo)",
        "4. Validar (schema + eval + safety)",
        "5. Instanciar (criar agente executável)",
        "6. Deploy (confiança incremental)",
        "Falha → feedback → iterar",
    ],
    "Passo 4 (validação) é o gargalo. Geração é rápida; eval é lento.",
    18, T,
, acronyms="LLM = Large Language Model")

content_slide(
    "Validação do Agente Gerado",
    [
        "Confie, mas verifique — gerado ≠ bom",
        "Schema check (config válida?)",
        "Eval em subset (passa nos critérios?)",
        "Safety check (não viola políticas?)",
        "Review humano para mudanças críticas",
        "Sem validação: agente pode ser pior que manual",
    ],
    "Sem validação, você está deployando código não-testado em produção.",
    19, T,
)

content_slide(
    "Caso Prático — Suporte por Produto",
    [
        "Empresa com 10 produtos, cada um precisa de agente de suporte",
        "Manual: 10 engenheiros × 2 semanas = 20 semanas",
        "Com meta-agente: 1 engenheiro × 2 dias = 2 dias",
        "Gera: prompt + tools (search_kb, escalate) + config",
        "Validação: 50 casos/produto, aprova se >90%",
        "Resultado: 10 agentes, 92% precisão média",
    ],
    "Caso concreto. Ganho não é só velocidade — é consistência.",
    20, T,
)

code_slide(
    "DEMO — Agente que Escreve Agente",
    """$ python meta_agent.py

[INPUT]  "preciso de um agente que classifique tickets"

[META-AGENT] Interpretando tarefa...
[META-AGENT] Selecionando primitivas: ClassifierPersona, ClassifyTool
[META-AGENT] Gerando config JSON...

config = {
  "system_prompt": "Você é um classificador de tickets...",
  "tools": ["classify_ticket"],
  "model": "gpt-4o-mini",
  "temperature": 0.1
}

[VALIDATE] Schema check... PASS
[VALIDATE] Eval em 3 casos... 3/3 PASS
[DEPLOY]  Agente instanciado e pronto.""",
    "DEMO central. Se API falhar, mostrar este screenshot. Destacar loop eval.",
    21, T,
    obj="DEMO (5 min)",
)

exercise_slide(
    "Perguntas sobre a DEMO",
    [
        "O agente gerado é pior que um manual?",
        "Como sabemos que o eval cobre casos suficientes?",
        "E se o meta-agente gerar config que funciona no eval mas falha em produção?",
        "Discussão em duplas (2 min)",
    ],
    "Pergunta mais importante: overfitting ao eval. Holdout set é a resposta.",
    22, T,
)

exercise_slide(
    "Critérios de Qualidade para Agentes Gerados",
    [
        "Em duplas: listar 5 critérios de qualidade",
        "Exemplos: success rate, latência, custo, safety, robustez",
        "Bônus: como detectar overfitting ao eval?",
        "2 min listar, 1 min compartilhar",
    ],
    "Exercício rápido para fixar critérios além de 'acertou?'.",
    23, T,
)

content_slide(
    "Lições da Geração de Agentes",
    [
        "Geração é rápida, validação é o gargalo",
        "Eval de qualidade > velocidade de geração",
        "Templates reduzem variabilidade e aumentam reuso",
        "Comece com domínio estreito, expanda gradualmente",
        "HITL para mudanças críticas",
        "Config gerada deve ser auditável e versionada",
    ],
    "Seis lições para levar. Eval é o novo bottleneck.",
    24, T,
)

section_slide(3, "Otimização Automatizada")

content_slide(
    "Por Que Otimizar Automaticamente?",
    [
        "Prompts manuais: arte, não ciência",
        "Humanos exploram 5-10 variações; espaço tem milhares",
        "Otimização automática explora centenas em horas",
        "Reprodutível: mesma otimização → mesmo resultado",
        "Escala: 100 agentes em paralelo",
    ],
    "Três razões: exploração sistemática, reprodutibilidade, escala.",
    26, T,
)

content_slide(
    "DSPy — Compilação Declarativa",
    [
        "DSPy: 'compilar' assinaturas declarativas em prompts otimizados",
        "Analogia: DSPy é para prompts o que compilador é para código",
        "Você escreve o QUE (assinatura), DSPy gera o COMO (prompt)",
        "Assinatura: question -> answer",
        "DSPy escolhe: few-shot examples, formato, instruções",
        "Fonte: arXiv:2310.03714",
    ],
    "Mudança de paradigma: imperativo → declarativo.",
    27, T,
)

content_slide(
    "DSPy — Teleprompters",
    [
        "Teleprompter = otimizador do DSPy",
        "BootstrapFewShot: seleciona melhores exemplos",
        "MIPRO: otimiza instruções + exemplos via busca bayesiana",
        "Processo: assinatura + métrica + dados → teleprompter → config",
        "Resultado: prompt otimizado sem intervenção manual",
    ],
    "Teleprompter é o otimizador. Diferente de fine-tuning (treina pesos).",
    28, T,
)

content_slide(
    "Promptbreeder — Evolução de Prompts",
    [
        "Prompts 'reproduzem' e 'mutam' (algoritmo genético)",
        "População inicial → mutação (LLM reescreve) → seleção (eval)",
        "Gerações sucessivas: prompts melhoram com evolução",
        "Mutação guiada por LLM, não aleatória",
        "Fonte: arXiv:2309.16797",
    ],
    "Evolução direcionada: mutação por LLM + seleção por métrica.",
    29, T,
)

content_slide(
    "Atlas e Outras Abordagens",
    [
        "Atlas: otimização guiada por árvore de pensamentos",
        "OPRO (Google): meta-prompting gera candidatos",
        "TextGrad: gradientes de texto no pipeline",
        "Padrão comum: gerar → avaliar → selecionar → iterar",
        "Diferença: estratégia de busca (bayesiana, evolutiva, LLM-guided)",
    ],
    "Não há 'melhor' — depende do problema. Teste 2-3 e escolha com eval.",
    30, T,
)

content_slide(
    "Otimização de Tools",
    [
        "Tool description afeta quando e como o modelo a chama",
        "Descrição ruim → modelo não chama ou chama errado",
        "Otimização: reescrever descrições baseada em taxa de erro",
        "Ex: 'Busca documentos' → 'Busca na base interna, use para produtos'",
        "Métrica: tool call accuracy",
        "Loop: medir erro → reescrever → re-avaliar",
    ],
    "Tools também são otimizáveis. Às vezes o problema é a tool, não o prompt.",
    31, T,
)

content_slide(
    "Otimização de Topologia",
    [
        "Topologia = como agentes estão conectados",
        "Perguntas: N workers ou N-1? Orchestrator ou peer-to-peer?",
        "Ex: 5 workers → 3 + 2 especializados = mesma qualidade, menor custo",
        "Desafio: espaço de topologias é combinatório",
        "Regra: comece simples (1-2 agentes), adicione só se eval mostrar necessidade",
    ],
    "Otimização mais estrutural. Heurísticas > busca exaustiva.",
    32, T,
)

content_slide(
    "Loop de Otimização (Strategy Evolver)",
    [
        "estratégia atual → gerar variações → avaliar em subset",
        "melhor que atual? sim: substituir · não: manter",
        "substituir → estratégia atual (loop)",
        "HITL: audit periódico para prevenir drift",
        "Padrão: mutação + seleção + retenção = evolução",
        "Diagrama: 12-Diagrams/ETHAGT15/evolution-loop.mmd",
    ],
    "HITL é o que mantém alinhado com objetivo real. Sem audit, drift.",
    33, T,
)

content_slide(
    "Manual vs Automatizado",
    [
        "Manual: lento, não-reprodutível, explora pouco, controlável, transparente",
        "Automatizado: rápido, reprodutível, explora muito, opaco",
        "Recomendação: automatizar para volume, manual para edge cases",
        "Custo de eval: baixo (manual) vs alto (automatizado)",
    ],
    "Não é um ou outro — é ambos. Volume → automático, edge → manual.",
    34, T,
)

exercise_slide(
    "Quando Otimizar é Melhor?",
    [
        "Quando otimizar: volume alto, métrica clara, espaço grande",
        "Quando manual: volume baixo, métrica subjetiva, domínio novo",
        "Como definir a métrica de avaliação?",
        "Discussão aberta (2 min)",
    ],
    "Métrica é o árbitro. Métrica ruim → otimização errada.",
    35, T,
)

section_slide(4, "Auto-Aprendizado Contínuo")

content_slide(
    "Memória de Sucesso e Falha",
    [
        "Memória de sucesso: 'esta estratégia funcionou para este tipo'",
        "Memória de falha: 'esta estratégia falhou para este tipo'",
        "Armazenamento: vetor + metadata (tipo, estratégia, outcome, timestamp)",
        "Recuperação: nova tarefa → buscar similar → aplicar bem-sucedida",
        "Sem memória: agente repete os mesmos erros",
    ],
    "Memória é retrieval em runtime, não fine-tuning. São complementares.",
    37, T,
)

content_slide(
    "Reflexion Sistêmico",
    [
        "Individual: agente reflete sobre SUA execução",
        "Sistêmico: meta-agente reflete sobre o SISTEMA",
        "Perguntas: qual agente falha mais? qual tool é subutilizada?",
        "Output: recomendações de mudança estrutural",
        "Na prática, 2 níveis bastam (não recursão infinita)",
    ],
    "Sistêmico vê padrões que individual não vê. Cuidado com meta-meta-meta.",
    38, T,
)

content_slide(
    "Strategy Evolver",
    [
        "Componente que evolui estratégias ao longo do tempo",
        "Estratégia = config de agente (prompt + tools + parâmetros)",
        "Processo: população → mutação → avaliar → selecionar → substituir",
        "Diferença vs DSPy: evolver no nível de sistema, DSPy no nível de prompt",
        "HITL: auditoria periódica para prevenir drift",
    ],
    "Evolução aplicada a config. HITL é obrigatório — sem audit, drift.",
    39, T,
)

content_slide(
    "Voyager — Aprendendo Skills no Minecraft",
    [
        "Voyager (arXiv:2305.16291): agente que aprende Minecraft sozinho",
        "3 componentes: automatic curriculum, skill library (JS), iterative prompting",
        "Resultado: skills complexas em horas, sem intervenção",
        "Lição: ambiente fechado + feedback automático = viável",
        "Limitação: Minecraft é determinístico e seguro — produção não é",
    ],
    "Caso canônico. Por que é seguro? Ambiente controlado. Produção não é.",
    40, T,
)

content_slide(
    "Quando Esquecer — Drift de Dados",
    [
        "Drift: ambiente muda, conhecimento antigo vira obsoleto",
        "Conceito: o que era correto não é mais",
        "Dados: distribuição de inputs mudou",
        "Sintomas: success rate cai, erros aumentam",
        "Estratégias: TTL, janela deslizante, detecção ativa",
        "Esquecer é tão importante quanto lembrar",
    ],
    "Memória obsoleta é pior que nenhuma. Mais memória ≠ sempre melhor.",
    41, T,
, acronyms="TTL = Time To Live")

exercise_slide(
    "O Que Acontece se o Ambiente Muda?",
    [
        "Resposta: overfitting ao antigo → performance cai",
        "Como detectar drift automaticamente?",
        "Esquecer tudo ou apenas o relevante?",
        "Discussão em duplas (2 min)",
    ],
    "TTL agressivo + detecção ativa é mais robusto que esquecer tudo.",
    42, T,
)

exercise_slide(
    "Exercício — Detectar Drift",
    [
        "Cenário: agente de suporte, 6 meses de memória",
        "Success rate caiu 92% → 78% em 30 dias",
        "Em duplas: propor 3 estratégias para detectar e mitigar",
        "Ex: comparar distribuição de queries, TTL, re-treinar",
        "2 min propor, 1 min compartilhar",
    ],
    "Diagnosticar antes de remediar. Sintoma é claro, causas são várias.",
    43, T,
)

exercise_slide(
    "Auto-Aprendizado Sempre Melhora?",
    [
        "V/F: 'Auto-aprendizado contínuo sempre melhora.'",
        "Resposta: FALSO",
        "Razões: overfitting, drift, jogar a métrica, catastrophic forgetting",
        "Precisa: monitoramento + validação + reset",
    ],
    "Mito para quebrar. Não é 'set and forget'.",
    44, T,
    obj="V/F",
)

content_slide(
    "Lições do Auto-Aprendizado",
    [
        "Memória é poder e risco",
        "Reflexão sistêmica > individual para mudanças estruturais",
        "Strategy evolver = evolução aplicada a configs",
        "Voyager prova que funciona em ambiente controlado",
        "Drift é inevitável — detecção e esquecimento obrigatórios",
        "Não é bala de prata: precisa governança",
    ],
    "Fogo: controlado aquece, sem controle queima. Mesma coisa.",
    45, T,
)

section_slide(5, "Riscos e Governança")

content_slide(
    "Recursão e Auto-Modificação",
    [
        "Meta-agente modifica a si mesmo → loop infinito possível",
        "Cada iteração introduz mudanças cumulativas",
        "Risco: comportamento emergente imprevisível",
        "Mitigação 1: max_iterations (hard limit)",
        "Mitigação 2: diff check (>N% requer aprovação)",
        "Mitigação 3: versionamento + rollback automático",
    ],
    "max_iter limita risco, não inteligência. Sem ele, bug vira loop infinito.",
    47, T,
)

content_slide(
    "Goal Drift",
    [
        "Goal drift: agente otimiza métrica que diverge do objetivo real",
        "Ex: objetivo 'resolver tickets' vs métrica 'tickets fechados/hora'",
        "Agente fecha sem resolver para maximizar métrica",
        "Causa: métrica é proxy imperfeito",
        "Detecção: divergência entre métrica e satisfação do usuário",
        "Mitigação: múltiplas métricas, nunca uma única",
    ],
    "Risco mais insidioso. Toda métrica é proxy. Múltiplas métricas + auditoria.",
    48, T,
)

content_slide(
    "Meta-Governor Pattern",
    [
        "Meta-governor: guardião que avalia mudanças propostas",
        "Separado do meta-agente (não pode auto-aprovar)",
        "Regras: policy-as-code (vetos), sandbox, shadow, canary, rollback",
        "HITL: aprovação humana para mudanças críticas",
        "Governor crítico deve ser determinístico, não LLM",
    ],
    "Separação de poderes é crítica. Lobo não guarda galinheiro.",
    49, T,
)

content_slide(
    "Safety Fences — Camadas de Defesa",
    [
        "Gate 1: policy-as-code (vetos) — viola? bloqueia",
        "Gate 2: sandbox test — falha? bloqueia",
        "Gate 3: shadow run — não é melhor? bloqueia",
        "Gate 4: canary 5% — degrada? rollback",
        "4 camadas = defesa em profundidade",
        "Diagrama: 12-Diagrams/ETHAGT15/safety-fences.mmd",
    ],
    "Blueprint de segurança. Cada camada filtra um tipo de problema.",
    50, T,
)

content_slide(
    "Confiança Incremental",
    [
        "Nível 0: Sandbox (dados sintéticos, sem produção)",
        "Nível 1: Shadow run (paralelo, não afeta usuários)",
        "Nível 2: Canary (5% tráfego, monitorado)",
        "Nível 3: Produção gradual (10→50→100%)",
        "Nível 4: Produção total",
        "Princípio: confiança se ganha, não se dá",
    ],
    "Eval ≠ produção. Shadow e canary pegam problemas que eval não pega.",
    51, T,
)

code_slide(
    "Vetos — Regras Inegociáveis",
    """veto(change) if change.removes_safety_constraint == true
veto(change) if change.cost_increase > 0.10
veto(change) if change.removes_logging == true
veto(change) if change.new_model not in APPROVED_MODELS

# Veto é binário: bloqueia ou não.
# Não há "quase veto" ou "veto com exceção".""",
    "Veto = hard block em código, não warning em prompt. Binariedade = auditável.",
    52, T,
)

exercise_slide(
    "Exercício — Projetar um Meta-Governor",
    [
        "Cenário: meta-agente modifica o próprio system prompt",
        "Em trios: projetar meta-governor com regras de veto",
        "Mínimo 3 regras de veto + pipeline de gates",
        "Quando HITL é obrigatório?",
        "Bônus: e se o governor também for agente?",
        "3 min design, 2 min compartilhar",
    ],
    "Exercício central. Governors críticos devem ser determinísticos.",
    53, T,
    obj="Exercício (5 min)",
)

exercise_slide(
    "Quando um Meta-Agente Se Torna Perigoso?",
    [
        "Pode modificar suas próprias regras de segurança",
        "Pode escolher sua própria métrica",
        "Não tem HITL para mudanças críticas",
        "Pode escalar sem canary",
        "Fórmula: perigo = capacidade × autonomia × falta de oversight",
        "Discussão aberta (2 min)",
    ],
    "Pergunta mais profunda. Objetivo é equilíbrio, não eliminação.",
    54, T,
)

content_slide(
    "Caso — Voyager (Segurança em Ambiente Fechado)",
    [
        "Voyager é seguro porque: ambiente fechado (sandbox)",
        "Feedback determinístico (crafteou o item?)",
        "Sem consequências reais (erro não causa dano)",
        "Skills são código verificável (auditáveis)",
        "Lição: replicar estas condições em produção",
    ],
    "Voyager é o contra-exemplo — quando auto-aprendizado É seguro.",
    55, T,
)

section_slide(6, "Boas Práticas e Anti-Patterns")

content_slide(
    "Boas Práticas — DO",
    [
        "Comece com domínio estreito",
        "Sempre valide (eval antes de deploy)",
        "Use confiança incremental",
        "Implemente meta-governor com vetos desde o dia 1",
        "Versionamento de todas as configs",
        "Monitore drift continuamente",
        "HITL para mudanças críticas",
        "Audite memória/skills periodicamente",
    ],
    "Checklist de pré-voo. Pular um item pode ser fatal.",
    57, T,
)

content_slide(
    "Anti-Patterns — DON'T",
    [
        "Auto-modificação sem limites (recursão)",
        "Métrica única (goal drift garantido)",
        "Deployar sem eval",
        "Deixar meta-agente escolher própria métrica",
        "Sem rollback automático",
        "Confiança total sem canary/shadow",
        "Ignorar drift",
        "Meta-agente sem governor (auto-aprovação)",
    ],
    "Os 3 mais perigosos: sem governor, métrica única, deployar sem eval.",
    58, T,
)

content_slide(
    "Resumo da Aula",
    [
        "Meta-agência = agentes sobre agentes (synthesis/evolution/optimization)",
        "Geração = meta-agente + primitivas + validação",
        "Otimização = DSPy + Promptbreeder + Atlas",
        "Auto-aprendizado = memória + reflexion + evolver + drift detection",
        "Riscos = recursão, goal drift, perda de controle",
        "Governança = governor + safety fences + confiança incremental + vetos",
    ],
    "Sete pontos para levar. Governança é o que separa útil de perigoso.",
    59, T,
)

content_slide(
    "Checklist da Aula",
    [
        "[ ] Definiu meta-agência e 3 estratégias",
        "[ ] Implementou geração de agente com validação",
        "[ ] Aplicou DSPy ou Promptbreeder",
        "[ ] Discutiu auto-aprendizado e drift",
        "[ ] Projetou um meta-governor com vetos",
        "[ ] Identificou riscos e mitigações",
    ],
    "Item mais importante: projetar meta-governor.",
    60, T,
)

exercise_slide(
    "Quiz — Pergunta 1",
    [
        "Diferença fundamental entre agente e meta-agente?",
        "A) Meta-agente usa modelo maior",
        "B) Meta-agente opera sobre agentes",
        "C) Meta-agente tem mais tools",
        "D) Meta-agente é mais rápido",
        "Resposta: B",
    ],
    "É sobre alvo, não capacidade.",
    61, T,
    obj="Quiz",
)

exercise_slide(
    "Quiz — Pergunta 2",
    [
        "O que é goal drift?",
        "A) Agente muda de tarefa",
        "B) Otimiza métrica que diverge do objetivo real",
        "C) Agente esquece o objetivo",
        "D) Agente muda de modelo",
        "Resposta: B",
    ],
    "Métrica sobe, objetivo real cai.",
    62, T,
    obj="Quiz",
)

exercise_slide(
    "Quiz — Pergunta 3",
    [
        "Primeira camada de defesa em safety fences?",
        "A) Canary deployment",
        "B) Shadow run",
        "C) Policy-as-code (vetos)",
        "D) Rollback automático",
        "Resposta: C",
    ],
    "Policy-as-code é o gate mais barato e rápido.",
    63, T,
    obj="Quiz",
)

exercise_slide(
    "Perguntas para Discussão",
    [
        "1. Quando otimizar prompts é melhor que reescrever manualmente?",
        "2. Defina goal drift e proponha detecção automática",
        "3. Por que meta-governor é necessário?",
        "4. Em que ponto um meta-agente se torna perigoso?",
    ],
    "Usar o tempo restante para discussão em grupo.",
    64, T,
)

content_slide(
    "Conexão com Próximos Módulos",
    [
        "ETHAGT16 — Sociedades de Agentes: meta-agentes em ecossistemas",
        "ETHAGT90 — Capstone: meta-agente como orquestração",
        "ETHAGT14 — Escalabilidade: custo aplicado a meta-agentes",
        "ETHAGT12 — AgentOps: avaliação e observabilidade",
    ],
    "ETHAGT15 é pré-requisito de ETHAGT16 e ETHAGT90.",
    65, T,
, acronyms="AgentOps = Agent Operations")

content_slide(
    "Referências, Projeto e Labs",
    [
        "Obrigatório: DSPy (arXiv:2310.03714)",
        "Obrigatório: Promptbreeder (arXiv:2309.16797)",
        "Recomendado: Meta-Prompting (arXiv:2311.11402)",
        "Recomendado: Voyager (arXiv:2305.16291)",
        "Recomendado: AI Scientist (arXiv:2408.06292)",
        "Projeto: otimização automática (DSPy) + benchmark",
        "Lab 1 (4h): Agente que escreve agente",
        "Próxima aula: ETHAGT16 — Sociedades de Agentes",
    ],
    "Lembrar prazos: Lab 1 (1 semana), Projeto (2 semanas).",
    66, T,
)

s = title_slide(
    "Q&A / Encerramento",
    "Próxima aula: ETHAGT16 — Sociedades de Agentes",
    CODE,
)
add_notes(s, "Se nenhuma pergunta: 'Qual parte foi menos clara?' ou 'Qual técnica vão aplicar?'")

out = os.path.join(os.path.dirname(__file__), "ETHAGT15-slides.pptx")
prs.save(out)
print(f"PPTX gerado: {out} ({len(prs.slides)} slides)")
