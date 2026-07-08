#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT11: Event-Driven Agents & Workflow Orchestration
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj=""):
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT11"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

T = 76

s = title_slide(
    "Event-Driven Agents & Workflow Orchestration",
    "Universidade Etho · Especialização em Programação Agêntica\nFase C — Orquestração & Produção · 25 h",
    "ETHAGT11"
)
add_notes(s, "📖 Bem-vindos à Fase C. Hoje vamos construir agentes que sobrevivem a falhas, rodem por dias, e coordenem via eventos. Se seus agentes morriam no primeiro crash, isso muda hoje.\n💡 Analogia: email (assíncrono, resiliente) vs ligação telefônica (síncrono, ambos presentes). Hoje vocês aprendem o 'email' da orquestração.\n❓ 'Quantos já tiveram agente que caiu no meio de processamento longo?'\n⚠️ Alunos acham que Kafka resolve tudo. Kafka é uma peça — Temporal é o complemento.\n➡️ Objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: construir sistemas de agentes orientados a eventos com orquestração robusta",
    "",
    "Objetivos específicos:",
    "1. Diferenciar orquestração síncrona de event-driven assíncrona",
    "2. Aplicar filas (Kafka, RabbitMQ, NATS) para coordenação de agentes",
    "3. Implementar durable execution (Temporal, Prefect) para workflows longos",
    "4. Lidar com falhas, retries, idempotência, compensação",
    "5. Avaliar escala, ordering, exactly-once vs at-least-once",
], "📖 Cada objetivo é mensurável: diferenciar, aplicar, implementar, lidar, avaliar.\n❓ 'Qual objetivo vocês acham mais desafiador?' (geralmente #3 ou #5)\n⚠️ 'Aprender Kafka' não é 'aprender event-driven'. Kafka é ferramenta; event-driven é paradigma.\n➡️ Competências.", 2, T, "Objetivos mensuráveis")

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → A (Avançado)",
    "C2 Multi-Agent Systems → A (Avançado)",
    "C3 MCP & Tool Use → B (Básico)",
    "C4 Agent Memory → B (Básico)",
    "C5 AgentOps & Avaliação → I (Intermediário)",
    "",
    "C1 e C2 atingem Avançado: arquitetar multi-agente orientado a eventos em produção",
    "C5 atinge Intermediário: tracing distribuído (aprofundamento em ETHAGT12)",
], "📖 C1 e C2 chegam ao Avançado — vocês vão arquitetar sistemas multi-agente orientados a eventos em produção.\n💡 Analogia: vocês já dirigem na cidade (Básico). Agora vão dirigir em tráfego intenso com GPS e sistemas de segurança.\n⚠️ 'Avançado' não é 'domínio total' — é capacidade de arquitetar, justificar e operar.\n➡️ Agenda.", 3, T, "Conectar ao Framework Etho")

s = content_slide("Agenda da Aula", [
    "Bloco 1 (45 min):",
    "  Abertura (8 min) — objetivos, motivação (crash no ticket 5.000), contexto",
    "  Por que Event-Driven (9 min) — limites síncrono, paradigma, trade-offs",
    "  Mensageria (15 min) — Kafka, RabbitMQ, NATS, CQRS, saga, CloudEvents, DEMO",
    "  Workflows (10 min) — Temporal, Prefect, Airflow, decisão",
    "  Intervalo (3 min)",
    "",
    "Bloco 2 (45 min):",
    "  Durable Execution (12 min) — crashes, long-running, HITL, replays, DEMO",
    "  Resiliência & Produção (12 min) — retries, idempotência, saga, circuit breakers",
    "  Fechamento (18 min) — boas práticas, caso de estudo, quiz, projeto, Q&A",
], "📖 Dois blocos. Primeiro: fundamentos + mensageria + workflows. Segundo: durable execution + resiliência + fechamento. Duas DEMOS ao vivo.\n➡️ Por que estamos aqui?", 4, T, "Roteiro da aula")

s = content_slide("O Crash no Ticket 5.000", [
    "Cenário: agente processando 10.000 tickets de suporte em série",
    "",
    "O processo cai no ticket 5.000:",
    "  • Sem durability: recomeça do ticket 1 (horas perdidas, custo desperdiçado)",
    "  • Com durability: retoma do ticket 5.001 (zero perda)",
    "",
    "Pergunta: Qual o pior crash que você já viu em produção?",
], "📖 O problema central da aula. Agente síncrono processando 10.000 tickets — se cai no 5.000, perde metade. Sem durability, restart recomeça do zero. Com durable execution, retoma do último checkpoint.\n💡 Analogia: escrever 100 páginas sem salvar vs com auto-save.\n❓ 'Qual o pior crash que vocês já viram em produção?' (2-3 alunos compartilharem)\n⚠️ Alunos acham que 'retry resolve'. Retry sem durability recomeça do início.\n➡️ Como chegamos aqui?", 5, T, "Criar tensão — agentes síncronos perdem tudo em crash")

s = content_slide("Evolução para Event-Driven", [
    "Linha do tempo:",
    "  2011: Kafka na LinkedIn (processamento de logs em escala)",
    "  2016: Predecessores do Temporal (Cadence, SWF)",
    "  2019: Temporal open source",
    "  2023: Agentes de longa duração entram em pauta",
    "  2024: Durable execution para LLM agents amadurece",
    "",
    "Confluência: agentes longos + custo de recomputação + necessidade de resiliência",
    "Fundamento: Kreps 'The Log' · Padrão: CloudEvents (CNCF)",
], "📖 Event-driven não é novo — Kafka existe desde 2011. Mas a combinação com agentes de LLM é recente: agentes executam por minutos/horas, custo de recomputação é alto, e HITL exige pausas duráveis. Tudo convergiu em 2024.\n💡 Analogia: avião — princípios existiam há séculos, mas precisou de motor leve + materiais + controle.\n⚠️ Alunos acham que event-driven é 'modinha'. O conceito tem 15 anos.\n➡️ Fundamentos.", 6, T, "Confluência histórica")

s = section_slide(1, "Por que Event-Driven")
add_notes(s, "Início do bloco de fundamentos. Antes de ferramentas, precisamos entender o PORQUÊ.\n➡️ Limites do síncrono.")

s = content_slide("Limites da Orquestração Síncrona", [
    "Acoplamento temporal: caller e callee devem estar vivos ao mesmo tempo",
    "Falta de resiliência: falha no meio = perda de progresso",
    "Difícil escalar: bottleneck no orquestrador síncrono",
    "Custo de recomputação: cada retry recomeça do zero",
    "",
    "Exemplo: agente que processa 10 tools em série",
    "  → se cai na tool 7, perde tudo",
], "📖 Síncrono funciona para tarefas curtas (segundos). Em pipelines longos (minutos, horas), cada limite se manifesta: acoplamento temporal, falta de resiliência, bottleneck de escala, custo de recomputação.\n💡 Analogia: linha de montagem que para toda se uma estação para.\n❓ 'Vocês têm pipeline síncrono hoje que já quebrou por essas razões?'\n⚠️ Retry sem durability recomeça do início. Retry sem idempotência duplica.\n➡️ Qual é a alternativa?", 8, T, "Onde síncrono quebra")

s = content_slide("O que é Event-Driven", [
    "Componentes se comunicam via eventos (mensagens assíncronas)",
    "Produtor publica → broker roteia → consumidor processa",
    "",
    "Desacoplamento temporal: não precisam coincidir no tempo",
    "Desacoplamento espacial: não precisam se conhecer",
    "",
    "Evento = fato que ocorreu (passado), não comando (futuro)",
    "  'PedidoCriado' (evento) vs 'CriarPedido' (comando)",
], "📖 Definição canônica: componentes se comunicam via eventos assíncronos. Evento é FATO (passado); comando é INSTRUÇÃO (futuro). Eventos são imutáveis; comandos podem falhar.\n💡 Analogia: jornal notifica 'Choveu ontem' (evento). Bilhete diz 'Regue as plantas' (comando).\n❓ 'Em seus sistemas, usam eventos ou comandos?'\n⚠️ Modelar tudo como comando perde o desacoplamento do event-driven.\n➡️ Benefícios concretos.", 9, T, "Definir paradigma event-driven")

s = content_slide("Benefícios do Event-Driven", [
    "Desacoplamento: componentes evoluem independentemente",
    "Escala horizontal: consumidores paralelos por partição",
    "Resiliência: broker persiste eventos; consumers recuperam do offset",
    "Replay: reprocessar histórico de eventos (corrigir bugs retroativamente)",
    "Backpressure natural: consumer puxa no seu ritmo",
], "📖 5 benefícios concretos. Desacoplamento (novo consumer sem mudar produtor). Escala (adicionar consumers ao group). Resiliência (broker persiste). Replay (reprocessar com nova lógica). Backpressure (consumer puxa, não é empurrado).\n💡 Analogia: biblioteca. Cada leitor lê independente. Adicione mais leitores. O livro fica na prateleira.\n⚠️ Replay não é 'recalcular tudo' — é reprocessar eventos com nova lógica.\n➡️ Mas não é grátis.", 10, T, "Ganhos concretos")

s = content_slide("Trade-offs do Event-Driven", [
    "Complexidade: mais componentes móveis (broker, consumers, DLQ)",
    "Eventual consistency: estado converge apenas após processar eventos",
    "Debugging distribuído: trace espalhado por serviços",
    "Ordering: ordem global é cara ou impossível",
    "",
    "'Event-driven não é grátis — é uma troca consciente'",
], "📖 4 trade-offs reais. Complexidade operacional. Eventual consistency (delay até convergir). Debugging distribuído. Ordering limitado.\n💡 Analogia: trocar reunião presencial (síncrono) por email (assíncrono, flexível, mas perde certeza de ordem).\n❓ 'Qual trade-off mais preocupa em produção?'\n⚠️ Adotar event-driven sem custo operacional = pesadelo.\n➡️ Quando síncrono, quando assíncrono?", 11, T, "Honestidade sobre o preço")

s = comparison_slide("Síncrono vs Assíncrono",
    "Síncrono",
    ["Latência: baixa (ms)", "Acoplamento: alto", "Resiliência: frágil", "Escala: vertical", "Complexidade: baixa", "Consistência: imediata"],
    "Assíncrono (Event-Driven)",
    ["Latência: alta (segundos+)", "Acoplamento: baixo", "Resiliência: resiliente", "Escala: horizontal", "Complexidade: alta", "Consistência: eventual"],
    "📖 Nenhum paradigma é universalmente melhor. Síncrono para UX (usuário esperando). Assíncrono para processamento pesado. Híbrido é o mais comum: API síncrona dispara processamento assíncrono.\n💡 Analogia: fast-food (síncrono) vs restaurante com garçom (assíncrono) vs iFood (híbrido — pede, recebe número, consulta status).\n❓ 'Suas tools de agente são síncronas ou assíncronas?'\n⚠️ Tentar fazer tudo assíncrono piora UX.\n➡️ Vamos praticar.", 12, T, "Quando cada paradigma")

s = exercise_slide("Exercício: Síncrono ou Assíncrono?", [
    "Para cada cenário, vote: Síncrono (S) ou Assíncrono (A)?",
    "",
    "1. Responder pergunta do usuário em chat",
    "2. Processar 10.000 documentos",
    "3. Buscar em RAG e responder",
    "4. Coordenar 3 agentes em pipeline longo",
    "",
    "Gabarito: 1-S, 2-A, 3-S(timeout), 4-A",
], "📖 Respostas: (1) Síncrono — UX espera. (2) Assíncrono — massa, longa duração. (3) Síncrono com timeout. (4) Assíncrono — coordenação complexa.\n❓ Votar em cada cenário. Anotar distribuição.\n⚠️ Responder 'assíncrono' para tudo é erro comum — UX NÃO pode ser assíncrona.\n➡️ Motor do event-driven: mensageria.", 13, T, "Votação rápida")

s = section_slide(2, "Mensageria: Kafka, RabbitMQ, NATS")
add_notes(s, "O coração do event-driven é a mensageria. Vamos cobrir os três brokers mais relevantes.\n➡️ Conceito fundacional: o Log.")

s = content_slide("O Log como Abstração Central", [
    "'O Log' (Kreps, LinkedIn) — append-only, ordenado, imutável",
    "Tópico = log particionado; cada partição é uma sequência ordenada",
    "Consumer mantém offset (posição no log)",
    "Replay: voltar o offset e reprocessar",
    "",
    "Fonte: Kreps, 'The Log' (LinkedIn Engineering)",
], "📖 Jay Kreps escreveu 'The Log' — essay seminal. O log é a abstração mais poderosa em sistemas distribuídos: append-only, ordenado, imutável. Consumer mantém offset. Replay = voltar o offset.\n💡 Analogia: diário de bordo do navio. Append-only, ordenado, imutável. Capitão lê a partir de onde parou (offset).\n❓ 'Por que o log é imutável?'\n⚠️ Log aqui não é 'log de debug' — é estrutura de dados distribuída, fonte de verdade.\n➡️ Kafka é a implementação mais famosa.", 15, T, "Conceito fundacional")

s = content_slide("Kafka: Arquitetura", [
    "Topic → N partitions → cada partition é um log imutável",
    "Producer escreve em partition (por key ou round-robin)",
    "Consumer Group: um consumer por partition",
    "Replicação: leader + followers por partition",
    "Retenção por tempo ou tamanho",
], "📖 Kafka organiza em tópicos. Cada tópico é dividido em partições — logs independentes. Produtores escolhem partição (key ou round-robin). Consumers em groups — cada partição vai para um consumer. Replicação garante durabilidade.\n💡 Analogia: biblioteca com estantes (partições). Bibliotecários (consumers) leem estantes diferentes em paralelo.\n⚠️ Partição ≠ fila. Partição é log (mensagens ficam); fila é consumida (somem).\n➡️ Particionamento e ordering.", 16, T, "Anatomia do Kafka")

s = content_slide("Kafka: Particionamento e Ordering", [
    "Ordering garantido DENTRO de uma partição, não global",
    "Key = agent_id → mesmo agente → mesma partição → ordering preservado",
    "Sem key → round-robin → sem ordering garantido",
    "",
    "Como garantir? hash(key) % num_partitions",
    "",
    "Pergunta: Como garantir ordering de mensagens do mesmo agente?",
], "📖 Regra mais importante do Kafka: ordering é por partição, não global. Para ordering do 'agente 42', use agent_id como key. hash(key) % num_partitions → mesma partição.\n💡 Analogia: filas de banco. Múltiplos caixas em paralelo. Para suas transações em ordem, vá sempre ao mesmo caixa (mesma key = mesma partição).\n❓ 'Como garantir ordering do mesmo agente?' (resposta: hash da key)\n⚠️ Ordering global não existe em Kafka. Se precisa, use 1 partição (perde paralelismo).\n➡️ Escalar consumidores.", 17, T, "Ordering por partição")

s = content_slide("Kafka: Consumer Groups e Escala Horizontal", [
    "Consumer Group: conjunto de consumers que dividem as partições",
    "1 consumer por partição dentro do grupo",
    "Adicionar consumers → rebalanceamento → mais paralelismo",
    "Limite: # consumers ≤ # partições",
    "",
    "Offset commit: at-least-once (após processar) vs at-most-once (antes)",
], "📖 Consumer Group é como Kafka escala. Grupo divide partições — cada uma vai para um consumer. Se 3 partições e 4 consumers, o 4º fica ocioso. Para escalar, adicione partições. Offset commit define semântica: commit após = at-least-once; commit antes = at-most-once.\n💡 Analogia: equipe de entregadores dividindo zonas da cidade.\n⚠️ Adicionar consumers sem partições = sem ganho. Limite é o número de partições.\n➡️ Alternativa: RabbitMQ.", 18, T, "Escala horizontal de consumers")

s = content_slide("RabbitMQ: Filas, Exchanges, Routing", [
    "Exchange → binding → queue",
    "Tipos de exchange: direct, topic, fanout, headers",
    "Routing rico: mensagens vão para filas baseadas em routing key",
    "ACK manual: consumer ack após processar",
    "Dead letter exchange para mensagens que falham",
], "📖 RabbitMQ segue AMQP. Diferença vs Kafka: RabbitMQ é orientado a filas (mensagem consumida some); Kafka é log (mensagem fica até retenção expirar). RabbitMQ brilha em routing rico: exchanges decidem para quais filas a mensagem vai.\n💡 Analogia: Kafka é jornal (todos leem a mesma edição, fica arquivada). RabbitMQ é central de taxistas (chamada vai para taxista específico).\n❓ 'Em que cenário RabbitMQ é melhor que Kafka?'\n⚠️ Usar Kafka quando precisavam de RabbitMQ (routing) é erro comum.\n➡️ Para leveza: NATS.", 19, T, "Modelo AMQP")

s = content_slide("NATS: JetStream, Leveza, Performance", [
    "NATS: mensagem simples, baixíssima latência, single binary",
    "JetStream: persistência, streams, consumers duráveis",
    "Modelos: at-least-once, exactly-once (com dedup window)",
    "Ideal para edge, IoT, sistemas embarcados",
    "Throughput: NATS > RabbitMQ > Kafka (mensagens pequenas)",
], "📖 NATS é o 'canivete suíço' da mensageria. Single binary, latência sub-milissegundo. JetStream adiciona persistência. Ideal para edge/IOT. Trade-off: menos features que Kafka/RabbitMQ.\n💡 Analogia: Kafka é caminhão (carrega muito, pesado). RabbitMQ é van (flexível). NATS é moto (rápido, leve).\n⚠️ Subestimar NATS por ser menos famoso é erro. Em muitos casos, resolve com 1/10 da complexidade do Kafka.\n➡️ Comparar os três.", 20, T, "Alternativa leve")

s = comparison_slide("Kafka vs RabbitMQ vs NATS",
    "Kafka (Log)",
    ["Ordering: por partição", "Retenção: dias (configurável)", "Throughput: altíssimo", "Latência: média", "Complexidade: alta", "Melhor para: volume + replay"],
    "RabbitMQ / NATS",
    ["RabbitMQ: routing rico, baixa latência", "RabbitMQ: ACK-based, médio throughput", "NATS: baixíssima latência, single binary", "NATS: stream-based, simplicidade", "Melhor routing: RabbitMQ", "Melhor leveza: NATS"],
    "📖 Kafka ganha em throughput e retenção. RabbitMQ em routing. NATS em simplicidade e latência. A escolha depende do problema.\n💡 Analogia: trem de carga (Kafka), táxi (RabbitMQ), bicicleta (NATS).\n⚠️ Começar com Kafka 'porque a empresa usa' sem avaliar é over-engineering.\n➡️ Critério prático.", 21, T, "Trade-offs lado a lado")

s = content_slide("Quando Usar Cada Um", [
    "Kafka: volume altíssimo, replay necessário, log de eventos, event sourcing",
    "RabbitMQ: routing complexo, filas de trabalho, request/reply, legados",
    "NATS: baixa latência, edge, simplicidade operacional, mensagens pequenas",
    "",
    "Regra: comece com o mais simples que resolve; Kafka é overkill para muitos",
], "📖 Regra de ouro: comece simples. Se RabbitMQ resolve, não use Kafka. Se NATS resolve, não use RabbitMQ. Kafka é mais poderoso mas mais complexo de operar.\n💡 Analogia: não use martelo pneumático para pregar quadro.\n⚠️ FOMO tecnológico — querer Kafka porque é 'moderno' sem necessidade.\n➡️ Padrões: CQRS e Saga.", 22, T, "Critério de escolha")

s = content_slide("Padrão: CQRS", [
    "Command Query Responsibility Segregation",
    "Write side (commands) → eventos → read side (queries) → views materializadas",
    "Para agentes: comandos geram eventos; agentes de consulta leem views",
    "",
    "Fonte: Microsoft Cloud Design Patterns",
], "📖 CQRS separa escrita (commands) de leitura (queries). Write side otimizado para consistência (eventos); read side otimizado para consulta (views materializadas).\n💡 Analogia: biblioteca. Catalogador (write) organiza; índices (read) são views otimizadas para busca.\n⚠️ CQRS adiciona complexidade — só use quando read e write têm perfis muito diferentes.\n➡️ Saga.", 23, T, "Separação comando/consulta")

s = content_slide("Padrão: Saga (Visão Geral)", [
    "Saga: sequência de transações locais com compensação",
    "Se passo N falha, executam compensações dos passos 1..N-1",
    "Dois estilos: orquestrada (central) vs coreografada (eventos)",
    "",
    "Aprofundamento na Seção F (Slide 51)",
], "📖 Saga resolve transações distribuídas sem 2PC. Transações locais com compensação. Não é rollback perfeito — é melhor esforço. Aprofundamos na Seção F com caso concreto.\n💡 Analogia: planejar viagem. Se voo cancela, cancela hotel e devolve carro (compensações).\n➡️ Padronização de eventos.", 24, T, "Transação compensatória")

s = content_slide("CloudEvents: Padronização de Eventos", [
    "CloudEvents: spec CNCF para descrição de eventos",
    "Campos: id, source, type, specversion, time, data, subject",
    "Interoperabilidade: mesmo formato entre Kafka, RabbitMQ, NATS, HTTP",
    "Para agentes: tool outputs como CloudEvents padroniza observabilidade",
], "📖 CloudEvents é padrão CNCF. Em vez de cada sistema ter seu formato, todos seguem o mesmo schema. Para agentes: tool outputs como CloudEvents = interoperabilidade + observabilidade.\n💡 Analogia: padrão de tomada elétrica. Um formato, qualquer aparelho.\n⚠️ Não é 'mais um formato' — é O formato que unifica.\n➡️ DEMO.", 25, T, "Padrão CNCF")

s = code_slide("DEMO: Agente com Kafka",
    "Referência: 05-Labs/ETHAGT11/Lab1-Agente-Kafka\n\n"
    "# Agente A publica tarefa\n"
    "producer.produce('agent-tasks',\n"
    "    key=agent_id,           # ordering por key\n"
    "    value=json.dumps(task))\n\n"
    "# Agente B consome\n"
    "def consume(msg):\n"
    "    result = agent.process(msg.value())\n"
    "    producer.produce('agent-results',\n"
    "        key=msg.key(),\n"
    "        value=json.dumps(result))\n"
    "    consumer.commit(msg)    # at-least-once",
    "📖 Demo ao vivo. Agente A publica no tópico agent-tasks com key=agent_id. Agente B consome, processa (chama LLM), publica resultado. Destacar: (1) ordering por key; (2) resiliência — matar consumer B e reiniciar, recupera do último offset.\n💡 Analogia: dois colegas via Slack. Um posta em #tasks; outro lê, faz, posta em #results.\n⚠️ Se demo falhar, screenshots/logs prontos. Não improvisar.\n➡️ Pergunta da demo.", 26, T, "Demo ao vivo")

s = exercise_slide("Pergunta da DEMO", [
    "Discutam em duplas (2 min):",
    "",
    "1. O que acontece se Agente B processar duas vezes a mesma mensagem?",
    "2. Como garantir que o resultado não seja duplicado?",
    "",
    "Resposta: idempotência (chaves de idempotência)",
    "At-least-once delivery exige consumidor idempotente",
], "📖 Resposta: idempotência. Kafka é at-least-once — mensagens podem chegar mais de uma vez. Consumer precisa ser idempotente: processar 2x = processar 1x. Solução: chaves de idempotência (Slide 49). Motiva a Seção F.\n💡 Analogia: garçom que anota pedido 2x. Cozinha idempotente verifica e faz só 1 prato.\n⚠️ Kafka não é exactly-once por default — é at-least-once.\n➡️ Workflow engines.", 27, T, "Discussão em duplas")

s = section_slide(3, "Orquestração de Workflows")
add_notes(s, "Vimos mensageria — como agentes se comunicam. Agora: workflow engines — como orquestrar execução com estado durável.\n➡️ Distinção importante.")

s = content_slide("Workflow Engine vs Agentes que Decidem", [
    "Workflow engine: passos predefinidos, determinísticos, duráveis",
    "Agente que decide: rota dinâmica, flexível, menos previsível",
    "Híbrido: agente decide dentro de activity; workflow garante durability",
    "",
    "Pergunta: Você orquestra via código ou via agente supervisor?",
], "📖 Não é 'workflow OU agente' — é 'workflow E agente'. Workflow (Temporal) lida com orquestração (sequência, retries, durability, HITL). Agente lida com decisões (qual tool, como interpretar). Híbrido é o mais poderoso.\n💡 Analogia: workflow é trilho do trem (determinístico). Agente é maquinista (decide velocidade, para em emergência).\n❓ 'Você orquestra via código ou via agente supervisor?'\n⚠️ Agente não é confiável para durability — pode alucinar, esquecer passos.\n➡️ Temporal.", 29, T, "São complementares")

s = content_slide("Temporal: Conceitos", [
    "Workflow: função que descreve a lógica (determinística, sem I/O direto)",
    "Activity: unidade de trabalho (chamada de API, tool, I/O)",
    "Worker: processo que executa workflows e activities",
    "Temporal Server: armazena estado e histórico de execução",
    "Task Queue: workflow/activity tasks esperando workers",
], "📖 5 conceitos-chave. Workflow é DETERMINÍSTICA (sem I/O, random, datetime). Activity é onde trabalho real acontece (API, tools). Worker executa. Server armazena histórico. Task Queue é onde tasks esperam.\n💡 Analogia: workflow é receita (determinística). Activity é ação física (cortar, cozinhar). Worker é cozinheiro. Server é gerente que anota.\n⚠️ I/O no workflow quebra replay. I/O sempre em activities.\n➡️ Durable execution.", 30, T, "Workflow, Activity, Worker, Server")

s = content_slide("Temporal: Durable Execution", [
    "Estado da execução é persistido a cada step",
    "Se o worker cai, outro worker retoma do último step completado",
    "Não há perda de progresso — o histórico é a fonte de verdade",
    "Replay: re-executar o workflow usando o histórico gravado",
    "",
    "Fonte: Temporal Durable Execution primer",
], "📖 A mágica do Temporal. Cada activity é persistida. Se worker cai, outro pega o histórico, faz replay (sem re-executar activities), e continua do ponto onde parou. Zero perda.\n💡 Analogia: videogame com auto-save a cada nível. Console desliga? Recomeça do último nível salvo.\n⚠️ Durable execution ≠ database. É capacidade de retomar execução de código do ponto exato — loops, variáveis, stack.\n➡️ Workflow é código comum.", 31, T, "O conceito central")

s = code_slide("Temporal: Código como Workflow",
    "import temporalio.workflow as wf\n\n"
    "@wf.workflow\n"
    "class ProcessTicketsWorkflow:\n"
    "    async def run(self, tickets: list):\n"
    "        for ticket in tickets:\n"
    "            # activity = checkpoint automático\n"
    "            result = await wf.execute_activity(\n"
    "                process_ticket, ticket,\n"
    "                retry_policy=RetryPolicy(\n"
    "                    maximum_attempts=5))\n"
    "\n"
    "        # timer durável (sobrevive a restart)\n"
    "        await wf.sleep(seconds=3600)\n\n"
    "        # signal: HITL\n"
    "        await wf.wait_condition(...)",
    "📖 Workflow é função Python comum. await activity(...). Não pode usar datetime.now() (usa workflow.now()), random() (usa workflow.random()), time.sleep() (usa workflow.sleep()).\n💡 Analogia: receita não pode dizer 'espere o relógio marcar 14h' (não-determinístico). Diz 'asse por 30 min' (determinístico).\n⚠️ time.sleep() no workflow trava o worker sem durabilidade. Use workflow.sleep().\n➡️ Alternativas mais simples.", 32, T, "Workflow é código Python")

s = code_slide("Prefect: Python Puro, Mais Simples",
    "from prefect import flow, task\n\n"
    "@task(retries=3, retry_delay_seconds=60)\n"
    "def process_ticket(ticket):\n"
    "    return agent.process(ticket)\n\n"
    "@flow\n"
    "def process_tickets(tickets):\n"
    "    for t in tickets:\n"
    "        process_ticket(t)\n\n"
    "# Menos verboso que Temporal\n"
    "# Menos durável: estado no Prefect server\n"
    "# Bom para data pipelines e protótipos",
    "📖 Prefect é Python-nativo. @flow e @task decorators. Muito mais simples que Temporal. Trade-off: durabilidade menor — se Prefect server cai, pode perder estado. Temporal sobrevive a qualquer falha.\n💡 Analogia: Prefect é Apple Watch (simples, out-of-the-box). Temporal é equipamento médico profissional (potente, precisa expertise).\n⚠️ Prefect para protótipos; Temporal para produção crítica.\n➡️ Airflow.", 33, T, "Alternativa simples")

s = content_slide("Airflow: Agendamento vs Orquestração", [
    "Airflow: DAGs agendadas (batch, cron-like)",
    "NÃO é orquestração em tempo real nem durable execution",
    "Ótimo para ETL e pipelines de dados agendados",
    "Para agentes de longa duração com HITL: Temporal é mais adequado",
], "📖 Airflow é excelente para DAGs agendadas (toda meia-noite, rode ETL). Mas NÃO é durable execution. Se task falha, retenta o task — não retoma execução parcial. Sem HITL em tempo real. Para agentes longos com HITL, Temporal.\n💡 Analogia: Airflow é despertador (dispara em horário). Temporal é assistente pessoal (executa, pausa, retoma).\n⚠️ Usar Airflow para orquestração de agentes funciona em casos simples, dói em long-running e HITL.\n➡️ Comparar os três.", 34, T, "Posicionamento correto")

s = comparison_slide("Temporal vs Prefect vs Airflow",
    "Temporal (Durabilidade Máxima)",
    ["Durable execution: máximo", "HITL: nativo (signals)", "Timers duráveis: sim", "Replay/debug: sim", "Complexidade: média", "Para: agentes longos, HITL"],
    "Prefect / Airflow",
    ["Prefect: durabilidade média, simplicidade alta", "Prefect: HITL limitado", "Airflow: durabilidade baixa", "Airflow: agendamento batch excelente", "Airflow: sem HITL real-time", "Para: ETL, protótipos"],
    "📖 Temporal ganha em durability, HITL e replay. Prefect em simplicidade. Airflow em agendamento batch. Para agentes de longa duração com HITL, Temporal.\n💡 Analogia: Temporal é tanque (resistente). Prefect é pickup (versátil). Airflow é ônibus (agenda rotas).\n⚠️ Escolher por familiaridade ('já sei Airflow') dói em long-running.\n➡️ Decisão final.", 35, T, "Trade-offs sistematizados")

s = content_slide("Código vs Agente Supervisor", [
    "Workflow engine (Temporal): passos previsíveis, determinismo, auditoria",
    "Agente supervisor: rota dinâmica, decisões contextuais, flexibilidade",
    "Híbrido: agente dentro de activity Temporal — flexibilidade + durability",
    "",
    "Regra: orquestração determinística no workflow; decisões inteligentes no agente",
], "📖 Regra de ouro: orquestração (sequência, retries, HITL) no workflow. Decisão inteligente (qual tool, interpretar) no agente. Híbrido é padrão de produção: workflow chama agente como activity.\n💡 Analogia: workflow é processo (checklist de voo). Agente é piloto (decisões contextuais). Auditoria quer checklist; flexibilidade quer piloto.\n⚠️ Fazer tudo no agente = sem reliability/auditoria.\n➡️ Durable execution para agentes.", 36, T, "Decisão arquitetural")

s = section_slide(4, "Durable Execution para Agentes")
add_notes(s, "O coração da aula. Durable execution = sobreviver a crashes, rodar por dias, pausar para HITL.\n➡️ Sobreviver a crashes.")

s = content_slide("Sobreviver a Crashes", [
    "Estado persistido a cada step (não a cada minuto — a cada chamada)",
    "Kill do processo → novo worker pega histórico → replay → continua",
    "Sem durability = recomeça do zero; com durability = retoma do checkpoint",
    "Para agentes: cada tool call é um activity → checkpoint automático",
], "📖 Diferença drástica. Sem: crash no 5.000 = recomeça do 1. Com: crash no 5.000 = retoma do 5.001. Cada activity é checkpoint. Para agentes, cada tool call (cara) é persistida.\n💡 Analogia: jogo com auto-save a cada ação. Sem save, desliga = recomeça. Com save, perde no máximo última ação.\n❓ 'Já calcularam quanto custa restart do zero?'\n⚠️ Checkpoint por step não é lento — é operação de DB (ms). Reprocessar é ordens de magnitude mais caro.\n➡️ Long-running.", 38, T, "Checkpoint a cada step")

s = content_slide("Long-Running Agents (horas, dias)", [
    "Agente tradicional: segundos a minutos (context window limita)",
    "Agente durável: horas, dias, semanas",
    "Padrão: workflow chama agente em loop, persistindo estado entre iterações",
    "Casos: pesquisa longitudinal, monitoramento contínuo, processamento em massa",
    "Cuidado: custo acumulado de tokens — orçamento por execução",
], "📖 Sem durable execution, agentes limitam-se a segundos/minutos. Com durable, rodam por dias: workflow persiste estado entre iterações, agente 'acorda' com contexto fresco. Casos: pesquisa longitudinal, massa, monitoramento.\n💡 Analogia: trabalhador que vai embora às 18h e volta às 9h. Sem 'memória' (state), esquece tudo. Com caderno (checkpoint), retoma.\n⚠️ Esquecer orçamento = agente rodando dias pode custar milhares.\n➡️ HITL.", 39, T, "Expandir horizonte de duração")

s = content_slide("Human-in-the-Loop via Timers e Signals", [
    "Signal: mensagem externa entregue ao workflow em execução",
    "Timer: workflow.sleep(dias) — pausa durável, não consome recursos",
    "Padrão HITL: workflow pausa em await signal → humano aprova → continua",
    "Timeout: se humano não responde em X → caminho alternativo",
    "Para agentes: agente propõe → workflow pausa → humano aprova → executa",
], "📖 HITL em workflows é elegante. Workflow pausa em await signal — fica pausado por DIAS sem consumir recursos (não é thread dormindo; é estado persistido). Humano aprova (envia signal) → workflow acorda. Timeout: se ninguém aprova, caminho alternativo.\n💡 Analogia: email esperando resposta. Fica na caixa sem consumir recursos. Responde → processo continua.\n❓ 'Em quais ações exigiria HITL?'\n⚠️ HITL com polling (loop checando a cada segundo) consome recursos. Signals são event-driven.\n➡️ Replays.", 40, T, "Pausa durável sem consumo")

s = content_slide("Replays e Debug Temporal", [
    "Replay: re-executar workflow usando histórico gravado (sem re-executar activities)",
    "Debug temporal: voltar no tempo, inspecionar estado em qualquer step",
    "'Time travel debugging' — ver o que o agente decidiu em cada ponto",
    "Ferramenta: Temporal Web UI mostra histórico completo",
], "📖 Replay é superpotência. Workflow falhou? Não precisa reproduzir (caro). Replay: Temporal re-executa código do workflow usando resultados gravados das activities. Inspecione variáveis em cada linha. Time travel debugging.\n💡 Analogia: gravação de DVR. Pausa, volta, avança. Frames são results das activities.\n⚠️ Replay NÃO re-executa activities — usa resultados gravados. Instantâneo, gratuito.\n➡️ Armadilha: non-determinism.", 41, T, "Time travel debugging")

s = content_slide("O Problema do Non-Determinism", [
    "Replay pressupõe que o workflow re-executa de forma idêntica",
    "Non-determinism quebra o replay:",
    "  • datetime.now() → valor diferente a cada replay",
    "  • random() → resultado diferente",
    "  • I/O direto → estado externo mudou",
    "  • Código alterado entre execução e replay",
    "Solução: todo I/O e não-determinismo vão em activities",
], "📖 Non-determinism é inimigo nº 1 do replay. Replay re-executa código esperando mesma sequência de activities. Se workflow usa datetime.now(), cada replay gera timestamp diferente — diverge — quebra.\n💡 Analogia: reassistir filme esperando mesmo final, mas alguém trocou roteiro.\n⚠️ Deployar nova versão sem versionar = replay quebra. Version via Task Queue.\n➡️ Regras práticas.", 42, T, "A armadilha do replay")

s = code_slide("Determinismo em Código de Workflow",
    "✗ Quebrado (non-determinístico):\n"
    "  now = datetime.now()      # NON-DETERMINISTIC\n"
    "  time.sleep(5)              # BLOQUEIA, não durable\n"
    "  result = requests.get(...)  # I/O direto\n\n"
    "✓ Correto (determinístico):\n"
    "  now = workflow.now()       # seeded do histórico\n"
    "  await workflow.sleep(5)    # durable\n"
    "  result = await workflow.execute_activity(\n"
    "      fetch_data)            # I/O em activity",
    "📖 Regras simples. No workflow: APIs determinísticas (workflow.now, workflow.random, workflow.sleep). Nunca I/O direto. Em activities: tudo permitido (HTTP, DB, API).\n💡 Analogia: workflow é contrato (reproduzível). Activity é execução (efeitos reais).\n⚠️ 'datetime.now() só uma vez' = tem. No primeiro replay, quebra.\n➡️ DEMO Temporal.", 43, T, "Regras de determinismo")

s = code_slide("DEMO: Workflow Durável em Temporal",
    "Referência: 05-Labs/ETHAGT11/Lab2-Workflow-Temporal\n\n"
    "1. Iniciar workflow:\n"
    "   temporal workflow start --task-queue agents \\\n"
    "     --type ProcessTicketsWorkflow\n\n"
    "2. Worker processando tickets...\n"
    "3. Ctrl+C (kill do worker)\n\n"
    "4. Reiniciar worker:\n"
    "   → Temporal faz replay do histórico\n"
    "   → Retoma do último ticket completado\n"
    "   → ZERO perda de progresso",
    "📖 Demo ao vivo. Workflow processa tickets em loop, chamando agente como activity. Na Temporal Web UI (localhost:8080). Ctrl+C no worker após alguns tickets. Workflow pausa. Reinicio worker — pega histórico, replay, continua do último ticket.\n💡 Analogia: desligar/religar computador durante compilação. Sem durable = recompila do zero. Com Temporal = continua do arquivo que parou.\n⚠️ Se demo falhar, vídeo/screenshot prontos. Demo é ponto alto.\n➡️ Pergunta da demo.", 44, T, "Demo ao vivo")

s = exercise_slide("Pergunta da DEMO", [
    "Discussão aberta (2 min):",
    "",
    "1. Replay: se uma tool externa mudou de comportamento, o replay quebra?",
    "2. O que acontece se o agente tomar decisão diferente no replay?",
    "",
    "Respostas:",
    "  • Replay NÃO re-executa activities — usa resultados gravados",
    "  • Se agente é determinístico (temp=0), reproduz mesma decisão",
], "📖 Respostas: (1) Tool mudou → replay NÃO quebra (usa resultados gravados). Mas próxima execução após replay pode diferir. (2) Agente determinístico (mesmo input = mesma output) = mesma decisão. Se temp>0, pode divergir — use temp=0 ou persista decisão como activity.\n💡 Analogia: reassistir filme. Atores não refilmam — gravação é replayed. Mas refazer a cena pode improvisar diferente.\n⚠️ Replay não re-executa tools. Usa resultados gravados.\n➡️ Exercício HITL.", 45, T, "Discussão aberta")

s = exercise_slide("Exercício: HITL com Signal", [
    "Cenário: agente propõe enviar email para cliente",
    "",
    "Em duplas, esboçar o workflow:",
    "1. Signal de aprovação (humano aprova/rejeita)",
    "2. Timeout de 24h",
    "3. O que acontece se humano rejeitar? E se não responder?",
    "",
    "2 min discussão, 1 min compartilhar",
], "📖 Exercício em duplas. Esboçar: (1) agente gera email; (2) workflow pausa em await signal com timeout 24h; (3) aprovado → envia; (4) rejeitado → volta ao agente; (5) timeout → escalada. Ponto é pensar nas ramificações.\n💡 Analogia: processo de compra com aprovação do gerente. Aprova → segue. Rejeita → revisão. Sem resposta → escala.\n⚠️ Esquecer timeout = workflow pausado para sempre se humano ignorar.\n➡️ Patterns de resiliência.", 46, T, "Exercício em duplas")

s = section_slide(5, "Patterns de Resiliência")
add_notes(s, "Event-driven traz resiliência, MAS falhas acontecem. Patterns: retries, idempotência, saga, circuit breakers.\n➡️ Retries.")

s = content_slide("Retries com Backoff", [
    "Retry simples: tentar N vezes",
    "Backoff exponencial: 1s → 2s → 4s → 8s → 16s",
    "Jitter: aleatoriedade para evitar thundering herd",
    "Limite de tentativas + circuit breaker quando excedido",
    "Para agentes: retry em tool calls com timeout",
    "Temporal oferece retry com backoff embutido (RetryPolicy)",
], "📖 Retry é o pattern mais fundamental. Se chamada falha, tentar de novo. Mas retries ingênuos causam thundering herd (1000 clientes retentando = derruba servidor). Solução: backoff exponencial + jitter.\n💡 Analogia: ligar para alguém ocupado. Em vez de discar a cada segundo (piora), espere crescentemente. Jitter dessincroniza.\n⚠️ Retry sem backoff = thundering herd. Sempre backoff exponencial + jitter.\n➡️ Mas retry só é seguro se idempotente.", 48, T, "Backoff exponencial + jitter")

s = content_slide("Idempotência (Chaves de Idempotência)", [
    "At-least-once delivery: mensagens podem chegar mais de uma vez",
    "Idempotência: processar 2x = processar 1x (mesmo resultado)",
    "Chave de idempotência: ID único por operação",
    "Antes de processar: verificar se já processou aquela chave",
    "Armazenar chaves processadas (Redis, DB)",
], "📖 Idempotência é complemento essencial do retry. Se retenta, pode processar 2x. Se operação é 'enviar email', manda 2. Solução: chave de idempotência única. Antes de processar, verifique se key já foi. Se sim, retorne resultado anterior.\n💡 Analogia: carimbo 'PAGO' na fatura. Se tentar pagar de novo, carimbo impede. Key é o carimbo digital.\n❓ 'Quais das suas tools são idempotentes?'\n⚠️ 'Se não retento, não preciso idempotência' = errado. At-least-once pode duplicar sem retry seu.\n➡️ Caso concreto.", 49, T, "Complemento essencial do retry")

s = code_slide("Idempotência: Tool 'Enviar Email'",
    "async def send_email(to, subject, body, idempotency_key):\n"
    "    # Verificar se já processou\n"
    "    existing = await db.find(\n"
    "        'sent_emails',\n"
    "        key=idempotency_key)\n\n"
    "    if existing:\n"
    "        return existing.result  # não reenvia\n\n"
    "    # Processar e gravar\n"
    "    result = await email_provider.send(\n"
    "        to, subject, body)\n"
    "    await db.insert('sent_emails', {\n"
    "        'key': idempotency_key,\n"
    "        'result': result})\n"
    "    return result",
    "📖 Caso concreto. send_email recebe idempotency_key (gerada pelo chamador). Antes de enviar, consulta DB. Se existe, retorna resultado anterior. Se não, envia, grava key+resultado. Resultado: mesmo com 10 retries, só 1 email enviado.\n💡 Analogia: cofre com código. Abre uma vez com aquele código. Tentar de novo não abre nada novo.\n⚠️ idempotency_key = hash do conteúdo = errado. Se conteúdo muda (timestamp), hash muda. Use keys estáveis (request_id, message_id).\n➡️ Transações distribuídas: saga.", 50, T, "Caso concreto")

s = content_slide("Compensação: Saga Pattern", [
    "Saga: sequência de transações locais com compensação",
    "Passo 1 (execute) → Passo 2 (execute) → Passo 3 (FALHA)",
    "Compensar Passo 2 → Compensar Passo 1",
    "Não é rollback — é ação reversa (nem sempre perfeita)",
    "Dois estilos: orquestrada (central) vs coreografada (eventos)",
], "📖 Saga resolve transações distribuídas sem 2PC. Transações locais com compensação. Compensação NÃO é rollback — é ação reversa. Nem sempre perfeita (email enviado não 'desenvia').\n💡 Analogia: planejar festa. Buffet cancela? Cancela salão, devolve convites. Não é 'desfazer' — é reverter efeito. Alguns convidados já leram convite.\n⚠️ Saga não é rollback. Compensação é melhor esforço. Algumas operações não têm compensação perfeita.\n➡️ Caso concreto.", 51, T, "Transação compensatória")

s = content_slide("Saga: Transferência entre Contas", [
    "Transferência: debita A → credita B → notifica usuário",
    "Se 'notifica' falha:",
    "  Compensar 'credita B' (estorna) → compensar 'debita A' (devolve)",
    "Saga orquestrada: orquestrador coordena passos e compensações",
    "Saga coreografada: cada serviço publica evento e reage",
], "📖 Caso clássico. (1) debita A, (2) credita B, (3) notifica. Se notifica falha: estorna B, devolve A. Usuário não notificado, mas dinheiro voltou. Orquestrada: orquestrador coordena. Coreografada: serviços reagem a eventos.\n💡 Analogia: devolução na loja. Não 'desfaz' compra — faz estorno. Dinheiro volta, transação original existe.\n❓ 'Qual compensação não é perfeita?' (notificação — email não 'desenvia')\n⚠️ Compensações também podem falhar. Saga precisa retry nas compensações.\n➡️ Circuit breaker.", 52, T, "Caso de transferência")

s = content_slide("Circuit Breakers", [
    "Estados: CLOSED (normal) → OPEN (falhando, rejeita) → HALF-OPEN (testa)",
    "Após N falhas consecutivas: OPEN (para de chamar a tool)",
    "Após timeout: HALF-OPEN (permite 1 tentativa)",
    "Se sucesso: CLOSED; se falha: OPEN",
    "Para agentes: se tool de API está down, breaker evita desperdício de tokens",
], "📖 Circuit breaker protege de falhas em cascata. CLOSED: normal. Após N falhas: OPEN (rejeita imediatamente). Cooldown: HALF-OPEN (1 teste). Sucesso: CLOSED. Falha: OPEN. Para agentes: se LLM down, breaker evita gastar tokens em chamadas que falham.\n💡 Analogia: disjuntor de casa. Sobrecarga → abre (corta energia). Não tenta ligar aparelhos. Reseta (HALF-OPEN) e se não há sobrecarga, volta (CLOSED).\n⚠️ Retry infinito sem breaker = desperdiça recursos e tokens se API realmente down.\n➡️ Composição de patterns.", 53, T, "Proteger de falhas em cascata")

s = content_slide("Composição de Patterns", [
    "Retry + Idempotência: retry é seguro porque idempotência garante não-duplicação",
    "Saga + Idempotência: compensações são idempotentes",
    "Circuit Breaker + Retry: breaker protege; retry recupera",
    "Durable Execution + tudo: Temporal oferece retry, timeout, idempotência embutidos",
    "",
    "Regra: não reinvente — Temporal já tem a maioria",
], "📖 Patterns se combinam. Retry+Idempotência: retry duplica, idempotência garante não-efeito. Saga+Idempotência: compensações idempotentes. Breaker+Retry: breaker para, retry recupera. Temporal oferece tudo embutido.\n💡 Analogia: equipamentos de segurança do carro. Cinto (idempotência), airbag (compensação), ABS (breaker). Juntos = máxima segurança.\n⚠️ Reinventar patterns que Temporal oferece é desperdício. Retry, timeout, idempotência — tudo configurável.\n➡️ Tópicos de produção.", 54, T, "Patterns se combinam")

s = section_slide(6, "Produção: Ordering, Escala, Observabilidade")
add_notes(s, "Dos fundamentos para produção. How colocar event-driven em escala.\n➡️ Debate mais quente: exactly-once.")

s = content_slide("Exactly-once vs At-least-once", [
    "At-least-once: mensagem pode chegar 1+ vezes (default Kafka)",
    "At-most-once: mensagem pode se perder (fire and forget)",
    "Exactly-once: cada mensagem processada exatamente 1 vez",
    "Realidade: exactly-once verdadeiro não existe na presença de falhas",
    "Prática: at-least-once + idempotência = 'effectively once'",
    "",
    "Pergunta: Exactly-once existe de verdade ou é mito?",
], "📖 Debate clássico. Exactly-once delivery é mito na presença de falhas (ACK pode falhar). Prática consagrada: at-least-once (pode duplicar) + idempotência (duplicar não causa efeito) = 'effectively once'. É assim que Kafka, NATS, Temporal funcionam.\n💡 Analogia: carteiro pode entregar carta 2x (at-least-once). Destinatário idempotente (arquiva por ID) = 'uma vez'. Exactly-once exigiria carteiro que nunca erra.\n❓ 'Exactly-once existe ou é mito?'\n⚠️ Perder tempo em exactly-once é desperdício. Idempotência é mais simples e robusta.\n➡️ Escalar consumidores.", 56, T, "Desconstruir o mito")

s = content_slide("Sharding de Consumidores", [
    "Consumer Group: partições divididas entre consumers",
    "Sharding por chave: agent_id → partition → consumer dedicado",
    "Adicionar consumers → rebalanceamento automático",
    "Limite: # consumers ativos ≤ # partições",
    "Para agentes: cada shard processa um conjunto de agentes",
], "📖 Sharding é como Kafka escala. Partições divididas entre consumers do mesmo group. 100 partições, 10 consumers = cada um pega 10. Sharding por agent_id garante mesmo agente sempre no mesmo consumer (ordering). Para escalar, adicione partições.\n💡 Analogia: divisão de trabalho em fábrica. Cada operário cuida de um conjunto de máquinas. Mais operários = mais máquinas necessárias.\n⚠️ Adicionar consumers sem partições = ociosos.\n➡️ Observabilidade.", 57, T, "Escala horizontal")

s = content_slide("Distributed Tracing em Pipelines de Agentes", [
    "Em event-driven: trace espalhado por producer, broker, consumer",
    "OpenTelemetry: spans distribuídos com context propagation",
    "Trace ID propagado via headers de mensagem",
    "Para agentes: cada tool call = span; cada agente = service",
    "Aprofundamento em ETHAGT12 (AgentOps)",
], "📖 Em event-driven, request passa por múltiplos serviços. Sem distributed tracing, não sabe onde latência está. OpenTelemetry propaga Trace ID via headers — cada serviço cria spans filhos. Para agentes: tool call = span, agente = service.\n💡 Analogia: rastrear encomenda. Cada parada registrada com timestamp. Vê onde passou e quanto ficou.\n⚠️ Tracing por serviço sem propagation = traces fragmentados. Sempre propagar Trace ID.\n➡️ Custo.", 58, T, "Observabilidade distribuída")

s = content_slide("Custo da Mensageria", [
    "Infraestrutura: brokers, ZooKeeper/KRaft, monitoramento",
    "Storage: retenção de logs consome disco",
    "Network: replicação entre datacenters",
    "Operação: equipe de platform engineering",
    "Para agentes: cada evento tem custo de storage + processamento",
    "Regra: mensure custo por evento antes de escalar",
], "📖 Mensageria não é grátis. Infraestrutura, storage (logs retidos), network (replicação), operação (platform engineering). Para agentes: cada evento tem custo. 1M eventos/dia com retenção 7 dias = 7M armazenados. Mensure antes de escalar.\n💡 Analogia: armazém. Cada caixa (evento) ocupa espaço. Retenção longa = armazém grande = caro.\n⚠️ Retenção infinita 'por segurança' = disco enche, cluster cai.\n➡️ Quebrar mito.", 59, T, "Mensageria não é grátis")

s = exercise_slide("V/F: Event-driven é sempre mais escalável", [
    "Verdadeiro ou Falso:",
    "",
    "'Event-driven é sempre mais escalável'",
    "",
    "Resposta: FALSO",
    "",
    "Event-driven escala horizontalmente, mas introduz overhead:",
    "latência de broker, complexidade de debugging, custo operacional",
    "Para tarefas simples e síncronas, não compensa",
], "📖 Resposta: Falso. Event-driven escala horizontalmente (mais consumers = mais throughput), MAS introduz overhead: latência de broker, complexidade, debugging distribuído. Para tarefas simples síncronas (responder chat), é overkill.\n💡 Analogia: 'caminhões são sempre melhores que carros'. Caminhão carrega mais, mas é lento, caro, difícil de estacionar. Para padaria, use carro.\n⚠️ Adotar event-driven para tudo 'porque escala' = over-engineering. Comece simples.\n➡️ Fechamento.", 60, T, "Quebrar mito")

s = section_slide(7, "Fechamento")
add_notes(s, "Última seção. Exercício saga, boas práticas, anti-patterns, caso real, quiz, Q&A.\n➡️ Exercício.")

s = exercise_slide("Exercício: Saga Compensatória", [
    "Cenário: transferência entre contas (debita A → credita B → notifica)",
    "",
    "Em duplas, escrever a compensação para cada passo falhar:",
    "  1. Se 'debita A' falha → ?",
    "  2. Se 'credita B' falha → ?",
    "  3. Se 'notifica' falha → ?",
    "",
    "Quais compensações não são perfeitas?",
    "3 min discussão, 2 min apresentar",
], "📖 Exercício em duplas. Cada dupla escreve compensação para cada falha. Discutir quais não são perfeitas (notificação — email não 'desenvia'). Pedir 2 duplas compartilharem.\n💡 Analogia: planejar viagem. Voo cancela? Hotel cancela? Carro indisponível? Para cada, qual compensação?\n⚠️ Compensação sem considerar que também pode falhar = lacuna.\n➡️ Boas práticas.", 62, T, "Saga em duplas")

s = content_slide("Boas Práticas (DO)", [
    "Comece com at-least-once + idempotência (não persiga exactly-once)",
    "Use chaves de partição para preservar ordering por entidade",
    "Coloque I/O em activities; mantenha workflows determinísticos",
    "Defina DLQ (dead letter queue) desde o dia 1",
    "Monitore lag de consumer (Kafka) / depth de fila (RabbitMQ)",
    "Use circuit breakers em tools externas",
    "Version seus eventos (schema registry)",
], "📖 7 boas práticas cumulativas. (1) at-least-once + idempotência é base. (2) chaves de partição preservam ordering. (3) I/O em activities. (4) DLQ desde dia 1. (5) Monitore lag/depth. (6) Circuit breakers. (7) Version eventos.\n💡 Analogia: checklist de pré-voo. Cada item evita classe de problema.\n⚠️ Esquecer DLQ = mensagens falhadas somem ou bloqueiam fila. Sempre defina.\n➡️ Anti-patterns.", 63, T, "Checklist verde")

s = content_slide("Anti-Patterns (DON'T)", [
    "Perseguir exactly-once sem idempotência",
    "I/O direto em código de workflow (non-determinism)",
    "Sem DLQ (mensagens falhadas somem)",
    "Sem schema versioning (quebra consumidores ao evoluir)",
    "Ordering global esperado em Kafka (só existe por partição)",
    "Sem monitoring de lag (consumer atrasa silenciosamente)",
    "Usar Airflow para orquestração em tempo real",
    "Começar com Kafka quando RabbitMQ bastava",
], "📖 8 anti-patterns comuns. (1) exactly-once sem idempotência = perda de tempo. (2) I/O no workflow = quebra replay. (3) Sem DLQ = somem. (4) Sem versioning = quebra consumers. (5) Ordering global em Kafka = não existe. (6) Sem monitoring de lag = atraso silencioso. (7) Airflow para tempo real = não foi feito. (8) Kafka quando RabbitMQ bastava = overkill.\n💡 Analogia: lista de 'não fazer na cozinha'. Cada item evita desastre.\n⚠️ 'Anti-patterns não acontecem comigo' = acontece com todo mundo na primeira vez.\n➡️ Caso real.", 64, T, "Checklist vermelho")

s = content_slide("Caso de Estudo: Processamento de Documentos", [
    "Pipeline: ingestão (Kafka) → classificação (agente) → extração (agente) → validação (HITL) → publicação",
    "Temporal orquestra; agentes são activities",
    "Saga compensatória se validação falha",
    "Circuit breaker no agente de extração (API de LLM)",
    "DLQ para documentos irreparáveis",
    "Distributed tracing via OpenTelemetry",
    "",
    "Lição: event-driven + durable execution + resiliência = produção",
], "📖 Caso real junta tudo. Ingestão via Kafka, classificação/extração como activities do Temporal, HITL para validação, saga se falha, circuit breaker no LLM, DLQ para irreparáveis, tracing. Cada pattern tem papel.\n💡 Analogia: fábrica de alimentos. Esteira (Kafka), classificador, extrator, inspetor (HITL), empacotador. Linha de reversão (saga), descarte controlado (DLQ).\n⚠️ 'Começo simples e adiciono resiliência depois' = refatoração cara. DLQ e breaker desde dia 1.\n➡️ Resumo.", 65, T, "Todos os patterns em um caso real")

s = content_slide("Resumo da Aula", [
    "Event-driven = desacoplamento, escala, resiliência (com complexidade)",
    "Mensageria: Kafka (log/volume), RabbitMQ (routing), NATS (leveza)",
    "Temporal = durable execution: sobreviver a crashes, long-running, HITL",
    "Resiliência: retry + idempotência + saga + circuit breaker",
    "Produção: at-least-once + idempotência = 'effectively once'",
    "Determinismo é a regra de ouro do replay",
], "📖 6 pontos para lembrar: (1) event-driven traz desacoplamento/escala/resiliência com complexidade. (2) Kafka=log, RabbitMQ=routing, NATS=leveza. (3) Temporal=durable execution. (4) Resiliência é composição. (5) Exactly-once é mito; effectively once é prática. (6) Determinismo é regra de ouro.\n💡 Analogia: 6 mandamentos da cozinha. Seguir = prato sai bom.\n➡️ Checklist.", 66, T, "6 pontos-chave")

s = content_slide("Checklist da Aula", [
    "[ ] Diferenciou síncrono de event-driven",
    "[ ] Comparou Kafka, RabbitMQ e NATS",
    "[ ] Explicou durable execution e Temporal",
    "[ ] Implementou retries e idempotência",
    "[ ] Escreveu uma saga compensatória",
    "[ ] Discutiu exactly-once vs at-least-once",
], "📖 Checklist final. Se 6 itens marcados, aula cumpriu objetivos. Se algum não, indicar slides para revisar.\n➡️ Quiz.", 67, T, "Confirmar cobertura")

s = exercise_slide("Quiz 1: Temporal vs Kafka", [
    "'Quando Temporal é preferível a Kafka puro?'",
    "",
    "A) Quando você precisa de máximo throughput",
    "B) Quando você precisa de durable execution, HITL e orquestração",
    "C) Quando você quer o broker mais leve",
    "D) Quando você não precisa de ordering",
    "",
    "Resposta: B",
], "📖 Resposta B. Temporal é durable execution, HITL, orquestração. Kafka é mensageria. Não são mutuamente exclusivos. A é Kafka. C é NATS. D é irrelevante.\n➡️ Quiz 2.", 68, T, "Quiz — múltipla escolha")

s = exercise_slide("Quiz 2: Ordering no Kafka", [
    "'Em Kafka, como garantir ordering de mensagens do mesmo agente?'",
    "",
    "A) Usar ordering global (não existe)",
    "B) Usar a mesma chave de partição (agent_id)",
    "C) Aumentar o número de partições",
    "D) Usar RabbitMQ no lugar",
    "",
    "Resposta: B",
], "📖 Resposta B. Mesma key = mesma partição = ordering. A é falso. C não ajuda (pode piorar sem key). D é trocar ferramenta.\n➡️ Quiz 3.", 69, T, "Quiz — múltipla escolha")

s = exercise_slide("Quiz 3: Non-Determinism", [
    "'O que quebra um replay no Temporal?'",
    "",
    "A) Usar workflow.sleep() no workflow",
    "B) Non-determinism (I/O direto, datetime.now(), random)",
    "C) Chamar activities",
    "D) Usar signals",
    "",
    "Resposta: B",
], "📖 Resposta B. Non-determinism quebra replay (sequência diverge do histórico). A é determinístico. C é uso correto. D é suportado.\n➡️ Quiz 4.", 70, T, "Quiz — múltipla escolha")

s = exercise_slide("Quiz 4: Effectively Once", [
    "'Qual é a estratégia prática para effectively once?'",
    "",
    "A) Implementar exactly-once delivery no broker",
    "B) At-least-once delivery + idempotência no consumidor",
    "C) At-most-once + retry infinito",
    "D) Usar transações distribuídas (2PC)",
    "",
    "Resposta: B",
], "📖 Resposta B. At-least-once + idempotência = effectively once. É prática consagrada. A é impossível. C perde mensagens. D é caro/frágil.\n➡️ Quiz 5.", 71, T, "Quiz — múltipla escolha")

s = exercise_slide("Quiz 5: Saga", [
    "'Em uma saga de transferência, se notificar falha, o que se faz?'",
    "",
    "A) Nada — a transferência já foi feita",
    "B) Compensar passos anteriores (estornar crédito e débito)",
    "C) Retentar infinitamente até notificar",
    "D) Cancelar via 2PC",
    "",
    "Resposta: B",
], "📖 Resposta B. Compensar: estornar B, devolver A. Usuário não notificado, mas dinheiro voltou. A = inconsistente. C = trava. D não é saga.\n➡️ Conexão.", 72, T, "Quiz — múltipla escolha")

s = content_slide("Conexão com Próximos Módulos", [
    "ETHAGT12 — AgentOps: observabilidade + avaliação (aprofunda tracing)",
    "ETHAGT14 — Escalabilidade: event-driven em produção em larga escala",
    "ETHAGT90 — Projeto final: pipeline event-driven como componente",
], "📖 ETHAGT11 conecta com ETHAGT12 (AgentOps — aprofunda tracing), ETHAGT14 (Escalabilidade — event-driven em larga escala) e ETHAGT90 (Projeto final). Base de hoje é usada em todos.\n➡️ Projeto e labs.", 73, T, "Mapa da especialização")

s = content_slide("Projeto do Módulo e Labs", [
    "Projeto: pipeline event-driven para tickets em massa",
    "  Durable execution, retries e compensação",
    "  Entrega: código + ADR + análise de falhas injetadas (chaos)",
    "  Critério: sobrevive a ≥2 falhas injetadas sem perda de dados",
    "",
    "Lab 1 (4h): 'Agente com Kafka' — dois agentes coordenados via tópicos",
    "Lab 2 (5h): 'Workflow durável em Temporal' — agente sobrevivendo a kill",
], "📖 Projeto é avaliação principal. Pipeline que sobrevive a falhas injetadas (chaos). Entrega inclui ADR. Critério: ≥2 falhas sem perda. Labs são preparatórios: Lab 1 (Kafka), Lab 2 (Temporal).\n💡 Analogia: teste de colisão de carro. Prove que sobrevive a 'colisões' sem 'ferir motorista' (perder dados).\n➡️ Leitura.", 74, T, "Avaliação + labs")

s = content_slide("Leitura Recomendada e Referências", [
    "Obrigatório:",
    "  • Temporal.io Durable Execution primer",
    "  • Kreps, The Log (LinkedIn Engineering)",
    "Recomendado:",
    "  • Narkhede et al., Kafka: The Definitive Guide",
    "  • Microsoft Cloud Design Patterns (saga, CQRS)",
    "Complementar:",
    "  • CloudEvents spec (CNCF); NATS docs; RabbitMQ docs",
    "  • Vídeo: Temporal Durable Execution talks (YouTube)",
], "📖 Obrigatórias: Temporal primer (durable execution profundo) e Kreps The Log (fundamento conceitual). Kafka: The Definitive Guide é referência canônica. Microsoft Cloud Design Patterns cobre saga/CQRS. CloudEvents é spec de padronização.\n➡️ Q&A.", 75, T, "Leitura obrigatória")

s = title_slide(
    "Perguntas?",
    "ETHAGT11 — Event-Driven Agents & Workflow Orchestration\nPróxima aula: ETHAGT12 — AgentOps: Observabilidade & Avaliação",
    "Universidade Etho"
)
add_notes(s, "📖 Abrir Q&A. Se não houver perguntas: 'Qual parte foi menos clara?'\nMensagem final: 'Event-driven não é sobre Kafka ou Temporal — é sobre desacoplamento, resiliência e durabilidade. As ferramentas mudam; os princípios permanecem. Se entenderam os princípios, qualquer ferramenta será uma instanciação.'\n💡 Analogia final: 'Vocês aprenderam a cozinhar. Kafka, Temporal, RabbitMQ são as panelas. O prato depende do cozinheiro.'\n➡️ Fim da aula.")

output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ETHAGT11-Apresentacao.pptx")
prs.save(output)
print(f"PPTX gerado: {output}")
print(f"Total de slides: {len(prs.slides)}")
