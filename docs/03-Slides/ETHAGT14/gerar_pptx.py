#!/usr/bin/env python3
"""
Gerador de PPTX — ETHAGT14: Escalabilidade & Sistemas Distribuídos de Agentes
Universidade Etho · Especialização em Programação Agêntica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, num, total, obj=""):
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_header(slide, code="ETHAGT14"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

def title_slide(title, subtitle, code):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return s

def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s

def content_slide(title, bullets, notes, num, total, obj="", subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def code_slide(title, code, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

def exercise_slide(title, items, notes, num, total, obj="Exercício"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj)
    add_notes(s, notes)
    return s

T = 75

s = title_slide(
    "Escalabilidade & Sistemas Distribuídos de Agentes",
    "Universidade Etho · Especialização em Programação Agêntica\nFase D — Sistemas Distribuídos · 30 h",
    "ETHAGT14"
)
add_notes(s, "📖 Bem-vindos à Fase D. Onde protótipos viram produção. Um agente que funciona para 1 usuário não é um sistema escalável.\n💡 Analogia: cozinhar para a família vs abrir um restaurante. Mesmos ingredientes, engenharia diferente.\n❓ 'Quantos já colocaram agente em produção para +100 usuários?'\n⚠️ Escala ≠ mais servidores. Escala é reduzir custo, latência e contenção sistemicamente.\n➡️ Vamos aos objetivos.")

s = content_slide("Objetivos do Módulo", [
    "Objetivo geral: levar sistemas de agentes a escala de produção",
    "",
    "Objetivos específicos:",
    "1. Identificar gargalos em sistemas multi-agente (LLM, tools, memória, rede)",
    "2. Aplicar caching (semântico, de prompts, de embeddings)",
    "3. Distribuir agentes (sharding, replica, partitioning)",
    "4. Otimizar custo e latência (model routing, batching, speculative)",
    "5. Operar FinOps de agentes (orçamento, observabilidade de custo)",
], "📖 Objetivos mensuráveis: identificar, aplicar, distribuir, otimizar, operar.\n❓ 'Qual objetivo é mais urgente no trabalho de vocês?' (geralmente caching ou FinOps)\n⚠️ Priorizem caching antes de distribuição (ordem errada: sharding antes de cache).\n➡️ Competências.", 2, T, "Objetivos mensuráveis")

s = content_slide("Competências Desenvolvidas", [
    "C1 Programação Agêntica → A (Avançado)",
    "C2 Multi-Agent Systems → A (Avançado)",
    "C3 MCP & Tool Use → B (Básico)",
    "C4 Agent Memory → A (Avançado)",
    "C5 AgentOps & Avaliação → A (Avançado)",
    "",
    "Módulo terminal: 4 de 5 competências em Avançado",
    "C4 ganha dimensão distribuída: checkpointer, estado entre réplicas",
], "📖 Módulo terminal — 4 competências em Avançado. Avançado = arquitetar, defender, operar.\n⚠️ Alunos subestimam C4 aqui. Memória distribuída é o calcanhar de Aquiles da escala.\n➡️ Agenda.", 3, T, "Conectar ao Framework Etho")

s = content_slide("Agenda da Aula", [
    "Bloco 1 (45 min):",
    "  Abertura (8 min) — objetivos, custo da escala, contexto",
    "  Gargalos (10 min) — latência, custo, rate limits, estado",
    "  Caching (22 min) — exact, semântico, embeddings, tools, DEMO",
    "  Intervalo (5 min)",
    "",
    "Bloco 2 (45 min):",
    "  Routing & Otimização (12 min) — Haiku/Sonnet, batching, speculative",
    "  Distribuição (10 min) — stateless, sharding, replica, consensus",
    "  Infra (5 min) — K8s, serverless, GPU, service mesh",
    "  FinOps (8 min) — orçamento, trade-offs, circuit breaker",
    "  Fechamento (10 min) — boas práticas, caso, quiz, Q&A",
], "📖 Bloco 1: diagnóstico + caching (estrela). Bloco 2: arquitetura + FinOps.\n1 DEMO ao vivo (Slide 26), 3 exercícios rápidos, quiz final.\n➡️ Por que estamos aqui?", 4, T, "Roteiro da aula")

s = content_slide("O Custo da Escala", [
    "Cenário: agente funciona para 1 usuário custa R$0,50/consulta",
    "",
    "Para 10.000 usuários simultâneos:",
    "  • Latência explode: 2s → 30s (rate limit, contenção)",
    "  • Custo mensal: R$50k (linear = explosão)",
    "  • Estado distribuído entre réplicas: race conditions",
    "",
    "Sem otimização: escala linear = explosão de custo",
    "",
    "Pergunta: Qual o maior custo oculto que já viram em LLMs?",
], "📖 Escala linear é o pesmaro. 10k usuários × 10 consultas/dia × R$0,50 = R$1,5M/mês. Sem caching/routing/FinOps, quebra em 30 dias.\n💡 Analogia: restaurante servindo 10 vs 1.000 clientes no mesmo horário.\n❓ 'Qual o maior custo oculto?' (deixar 2-3 responderem — contexto crescente, retries, logs)\n⚠️ Custo de LLM é só 60-70% do total. Infra é 30-40%.\n➡️ Por que só agora escala é urgente?", 5, T, "Criar tensão — escala custa caro")

s = content_slide("Por Que Escalabilidade Agora", [
    "Timeline:",
    "  2023: protótipos (todo mundo experimentando)",
    "  2024: frameworks maduros (LangGraph, CrewAI, OpenAI SDK)",
    "  2025: sistemas em produção",
    "  2026: escala como urgência financeira",
    "",
    "Confluência:",
    "  1. Agentes = múltiplas chamadas de LLM por execução",
    "  2. Contexto cresce dentro de uma sessão",
    "  3. APIs impõem rate limits (RPM/TPM)",
    "  4. Estado de sessão precisa ser distribuído",
], "📖 Em 2026 a pergunta virou 'como escalo sem quebrar o caixa?'. Empresas gastando 6 dígitos/mês percebem que sem engenharia de escala, não dá.\n💡 Analogia: web em 1999. Primeiro vieram os sites, depois a pergunta: como servir 1M de usuários?\n⚠️ Escala não é só problema do DevOps. Começa no design do agente.\n➡️ Vamos diagnosticar gargalos.", 6, T, "Confluência histórica")

s = section_slide(1, "Onde Agentes Esbarram em Escala")
add_notes(s, "Início do bloco de gargalos. Os próximos 8 slides identificam os 5 gargalos principais. Cada um tem uma técnica de mitigação que veremos adiante.")

s = content_slide("Latência de LLM: O Gargalo Dominante", [
    "TTFT (Time To First Token): ~0.3-1s antes do primeiro token",
    "TPOT (Time Per Output Token): ~30-80ms por token gerado",
    "",
    "Loop ReAct de 5 steps:",
    "  5 × (TTFT + TPOT × tokens) = latência cumulativa",
    "  Latência serial: cada step depende do anterior",
    "",
    "Comparação:",
    "  Chat simples: ~2s",
    "  Agente 3-step: ~6s",
    "  Agente 5-step: ~15s",
    "  Agente 10-step: ~60s",
], "📖 LLM é gargalo serial. Cada chamada leva segundos; em loop ReAct, são sequenciais. 1.000 usuários simultâneos sem paralelização = fila cresce.\n💡 Analogia: restaurante com 1 cozinheiro. Adicionar cozinheiro (réplica) ou pré-preparar (caching) reduz espera.\n❓ 'Qual a latência p95 do agente mais crítico de vocês?'\n⚠️ Confundir TTFT com latência total. TTFT é só até o primeiro token.\n➡️ Latência não é o único problema. Custo cresce mais rápido.", 8, T, "LLM domina o tempo total")

s = content_slide("Custo: Crescimento Quadrático", [
    "Cada step reenvia TODO o histórico anterior",
    "",
    "Step 1: 1k tokens",
    "Step 5: 15k tokens",
    "Step 10: 40k tokens",
    "",
    "Custo de input cresce QUADRATICAMENTE com profundidade",
    "Soma triangular ≈ n²/2 (não é n × tokens-por-step)",
    "",
    "KV cache do provider reduz mas não elimina",
    "Implicação: agente profundo = custo explosivo",
], "📖 O custo oculto que mais surpreende. Em loop ReAct, cada step reenvia histórico. Não é 5 × 1k = 5k — é soma triangular ≈ n²/2.\n💡 Analogia: reunião onde cada participante ouve todo o histórico. Custo de comunicação cresce quadraticamente.\n⚠️ Calcular custo como 'steps × tokens-por-step'. Errado. Use soma triangular.\n➡️ E quando mil usuários fazem isso ao mesmo tempo?", 9, T, "Custo cresce quadraticamente")

s = content_slide("Concorrência e Rate Limits", [
    "APIs limitam RPM (requests/min) e TPM (tokens/min)",
    "Anthropic/OpenAI: tiers com limites crescentes",
    "",
    "1.000 usuários × 5 steps = 5.000 RPM necessários",
    "Sem gestão: 429 Too Many Requests → falhas em cascata",
    "",
    "Estratégias:",
    "  • Filas com prioridade",
    "  • Backoff exponencial com jitter",
    "  • Multi-provider (failover Anthropic→OpenAI→Bedrock)",
    "  • Batching de requests",
], "📖 Rate limit é o teto duro. Por mais que tenham orçamento, a API limita por minuto.\n💡 Analogia: porta giratória. Só passa N pessoas/min. Empurrar mais gente não acelera — trava.\n❓ 'Vocês já tomaram 429 em produção?' (mãos levantadas)\n⚠️ Retry sem backoff = tempestade de retries. Sempre exponencial com jitter.\n➡️ Mesmo com rate limit ok, há outro problema: estado.", 10, T, "Rate limits impedem escala")

s = content_slide("Estado Distribuído", [
    "Agente stateful: sessão, memória, contexto acumulado",
    "",
    "Com múltiplas réplicas: usuário pode cair em réplica diferente",
    "",
    "Opções:",
    "  • Sticky sessions (simples mas frágil)",
    "  • Estado externo (Redis/Postgres checkpoint)",
    "  • Stateless + checkpoint restore on-demand",
    "",
    "Tensão: stateless é fácil de escalar, mas agentes são stateful",
    "Aprofundamento: ETHAGT05 (Memória de Agentes)",
], "📖 Agentes têm estado. Em 1 instância, mora em memória. Com réplicas, usuário pode cair em réplica diferente. Solução: externalizar (Redis/Postgres). Transforma stateful em 'stateless-friendly'.\n💡 Analogia: hotel. Troca de quarto a cada dia → guarda coisas no cofre da recepção (Redis).\n⚠️ Estado em memória sem checkpoint. Deploy/crash = perda total.\n➡️ Vamos visualizar todos os gargalos.", 11, T, "Estado entre réplicas")

s = content_slide("Anatomia de um Gargalo", [
    "Fluxo: request → LLM (2s) → tools (0.5s) → contexto cresce ($$) → loop?",
    "",
    "Cada componente é um potencial gargalo",
    "Latência total = soma serial de todos os steps",
    "Custo total = soma de tokens × preço por step",
    "",
    "Mitigação: atacar o maior gargalo primeiro (maior ROI)",
    "Ver diagrama: 12-Diagrams/ETHAGT14/bottleneck-analysis.mmd",
], "📖 Diagrama-chave da aula. Cada nó é gargalo. LLM é dominante. Tools somam latência. Contexto multiplica custo. Loop multiplica tudo. Atacar nó vermelho de maior impacto (geralmente caching na LLM).\n💡 Analogia: via expressa. Cada semáforo/obra/pedágio é gargalo. Otimizar um semáforo não adianta se o pedágio é o gargalo real. Measure first.\n❓ 'No sistema de vocês, qual nó é o mais vermelho?'\n➡️ Para otimizar, precisamos medir.", 12, T, "Visualizar gargalos")

s = content_slide("Métricas-Chave de Escalabilidade", [
    "Latência: TTFT, TPOT, p50/p95/p99 de tempo total",
    "Custo: $ por execução, $ por usuário, $ por tenant",
    "Throughput: requisições concorrentes, tokens/min",
    "Taxa de erro: 429s, timeouts, falhas de tool",
    "Cache hit rate: % de requisições servidas do cache",
    "",
    "Dashboard mínimo: latência, custo, throughput, cache hit, erros",
    "Sempre olhar p95 e p99, não só média",
], "📖 Sem métricas, vocês otimizam no escuro. 5 dimensões formam o dashboard mínimo. p99 revela outliers (usuário esperando 30s). Custo deve ser por execução (granular).\n💡 Analogia: painel de carro. Não basta velocidade média — precisa combustível, temperatura, RPM.\n⚠️ 'Premature optimization is the root of all evil' — Knuth. Meçam primeiro.\n➡️ Pergunta rápida.", 13, T, "Métricas desde o dia 1")

s = content_slide("Pergunta: O Maior Custo Oculto", [
    "Qual o maior custo oculto que vocês já viram em sistemas de LLM?",
    "",
    "Opções:",
    "  • Contexto crescente",
    "  • Retries silenciosos",
    "  • Tools (APIs externas)",
    "  • Logs/traces",
    "  • Infra (Redis, K8s, observabilidade)",
    "",
    "Discussão aberta (3 min)",
    "Resposta típica: contexto crescente + retries silenciosos",
], "📖 Deixar a turma falar. Respostas comuns: contexto crescente, retries sem log, tools que chamam APIs pagas, logs verbosos, observabilidade cara. Lição: custo de LLM é 60-70%. Os 30-40% são 'invisíveis'.\n❓ Deixar 3-4 alunos compartilharem.\n⚠️ Calcular TCO esquecendo observabilidade e rede.\n➡️ Agora o tratamento. Caching é a primeira linha.", 14, T, "Engajamento")

s = section_slide(2, "Caching: A Primeira Linha de Defesa")
add_notes(s, "Caching é a otimização de maior ROI. Custo baixo, impacto alto, risco baixo. Se só vão implementar uma técnica, implementem caching.")

s = content_slide("Por Que Caching?", [
    "Caching elimina chamadas de LLM redundantes",
    "Impacto: reduz custo E latência simultaneamente",
    "Cache hit = resposta instantânea + custo zero",
    "",
    "Em sistemas reais: 30-60% das queries são repetidas ou similares",
    "Em FAQs/documentação técnica: até 70-80%",
    "",
    "Analogia: 'Cache é o único caso onde você ganha algo por nada'",
], "📖 30-60% das queries em produção são repetidas ou similares. Cache hit nessas = custo zero, latência ~50ms vs 2-3s do LLM. Única otimização que reduz custo E latência simultaneamente.\n💡 Analogia: biblioteca. Sem cache, toda pergunta vai ao arquivo morto (LLM). Com cache, frequentes ficam no balcão.\n⚠️ Achar que cache é só 'memória'. É estratégia: o quê, quanto tempo, como invalidar.\n➡️ Do mais simples ao mais sofisticado.", 16, T, "Maior ROI")

s = content_slide("Cache de Prompts / Exact Match", [
    "Mesmo prompt (exato) → mesma resposta",
    "Implementação: hash (SHA-256) do prompt → lookup em Redis",
    "",
    "Prompt caching nativo (Anthropic, OpenAI):",
    "  • Cache de prefixo no provider-side",
    "  • Reduz custo de system prompt repetido",
    "",
    "Vantagem: trivial de implementar",
    "Limitação: só prompts idênticos (case/whitespace-sensitive)",
], "📖 Exact match é o cache mais simples. Hash do prompt = chave. Se existe, retorna. Senão, chama LLM, armazena.\n💡 Analogia: reconhecer cliente pelo nome exato. 'João da Silva' = hit. 'João Silva' = miss.\n❓ 'Em qual caso exact match falha?' (qualquer variação)\n⚠️ Esquecer de normalizar antes de hashar (lowercase, trim).\n➡️ E se perguntas forem semanticamente iguais com palavras diferentes?", 17, T, "Forma mais simples")

s = content_slide("Cache Semântico: Conceito", [
    "Perguntas semanticamente similares → mesma resposta",
    "",
    "Exemplos equivalentes:",
    "  'Qual a capital do Brasil?'",
    "  'Capital do Brasil?'",
    "  'Brasil capital?'",
    "",
    "Implementação:",
    "  embedding da query → busca por similaridade → threshold",
    "  Similaridade > threshold: retorna cacheada",
    "  Senão: chama LLM, armazena resultado",
], "📖 Onde a mágica acontece. Compara SIGNIFICADO, não strings. Embeddings convertem query em vetor; busca vetorial encontra a mais similar. Se > threshold (e.g., 0.92), hit. Captura variações linguísticas, typos, sinônimos. Hit rate tipicamente dobra vs exact match.\n💡 Analogia: balconista experiente. Entende que variações da mesma pergunta são a mesma pergunta.\n⚠️ Usar embedding caro para cache. Use text-embedding-3-small (~1000x mais barato que LLM).\n➡️ Como implementar?", 18, T, "Cache por significado")

s = code_slide("Cache Semântico: Implementação", (
    "def semantic_cache(query, threshold=0.92):\n"
    "    # Passo 1: embeddar com modelo leve\n"
    "    emb = embed(query, model='text-embedding-3-small')\n"
    "\n"
    "    # Passo 2: busca vetorial no Redis\n"
    "    result = vector_db.search(emb, top_k=1)\n"
    "\n"
    "    # Passo 3: se similaridade > threshold, hit\n"
    "    if result and result.similarity > threshold:\n"
    "        return result.response  # instantâneo, custo ~zero\n"
    "\n"
    "    # Passo 4: miss → chamar LLM, armazenar\n"
    "    response = llm(query)\n"
    "    vector_db.store(emb, query, response)\n"
    "    return response\n"
    "\n"
    "# Custo do embedding ~1000x menor que custo do LLM"
), "📖 Implementação direta. Embedding é ~1000x mais barato que LLM. Mesmo com hit rate 20%, já vale a pena. TTL obrigatório para evitar cache stale.\n💡 Analogia: sugestão do Google. Calcula embedding da query parcial, busca similares.\n⚠️ Esquecer TTL. Sem TTL, cache acumula queries obsoletas.\n➡️ Threshold é o parâmetro crítico.", 19, T, "Código real")

s = content_slide("Threshold de Similaridade", [
    "Threshold alto (0.95): poucos falsos positivos, baixa hit rate",
    "Threshold baixo (0.80): alta hit rate, falsos positivos (erro!)",
    "",
    "Falso positivo = resposta semanticamente próxima mas ERRADA",
    "Exemplo: 'Preço iPhone 15' vs 'Preço iPhone 14' → 0.89 (parecido mas diferente)",
    "",
    "Recomendação: começar com 0.92, ajustar com eval",
    "",
    "Trade-off central: hit rate vs precisão",
], "📖 Threshold é delicado. Alto = só captura quase idênticas. Baixo = captura tudo, inclusive o que não deveria. Falso positivo é grave — serve resposta errada como certa.\n💡 Analogia: controle de fronteira. Muito rigoroso = poucos entram. Muito frouxo = entra gente errada.\n❓ Duplas (2 min): 'Em que tipo de sistema threshold baixo? E alto?' (baixo: FAQs; alto: dados numéricos)\n⚠️ Configurar uma vez e nunca revisar. Threshold ótimo muda.\n➡️ Outro cache: embeddings.", 20, T, "Trade-off central")

s = content_slide("Cache de Embeddings", [
    "Embedding de um texto é determinístico (mesmo modelo → mesmo vetor)",
    "Re-embeddar mesmo documento a cada query = desperdício",
    "",
    "Solução: cache de embeddings (hash do texto → vetor)",
    "Útil em RAG: documentos indexados não precisam re-embedding",
    "",
    "TTL: invalidar quando o modelo de embedding muda",
    "Troca de modelo = vetores incompatíveis = busca retorna lixo",
], "📖 Embeddings são determinísticos. Re-embeddar documentos em cada query é desperdício. Cachear por hash do texto. Em RAG com milhões de docs, economiza muita computação.\n💡 Analogia: biblioteca que re-cataloga o mesmo livro toda vez. Bobo. Cataloga uma vez.\n⚠️ Trocar modelo de embedding sem invalidar cache. Vetores incompatíveis = lixo.\n➡️ E os resultados de tools?", 21, T, "Evitar re-embedding")

s = content_slide("Cache de Tool Results", [
    "Tools idempotentes (GET, search, lookup) → não mudam frequentemente",
    "",
    "Exemplos com TTL:",
    "  • get_weather('Rio de Janeiro') → cache 30 min",
    "  • search_docs('API reference') → cache 24h",
    "  • get_price('AAPL') → cache 5 min",
    "",
    "Cuidado: tools não-idempotentes (POST, mutate) NUNCA cacheadas",
    "Sempre verificar idempotência antes de cachear",
], "📖 Tool results são cache subaproveitado. search_docs('configurar Redis') não muda → cachear por horas. Reduz latência e custo de API externa. MAS só tools idempotentes.\n💡 Analogia: cache do navegador. Imagens/CSS cacheiam (estáticos). POST de formulário nunca (muta estado).\n❓ 'Quais tools de vocês são idempotentes e cacheáveis?'\n⚠️ Cachear 'enviar email'. Disaster. SEMPRE verificar idempotência.\n➡️ Cache tem um inimigo: tempo.", 22, T, "Cache de tools")

s = content_slide("Invalidação e Consistência", [
    "Cache stale = resposta desatualizada",
    "",
    "Estratégias de invalidação:",
    "  • TTL (time-to-live) — simples mas impreciso",
    "  • Event-driven — invalidar quando fonte muda",
    "  • Versionada — chave inclui versão da fonte",
    "",
    "Consistência eventual: cache diverge da fonte por um tempo",
    "Em agentes: contexto muda a cada step → cache por step",
], "📖 'There are only two hard things in CS: cache invalidation and naming things.' TTL é simples mas impreciso. Event-driven é preciso mas complexo. Em agentes, contexto muda a cada step → cache por step.\n💡 Analogia: cardápio de restaurante. TTL = impresso semanalmente (pode estar desatualizado). Event-driven = garçom pergunta cozinha.\n⚠️ Cachear sem invalidar quando fonte muda. Cliente recebe preço antigo.\n➡️ E se alguém injetar dados errados?", 23, T, "Cache invalidation é difícil")

s = content_slide("Cache Poisoning e Mitigação", [
    "Cache poisoning: usuário malicioso injeta resposta errada no cache",
    "",
    "Cenário de ataque:",
    "  query maliciosa → resposta manipulada → cacheada → servida a outros",
    "",
    "Mitigações:",
    "  • Separar cache por tenant/usuário",
    "  • Não cachear respostas de usuários não confiáveis",
    "  • Validar resposta antes de cachear",
    "  • Rate limit por tenant",
], "📖 Risco real e subestimado. Sistema multi-tenant com cache compartilhado: tenant malicioso injeta resposta errada ('capital do Brasil = São Paulo'), todos recebem.\n💡 Analogia: placa falsa no museu ('Mona Lisa — 1990'). Outros leem e acreditam. Solução: placas só por curadores.\n⚠️ Cache compartilhado entre tiers gratuito e pago. Gratuito envenena pago.\n➡️ Vamos arquitetar múltiplos caches em camadas.", 24, T, "Segurança de cache")

s = content_slide("Estratégia de Camadas (L1/L2/L3)", [
    "L1: Cache local em memória (~1ms, por instância)",
    "L2: Cache compartilhado em Redis (~5ms, por cluster)",
    "L3: Prompt caching do provider (~50ms, Anthropic/OpenAI)",
    "",
    "Fluxo: L1 → L2 → L3 → LLM",
    "Cada camada: latência incremental mas reduz custo",
    "Para no primeiro hit",
], "📖 Camadas combinam velocidade e cobertura. L1 (memória da instância) mais rápida mas só serve aquela instância. L2 (Redis) compartilhada entre réplicas. L3 (prompt cache do provider) cacheia prefixo.\n💡 Analogia: memória de computador. L1 = registrador (instantâneo), L2 = RAM, L3 = SSD, LLM = internet (lenta).\n➡️ Hora da DEMO.", 25, T, "Camadas de cache")

s = code_slide("DEMO: Cache Semântico em Ação", (
    "$ python main.py --no-cache\n"
    "Query 1: 'Qual a capital do Brasil?'\n"
    "  → LLM call... 2.3s · custo: R$0.05\n"
    "Query 2: 'Capital brasileira?'\n"
    "  → LLM call... 2.1s · custo: R$0.05\n"
    "...\n"
    "Total (10 queries): 24.8s · R$0.50\n"
    "\n"
    "$ python main.py --semantic-cache\n"
    "Query 1: 'Qual a capital do Brasil?'\n"
    "  → cache MISS · LLM call... 2.3s · R$0.05\n"
    "Query 2: 'Capital brasileira?'\n"
    "  → cache HIT (sim=0.94) · 0.05s · R$0.001\n"
    "...\n"
    "Total (10 queries): 10.2s · R$0.21\n"
    "\n"
    "Redução: 58% custo · 60% latência · 60% hit rate"
), "📖 Demo ao vivo. Primeira execução sem cache: cada query demora ~2-3s e custa R$0,05. Segunda com cache semântico: primeira vai ao LLM (miss), similares são hits (~50ms, ~R$0,001).\n💡 Analogia: ligar ar-condicionado pela primeira vez vs deixá-lo ligado. Primeiro ciclo gasta energia, depois só mantém.\n❓ 'Se a API falhar, qual o plano B?' (screenshot pré-gravado)\n⚠️ Demo falha por rate limit. Sempre pré-popular cache.\n➡️ Vamos quantificar o ganho.", 26, T, "DEMO ao vivo")

s = content_slide("Medindo o Impacto (Antes/Depois)", [
    "Antes (sem cache):",
    "  10 queries × R$0,05 = R$0,50",
    "  Latência média: 3s",
    "",
    "Depois (com cache semântico):",
    "  4 queries × R$0,05 + 6 hits × R$0,001 = R$0,21",
    "  Latência média: 1.2s",
    "",
    "Redução de custo: 58%",
    "Redução de latência: 60%",
    "Cache hit rate: 60%",
], "📖 Os números falam. 58% custo, 60% latência, hit rate 60%. Em produção: 30-70% dependendo do domínio. FAQs/doc técnica: até 80%. Monitorem hit rate continuamente — se cair, investiguem.\n💡 Analogia: medir economia de combustível após ajustar motor. Precisa do before/after.\n⚠️ Medir uma vez e assumir que fica estável. Varia com usuários, produto, sazonalidade.\n➡️ Pergunta rápida sobre a demo.", 27, T, "Quantificar ganho")

s = content_slide("Pergunta da DEMO", [
    "Discussão em duplas (2 min):",
    "",
    "1. O que acontece se o threshold for muito baixo?",
    "2. E se o usuário fizer pergunta factual que mudou ('quem é o presidente')?",
    "3. Como detectar que o cache está stale?",
    "",
    "Respostas esperadas:",
    "  1. Falsos positivos (resposta errada servida)",
    "  2. Cache stale — daí importância do TTL por categoria",
    "  3. A/B testing, feedback do usuário, queda de qualidade",
], "📖 Deixar duplas discutirem. Lição: cache semântico não é 'set and forget'. Precisa de governança contínua.\n❓ Deixar 2-3 duplas compartilharem.\n➡️ Exercício rápido.", 28, T, "Discussão em duplas")

s = exercise_slide("Exercício: Quando Cache Semântico Falha", [
    "Para cada cenário, vote: cache semântico FUNCIONA (F) ou FALHA (X)?",
    "",
    "1. 'Preço das ações da Apple hoje'  →  X (temporal)",
    "2. 'Qual a capital da França?'  →  F (estável)",
    "3. 'Resuma o último email recebido'  →  X (contextual)",
    "4. 'Como fazer um loop em Python?'  →  F (estável)",
    "5. 'Qual a previsão do tempo amanhã?'  →  X (temporal)",
    "",
    "Lição: cache funciona para conhecimento ESTÁVEL.",
    "Falha para TEMPORAL (preço) ou CONTEXTUAL (último email).",
    "Para esses: TTL curto ou não cachear.",
], "📖 Votação rápida. Cache semântico funciona para ESTÁVEL (geografia, programação, conceitos). Falha para TEMPORAL/CONTEXTUAL.\n💡 Analogia: enciclopédia. Excelente para 'capital da França'. Péssima para 'preço da Apple hoje'.\n⚠️ Aplicar cache sem categorizar. Configure TTL por categoria.\n➡️ Intervalo. Voltamos com routing.", 29, T, "Votação rápida")

s = section_slide(3, "Model Routing & Otimização")
add_notes(s, "Depois do caching, a segunda maior otimização: routing. Nem toda query precisa do modelo mais caro.")

s = content_slide("Roteamento por Complexidade", [
    "Nem toda query precisa do modelo mais caro",
    "",
    "Haiku ($0.25/M): classificação, extração, Q&A simples",
    "Sonnet ($3/M): raciocínio moderado, tool use, análise",
    "Opus ($15/M): raciocínio complexo, multi-step, crítico",
    "",
    "Roteamento automático: classifier → modelo apropriado",
    "Economia: 70% das queries podem usar Haiku",
    "Haiku é 60x mais barato que Opus",
], "📖 70% das queries em produção são 'simples'. Podem ir para Haiku ($0.25/M) vs Opus ($15/M) — 60x mais barato. Classifier decide. Reduz custo médio 40-60% sem perda mensurável.\n💡 Analogia: pronto-socorro. Triagem: enfermeiro (Haiku), geral (Sonnet), especialista (Opus).\n⚠️ Usar Opus para tudo. Erro mais comum e mais caro.\n➡️ Mas como classificar complexidade SEM gastar tokens?", 31, T, "Rotear por complexidade")

s = content_slide("Como Medir Complexidade", [
    "Heurísticas baratas: tamanho do input, nº de tools, palavra-chave",
    "Classifier leve (Haiku): 'esta query é simples ou complexa?'",
    "Histórico: queries similares já classificadas (cache de roteamento)",
    "",
    "Cascata: tentar Haiku primeiro, escalar para Sonnet se falhar",
    "  Captura melhor dos dois mundos: barato para simples, capaz para complexas",
    "",
    "Trade-off: classifier errado = resposta ruim ou custo desnecessário",
], "📖 Classifier não pode custar mais que a economia. Heurísticas são de graça. Haiku classifier é quase grátis. Cascata é elegante: tenta Haiku, escala se baixa confiança.\n💡 Analogia: resolver problema matemático. Tente fórmula simples primeiro. Se não, chame o professor.\n❓ 'Que heurística você usaria para classificar queries de vocês?'\n⚠️ Classifier muito complexo (caro) ou muito simples (erra muito).\n➡️ Outra otimização: batching.", 32, T, "Classificar sem gastar tokens")

s = content_slide("Batching de Requests", [
    "N requisições independentes → 1 chamada de API com N prompts",
    "",
    "Batch de embeddings: 100 textos em 1 chamada",
    "Batch de classificações: 10 queries em 1 prompt",
    "",
    "Reduz: nº de requests (rate limit), overhead de rede, custo de TTFT",
    "",
    "Limitação: só funciona para requisições INDEPENDENTES",
    "Em loop serial, cada step depende do anterior — não batchear",
], "📖 Batching em agentes: (1) embeddings — 100 textos em 1 chamada; (2) classificações — 10 tickets em 1 prompt. Reduz latência, custo, respeita rate limit. MAS só para independentes.\n💡 Analogia: ir ao mercado. 10 viagens de 1 item vs 1 viagem de 10 itens.\n⚠️ Tentar batchear steps de loop ReAct. Não funciona — cada step depende do anterior.\n➡️ Para reduzir latência real: speculative decoding.", 33, T, "Agrupar requests")

s = content_slide("Speculative Decoding", [
    "Modelo pequeno (draft) gera rascunho rápido",
    "Modelo grande (target) verifica rascunho em paralelo",
    "",
    "Se rascunho certo: enorme economia de latência",
    "Se errado: target corrige (pequeno overhead)",
    "",
    "Em média: latência cai 2-3x",
    "Variação para agentes: prever próxima action com modelo leve",
    "Fonte: arXiv:2211.17192",
], "📖 Draft gera K tokens. Target verifica em paralelo. Se corretos, economia. Se não, corrige a partir do primeiro erro. Latência cai 2-3x.\n💡 Analogia: estagiário escreve rascunho, sênior revisa. Se bom, aprova rápido. Se não, reescreve.\n⚠️ Achar que é simples de implementar. Exige vLLM/TGI. Não faça do zero.\n➡️ Streaming é mais simples e impacta UX.", 34, T, "Reduz latência de geração")

s = content_slide("Streaming para Latência Percebida", [
    "Sem streaming: usuário espera 10s, recebe tudo de uma vez",
    "Com streaming: usuário vê tokens em ~0.5s, recebe gradualmente",
    "",
    "Latência TOTAL é a mesma",
    "Latência PERCEBIDA cai drasticamente",
    "",
    "Em agentes: streamar thought/answer, não tool calls",
    "",
    "V/F: 'Streaming reduz latência total.' → FALSO",
    "Reduz percebida, não real",
], "📖 Pegadinha clássica. Streaming NÃO reduz latência total — LLM gera mesmos tokens no mesmo tempo. Mas reduz PERCEBIDA: usuário vê algo em 500ms vs 10s de tela vazia. Em agentes, streamar resposta final.\n💡 Analogia: carregamento de página. Barra vazia por 5s parece eterna. Barra preenchendo parece rápida.\n❓ 'Vocês streamam respostas em produção?'\n⚠️ Implementar streaming mas não para tool calls. Usuário espera sem feedback.\n➡️ Panorama: distilação e fine-tuning.", 35, T, "UX perceived latency")

s = content_slide("Distilação e Fine-Tuning (Panorama)", [
    "Distilação: treinar modelo menor com saídas de modelo maior",
    "Fine-tuning: especializar modelo para domínio (reduz prompt size)",
    "",
    "Quando vale:",
    "  • Volume alto (>100k requests/mês)",
    "  • Domínio estável (não muda toda semana)",
    "",
    "Quando NÃO vale:",
    "  • Volume baixo (custo de treino supera economia)",
    "  • Domínio muda frequentemente",
    "",
    "Aprofundamento: ETHAGT08 (Fine-tuning), ETHAGT15 (DSPy)",
], "📖 Otimizações profundas. Distilação: treine pequeno com saídas do grande. Fine-tuning: especialize para domínio. Ambos valem quando volume alto E domínio estável.\n💡 Analogia: distilar = aprendiz copia mestre. Fine-tune = generalista vira cardiologista (menos contexto necessário).\n⚠️ Fine-tunar com volume baixo. Custo de treino supera economia. Mínimo 100k/mês.\n➡️ Vamos sistematizar tudo.", 36, T, "Otimizações profundas")

s = content_slide("Matriz de Otimização", [
    "Caching: alto impacto custo/latência, baixo esforço, baixo risco",
    "Routing: médio impacto custo, médio esforço, médio risco",
    "Batching: médio impacto latência, baixo esforço, baixo risco",
    "Speculative: alto impacto latência, alto esforço, alto risco",
    "Fine-tuning: alto impacto custo, alto esforço, alto risco",
    "",
    "Regra: comece barato e de baixo risco (caching)",
    "Escale para caro/risco conforme ROI justifique",
], "📖 Matriz ajuda a priorizar. Canto sup. esquerdo (alto impacto, baixo esforço): caching — comece aqui. Routing/batching em seguida. Speculative/fine-tuning só depois de esgotar as anteriores.\n💡 Analogia: investir. ETFs (caching) → ações (routing) → derivativos (speculative).\n⚠️ Pular caching e ir para fine-tuning. Desperdício de risco.\n➡️ Pergunta de reflexão.", 37, T, "Impacto vs esforço")

s = content_slide("Pergunta: Como Decidir?", [
    "Se você tivesse 1 semana para otimizar um agente em produção,",
    "por onde começaria?",
    "",
    "Resposta esperada: caching primeiro (maior ROI, menor esforço)",
    "",
    "Em que cenário routing é melhor que caching?",
    "Resposta: queries majoritariamente únicas (baixa repetição)",
    "mas com clara separação de complexidade",
    "",
    "Discussão aberta (2 min)",
], "📖 Resposta esperada: caching primeiro. Em 1 semana, dá para implementar cache semântico, medir hit rate, ajustar threshold. Routing é segunda prioridade.\n❓ Deixar 2-3 alunos compartilharem prioridades.\n➡️ Exercício rápido de routing.", 38, T, "Reflexão sobre priorização")

s = exercise_slide("Exercício: Routing por Complexidade", [
    "Para cada tarefa, vote: Haiku (H), Sonnet (S) ou Opus (O)?",
    "",
    "1. 'Classificar este ticket como bug/feature/question'  →  H",
    "2. 'Resolver bug complexo de concorrência'  →  O",
    "3. 'Resumir documento de 5 páginas'  →  H",
    "4. 'Planejar arquitetura de sistema multi-agente'  →  O",
    "5. 'Extrair entidades deste email'  →  H",
    "6. 'Decidir qual tool chamar e por quê'  →  S",
    "",
    "Lição: maioria das tarefas é Haiku. Opus com parcimônia.",
], "📖 Votação rápida. 1 Haiku (classificação), 2 Opus (concorrência), 3 Haiku (sumarização), 4 Opus (arquitetura), 5 Haiku (extração), 6 Sonnet (decisão com tools).\n❓ Votação de mãos para cada tarefa.\n⚠️ Usar Opus para tudo por 'medo de errar'. Meça % de Haiku rejeitada — geralmente baixo.\n➡️ Resumo da seção.", 39, T, "Votação")

s = content_slide("Resumo — Routing & Otimização", [
    "Routing: Haiku para 70% das queries → reduz custo 40-60%",
    "Batching: agrupar requests independentes → reduz latência e rate limit",
    "Speculative: draft + target → reduz latência 2-3x",
    "Streaming: reduz latência percebida (não real)",
    "Fine-tuning/Distilação: alto esforço, alto ganho, só volume alto",
    "",
    "Ordem de implementação:",
    "  caching → routing → batching → speculative → fine-tuning",
], "📖 Ordem de implementação é crucial: caching primeiro (maior ROI), depois routing, depois batching. Speculative e fine-tuning só depois.\n➡️ Agora vamos além de otimizar chamadas. Vamos DISTRIBUIR agentes.", 40, T, "Síntese da seção D")

s = section_slide(4, "Distribuição de Agentes")
add_notes(s, "Agora vamos além de otimizar chamadas individuais. Vamos DISTRIBUIR: réplicas, sharding, balanceamento, coordenação.")

s = comparison_slide("Stateless vs Stateful Workers",
    "Stateless", [
        "Cada request é independente",
        "Sem sessão entre chamadas",
        "Escala trivialmente (qualquer réplica atende)",
        "Ex: API REST simples",
    ],
    "Stateful", [
        "Mantém sessão/contexto entre chamadas",
        "Necessário para agentes",
        "Mais complexo de escalar",
        "Ex: agente com loop, memória",
    ],
    "📖 Stateless = qualquer réplica atende. Stateful = precisa da réplica com estado. Agentes são naturalmente stateful. Solução moderna: 'stateless-friendly' — worker stateless + checkpoint externo (Redis/Postgres).\n💡 Analogia: caixa de banco (stateless, atende qualquer cliente) vs gerente de conta (stateful, só 'seus' clientes). Híbrido: qualquer gerente pode atender consultando arquivo central.\n⚠️ Estado em memória sem checkpoint. Deploy/crash = perda total.\n➡️ Para distribuir: sharding.", 41, T, "Distinção fundamental")

s = content_slide("Sharding por Tenant", [
    "Sharding: dividir carga por chave (tenant_id, user_id, domínio)",
    "",
    "Shard 1: tenants A, B",
    "Shard 2: tenants C, D",
    "Shard 3: tenant 'hot' Z (dedicado)",
    "",
    "Cada shard tem DB isolado → isolamento de dados e custo",
    "Vantagem: isolamento, custo rastreável, sem contenção",
    "Desafio: rebalanceamento quando shards desequilibram",
    "",
    "Ver diagrama: 12-Diagrams/ETHAGT14/sharding.mmd",
], "📖 Sharding é particionamento por chave. Múltiplos DBs menores, cada um responsável por tenants. Vantagens: isolamento, custo rastreável, escalabilidade. Desafio: hot shard (tenant gigante), rebalanceamento. Kleppmann cap. 6.\n💡 Analogia: banco com múltiplas agências. Cada uma atende região. Agência maior no centro (hot shard para premium).\n❓ 'Em qual chave vocês fariam sharding?'\n⚠️ Sharding mal escolhido (chave não uniforme). Tenant gigante sobrecarrega 1 shard.\n➡️ Com shards, precisamos de réplicas dentro de cada shard.", 42, T, "Particionamento")

s = content_slide("Replica e Balanceamento de Carga", [
    "Réplica = cópia idêntica do worker para atender mais requisições",
    "Load balancer distribui requisições entre réplicas",
    "",
    "Com N réplicas: throughput × N (se stateless)",
    "Com stateful: sticky sessions ou estado externo",
    "",
    "Health checks: réplica doente é removida do pool",
    "Auto-scaling: adiciona/remove réplicas com base na carga",
], "📖 Réplicas aumentam throughput horizontalmente. 1 réplica = X usuários. 10 réplicas = 10X (se stateless). LB é o maestro: distribui, monitora, remove. Auto-scaling adiciona/remove. Importante: DB compartilhado ainda é gargalo — sharding é complementar.\n💡 Analogia: call center. 10 atendentes em paralelo. PBX distribui chamadas.\n⚠️ Adicionar réplicas sem checar gargalo real. Se DB é gargalo, réplicas não ajudam.\n➡️ Quais algoritmos de balanceamento?", 43, T, "Throughput horizontal")

s = content_slide("Estratégias de Balanceamento", [
    "Round-robin: distribuição circular, simples",
    "Least-connections: envia para réplica com menos conexões ativas",
    "Weighted: réplicas mais potentes recebem mais carga",
    "Latency-based: envia para réplica com menor latência",
    "Hash-based (consistent hashing): mesmo tenant → mesma réplica",
    "",
    "Para agentes: least-connections ou consistent hashing",
], "📖 Para agentes, mais úteis: Least-connections (envia para menos ocupada) e Consistent hashing (mesmo tenant → mesma réplica, preserva estado em memória).\n💡 Analogia: round-robin = fila única. Least-connections = fila mais curta. Consistent hashing = seu caixa de sempre.\n⚠️ Round-robin com requisições de duração variável. Algumas sobrecarregadas, outras ociosas.\n➡️ Quando precisamos de coordenação?", 44, T, "Algoritmos de LB")

s = content_slide("Coordenação: Consensus e Leader Election", [
    "Quando múltiplos workers precisam concordar: consensus",
    "Raft/Paxos: algoritmos de consenso distribuído",
    "Leader election: um worker é 'líder', outros são 'followers'",
    "",
    "Casos de uso em agentes:",
    "  • Orquestrador (leader) → workers (followers)",
    "  • Distribuição de tarefas sem duplicação",
    "  • Locks distribuídos (apenas um edita recurso por vez)",
    "",
    "Ferramentas: etcd, ZooKeeper, Redis Redlock",
], "📖 Em sistemas distribuídos, às vezes workers precisam concordar. Exemplos: (1) quem é orquestrador principal (leader election); (2) distribuir tarefas sem duplicação; (3) locks distribuídos — apenas um agente edita recurso. Ferramentas: etcd (padrão K8s), ZooKeeper, Redis Redlock.\n💡 Analogia: leader election = escolher líder de equipe. Raft/Paxos = regras de votação.\n⚠️ Implementar consenso do zero. Use etcd/ZooKeeper/Redlock — testados em batalha.\n➡️ E para usuários globais?", 45, T, "Consenso distribuído")

s = content_slide("Multi-Region e Autoscaling", [
    "Multi-region: reduzir latência para usuários globais",
    "Latência inter-região: ~50-100ms (vs <5ms intra-região)",
    "Desafio: replicação de estado entre regiões",
    "",
    "Autoscaling: HPA (Horizontal Pod Autoscaler) no Kubernetes",
    "KEDA: autoscaling baseado em eventos (fila, métricas custom)",
    "Escala para zero: serverless quando não há tráfego",
    "",
    "Métricas de escala: CPU, memória, tamanho de fila, latência",
], "📖 Multi-region para latência global. Usuário na Ásia bate em região asiática. Desafio: replicar estado (50-100ms por sync). HPA adiciona pods com CPU. KEDA escala por eventos (fila cresceu?). Importante: escalonar por métricas de negócio (fila, latência) melhor que só CPU.\n💡 Analogia: franquia global. McDonald's SP atende brasileiros, Tóquio atende japoneses. Cada um com estoque próprio.\n⚠️ Autoscaling só por CPU. Pode ter CPU baixa mas fila crescendo (I/O bound).\n➡️ Pergunta para discussão.", 46, T, "Escala geográfica e automática")

s = content_slide("Pergunta: Stateless é Sempre Preferível?", [
    "Stateless é sempre preferível? Justifique.",
    "",
    "Casos onde stateful é necessário:",
    "  • Sessões longas com contexto acumulado",
    "  • WebSockets / streaming",
    "  • Estado em memória para baixa latência",
    "",
    "Como tornar stateful mais 'stateless-friendly'?",
    "  checkpoint externo + restore on-demand",
    "",
    "Discussão aberta (2 min)",
], "📖 Falso. Há casos legítimos para stateful. A pergunta certa: 'como tornar stateful mais resiliente?'. Resposta: checkpoint externo + restore on-demand.\n❓ Duplas (2 min): 'Pensem em um caso onde stateful é claramente melhor.'\n⚠️ Dogma 'stateless sempre'. Avalie caso a caso.\n➡️ Resumo da seção.", 47, T, "Trade-off stateless/stateful")

s = content_slide("Resumo — Distribuição", [
    "Stateless-friendly: worker stateless + checkpoint externo",
    "Sharding: particionar por tenant/user/domínio",
    "Réplicas: throughput × N (com load balancer)",
    "Balanceamento: least-connections ou consistent hashing",
    "Coordenação: leader election (Raft/Paxos), locks distribuídos",
    "Multi-region + autoscaling: latência global + escala automática",
    "",
    "Arquitetura típica: clients → LB → shards → DBs isolados",
], "📖 Recapitulando distribuição. Arquitetura típica: clients → load balancer → shards (cada um com N réplicas) → DBs isolados. Tudo monitorado, com autoscaling. Coordenado por etcd/Redis. Estado externalizado.\n➡️ Vamos à infraestrutura que sustenta tudo.", 48, T, "Síntese da seção E")

s = section_slide(5, "Infraestrutura para Agentes")
add_notes(s, "5 slides de panorama: Kubernetes, serverless vs dedicado, GPUs, service mesh, custo total. Rápido — ideia é panorama.")

s = content_slide("Kubernetes para Agentes", [
    "Pods: unidade mínima (1 agente worker por pod)",
    "Deployments: réplicas + rolling updates",
    "HPA: autoscaling horizontal baseado em CPU/métricas",
    "Services: load balancing interno",
    "Ingress: entrada de tráfego externo",
    "ConfigMaps/Secrets: configuração e secrets",
    "",
    "Vantagem: padronização, portabilidade entre clouds",
], "📖 K8s é padrão para orquestração. Cada pod roda 1 worker. Deployments gerenciam réplicas + rolling updates. HPA adiciona/remove pods. Services fazem LB interno. Vantagem: portabilidade (AWS, GCP, Azure, on-prem).\n💡 Analogia: K8s é orquestra. Cada pod = músico. Deployment = maestro. HPA contrata/demite conforme a peça.\n⚠️ Usar K8s sem necessidade. Para <100 usuários, Docker Compose basta.\n➡️ Alternativa: serverless vs dedicado.", 50, T, "Plataforma padrão")

s = comparison_slide("Serverless vs Dedicado",
    "Serverless (Lambda, Cloud Run)", [
        "Escala para zero",
        "Paga por execução",
        "Cold start: 1-5s de latência inicial",
        "Ideal: tráfego esporádico, baixo volume",
    ],
    "Dedicado (EC2, ECS, EKS)", [
        "Sempre ativo, sem cold start",
        "Custo fixo, independente de uso",
        "Sem latência inicial",
        "Ideal: tráfego constante, alto volume",
    ],
    "📖 Serverless ótimo para cargas esporádicas. Mas cold start (1-5s) mata UX de agentes. Em agente de 10s, 5s de cold start é 50% de overhead. Para agentes em produção, dedicado geralmente melhor.\n💡 Analogia: serverless = táxi (paga por viagem, mas espera motorista). Dedicado = carro próprio (sempre pronto, paga mesmo sem usar).\n❓ 'Vocês usam serverless para algum agente?'\n⚠️ Serverless para agente de chat ao vivo. Cold start irrita.\n➡️ E se quiser inferência local?", 51, T, "Modelos de deploy")

s = content_slide("GPUs para Inferência Local", [
    "vLLM, TGI: servidores de inferência otimizados",
    "Modelos open-source: Llama, Mistral, Qwen",
    "",
    "Vantagem: custo fixo, sem rate limit, baixa latência, privacidade",
    "Desafio: custo de GPU ($1-8/hora), manutenção, qualidade vs Claude/GPT",
    "",
    "Quando vale: alto volume (>1M tokens/dia), privacidade, latência crítica",
    "Quando não vale: baixo volume, qualidade prioritária",
], "📖 Self-hosting para alto volume. vLLM/TGI otimizam inferência. Llama/Mistral/Qwen aproximam-se de Claude/GPT. Vantagens: custo fixo, sem rate limit, dados ficam com você. Desafios: GPU cara, ops, qualidade geralmente abaixo dos fechados.\n💡 Analogia: alugar vs comprar casa. API = aluguel (flexível). Self-hosted = comprar (alto custo inicial, manutenção).\n⚠️ Self-hosting por 'economia' com volume baixo. GPU + ops supera API.\n➡️ Com múltiplos serviços, precisamos de service mesh.", 52, T, "Inferência local")

s = content_slide("Service Mesh e Custo Total", [
    "Service mesh (Istio, Linkerd): camada de comunicação transparente",
    "  Features: mTLS, retries, circuit breakers, observabilidade",
    "  Traffic splitting: canary deployments",
    "  Para agentes: tracing distribuído entre chamadas",
    "",
    "Custos escondidos (além do LLM):",
    "  Memória: Redis (~$50-500/mês)",
    "  Rede: transferência entre regiões (~$50-200/mês)",
    "  Storage: vetores, logs, traces (~$20-200/mês)",
    "  Compute: K8s nodes, GPUs (~$200-2000/mês)",
    "  Observabilidade: Datadog, LangSmith (~$50-300/mês)",
    "",
    "Regra: LLM é ~60-70% do total, infra é 30-40%",
], "📖 Service mesh resolve comunicação entre serviços. Para multi-agente: mTLS, retries, circuit breakers, tracing distribuído (request atravessa A→B→C). Quanto a custo: muitas empresas calculam só LLM e se surpreendem. LLM é 65%, infra é 35%. Calculem TCO.\n💡 Analogia: service mesh = sistema nervoso. Conecta órgãos com segurança e observação.\n⚠️ Esquecer custo de observabilidade. Pode chegar a $300+/mês.\n➡️ Agora o tema mais importante: FinOps.", 53, T, "Comunicação + custo total")

s = section_slide(6, "FinOps de Agentes")
add_notes(s, "FinOps é a disciplina de gerenciar custo em escala. Sem FinOps, vocês otimizam técnicas mas perdem no orçamento.")

s = content_slide("Orçamento como Guardrail", [
    "Orçamento por execução: máximo de $X por run do agente",
    "Orçamento por usuário: máximo de $Y/mês por usuário",
    "Orçamento por tenant: máximo de $Z/mês por empresa/cliente",
    "",
    "Implementação: contador de custo acumulado → circuit breaker",
    "",
    "Sem orçamento: um usuário pode gerar $1000 em uma noite",
    "Com orçamento: agente para e avisa quando limite é atingido",
], "📖 Orçamento evita desastre financeiro. 3 níveis: por execução (evita loops infinitos), por usuário (evita abuso individual), por tenant (evita que cliente consuma tudo). Implementação: contador em Redis, circuit breaker verifica antes de cada chamada.\n💡 Analogia: limite de cartão. Sem limite, fraude te quebra. Com limite, contido.\n❓ 'Vocês têm orçamento configurado em produção?'\n⚠️ Orçamento só agregado mensal. Descobre estouro só no fim do mês. Precisa por execução.\n➡️ Para orçamento funcionar, precisa medir granularmente.", 55, T, "Guardrail financeiro")

s = content_slide("Medição Granular de Custo", [
    "Por step: cada iteração do loop tem custo próprio",
    "Por tool: qual tool gera mais tokens de contexto?",
    "Por token: input vs output (preços diferentes)",
    "Por modelo: se há routing, custo varia",
    "",
    "Tagging: cada execução carrega tags (usuário, tenant, categoria, feature)",
    "Dashboard: custo por feature, por tenant, por modelo, por tool",
], "📖 Medição agregada é insuficiente. Precisamos granular: qual step mais consome? Qual tool mais infla contexto? Qual feature mais custa? Exige tagging. Empresa real descobriu 80% do custo em 1 feature (RAG com contexto gigante); otimizaram e cortaram 60%.\n💡 Analogia: extrato do cartão. Total mensal não ajuda. Categorizado (alimentação, transporte) revela onde cortar.\n⚠️ Medir só por execução sem tags. Não sabe qual feature/tenant mais consome.\n➡️ O trade-off fundamental.", 56, T, "Custo granular")

s = content_slide("Trade-offs: Custo × Latência × Qualidade", [
    "Os 3 eixos estão em tensão:",
    "  Reduzir custo → pode aumentar latência ou reduzir qualidade",
    "  Reduzir latência → pode aumentar custo ou reduzir qualidade",
    "  Aumentar qualidade → pode aumentar custo e latência",
    "",
    "Não existe otimização gratuita (pick two)",
    "",
    "Haiku: barato e rápido, mas menos capaz",
    "Opus: capaz, mas caro e lento",
    "",
    "Decisão: qual eixo é mais importante para SEU caso?",
], "📖 Trade-off fundamental. Os 3 eixos em tensão permanente. Não existe almoço grátis — você escolhe 2, o terceiro sofre. Em saúde, qualidade vence. Em suporte de baixa complexidade, custo vence. Em chat ao vivo, latência vence.\n💡 Analogia: projeto de engenharia. Rápido, barato, bom — escolha dois.\n❓ 'No sistema de vocês, qual eixo é prioritário?'\n⚠️ Prometer os 3 eixos. Quebra a física do sistema.\n➡️ Como operacionalizar: fluxo FinOps.", 57, T, "Triângulo de trade-offs")

s = content_slide("Fluxo FinOps", [
    "execução → medir custo granular → tag (usuário/tenant/categoria)",
    "tag → FinOps dashboard",
    "",
    "dashboard → alerta: custo > orçamento?",
    "  sim: circuit breaker aciona",
    "  não: ok",
    "",
    "dashboard → otimização contínua (routing, cache)",
    "",
    "Loop fechado: medir → alertar → otimizar → medir",
    "Ver diagrama: 12-Diagrams/ETHAGT14/finops-flow.mmd",
], "📖 Ciclo fechado. Toda execução gera medição granular. Tagueada. Dashboard dispara alertas quando > orçamento. Informa otimização contínua. Sem este loop, vocês voam cegos.\n💡 Analogia: dieta. Sem balança (medir), sem plano (alertar), sem ajuste (otimizar), engorda sem perceber.\n➡️ Vamos ver o dashboard na prática.", 58, T, "Loop FinOps")

s = content_slide("Dashboards de Custo Contínuo", [
    "Dashboard mínimo (Grafana):",
    "  Custo por dia/semana/mês",
    "  Custo por feature (qual funcionalidade consome mais)",
    "  Custo por tenant (qual cliente consome mais)",
    "  Cache hit rate (efetividade do cache)",
    "  Top 10 execuções mais caras (outliers)",
    "",
    "Alertas: custo diário > threshold → Slack/email",
    "Otimização contínua: identificar anomalias e agir",
    "",
    "Ferramentas: Grafana+Prometheus, Datadog, LangSmith, Helicone",
], "📖 Dashboard não é opcional em produção. Mínimo: agregado, por feature, por tenant, cache hit, top 10 caros. Alertas em Slack. Revela anomalias: 'tenant X gastou 10x mais hoje' → investigar. 'Cache hit caiu 60%→30%' → threshold errado.\n💡 Analogia: painel de carro. Sem isso, dirige no escuro até o motor fundir.\n⚠️ Dashboard só mensal. Descobre estrago tarde demais. Diário com alertas.\n➡️ E se vocês estiverem COBRANDO de clientes?", 59, T, "Dashboard FinOps")

s = content_slide("Pricing para Clientes", [
    "Modelo per-call: cobrar por execução (transparência)",
    "Modelo per-token: repassar custo + margem",
    "Modelo assinatura: valor fixo mensal com limite (previsibilidade)",
    "Modelo freemium: tier gratuito limitado + tier pago",
    "Modelo por valor: cobrar por outcome (resolvido/não resolvido)",
    "",
    "Tensão: custo é variável, pricing precisa ser previsível",
    "Recomendação: assinatura com limite + excedente",
], "📖 Pricing de produto agêntico é desafiador. Custo variável, cliente quer previsibilidade. Modelos: per-call (transparente mas imprevisível), per-token (justo mas complexo), assinatura (previsível mas você absorve risco), freemium (atrai mas pode queimar orçamento), por valor (alinhado mas difícil de medir).\n💡 Analogia: plano de celular. Não paga por minuto. Plano fixo com franquia. Excedente se passar.\n❓ 'Vocês cobrariam por chamada ou assinatura?'\n⚠️ Freemium sem limite. Usuário/bug consome orçamento de aquisição.\n➡️ Implementação técnica: circuit breaker.", 60, T, "Modelos de pricing")

s = code_slide("Circuit Breaker de Custo", (
    "class CostCircuitBreaker:\n"
    "    def __init__(self, budget):\n"
    "        self.budget = budget\n"
    "        self.spent = 0.0\n"
    "\n"
    "    def check(self, estimated_cost):\n"
    "        # fail-safe: bloqueia se exceder\n"
    "        if self.spent + estimated_cost > self.budget:\n"
    "            logger.warning(f'Budget exceeded: {self.spent}')\n"
    "            alert_slack(f'Circuit breaker acionado')\n"
    "            raise BudgetExceeded(\n"
    "                f'Orçamento de R${self.budget} excedido'\n"
    "            )\n"
    "\n"
    "    def add(self, actual_cost):\n"
    "        self.spent += actual_cost\n"
    "\n"
    "# Uso:\n"
    "breaker = CostCircuitBreaker(budget=0.10)\n"
    "breaker.check(estimated_cost=0.02)\n"
    "response = llm(query)\n"
    "breaker.add(actual_cost=response.cost)"
), "📖 Antes de cada chamada, verifica se custo acumulado + estimado excede orçamento. Se exceder, exceção amigável. Após, adiciona custo real. Fail-safe: se não sabe, bloqueie.\n💡 Analogia: disjuntor elétrico. Corrente excede → desliga. Melhor escuro que incêndio.\n⚠️ Circuit breaker sem log. Não sabe quando nem por que acionou.\n➡️ Exercício prático.", 61, T, "Guardrail operacional")

s = exercise_slide("Exercício: Circuit Breaker", [
    "Cenário: sistema multi-tenant com orçamento de R$100/usuário/mês",
    "",
    "Em duplas (3 min código + 2 min compartilhar):",
    "  Escrevam lógica de circuit breaker que:",
    "  1. Verifica custo acumulado do mês para o usuário",
    "  2. Bloqueia execuções se custo mensal exceder",
    "  3. Inclui: alerta, log, exceção amigável",
    "",
    "Bônus: o que acontece se orçamento for compartilhado entre features?",
    "Dica: pesos por feature (chat=1x, RAG=2x, code=3x)",
], "📖 Exercício em duplas. Cenário: orçamento R$100/usuário/mês. A cada chamada, breaker verifica acumulado. Se exceder, bloqueia. Bônus: orçamento compartilhado entre features → pesos ou segregação.\n❓ Após exercício: 'Como vocês trataram a exceção amigável?'\n⚠️ Bloquear sem mensagem clara. Usuário frustrado sem entender.\n➡️ Vamos ao fechamento.", 62, T, "Duplas — circuit breaker")

s = section_slide(7, "Boas Práticas e Anti-Patterns")
add_notes(s, "Síntese em boas práticas (DO) e anti-patterns (DON'T). Depois caso de estudo, resumo, quiz, Q&A.")

s = content_slide("Boas Práticas (DO)", [
    "✓ Comece por caching (maior ROI, menor esforço)",
    "✓ Meça custo desde a primeira linha de código",
    "✓ Defina orçamento por execução desde o dia 1",
    "✓ Use routing para separar queries simples de complexas",
    "✓ Faça checkpoint do estado externamente (stateless-friendly)",
    "✓ Monitore cache hit rate continuamente",
    "✓ Use autoscaling com métricas de negócio (fila, latência), não só CPU",
    "✓ Teste sob carga antes de ir para produção",
], "📖 8 boas práticas essenciais. A ordem importa: caching primeiro, medição desde início, orçamento dia 1, routing para custo, checkpoint para resiliência, monitoramento cache hit, autoscaling inteligente, teste de carga.\n💡 Analogia: checklist de pré-voo. Cada item evita um desastre.\n⚠️ Implementar 7 e esquecer 1 (geralmente 'teste sob carga'). Item esquecido te derruba.\n➡️ Agora o que NÃO fazer.", 64, T, "Checklist verde")

s = content_slide("Anti-Patterns (DON'T)", [
    "✗ Não medir custo (surpresa no fim do mês)",
    "✗ Não definir orçamento (usuário gera $1000 em uma noite)",
    "✗ Usar sempre o modelo mais caro (desperdício)",
    "✗ Não cachear (todas as queries batem no LLM)",
    "✗ Manter estado em memória sem checkpoint (perda em restart)",
    "✗ Ignorar rate limits (429 em cascata)",
    "✗ Não monitorar cache hit rate (cache ineficaz sem saber)",
    "✗ Otimizar antes de medir (premature optimization)",
], "📖 Cada anti-pattern é caso real de desastre. Não medir: $50k surpresa. Sem orçamento: $5k em uma noite. Sempre Opus: 60% desperdício. Sem cache: latência 3x. Sem checkpoint: 100 conversas perdidas por deploy. Ignorando rate limit: indisponível 1h. Sem monitorar cache: 5% hit por meses. Otimizar sem medir: 2 semanas em fine-tuning sem efeito.\n💡 Analogia: dirigir bêbado. Cada anti-pattern é um copo. Eventualmente, acidente.\n❓ 'Qual anti-pattern vocês cometeram?'\n➡️ Caso real aplicando tudo.", 65, T, "Checklist vermelho")

s = content_slide("Caso de Estudo: Assistente Corporativo", [
    "Cenário: assistente corporativo para 5.000 funcionários",
    "",
    "Antes: 1 agente/usuário, sem cache, sempre Sonnet",
    "  Custo: R$15k/mês · Latência p95: 12s · Erro: 5%",
    "",
    "Depois: cache semântico + routing + sharding + FinOps",
    "  Custo: R$4k/mês (−73%) · Latência p95: 4s (−67%) · Erro: 1%",
    "",
    "Decomposição: cache (40%), routing (20%), batching (10%), infra (3%)",
    "Lição: escalabilidade é combinação de técnicas, não bala de prata",
], "📖 Caso real sintetizando tudo. Antes: 1 agente/usuário, sem cache, sempre Sonnet. R$15k/mês, 12s p95, 5% erro. Depois: cache semântico (40%), routing Haiku/Sonnet (20%), batching (10%), otimização de infra (3%). Total: 73% redução custo, 67% latência, erro 1%. Lição: nenhuma técnica sozinha resolve. Foi a COMBINAÇÃO.\n💡 Analogia: emagrecimento. Dieta + exercício + sono. A combinação transforma.\n⚠️ Buscar 'bala de prata'. Não existe. Combinem técnicas.\n➡️ Resumo da aula.", 66, T, "Caso real")

s = content_slide("Resumo da Aula", [
    "Gargalos: latência de LLM, custo de contexto, rate limits, estado distribuído",
    "Caching = primeira linha de defesa (exact, semântico, embeddings, tools)",
    "Routing = separar simples de complexo (Haiku/Sonnet/Opus)",
    "Distribuição = stateless + checkpoint, sharding, replica + balanceamento",
    "Infra = K8s + autoscaling + service mesh",
    "FinOps = orçamento, medição granular, circuit breaker, trade-offs",
    "",
    "Projeto: otimizar custo/latência ≥40% sem perder qualidade",
], "📖 Recapitulando os 6 blocos: gargalos (diagnóstico), caching (maior ROI), routing (separação por complexidade), distribuição (escala horizontal), infra (plataforma), FinOps (sustentabilidade). O projeto consolida tudo. Usem técnicas na ordem: caching → routing → distribuição → FinOps.\n➡️ Checklist final.", 67, T, "Síntese")

s = content_slide("Checklist da Aula", [
    "[ ] Identificou gargalos de escala",
    "[ ] Implementou cache semântico",
    "[ ] Configurou model routing",
    "[ ] Distribuiu agentes com sharding",
    "[ ] Configurou FinOps com circuit breaker",
    "[ ] Discutiu trade-offs custo × latência × qualidade",
    "",
    "Se marcou todos, a aula cumpriu o objetivo.",
    "Labs: Cache Semântico (4h), Sharding por Tenant (5h)",
], "📖 Checagem final. Se marcarem todos, aula cumpriu objetivo. Labs darão prática adicional.\n➡️ Quiz rápido.", 68, T, "Confirmação")

s = exercise_slide("Quiz — Pergunta 1", [
    "Qual é a PRIMEIRA otimização que você deve aplicar a um agente em produção?",
    "",
    "A) Fine-tuning do modelo",
    "B) Caching (semântico + exact match)",
    "C) Multi-region deployment",
    "D) Speculative decoding",
    "",
    "Resposta: B (maior ROI, menor esforço, menor risco)",
], "📖 Resposta B. Caching tem maior ROI (alto impacto, baixo esforço, baixo risco). Fine-tuning (A) é alto esforço. Multi-region (C) para latência global, não primeira otimização. Speculative (D) é complexo. Sempre comecem por caching.\n➡️ Pergunta 2.", 69, T, "Quiz — pergunta 1")

s = exercise_slide("Quiz — Pergunta 2", [
    "Cache semântico com threshold muito baixo (ex: 0.75) causa qual problema?",
    "",
    "A) Aumenta latência",
    "B) Falsos positivos (resposta errada servida como correta)",
    "C) Rate limit excedido",
    "D) Aumenta custo de embedding",
    "",
    "Resposta: B (falsos positivos — o pior tipo de erro, silencioso)",
], "📖 Resposta B. Threshold baixo captura semanticamente próximas mas factualmente diferentes ('iPhone 15' vs 'iPhone 14'). Serve resposta errada como certa. Pior erro — silencioso e confiável.\n➡️ Pergunta 3.", 70, T, "Quiz — pergunta 2")

s = exercise_slide("Quiz — Pergunta 3", [
    "Streaming reduz a latência TOTAL de uma resposta?",
    "",
    "A) Sim, significativamente",
    "B) Sim, marginalmente",
    "C) Não, reduz apenas a latência PERCEBIDA",
    "D) Não, aumenta a latência total",
    "",
    "Resposta: C (reduz percebida, não real — LLM gera mesmos tokens no mesmo tempo)",
], "📖 Resposta C. Streaming NÃO reduz latência total. LLM ainda gera mesmos tokens no mesmo tempo. Mas reduz PERCEBIDA: usuário vê algo em 500ms vs 10s de tela vazia.\n➡️ Perguntas para discussão profunda.", 71, T, "Quiz — pergunta 3")

s = content_slide("Perguntas para Discussão", [
    "1. Em que cenário cache semântico é mais perigoso que útil?",
    "",
    "2. Stateless é sempre preferível? Justifique com um contra-exemplo.",
    "",
    "3. Como você convenceria um CTO a investir em FinOps antes de escalar?",
    "",
    "4. Qual o maior risco de autoscaling baseado apenas em CPU?",
], "📖 Perguntas profundas. (1) Cache perigoso em domínios sensíveis a tempo/contexto; (2) Stateful necessário para sessões longas, WebSockets; (3) Argumento CTO: sem FinOps, bug/pico gera $50k numa noite, ROI imediato; (4) Autoscaling por CPU não captura I/O bound (fila crescendo com CPU baixa).\n❓ Deixar 2-3 alunos responderem cada.\n➡️ Conexão com próximos módulos.", 72, T, "Debate profundo")

s = content_slide("Conexão com Próximos Módulos", [
    "ETHAGT11 — Pré-requisito: K8s, Docker, deploy",
    "ETHAGT05 — Memória de Agentes: estado distribuído em detalhe",
    "ETHAGT12 — AgentOps: observabilidade em profundidade (traces, eval)",
    "ETHAGT15 — Meta-Agentes: auto-aprendizado que pode otimizar custos",
    "ETHAGT90 — Capstone: otimização de custo em sistema real",
], "📖 ETHAGT14 não é ilha. Conecta com ETHAGT15 (meta-agentes — agentes que otimizam outros, incluindo custo), ETHAGT90 (Capstone — projeto real), ETHAGT12 (AgentOps — observabilidade profunda), ETHAGT05 (Memória — estado distribuído em detalhe).\n➡️ Leitura recomendada.", 73, T, "Mapa da especialização")

s = content_slide("Referências e Leitura Recomendada", [
    "Obrigatório: Kleppmann — Designing Data-Intensive Applications",
    "Obrigatório: FinOps Foundation — FinOps for ML/AI",
    "",
    "Recomendado: Speculative Decoding (arXiv:2211.17192)",
    "Recomendado: vLLM docs (serving LLMs)",
    "Recomendado: Richards & Ford — Fundamentals of Software Architecture",
    "",
    "Docs: LiteLLM (roteamento), OpenTelemetry, Grafana",
    "",
    "Projeto: otimizar custo/latência ≥40% + ADR + FinOps dashboard",
    "Labs: Lab 1 (Cache Semântico, 4h), Lab 2 (Sharding, 5h)",
], "📖 Para aprofundar: Kleppmann (cap. 5, 6, 9 essenciais), FinOps Foundation, paper Speculative Decoding, vLLM. Projeto: ADR documentando otimizações + dashboard. Labs dão prática guiada.\n➡️ Q&A.", 74, T, "Leitura")

s = title_slide(
    "Perguntas?",
    "ETHAGT14 — Escalabilidade & Sistemas Distribuídos\nPróxima aula: ETHAGT15 — Meta-Agentes",
    "Obrigado!"
)
add_notes(s, "📖 Hora do Q&A. Deixar turma fazer perguntas. Se ninguém, pergunta inversa: 'Qual técnica vão aplicar primeiro?'\n❓ 'Perguntas?'\n⚠️ Se nenhuma: 'Qual parte foi menos clara?' ou 'Qual técnica vão aplicar primeiro?'")

out = "ETHAGT14-escalabilidade-sistemas-distribuidos.pptx"
prs.save(out)
print(f"Gerado: {out} ({len(prs.slides)} slides)")
