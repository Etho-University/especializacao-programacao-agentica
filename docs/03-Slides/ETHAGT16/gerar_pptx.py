#!/usr/bin/env python3
"""Gerador de PPTX - ETHAGT16: Sociedades de Agentes & Autonomous Research Systems
Universidade Etho - Especializacao em Programacao Agentica"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY  = RGBColor(0x1B, 0x3A, 0x5E)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2F)
DARK     = RGBColor(0x0F, 0x1E, 0x2D)
LIGHT    = RGBColor(0xF7, 0xF8, 0xFA)
MUTED    = RGBColor(0x6B, 0x7A, 0x8F)
SUCCESS  = RGBColor(0x2D, 0x86, 0x59)
DANGER   = RGBColor(0xC0, 0x39, 0x2B)
INFO     = RGBColor(0x29, 0x80, 0xB9)
WARNING  = RGBColor(0xD4, 0xA0, 0x17)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
    return txBox


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_footer(slide, num, total, obj="", acronyms=""):
    if acronyms:
        add_textbox(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.35), acronyms, size=9, color=INFO)
    add_textbox(slide, Inches(0.3), Inches(7.0), Inches(7), Inches(0.4), obj, size=10, color=MUTED)
    add_textbox(slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.4), f"Slide {num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_header(slide, code="ETHAGT16"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.3), Inches(0.02), Inches(3), Inches(0.3), "Universidade Etho", size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(10.5), Inches(0.02), Inches(2.5), Inches(0.3), code, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)


def title_slide(title, subtitle, code, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_textbox(s, Inches(1), Inches(2.3), Inches(11), Inches(1.6), title, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(3.9), Inches(11), Inches(0.8), subtitle, size=20, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.6), Inches(11), Inches(0.5), code, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    if acronyms:
        add_textbox(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4), acronyms, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    return s


def section_slide(num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_textbox(s, Inches(1), Inches(2.5), Inches(2), Inches(2), str(num), size=80, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(3.5), Inches(3.0), Inches(9), Inches(1.5), title, size=32, color=WHITE, bold=True)
    return s


def content_slide(title, bullets, notes, num, total, obj="", subtitle=None, acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=32, color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5), subtitle, size=18, color=MUTED)
    add_bullets(s, Inches(0.7), Inches(1.8 if subtitle else 1.5), Inches(12), Inches(5), bullets, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s


def code_slide(title, code, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, DARK)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=WHITE, bold=True)
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0xFC)
    p.font.name = "Consolas"
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s


def comparison_slide(title, lt, li, rt, ri, notes, num, total, obj="", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=PRIMARY, bold=True)
    add_textbox(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.5), lt, size=20, color=DANGER, bold=True)
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), li, size=15)
    add_textbox(s, Inches(7), Inches(1.5), Inches(6), Inches(0.5), rt, size=20, color=SUCCESS, bold=True)
    add_bullets(s, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5), ri, size=15)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s


def exercise_slide(title, items, notes, num, total, obj="Exercicio", acronyms=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_header(s)
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.7), title, size=28, color=ACCENT, bold=True)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(12), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    box.line.color.rgb = WARNING
    box.line.width = Pt(2)
    add_bullets(s, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.5), items, size=16)
    add_footer(s, num, total, obj, acronyms)
    add_notes(s, notes)
    return s


T = 50

s = title_slide(
    "Sociedades de Agentes & Autonomous Research Systems",
    "Universidade Etho - Especializacao em Programacao Agentica",
    "ETHAGT16 - Fase D - Fronteira"
)

content_slide(
    "Objetivos do Modulo",
    [
        "Objetivo geral: atingir a fronteira do estado da arte",
        "1. Modelar sociedades de agentes (papeis, instituicoes, normas)",
        "2. Aplicar simulacoes multi-agente (Smallville-like)",
        "3. Construir um sistema de pesquisa autonoma",
        "4. Discutir emergencia, convergencia e alinhamento",
        "5. Conhecer a fronteira (AI Scientist, AlphaEvolve, Swarm)",
    ],
    "Cada objetivo e mensuravel. O #3 e o Projeto do modulo e o esboco do Capstone. O #4 (emergencia) e o mais sutil.",
    2, T, "Objetivos"
)

content_slide(
    "Competencias Desenvolvidas",
    [
        "C1 Programacao Agentica -> A (Avancado)",
        "C2 Multi-Agent Systems -> A (Avancado)",
        "C3 MCP & Tool Use -> B (Basico)",
        "C4 Agent Memory -> A (Avancado)",
        "C5 AgentOps & Avaliacao -> I (Intermediario)",
        "C6 Agent Security -> I (Intermediario)",
    ],
    "C1, C2 e C4 chegam ao Avancado. C5 e C6 em Intermediario porque a propria area e imatura em sociedades.",
    3, T, "Competencias"
, acronyms="AgentOps = Agent Operations — operacao e monitoramento de agentes em producao  ·  MCP = Model Context Protocol — Protocolo de Contexto de Modelo")

content_slide(
    "Agenda da Aula",
    [
        "Bloco 1 (25 min): Abertura + Sociedades + Simulacoes + Intervalo",
        "Bloco 2 (25 min): Research Systems + Emergencia + Fechamento",
        "DEMO ao vivo: Mini Sociedade de 5 agentes (Slide 30)",
        "Quiz final: 3 perguntas rapidas (Slide 48)",
    ],
    "Dois blocos. O segundo e onde mora a fronteira. DEMO Mini Sociedade e o ponto alto.",
    4, T, "Agenda"
)

content_slide(
    "O Problema do Agente Isolado",
    [
        "Agente isolado: resolve uma tarefa bem",
        "Sociedade de agentes: resolve o que nenhum agente resolve sozinho",
        "Generative Agents (Smallville): 25 agentes organizam festa de Valentine's Day",
        "AI Scientist (Sakana): conduz pesquisa cientifica do inicio ao fim",
        "Pergunta: o que acontece quando 100 agentes interagem sem supervisao?",
    ],
    "Hook da aula. A resposta honesta: emergencia - pode ser cooperacao ou conluio. Deixar a pergunta no ar.",
    5, T, "Motivacao"
)

content_slide(
    "Por Que Sociedades Agora",
    [
        "2022 - ReAct (padrao de agente individual)",
        "2023 - tool calling nativo + frameworks multi-agent",
        "abr/2023 - Generative Agents / Smallville (arXiv:2304.03442)",
        "ago/2024 - AI Scientist (Sakana, arXiv:2408.06292) + AlphaEvolve (DeepMind)",
        "Confluencia: reasoning + tools + memory + context window + custo baixo",
        "Proxima fronteira: sociedades autonomas",
    ],
    "Tres fatores convergiram: reasoning estruturado, tool calling, custo baixo. Smallville (2023) e AI Scientist (2024) sao os marcos.",
    6, T, "Contexto"
, acronyms="ReAct = Reasoning and Acting — padrao de loop Thought / Action / Observation")

section_slide(1, "Sociedades de Agentes")

content_slide(
    "Do Agente Individual a Sociedade",
    [
        "Nivel 0: Agente individual (LLM + tools + memory)",
        "Nivel 1: Pequeno grupo (2-5 agentes, colaboracao direta)",
        "Nivel 2: Instituicao (papeis fixos, normas explicitas, hierarquia)",
        "Nivel 3: Sociedade (25+ agentes, normas emergentes)",
        "Cada nivel adiciona: complexidade, potencial e risco",
    ],
    "Escada cumulativa. Nivel 2 (instituicao) e pre-requisito para Nivel 3. Sem normas explicitas, sociedade vira caos.",
    8, T, "Sociedades"
, acronyms="LLM = Large Language Model — Modelo de Linguagem de Grande Escala")

content_slide(
    "Papeis em Sociedades de Agentes",
    [
        "Pesquisador: explora, levanta informacoes",
        "Critico: identifica falhas, questiona premissas",
        "Sintetizador: integra contribuicoes divergentes",
        "Revisor: valida qualidade e coerencia",
        "Editor: finaliza, formata, publica",
        "Papeis podem ser dinamicos (agentes trocam de papel)",
    ],
    "5 papeis canonicos. Serao vistos na DEMO (Slide 30). Pergunta: qual e mais dificil de automatizar? (costuma ser critico ou sintetizador)",
    9, T, "Papeis"
)

content_slide(
    "Normas e Instituicoes",
    [
        "Normas: regras compartilhadas (explicitas ou emergentes)",
        "Instituicoes: estruturas que mantem normas (ex: revisao por pares)",
        "Mecanismos: constitution, votacao, reputacao",
        "Tenso central: explicitas (controlaveis) vs emergentes (adaptaveis)",
    ],
    "Normas explicitas dao controle mas sao rigidas. Emergentes sao adaptaveis mas imprevisiveis. Sociedade precisa de espaco para emergencia.",
    10, T, "Normas"
)

content_slide(
    "Reputacao e Confianca",
    [
        "Reputacao: historico de interacoes de um agente",
        "Confianca: funcao da reputacao + contexto",
        "Analogia: PageRank - confianca propagada pela rede",
        "Sem reputacao: sociedade colapsa (free-riders, manipulacao)",
    ],
    "Sem reputacao, free-riders dominam. Reputacao e contextual e propagada - nao e metrica simples.",
    11, T, "Reputacao"
)

content_slide(
    "Modelos: Generative Agents, AgentVerse, ChatArena",
    [
        "Generative Agents (Stanford): 25 agentes, Smallville, vida social simulada",
        "AgentVerse (arXiv:2308.10848): framework multi-agente, papeis flexiveis",
        "ChatArena: ambientes de jogo/dialogo multi-agente",
        "Cada modelo resolve um problema diferente",
    ],
    "Generative Agents = paper/pesquisa. AgentVerse e ChatArena = frameworks. Nao sao equivalentes.",
    12, T, "Modelos"
)

content_slide(
    "Arquitetura de uma Sociedade (Diagrama)",
    [
        "Diagrama: 12-Diagrams/ETHAGT16/society.mmd",
        "Centro: sociedade com 5 papeis (Pesquisador/Critico/Sintetizador/Revisor/Editor)",
        "Normas governam as interacoes ('citar fontes', 'editor tem veto')",
        "Saida: comportamento emergente",
        "Normas podem vir de fora (designer) ou emergir de dentro",
    ],
    "Mostrar o .mmd. Destacar que normas sao externas (instituicao) ou emergem (sociedade). Emergencia nao e magica - e deterministica mas imprevisivel.",
    13, T, "Diagrama"
)

content_slide(
    "Grupo vs Instituicao vs Sociedade",
    [
        "Grupo: colaboracao ad-hoc, sem normas fixas - controle alto, adaptabilidade baixa",
        "Instituicao: papeis fixos, normas explicitas - meio-termo",
        "Sociedade: normas emergentes - controle baixo, adaptabilidade alta, risco alto",
        "80% dos problemas se resolvem em instituicao",
    ],
    "Sociedade e para quando se PRECISA de emergencia. Nao comecar pela sociedade - a maioria dos problemas nao precisa.",
    14, T, "Comparacao"
)

exercise_slide(
    "Exercicio - Comite Editorial",
    [
        "Cenario: simular um comite editorial de revista cientifica",
        "Em duplas (1 min): definir 5 papeis e responsabilidades",
        "Qual norma deve ser compartilhada por todos?",
        "Como a reputacao e acumulada?",
        "Dica: normas precisam ser operacionais (executaveis por LLM)",
    ],
    "1 min em duplas. Rebater normas vagas ('seja justo') - como um LLM verifica isso? Norma precisa ser operacional.",
    15, T
)

section_slide(2, "Simulacoes Sociais")

content_slide(
    "Sandbox Social - Smallville",
    [
        "Smallville: mundo virtual com 25 agentes",
        "Cada agente: perfil, rotina, memoria, objetivos",
        "Ambiente: mapas, locais, recursos compartilhados",
        "Resultado: agentes organizam festa de Valentine's Day sozinhos",
        "Fonte: Park et al., arXiv:2304.03442",
    ],
    "O experimento que abriu o campo. Festa de Valentine's Day emergiu - ninguem programou. Plausivel mas nao necessariamente preditivo.",
    17, T, "Smallville"
)

content_slide(
    "Como Smallville Funciona",
    [
        "Memory stream: log cronologico de experiencias",
        "Retrieval: relevancia + recencia + importancia",
        "Reflection: sintese de memorias em insights de alto nivel",
        "Planning: planos diarios (rotina + objetivos)",
        "Action: acao baseada em contexto + reflexao",
    ],
    "Motor de memoria. Reflection nao e summarization - e sintese semantica que produz novos insights. Analogia com memoria humana.",
    18, T, "Arquitetura"
)

content_slide(
    "Casos de Uso de Simulacoes Sociais",
    [
        "Policy simulation: impacto de lei (imposto, subsidio)",
        "Mercado: agentes compram/vendem, formacao de precos",
        "Opiniao publica: difusao de ideias, polarizacao",
        "Vantagem: testar hipoteses sem custo social",
        "Mas: validacao e dificil (sem ground truth)",
    ],
    "Em mercados com dados reais - tem valor preditivo. Em opiniao publica - so valor explicativo. Cuidado em alta stakes.",
    19, T, "Casos de Uso"
)

content_slide(
    "Limites e Criticas",
    [
        "Agentes muito coerentes vs humanos (menos ruido)",
        "Vieses dos LLMs se amplificam em sociedade",
        "Custo computacional (25+ agentes x N interacoes)",
        "Dificil validacao (sem ground truth)",
        "Pergunta: simulacao social com LLMs pode prever comportamento humano?",
    ],
    "Plausivel nao e correto. Sem validacao, simulacao gera hipoteses, nao predicoes. Honestidade e crucial.",
    20, T, "Limites"
)

exercise_slide(
    "Simulacao = Predicao?",
    [
        "Uma simulacao social com LLMs pode prever comportamento humano real?",
        "V/F: 'Sociedades de agentes sempre convergem.'",
        "Resposta: FALSO - podem polarizar, oscilar ou colapsar",
        "Depende de normas, reputacao e alinhamento",
    ],
    "Votacao rapida. Maioria erra para V. Usar como gancho: convergencia nao e garantida.",
    21, T
)

exercise_slide(
    "Exercicio - Home Office",
    [
        "Cenario: simular impacto de uma politica de home office",
        "Em duplas (1 min): definir 3 tipos de agentes, 2 normas, 1 metrica",
        "Qual resultado voce espera? Qual pode surpreender?",
        "Surpresa comum: junior sem mentoria presencial cai (emergencia)",
    ],
    "1 min em duplas. Rebater metricas vagas ('satisfacao') - como agente reporta de forma mensuravel?",
    22, T
)

section_slide(3, "Autonomous Research Systems")

content_slide(
    "Pipeline de Pesquisa Autonoma",
    [
        "Pipeline: pergunta -> literatura -> hipotese -> experimento -> analise -> relatorio",
        "Cada etapa pode ser um agente diferente",
        "Ou um agente que executa todas as etapas (AI Scientist)",
        "Desafio central: manter coerencia entre etapas",
        "HITL em pontos criticos",
    ],
    "Coerencia e o coracao. Sem loops de feedback, o sistema produz lixo plausivel. HITL nao e opcional em alta stakes.",
    24, T, "Pipeline"
, acronyms="HITL = Human-in-the-Loop — Humano no Ciclo")

content_slide(
    "Research Pipeline (Diagrama)",
    [
        "Diagrama: 12-Diagrams/ETHAGT16/research-pipeline.mmd",
        "Fluxo: pergunta -> revisao -> hipotese -> experimento -> execucao -> analise -> escrita",
        "Loop no fim: HITL revisa - aprovar ou revisar",
        "Loop de revisao e critico em producao",
    ],
    "Mostrar o .mmd. Destacar o loop HITL. Autonomia sem HITL em alta stakes e irresponsavel.",
    25, T, "Diagrama"
)

content_slide(
    "Sakana AI Scientist - Pesquisa ML End-to-End",
    [
        "Sakana AI: pesquisa ML do inicio ao fim",
        "Etapas: ideacao -> literatura -> codigo -> experimento -> paper",
        "Custo: ~$15 por paper",
        "Resultado: papers aceitos em workshop ICLR",
        "Fonte: Lu et al., arXiv:2408.06292",
    ],
    "Primeiro caso de pesquisa autonoma passando em peer review. Workshop (nao main conference). Auxilia, nao substitui pesquisadores.",
    26, T, "AI Scientist"
)

content_slide(
    "Como o AI Scientist Funciona",
    [
        "Stage 1: Ideacao (LLM gera ideias de pesquisa)",
        "Stage 2: Experimentacao (LLM escreve codigo, roda experimentos)",
        "Stage 3: Paper writing (LLM escreve paper em LaTeX)",
        "Stage 4: Review (LLM revisa como reviewer)",
        "Agente decide: continuar, refinar ou abandonar",
    ],
    "4 stages com loop de decisao. Review automatico (LLM revisando LLM) tem vieses compartilhados - erros sutis passam.",
    27, T, "Arquitetura"
)

content_slide(
    "DeepMind AlphaEvolve - Evolucao de Algoritmos",
    [
        "AlphaEvolve: evolucao de algoritmos via LLM + avaliacao automatica",
        "LLM propoe mutacoes -> avaliador testa -> mantem melhores",
        "Resultado: descobriu novo algoritmo de multiplicacao de matrizes",
        "Licao: evolucao + avaliacao automatica = descoberta",
    ],
    "Diferente do AI Scientist: evolucao em vez de pipeline. Sem avaliador automatico (fitness function), e so gerador aleatorio.",
    28, T, "AlphaEvolve"
)

content_slide(
    "Multi-Agent Research Labs",
    [
        "Em vez de 1 agente: time de agentes especializados",
        "Pesquisador, programador, revisor, escritor",
        "Vantagem: especializacao + divisao de trabalho",
        "Desafio: coordenacao, comunicacao, overhead",
        "Sweet spot: 3-5 papeis",
    ],
    "Mais agentes = mais overhead, nao necessariamente mais capacidade. Sweet spot costuma ser 3-5.",
    29, T, "Multi-Agent"
)

code_slide(
    "DEMO - Mini Sociedade de 5 Agentes",
    "5 agentes: pesquisador, critico, sintetizador, revisor, editor\n"
    "Tarefa: relatorio sobre 'impacto de agentes na educacao'\n\n"
    "$ python main.py --task \"impacto de agentes na educacao\"\n"
    "[pesquisador] buscando fontes...\n"
    "[critico]      fonte A e fraca - sem baseline\n"
    "[sintetizador] integrando 3 fontes validas\n"
    "[revisor]      coerencia OK, citacoes OK\n"
    "[editor]       versao final pronta",
    "DEMO principal. Observar: (1) divergencias, (2) consenso, (3) qualidade. Se API falhar, usar screenshot pre-gravado. Nao consertar ao vivo.",
    30, T, "DEMO"
)

exercise_slide(
    "Pergunta da DEMO",
    [
        "Qual papel foi mais critico para o resultado?",
        "O que aconteceria se removesssemos o critico?",
        "E se o editor tivesse poder de veto?",
        "1 min em duplas",
    ],
    "Analise contrafactual. Sem critico - qualidade cai. Editor com veto - convergencia rapida mas menos diversidade. Sintetizador e subestimado.",
    31, T
)

content_slide(
    "AI Scientist - O Que Funcionou",
    [
        "Papers aceitos em workshop ICLR (peer review real)",
        "Custo baixo (~$15/paper) viabiliza exploracao",
        "Pipeline estruturado garante reprodutibilidade",
        "Review automatico identifica falhas antes da submissao",
        "Licao: estrutura > brilho individual",
    ],
    "Sucessos reais mas precisos: workshop ICLR (nao main conference). Estrutura vence brilho individual.",
    32, T, "Estudo de Caso"
)

content_slide(
    "AI Scientist - O Que Falhou",
    [
        "Experimentos as vezes nao convergem (codigo com bugs)",
        "Papers com contribuicao marginal (novidade baixa)",
        "Alucinacao em revisao de literatura",
        "Dificil avaliar 'qualidade cientifica' sem humano",
        "Licao: autonomia != qualidade",
    ],
    "Licao critica: autonomia nao implica qualidade. Sistema autonomo que produz lixo plausivel e PIOR que assistido.",
    33, T, "Estudo de Caso"
)

section_slide(4, "Emergencia e Alinhamento")

content_slide(
    "Comportamento Emergente - A Soma != Partes",
    [
        "Emergencia: padroes que surgem da interacao, nao do individuo",
        "A soma e diferente das partes",
        "Positivos: cooperacao espontanea, formacao de normas, divisao de trabalho",
        "Negativos: conluio, discriminacao, corrida armamentista",
        "Pergunta: quando a soma e melhor? Quando e pior?",
    ],
    "Emergencia nao e magica - e deterministica dadas as regras. Mas imprevisivel para o designer em sistemas grandes.",
    35, T, "Emergencia"
)

content_slide(
    "Emergencia (Diagrama)",
    [
        "Diagrama: 12-Diagrams/ETHAGT16/emergence.mmd",
        "Agentes individuais -> interacoes locais",
        "Interacoes -> padroes emergentes (nao programados)",
        "Padroes -> propriedades coletivas: coordenacao, polarizacao, especializacao, echo chamber",
    ],
    "As propriedades coletivas nao existem nos agentes - so na coletividade. Esse e o ponto-chave.",
    36, T, "Diagrama"
)

content_slide(
    "Emergencia Desejada vs Indesejada",
    [
        "Desejada: cooperacao, divisao de trabalho, inovacao, formacao de normas",
        "Indesejada: conluio, discriminacao, echo chamber, corrida armamentista",
        "Fator-chave: alinhamento de valores + supervisao",
        "Sem supervisao: emergencia tende ao extremo",
    ],
    "Conluio = agentes maximizam retorno individual contra objetivo coletivo. Mais agentes sem alinhamento = emergencia EXTREMA, boa ou ruim.",
    37, T, "Emergencia"
)

content_slide(
    "Alinhamento de Valores em Populacoes",
    [
        "Alinhamento individual != alinhamento coletivo",
        "Constitution para sociedades: regras globais",
        "Votacao e consenso como mecanismos",
        "Tenso: normas rigidas (seguras) vs flexiveis (adaptaveis)",
        "Referencia: Constitutional AI",
    ],
    "Insight mais profundo: alinhar cada agente individualmente NAO alinha a sociedade. Constituicao e necessaria.",
    38, T, "Alinhamento"
)

content_slide(
    "Avaliacao de Emergencia",
    [
        "Metricas de divergencia: distancia do comportamento esperado",
        "Metricas de surpresa: comportamento nao previsto",
        "Metricas de convergencia: estabiliza ou oscila?",
        "Sem metricas: emergencia e caixa preta",
        "Em producao: monitoramento continuo",
    ],
    "Logs geram dados, nao insights. Precisa de metricas agregadas que capturem propriedades coletivas.",
    39, T, "Avaliacao"
)

exercise_slide(
    "Quando Emergencia e Indesejada?",
    [
        "Quando comportamento emergente e indesejado?",
        "3 cenarios rapidos: em qual voce interviria?",
        "Regra: stakes + irreversibilidade",
        "Discussao aberta (1 min)",
    ],
    "Depende das stakes. Em pesquisa - baixas stakes, deixa rolar. Em politica publica - altas stakes, intervem cedo.",
    40, T
)

exercise_slide(
    "Exercicio - Detectando Emergencia Indesejada",
    [
        "Trace de 5 agentes, 10 interacoes",
        "Todos convergem para 'nada e confiavel' - paralisia",
        "Padrao: echo chamber de critica",
        "Em duplas (1 min): propor 2 correcoes",
        "Dica: nao remover o critico - ajustar seu PODER",
    ],
    "1 min em duplas. Correcoes: norma (critica com alternativa), HITL (intervir apos 5 rodadas), peso (limitar veto).",
    41, T
)

section_slide(5, "Fronteira e Etica")

content_slide(
    "Onde a Pesquisa Esta Agora",
    [
        "AI Scientist (Sakana, 2024): pesquisa ML automatizada",
        "AlphaEvolve (DeepMind, 2024): evolucao de algoritmos",
        "AutoGen research: framework multi-agente para pesquisa",
        "Swarm research: coordenacao de grandes populacoes",
        "Tendencia: 1 agente -> times -> sociedades autonomas",
    ],
    "Estamos na transicao de times para sociedades. Em producao ainda e experimental.",
    43, T, "Fronteira"
)

content_slide(
    "Questoes Eticas",
    [
        "Automacao de pesquisa: autoria? revisao?",
        "Responsible AI: vies amplificado em sociedades",
        "Autopropagacao: agente cria copias de si mesmo",
        "Auto-modificacao: agente altera seu proprio codigo",
        "Pergunta: onde voce traca a linha entre util e perigoso?",
    ],
    "Autopropagacao e auto-modificacao sao os mais perigosos. Alunos vao construir esses sistemas - sao a linha de frente.",
    44, T, "Etica"
)

content_slide(
    "O Que NAO Fazer",
    [
        "Sistemas sem HITL em alta stakes",
        "Auto-modificacao irrestrita",
        "Deploy sem avaliacao de risco",
        "Autopropagacao sem limites",
        "Sociedades sem constitution ou normas",
        "Sem metricas de emergencia",
    ],
    "6 proibicoes. Cada uma corresponde a incidente real. Autopropagacao irrestrita ja consumiu orcamentos inteiros.",
    45, T, "Anti-padroes"
)

exercise_slide(
    "Exercicio - Avaliar AI Scientist",
    [
        "Cenario: AI Scientist gerou paper sobre 'otimizacao de prompts'",
        "Em grupos (2 min): definir 3 criterios de qualidade cientifica",
        "Exemplos: reprodutibilidade, novidade, corretude metodologica",
        "Dica: criterios precisam ser operacionais (mensuraveis)",
        "1 min compartilhar",
    ],
    "2 min em grupos de 3. Rebater 'qualidade da escrita' - escrita pode ser boa mas conteudo errado. Foco no CONTEUDO.",
    46, T
)

content_slide(
    "Resumo da Aula",
    [
        "Sociedades de agentes = papeis + normas + reputacao + emergencia",
        "Smallville = caso canonico de simulacao social",
        "AI Scientist = fronteira de pesquisa autonoma",
        "Emergencia pode ser desejada (cooperacao) ou indesejada (conluio)",
        "Alinhamento coletivo != alinhamento individual",
        "Etica: supervisao, avaliacao de risco, responsible AI",
    ],
    "6 pontos para memorizar. Fechar o arco aberto no Slide 5.",
    47, T, "Resumo"
)

content_slide(
    "Quiz - 3 Perguntas",
    [
        "P1: Qual pipeline o AI Scientist automatiza?",
        "  -> B) Ideacao -> literatura -> codigo -> experimento -> paper -> review",
        "P2: O que e comportamento emergente?",
        "  -> B) Padroes que surgem da interacao, nao do individuo",
        "P3: V/F 'Sociedades de agentes sempre convergem'",
        "  -> FALSO - podem polarizar, oscilar ou colapsar",
    ],
    "Individual, sem consulta, 1 min. >=2 acertos = compreensao basica. Maioria erra P3 - usar como gancho.",
    48, T, "Quiz"
)

content_slide(
    "Capstone e Referencias",
    [
        "ETHAGT90 - Capstone: projeto final integrador",
        "Espectro completo: ETHAGT01 (Augmented LLM) -> ETHAGT16 (Sociedades)",
        "Leitura obrigatoria: arXiv:2304.03442 (Generative Agents)",
        "Leitura obrigatoria: arXiv:2408.06292 (AI Scientist)",
        "Leitura obrigatoria: arXiv:2308.10848 (AgentVerse)",
        "Complementar: DeepMind AlphaEvolve (2024)",
        "Convite: pesquisa, produto ou impacto social?",
    ],
    "O Projeto deste modulo e o esboco do Capstone. Fechar com a pergunta: para onde voce quer levar isso?",
    49, T, "Referencias"
)

s = title_slide(
    "Perguntas?",
    "Proximo modulo: ETHAGT90 - Capstone\nVoces estao na fronteira. Agora e hora de construir.",
    "ETHAGT16 - Q&A"
)

out = os.path.dirname(os.path.abspath(__file__))
out_file = os.path.join(out, "ETHAGT16-slides.pptx")
prs.save(out_file)
print(f"OK: {out_file} ({len(prs.slides)} slides)")
